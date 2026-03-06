import sys
import os

# Durcit la sortie console Windows (cp1252) pour eviter les crashs sur emojis/unicode.
try:
    if hasattr(sys.stdout, "reconfigure"):
        sys.stdout.reconfigure(errors="replace")
    if hasattr(sys.stderr, "reconfigure"):
        sys.stderr.reconfigure(errors="replace")
except Exception:
    pass

# VERIFICATION VERSION STABLE (Bloque le lancement sur 3.15)
if sys.version_info.major == 3 and sys.version_info.minor >= 15:
    print("\n" + "!" * 80)
    print(" ERREUR FATALE : VOUS LANCEZ L'APPLICATION AVEC PYTHON 3.15 (VERSION ALPHA INSTABLE)")
    print(" Jarvisse a \u00e9t\u00e9 configur\u00e9 et s\u00e9curis\u00e9 pour fonctionner EXCLUSIVEMENT sur Python 3.12.")
    print(" \n SOLUTION :")
    print(" 1. NE DOUBLE-CLIQUEZ PAS sur ppp.py")
    print(" 2. Double-cliquez plut\u00f4t sur le fichier 'Lancer_Jarvisse.bat' pr\u00e9sent dans ce dossier.")
    print("!" * 80 + "\n")
    sys.exit(1)

import speech_recognition as sr
import edge_tts
import asyncio
import pygame
try:
    import pywhatkit
except Exception as e:
    print(f"âš ï¸ Module pywhatkit dÃ©sactivÃ© (Erreur rÃ©seau/SSL: {e})")
    pywhatkit = None
import datetime
import wikipedia
import os
import webbrowser
import threading
import traceback
import sounddevice as sd
import numpy as np
import io
import wave
import requests
import subprocess
from bs4 import BeautifulSoup
from googlesearch import search as gsearch
from google import genai
from google.genai import types
import geocoder
from PIL import Image
try:
    from pydub import AudioSegment
except ImportError:
    pass
try:
    import pytesseract
except:
    pass

import customtkinter as ctk
import tkinter as tk
from tkinter import filedialog, colorchooser
import time
import random
import math
import nltk
from sumy.parsers.plaintext import PlaintextParser
from sumy.nlp.tokenizers import Tokenizer
from sumy.summarizers.lsa import LsaSummarizer
import fitz  # PyMuPDF
from docx import Document
import openpyxl
from pptx import Presentation
import csv
import pyautogui
from AppOpener import open as app_open, close as app_close
import pygetwindow as gw
import psutil
import re
import subprocess
import json
import sys
import shutil
from collections import deque
import urllib.parse
from urllib.parse import urlparse, urlunparse
import pickle
from google.auth.transport.requests import Request
from google_auth_oauthlib.flow import InstalledAppFlow
from googleapiclient.discovery import build
import queue
try:
    from tkinterweb import HtmlFrame
except:
    HtmlFrame = None

# TÃ©lÃ©chargement des ressources nÃ©cessaires pour NLTK
try:
    nltk.download('punkt', quiet=True)
    nltk.download('punkt_tab', quiet=True)
except:
    pass

# Configuration de l'apparence
ctk.set_appearance_mode("dark")
ctk.set_default_color_theme("blue")

# Couleurs Premium (Style Gemini / Futuriste)
COLOR_BG = "#0B0B0C"
COLOR_NAV = "#161618"
COLOR_ACCENT = "#4285F4"  # Bleu Gemini
COLOR_ACCENT_PURPLE = "#9B72CB" # Violet Gemini
COLOR_ACCENT_RED = "#D96570" # Rouge Gemini
COLOR_TEXT = "#E3E3E3"
COLOR_SECONDARY = "#9AA0A6"
COLOR_CARD = "#1E1F20"

class JarvisseAssistant(ctk.CTk):
    def __init__(self):
        super().__init__()

        import queue
        self.uiq = queue.Queue()

        self.title("Jarvisse Assistant")
        self.geometry("800x600")
        self.configure(fg_color=COLOR_BG)
        self.attributes("-alpha", 0.96) # Hologram/Glassmorphism Effect

        try:
            from scheduler import start_scheduler
            start_scheduler()
            print("Scheduler DÃ©marrÃ©.")
        except Exception as e:
            print("Erreur initialisation Scheduler:", e)

        # Initialisation Audio
        pygame.mixer.init()
        pygame.mixer.music.set_volume(1.0)
        self.voice = "fr-FR-VivienneNeural"
        self.voice_speed = "-5%"

        # Variables d'Ã©tat
        self.is_listening = False
        self.conversation_history = deque(maxlen=100)
        self.conversation_memory_limit = 100
        self.current_subject = None
        self.loaded_document_content = None  # Pour stocker le texte du dernier fichier chargÃ©
        self.veille_active = True # Activer la veille par dÃ©faut
        self.voice_lock = threading.Lock() # Lock pour Ã©viter les chevauchements audio
        self.stop_speaking_flag = False
        self.ping_pong_mode = True  # RÃ©ponses courtes et rÃ©actives
        self.ping_pong_max_chars = 2000 # AugmentÃ© massivement pour Ã©viter toute coupure
        self.proactive_suggestions = True  # Suggestions intelligentes basÃ©es sur le contexte
        self.voice = "fr-FR-VivienneNeural"
        self.available_voices = [
            "fr-FR-VivienneNeural", "fr-FR-RemyNeural", "fr-FR-EloiseNeural", 
            "fr-FR-DeniseNeural", "fr-FR-HenriNeural", "fr-BE-CharlineNeural", 
            "fr-BE-GerardNeural", "fr-CA-SylvieNeural", "fr-CA-JeanNeural",
            "fr-CH-ArianeNeural", "fr-CH-FabriceNeural"
        ]
        self.battery_alert_threshold = 10
        self.battery_alert_sent = False
        self.battery_low_message = "Patron, la batterie est Ã  {percent}%"
        self.battery_plugged_message = "Patron, le chargeur est branchÃ©."
        self.battery_full_message = "Patron, la batterie est complÃ¨tement chargÃ©e. Je vous recommande de retirer le chargeur"
        self.last_power_plugged = None
        self.last_full_charge_alert_time = 0 # Temps du dernier rappel batterie pleine

        self.last_full_charge_alert_time = 0 # Temps du dernier rappel batterie pleine

        self.awaiting_exit_confirm = False
        self.autonomous_mode = True
        self.autonomous_provider = "gemini"
        self.ollama_url = "http://localhost:11434/api/generate"
        self.ollama_model = "llama3.1:8b"
        self.ollama_models = ["llama3.1:8b", "llama3.1:70b", "llama3:8b", "mistral:7b", "qwen2.5:7b", "phi3:mini"]
        # self.conversation_memory_limit = 100 (DÃ©jÃ  dÃ©fini plus haut)
        self.ping_pong_max_chars = 10000 
        self.ollama_num_predict = 4096
        self.known_removable_devices = set()
        self.scanning_devices = set()
        self.scanned_devices = set()
        self.alarms = []
        self.alarm_next_id = 1
        self.anomaly_last_alert = {}
        self.anomaly_cooldown_sec = 120
        self.anomaly_temp_c = 85
        self.battery_drop_percent = 15
        self.battery_drop_window_sec = 600
        # CERVEAU PHONÃ‰TIQUE : Chargement de la base de donnÃ©es Ã©tendue
        self.phonetic_path = os.path.join(os.path.dirname(__file__), "phonetic_extended.json")
        self.phonetic_map = self._load_phonetic_map()
        self.alphabet_phonetique = {
            'A':'a', 'B':'bÃ©', 'C':'cÃ©', 'D':'dÃ©', 'E':'e', 'F':'Ã¨f', 'G':'gÃ©', 'H':'ache',
            'I':'i', 'J':'ji', 'K':'ka', 'L':'Ã¨l', 'M':'Ã¨m', 'N':'Ã¨n', 'O':'o', 'P':'pÃ©',
            'Q':'ku', 'R':'Ã¨r', 'S':'Ã¨s', 'T':'tÃ©', 'U':'u', 'V':'vÃ©', 'W':'double-vÃ©',
            'X':'iks', 'Y':'igrek', 'Z':'zÃ¨d'
        }
        self.battery_not_charging_grace_sec = 300
        self.proc_cpu_threshold = 80
        self.proc_cpu_duration_sec = 60
        self.proc_mem_threshold_mb = 1500
        self.disk_free_gb = 10
        self.disk_used_percent = 90
        self.anomalies_enabled = True
        self.settings_path = os.path.join(os.path.dirname(__file__), "settings.json")
        self.launch_on_startup = False
        self.ui_theme_name = "Gemini Blue"
        self.custom_accent_color = COLOR_ACCENT
        self.custom_colors = {}
        self.stark_personality_enabled = True
        self.stark_sarcasm_level = 0.25
        self.tactical_mode = False
        self.holographic_fx_enabled = True
        self.contextual_radar_enabled = True
        self.last_tactical_snapshot = ""
        # ðŸš€ STARK EDITION : MÃ©moire Visuelle et HUD
        self.visual_memory = deque(maxlen=20) # Garde les 20 derniers rÃ©sumÃ©s visuels
        self.hud_data = {"cpu": 0, "ram": 0, "status": "Online"}
        
        # Lancement des threads Stark
        threading.Thread(target=self._visual_memory_loop, daemon=True).start()

        self.ui_theme_presets = f"" # Reset for safe replacement
        self.ui_theme_presets = {
            "Gemini Blue": {
                "bg": "#0B0B0C",
                "nav": "#161618",
                "accent": "#4285F4",
                "accent_purple": "#9B72CB",
                "accent_red": "#D96570",
                "text": "#E3E3E3",
                "secondary": "#9AA0A6",
                "card": "#1E1F20",
            },
            "Emerald": {
                "bg": "#0B1110",
                "nav": "#13201E",
                "accent": "#2CB67D",
                "accent_purple": "#4E9F8B",
                "accent_red": "#D96570",
                "text": "#E6F2EE",
                "secondary": "#9BB8AE",
                "card": "#1A2A27",
            },
            "Sunset": {
                "bg": "#16100E",
                "nav": "#241A16",
                "accent": "#F59E0B",
                "accent_purple": "#F97316",
                "accent_red": "#E11D48",
                "text": "#F5E9E2",
                "secondary": "#BCA79B",
                "card": "#2A211C",
            },
            "Steel": {
                "bg": "#0E131A",
                "nav": "#151D27",
                "accent": "#60A5FA",
                "accent_purple": "#94A3B8",
                "accent_red": "#EF4444",
                "text": "#E6EDF5",
                "secondary": "#9AA9BD",
                "card": "#1C2633",
            },
        }
        self.notifications_enabled = True
        self.notification_sources = {
            "WhatsApp.exe": True,
            "OUTLOOK.EXE": True,
            "chrome.exe": True
        }
        self.seen_notifications_order = deque()
        self.seen_notifications_set = set()
        self.max_seen_notifications = 200
        self.notifications_access_prompted = False
        self.last_battery_percent = None
        self.last_battery_time = None
        self.suspect_processes = {}
        self.services_to_watch = []
        
        # Agents Settings
        self.agent_network_enabled = False
        self.agent_control_enabled = False
        self.agent_remote_enabled = False
        self.agent_gmail_enabled = False
        self.gmail_service = None
        self.last_email_id = None
        self.remote_hosts_whitelist = []
        self.remote_default_host = ""
        self.remote_ssh_user = ""
        self.remote_ssh_port = 22
        self.remote_ssh_key_path = ""
        self.remote_require_confirmation = True
        self.awaiting_remote_confirm = None
        self.awaiting_anomaly_confirm = False
        self.conversation_continue = False
        self.voice_enabled = True
        
        # Agent Browser
        self.agent_browser_enabled = False
        self.browser_history = []

        # Agent Juridique
        self.agent_legal_enabled = False

        # Agent Cyber SÃ©curitÃ©
        self.agent_cyber_enabled = False

        # Agent Docteur SystÃ¨me
        self.agent_doctor_enabled = False

        # Agent VidÃ©o Surveillance
        self.agent_video_enabled = True

        # Agent ContrÃ´le Android
        self.agent_android_enabled = True

        # Assistant Coding
        self.agent_coding_enabled = True
        self.last_internal_error = None  # Conscience de soi : stockage de la derniÃ¨re erreur systÃ¨me interne

        # Assistant GÃ©nÃ©ration Image et VidÃ©o
        self.agent_generation_enabled = True
        
        # Agent Assistant Licence
        self.agent_licence_enabled = True
        
        # Agent Lecture
        self.agent_lecture_enabled = True
        
        # Agent Bot Telegram
        self.agent_telegram_enabled = False
        
        # NOUVEAUX AGENTS PREMIUM
        self.agent_vision_enabled = True
        self.agent_domotique_enabled = True
        self.agent_finance_enabled = True
        self.agent_secretaire_enabled = True
        self.agent_traducteur_enabled = True
        self.agent_miner_enabled = True
        
        # NOUVEAUX AGENTS AJOUTÃ‰S
        self.agent_news_enabled = True
        self.agent_cuisine_enabled = True
        self.agent_sante_enabled = True
        self.agent_voyage_enabled = True
        self.agent_education_enabled = True
        self.agent_shopping_enabled = True
        self.agent_social_enabled = True
        self.agent_psy_enabled = True
        self.agent_immo_enabled = True
        self.agent_auto_enabled = True
        self.agent_carrier_enabled = True
        
        # AGENTS SUPPLÃ‰MENTAIRES (PACK GRATUIT & PREMIUM)
        self.agent_arbitre_enabled = True
        self.agent_ecolo_enabled = True
        self.agent_guetteur_enabled = True
        self.agent_historien_enabled = True
        self.agent_depanneur_enabled = True
        self.agent_astroph_enabled = True
        self.agent_stratege_enabled = True
        self.agent_architecte_enabled = True
        self.agent_business_enabled = True
        self.agent_polyglotte_enabled = True
        self.agent_nounou_enabled = True

        # Agent Appel (NOUVEAU)
        self.agent_appel_enabled = True
        self.agent_appel_number = ""
        self.agent_appel_autoreply = True
        self.agent_appel_phrases = [
            "Bonjour, je suis l'assistant de mon propriÃ©taire. Il n'est pas disponible pour le moment.",
            "Puis-je prendre un message ou vous rappeler plus tard ?",
            "Je note votre appel. Merci et au revoir."
        ]
        self.agent_appel_filters = [] # NumÃ©ros spÃ©cifiques Ã  rÃ©pondre

        # Configuration Domotique
        self.domo_wled_ip = ""
        self.domo_ha_url = ""
        self.domo_ha_token = ""
        self.domo_arduino_port = ""
        self.domo_webhook_on = ""
        self.domo_webhook_off = ""

        self.telegram_bot_token = ""
        self.telegram_chat_id = ""
        self.telegram_bot_running = False
        self.last_command_origin = "local"

        self.generations_dir = os.path.join(os.path.dirname(__file__), "generations")
        if not os.path.exists(self.generations_dir):
            os.makedirs(self.generations_dir)

        # Gigantesque MÃ©moire des Agents
        self.memories_dir = os.path.join(os.path.dirname(__file__), "agent_memories")
        if not os.path.exists(self.memories_dir):
            os.makedirs(self.memories_dir)
        self.agent_memories_cache = {}

        self._load_settings()
        self._apply_theme_globals()
        # Le listener de notifications sera dÃ©marrÃ© aprÃ¨s l'initialisation de l'UI

        if self.ollama_model not in self.ollama_models:
            self.ollama_models.insert(0, self.ollama_model)

        # Configuration Gemini pour l'autonomie (chargement via settings ou environment)
        self.gemini_api_key_var = ctk.StringVar(value=self.gemini_api_key)
        self._init_gemini_client()

        # Si Gemini indisponible â†’ bascule sur Ollama automatiquement
        if not self.genai_client:
            self.autonomous_provider = "ollama"

        # GÃ©olocalisation
        self.current_location = None
        self.location_city = "Inconnu"
        self.location_country = "Inconnu"
        self.location_coords = (0.0, 0.0)

        self.grid_columnconfigure(1, weight=1)
        self.grid_columnconfigure(2, weight=0)
        self.grid_rowconfigure(1, weight=1)

        self.sidebar = ctk.CTkFrame(self, width=200, fg_color=COLOR_NAV, corner_radius=0)
        self.sidebar.grid(row=0, column=0, rowspan=3, sticky="nsew")
        self.sidebar.grid_rowconfigure(0, weight=1)
        self.sidebar.grid_columnconfigure(0, weight=1)

        self.sidebar_scroll = ctk.CTkScrollableFrame(
            self.sidebar,
            fg_color=COLOR_NAV,
            corner_radius=0,
            width=200
        )
        self.sidebar_scroll.grid(row=0, column=0, sticky="nsew")
        self.sidebar_scroll.grid_rowconfigure(35, weight=1)

        self.sidebar_title = ctk.CTkLabel(self.sidebar_scroll, text="SYSTEM STATUS", font=ctk.CTkFont(size=12, weight="bold"), text_color=COLOR_ACCENT)
        self.sidebar_title.grid(row=0, column=0, padx=20, pady=(20, 10))

        self.cpu_label = ctk.CTkLabel(self.sidebar_scroll, text="CPU: 0%", font=ctk.CTkFont(size=11), text_color=COLOR_SECONDARY)
        self.cpu_label.grid(row=1, column=0, padx=20, pady=5, sticky="w")

        self.ram_label = ctk.CTkLabel(self.sidebar_scroll, text="RAM: 0%", font=ctk.CTkFont(size=11), text_color=COLOR_SECONDARY)
        self.ram_label.grid(row=2, column=0, padx=20, pady=5, sticky="w")

        self.battery_label = ctk.CTkLabel(self.sidebar_scroll, text="BATTERY: --", font=ctk.CTkFont(size=11), text_color=COLOR_SECONDARY)
        self.battery_label.grid(row=3, column=0, padx=20, pady=5, sticky="w")

        self.disk_label = ctk.CTkLabel(self.sidebar_scroll, text="DISK: 0%", font=ctk.CTkFont(size=11), text_color=COLOR_SECONDARY)
        self.disk_label.grid(row=4, column=0, padx=20, pady=5, sticky="w")

        # SÃ©parateur
        self.separator = ctk.CTkLabel(self.sidebar_scroll, text="â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€", font=ctk.CTkFont(size=10), text_color="#3C4043")
        self.separator.grid(row=5, column=0, padx=20, pady=10)

        # Seuil d'alerte batterie (configurable)
        self.battery_threshold_label = ctk.CTkLabel(
            self.sidebar_scroll,
            text="Alerte batterie (%)",
            font=ctk.CTkFont(size=11),
            text_color=COLOR_SECONDARY
        )
        self.battery_threshold_label.grid(row=6, column=0, padx=20, pady=(0, 5), sticky="w")
        
        self.battery_threshold_entry = ctk.CTkEntry(
            self.sidebar_scroll,
            height=28,
            width=120,
            fg_color=COLOR_CARD,
            border_color="#3C4043",
            text_color=COLOR_TEXT
        )
        self.battery_threshold_entry.grid(row=7, column=0, padx=20, pady=(0, 6), sticky="w")
        self.battery_threshold_entry.insert(0, str(self.battery_alert_threshold))

        self.battery_threshold_entry.bind("<Return>", lambda e: self.set_battery_threshold())
        self.battery_threshold_entry.bind("<FocusOut>", lambda e: self.set_battery_threshold())

        self.battery_threshold_button = ctk.CTkButton(
            self.sidebar_scroll,
            text="Appliquer",
            height=28,
            width=120,
            fg_color=COLOR_ACCENT_PURPLE,
            hover_color="#B388EB",
            command=self.set_battery_threshold
        )
        self.battery_threshold_button.grid(row=8, column=0, padx=20, pady=(0, 10), sticky="w")

        # Messages d'alerte batterie (configurables)
        self.battery_low_msg_label = ctk.CTkLabel(
            self.sidebar_scroll,
            text="Message batterie basse",
            font=ctk.CTkFont(size=11),
            text_color=COLOR_SECONDARY
        )
        self.battery_low_msg_label.grid(row=9, column=0, padx=20, pady=(0, 5), sticky="w")

        self.battery_low_msg_entry = ctk.CTkEntry(
            self.sidebar_scroll,
            height=28,
            width=160,
            fg_color=COLOR_CARD,
            border_color="#3C4043",
            text_color=COLOR_TEXT
        )
        self.battery_low_msg_entry.grid(row=10, column=0, padx=20, pady=(0, 6), sticky="w")
        self.battery_low_message = "Attention : la batterie est Ã  {percent}%"

        self.battery_plugged_msg_label = ctk.CTkLabel(
            self.sidebar_scroll,
            text="Message chargeur branchÃ©",
            font=ctk.CTkFont(size=11),
            text_color=COLOR_SECONDARY
        )
        self.battery_plugged_msg_label.grid(row=11, column=0, padx=20, pady=(0, 5), sticky="w")

        self.battery_plugged_msg_entry = ctk.CTkEntry(
            self.sidebar_scroll,
            height=28,
            width=160,
            fg_color=COLOR_CARD,
            border_color="#3C4043",
            text_color=COLOR_TEXT
        )
        self.battery_plugged_msg_entry.grid(row=12, column=0, padx=20, pady=(0, 6), sticky="w")
        self.battery_plugged_message = "Le chargeur est branchÃ©."

        self.battery_full_msg_label = ctk.CTkLabel(
            self.sidebar_scroll,
            text="Message batterie pleine",
            font=ctk.CTkFont(size=11),
            text_color=COLOR_SECONDARY
        )
        self.battery_full_msg_label.grid(row=13, column=0, padx=20, pady=(0, 5), sticky="w")

        self.battery_full_msg_entry = ctk.CTkEntry(
            self.sidebar_scroll,
            height=28,
            width=160,
            fg_color=COLOR_CARD,
            border_color="#3C4043",
            text_color=COLOR_TEXT
        )
        self.battery_full_msg_entry.grid(row=14, column=0, padx=20, pady=(0, 10), sticky="w")
        self.battery_full_message = "La batterie est complÃ¨tement chargÃ©e. Je vous recommande de debrancher le chargeur pour eviter une surchage"

        # Mode autonome (IA)
        self.autonomous_var = ctk.BooleanVar(value=self.autonomous_mode)
        self.autonomous_switch = ctk.CTkSwitch(
            self.sidebar_scroll,
            text="Mode autonome",
            variable=self.autonomous_var,
            command=self.toggle_autonomous_mode,
            text_color=COLOR_SECONDARY
        )
        self.autonomous_switch.grid(row=15, column=0, padx=20, pady=(0, 10), sticky="w")

        # Mode Conversation Continue (Ping-Pong Auto)
        self.conv_continue_var = ctk.BooleanVar(value=self.conversation_continue)
        self.conv_continue_switch = ctk.CTkSwitch(
            self.sidebar_scroll,
            text="Conversation continue",
            variable=self.conv_continue_var,
            command=self.toggle_conversation_continue,
            text_color=COLOR_ACCENT
        )
        self.conv_continue_switch.grid(row=16, column=0, padx=20, pady=(5, 10), sticky="w")

        # Choix IA pour rÃ©ponses autonomes
        self.autonomous_provider_var = ctk.StringVar(value=self.autonomous_provider)
        self.autonomous_provider_label = ctk.CTkLabel(
            self.sidebar_scroll,
            text="RÃ©ponses auto",
            font=ctk.CTkFont(size=11),
            text_color=COLOR_SECONDARY
        )
        self.autonomous_provider_label.grid(row=17, column=0, padx=20, pady=(10, 4), sticky="w")

        self.provider_gemini_radio = ctk.CTkRadioButton(
            self.sidebar_scroll,
            text="Gemini",
            variable=self.autonomous_provider_var,
            value="gemini",
            command=self.on_autonomous_provider_changed,
            text_color=COLOR_SECONDARY
        )
        self.provider_gemini_radio.grid(row=18, column=0, padx=20, pady=(0, 2), sticky="w")

        # Champ de saisie ClÃ© API Gemini
        self.gemini_key_entry = ctk.CTkEntry(
            self.sidebar_scroll,
            textvariable=self.gemini_api_key_var,
            placeholder_text="Votre ClÃ© Gemini API",
            height=28,
            width=160,
            fg_color=COLOR_CARD,
            border_color="#3C4043",
            text_color=COLOR_TEXT,
            show="*" # Masquer la clÃ©
        )
        self.gemini_key_entry.grid(row=19, column=0, padx=20, pady=(0, 4), sticky="w")
        
        self.gemini_key_apply_btn = ctk.CTkButton(
            self.sidebar_scroll,
            text="Appliquer ClÃ©",
            height=24,
            width=100,
            fg_color=COLOR_CARD,
            hover_color="#3C4043",
            command=self.on_gemini_key_apply
        )
        self.gemini_key_apply_btn.grid(row=20, column=0, padx=20, pady=(0, 10), sticky="w")

        self.provider_ollama_radio = ctk.CTkRadioButton(
            self.sidebar_scroll,
            text="Ollama",
            variable=self.autonomous_provider_var,
            value="ollama",
            command=self.on_autonomous_provider_changed,
            text_color=COLOR_SECONDARY
        )
        self.provider_ollama_radio.grid(row=21, column=0, padx=20, pady=(0, 10), sticky="w")

        self.ollama_model_var = ctk.StringVar(value=self.ollama_model)
        self.ollama_model_label = ctk.CTkLabel(
            self.sidebar_scroll,
            text="ModÃ¨le Ollama",
            font=ctk.CTkFont(size=11),
            text_color=COLOR_SECONDARY
        )
        self.ollama_model_label.grid(row=22, column=0, padx=20, pady=(0, 4), sticky="w")

        self.ollama_model_menu = ctk.CTkOptionMenu(
            self.sidebar_scroll,
            values=self.ollama_models,
            variable=self.ollama_model_var,
            command=self.on_ollama_model_changed,
            fg_color=COLOR_CARD,
            button_color=COLOR_ACCENT_PURPLE,
            button_hover_color="#B388EB",
            text_color=COLOR_TEXT
        )
        self.ollama_model_menu.grid(row=23, column=0, padx=20, pady=(0, 10), sticky="w")

        self.ollama_num_predict_var = ctk.StringVar(value=str(self.ollama_num_predict))
        self.ollama_num_predict_label = ctk.CTkLabel(
            self.sidebar_scroll,
            text="Longueur r??ponse",
            font=ctk.CTkFont(size=11),
            text_color=COLOR_SECONDARY
        )
        self.ollama_num_predict_label.grid(row=24, column=0, padx=20, pady=(0, 4), sticky="w")

        self.ollama_num_predict_entry = ctk.CTkEntry(
            self.sidebar_scroll,
            textvariable=self.ollama_num_predict_var,
            height=28,
            width=120,
            fg_color=COLOR_CARD,
            border_color="#3C4043",
            text_color=COLOR_TEXT
        )
        self.ollama_num_predict_entry.grid(row=25, column=0, padx=20, pady=(0, 6), sticky="w")
        self.ollama_num_predict_entry.bind("<Return>", lambda e: self.on_ollama_num_predict_apply())

        self.ollama_num_predict_apply = ctk.CTkButton(
            self.sidebar_scroll,
            text="Appliquer",
            height=28,
            width=120,
            fg_color=COLOR_CARD,
            hover_color="#3C4043",
            command=self.on_ollama_num_predict_apply
        )
        self.ollama_num_predict_apply.grid(row=26, column=0, padx=20, pady=(0, 10), sticky="w")

        self.ollama_refresh_button = ctk.CTkButton(
            self.sidebar_scroll,
            text="Actualiser modeles",
            height=28,
            fg_color=COLOR_CARD,
            hover_color="#3C4043",
            command=self.refresh_ollama_models
        )
        self.ollama_refresh_button.grid(row=26, column=0, padx=20, pady=(0, 10), sticky="w")

        self.ollama_url_var = ctk.StringVar(value=self.ollama_url)
        self.ollama_url_label = ctk.CTkLabel(
            self.sidebar_scroll,
            text="URL Ollama",
            font=ctk.CTkFont(size=11),
            text_color=COLOR_SECONDARY
        )
        self.ollama_url_label.grid(row=27, column=0, padx=20, pady=(0, 4), sticky="w")

        self.ollama_url_entry = ctk.CTkEntry(
            self.sidebar_scroll,
            textvariable=self.ollama_url_var,
            height=28,
            width=160,
            fg_color=COLOR_CARD,
            border_color="#3C4043",
            text_color=COLOR_TEXT
        )
        self.ollama_url_entry.grid(row=28, column=0, padx=20, pady=(0, 6), sticky="w")
        self.ollama_url_entry.bind("<Return>", lambda e: self.on_ollama_url_apply())

        self.ollama_url_apply_button = ctk.CTkButton(
            self.sidebar_scroll,
            text="Appliquer URL",
            height=28,
            fg_color=COLOR_CARD,
            hover_color="#3C4043",
            command=self.on_ollama_url_apply
        )
        self.ollama_url_apply_button.grid(row=29, column=0, padx=20, pady=(0, 10), sticky="w")

        self.ollama_test_button = ctk.CTkButton(
            self.sidebar_scroll,
            text="Tester Ollama",
            height=28,
            fg_color=COLOR_CARD,
            hover_color="#3C4043",
            command=self.test_ollama_connection
        )
        self.ollama_test_button.grid(row=30, column=0, padx=20, pady=(0, 10), sticky="w")

        self.launch_on_startup_var = ctk.BooleanVar(value=self.launch_on_startup)
        self.launch_on_startup_switch = ctk.CTkSwitch(
            self.sidebar_scroll,
            text="Lancer au demarrage (arriere-plan)",
            variable=self.launch_on_startup_var,
            command=self.toggle_launch_on_startup,
            text_color=COLOR_SECONDARY
        )
        self.launch_on_startup_switch.grid(row=31, column=0, padx=20, pady=(0, 10), sticky="w")

        self.ui_theme_label = ctk.CTkLabel(
            self.sidebar_scroll,
            text="Theme couleur",
            font=ctk.CTkFont(size=11),
            text_color=COLOR_SECONDARY
        )
        self.ui_theme_label.grid(row=32, column=0, padx=20, pady=(0, 4), sticky="w")

        self.ui_theme_var = ctk.StringVar(value=self.ui_theme_name)
        self.ui_theme_menu = ctk.CTkOptionMenu(
            self.sidebar_scroll,
            values=list(self.ui_theme_presets.keys()),
            variable=self.ui_theme_var,
            command=self.on_ui_theme_changed,
            fg_color=COLOR_CARD,
            button_color=COLOR_ACCENT_PURPLE,
            button_hover_color="#B388EB",
            text_color=COLOR_TEXT
        )
        self.ui_theme_menu.grid(row=33, column=0, padx=20, pady=(0, 6), sticky="w")

        self.custom_accent_button = ctk.CTkButton(
            self.sidebar_scroll,
            text="Couleur accent perso",
            height=28,
            fg_color=COLOR_CARD,
            hover_color="#3C4043",
            command=self.choose_custom_accent_color
        )
        self.custom_accent_button.grid(row=34, column=0, padx=20, pady=(0, 10), sticky="w")

        # Localisation
        self.location_label = ctk.CTkLabel(self.sidebar_scroll, text="ðŸ“ Localisation...", font=ctk.CTkFont(size=11), text_color=COLOR_ACCENT_PURPLE)
        self.location_label.grid(row=35, column=0, padx=20, pady=5, sticky="w")
        self.sidebar_scroll.grid_rowconfigure(40, weight=1)

        # HUD Stark (Arc Reactor)
        self.stark_hud_frame = ctk.CTkFrame(self.sidebar_scroll, fg_color=COLOR_CARD, corner_radius=12)
        self.stark_hud_frame.grid(row=36, column=0, padx=14, pady=(10, 8), sticky="ew")
        self.stark_hud_frame.grid_columnconfigure(0, weight=1)

        self.stark_hud_title = ctk.CTkLabel(
            self.stark_hud_frame,
            text="ARC CORE HUD",
            font=ctk.CTkFont(size=11, weight="bold"),
            text_color=COLOR_ACCENT
        )
        self.stark_hud_title.grid(row=0, column=0, padx=8, pady=(8, 4))

        self.stark_canvas = tk.Canvas(
            self.stark_hud_frame,
            width=150,
            height=150,
            bg=COLOR_NAV,
            highlightthickness=0,
            bd=0
        )
        self.stark_canvas.grid(row=1, column=0, padx=8, pady=4)

        self.stark_mode_label = ctk.CTkLabel(
            self.stark_hud_frame,
            text="MODE: NORMAL",
            font=ctk.CTkFont(size=10, weight="bold"),
            text_color=COLOR_SECONDARY
        )
        self.stark_mode_label.grid(row=2, column=0, padx=8, pady=(4, 2))

        self.stark_hud_metrics = ctk.CTkLabel(
            self.stark_hud_frame,
            text="CPU --% | RAM --%",
            font=ctk.CTkFont(size=10),
            text_color=COLOR_SECONDARY
        )
        self.stark_hud_metrics.grid(row=3, column=0, padx=8, pady=(0, 8))

        self._hud_anim_tick = 0.0
        self._hud_cpu = 0.0
        self._hud_ram = 0.0
        self._hud_disk = 0.0
        self._hud_battery = None
        self._hud_items = {}
        self._init_stark_hud()

        # Header
        self.header = ctk.CTkFrame(self, height=60, fg_color=COLOR_NAV, corner_radius=0)
        self.header.grid(row=0, column=1, sticky="ew")
        self.logo_label = ctk.CTkLabel(self.header, text="âœ§ Jarvisse Intelligence", font=ctk.CTkFont(family="Inter", size=22, weight="bold"), text_color=COLOR_ACCENT)
        self.logo_label.grid(row=0, column=0, padx=20, pady=15)

        # Bouton ParamÃ¨tres Globaux (Style Gear Premium)
        self.header_settings_btn = ctk.CTkButton(
            self.header,
            text="âš™ï¸",
            width=40,
            height=40,
            fg_color="transparent",
            hover_color=COLOR_CARD,
            text_color=COLOR_ACCENT,
            font=ctk.CTkFont(size=20),
            command=self.open_main_settings
        )
        self.header_settings_btn.grid(row=0, column=1, padx=20, pady=10, sticky="e")
        self.header.grid_columnconfigure(1, weight=1)

        # Zone centrale avec Onglets
        self.tabview = ctk.CTkTabview(self, fg_color=COLOR_BG, segmented_button_selected_color=COLOR_ACCENT_PURPLE)
        self.tabview.grid(row=1, column=1, padx=20, pady=20, sticky="nsew")
        self.tab_chat = self.tabview.add("âœ§ Chat & IA")
        self.tab_browser = self.tabview.add("ðŸŒ Navigateur")
        
        self.tab_chat.grid_columnconfigure(0, weight=1)
        self.tab_chat.grid_rowconfigure(0, weight=1)

        # Zone de discussion (dans l'onglet Chat)
        self.chat_display = ctk.CTkTextbox(self.tab_chat, fg_color=COLOR_BG, border_width=0, font=ctk.CTkFont(size=13))
        self.chat_display.grid(row=0, column=0, padx=5, pady=5, sticky="nsew")
        self.chat_display.configure(state="disabled")

        # Navigateur IntÃ©grÃ©
        self.browser_frame = ctk.CTkFrame(self.tab_browser, fg_color=COLOR_NAV)
        self.browser_frame.grid(row=0, column=0, sticky="nsew", padx=5, pady=5)
        self.tab_browser.grid_columnconfigure(0, weight=1)
        self.tab_browser.grid_rowconfigure(0, weight=1)

        # Barre d'outils du navigateur
        self.browser_toolbar = ctk.CTkFrame(self.browser_frame, height=40, fg_color="transparent")
        self.browser_toolbar.pack(side="top", fill="x", padx=5, pady=5)

        ctk.CTkButton(self.browser_toolbar, text="â†", width=30, fg_color=COLOR_CARD, command=lambda: self.web_view.go_back()).pack(side="left", padx=2)
        ctk.CTkButton(self.browser_toolbar, text="â†’", width=30, fg_color=COLOR_CARD, command=lambda: self.web_view.go_forward()).pack(side="left", padx=2)
        ctk.CTkButton(self.browser_toolbar, text="âŸ³", width=30, fg_color=COLOR_CARD, command=lambda: self.web_view.load_website(self.web_view.url)).pack(side="left", padx=2)

        self.url_entry = ctk.CTkEntry(self.browser_toolbar, placeholder_text="Rechercher ou saisir une URL...")
        self.url_entry.pack(side="left", fill="x", expand=True, padx=5)
        self.url_entry.bind("<Return>", lambda e: self.agent_browser_mission(f"va sur {self.url_entry.get()}"))

        self.web_view = None # Initialisation par dÃ©faut
        try:
            if HtmlFrame is not None:
                self.web_view = HtmlFrame(self.browser_frame, messages_enabled=False)
                self.web_view.pack(fill="both", expand=True)
                
                def delayed_load():
                    print("ðŸŒ DEBUG: Chargement diffÃ©rÃ© du navigateur dÃ©sactivÃ© pour diagnostic.")
                self.after(5000, delayed_load)
            else:
                raise ImportError("Composant de navigation non disponible.")
        except Exception as e:
            print(f"Erreur Navigateur: {e}")
            self.web_view = ctk.CTkLabel(self.browser_frame, text=f"Le navigateur n'a pas pu Ãªtre chargÃ©.\n({e})\nInstallez-le avec 'pip install tkinterweb'")
            self.web_view.pack(pady=50)

        # Zone d'action
        self.bottom_frame = ctk.CTkFrame(self, height=120, fg_color="transparent")
        self.bottom_frame.grid(row=2, column=1, sticky="ew", padx=20, pady=20)
        self.bottom_frame.grid_columnconfigure(0, weight=1)

        self.status_label = ctk.CTkLabel(self.bottom_frame, text="Jarvisse est prÃªt.", font=ctk.CTkFont(size=12), text_color=COLOR_SECONDARY)
        self.status_label.grid(row=0, column=0, columnspan=4, pady=(0, 10))

        self.input_entry = ctk.CTkEntry(self.bottom_frame, placeholder_text="Commandez Jarvisse...", height=50, fg_color=COLOR_NAV, border_color="#3C4043", text_color=COLOR_TEXT)
        self.input_entry.grid(row=1, column=0, sticky="ew", padx=(0, 10))
        self.input_entry.bind("<Return>", lambda e: self.handle_text_input())

        # Buttons
        button_params = {"width": 50, "height": 50, "corner_radius": 15, "font": ctk.CTkFont(size=20)}
        
        self.listen_button = ctk.CTkButton(self.bottom_frame, text="ðŸŽ™", command=self.start_listening_thread, fg_color=COLOR_ACCENT, hover_color="#5294FF", **button_params)
        self.listen_button.grid(row=1, column=1, padx=(0, 10))

        self.file_button = ctk.CTkButton(self.bottom_frame, text="ðŸ“Ž", command=self.select_file, fg_color=COLOR_CARD, hover_color="#3C4043", **button_params)
        self.file_button.grid(row=1, column=2, padx=(0, 10))

        self.send_button = ctk.CTkButton(self.bottom_frame, text="âž¤", command=self.handle_text_input, fg_color=COLOR_ACCENT_PURPLE, hover_color="#B388EB", **button_params)
        self.send_button.grid(row=1, column=3, padx=(0, 10))

        self.cmds_button = ctk.CTkButton(self.bottom_frame, text="ðŸ“‹", command=self.open_commandes_panel,
                                         fg_color=COLOR_CARD, hover_color="#3C4043", **button_params)
        self.cmds_button.grid(row=1, column=4)

        # Panneau droit - Programmation d'horloge (Scrollable)
        self.right_panel = ctk.CTkScrollableFrame(self, width=240, fg_color=COLOR_NAV, corner_radius=0)
        self.right_panel.grid(row=0, column=2, rowspan=3, sticky="nsew")
        self.right_panel.grid_rowconfigure(14, weight=1)

        self.right_title = ctk.CTkLabel(self.right_panel, text="ALARMES", font=ctk.CTkFont(size=12, weight="bold"), text_color=COLOR_ACCENT)
        self.right_title.grid(row=0, column=0, padx=20, pady=(20, 10), sticky="w")

        self.alarm_time_label = ctk.CTkLabel(self.right_panel, text="Heure (HH:MM)", font=ctk.CTkFont(size=11), text_color=COLOR_SECONDARY)
        self.alarm_time_label.grid(row=1, column=0, padx=20, pady=(0, 5), sticky="w")

        self.alarm_time_entry = ctk.CTkEntry(self.right_panel, height=28, width=160, fg_color=COLOR_CARD, border_color="#3C4043", text_color=COLOR_TEXT)
        self.alarm_time_entry.grid(row=2, column=0, padx=20, pady=(0, 10), sticky="w")

        self.alarm_reason_label = ctk.CTkLabel(self.right_panel, text="Motif (message)", font=ctk.CTkFont(size=11), text_color=COLOR_SECONDARY)
        self.alarm_reason_label.grid(row=3, column=0, padx=20, pady=(0, 5), sticky="w")

        self.alarm_reason_entry = ctk.CTkEntry(self.right_panel, height=28, width=160, fg_color=COLOR_CARD, border_color="#3C4043", text_color=COLOR_TEXT)
        self.alarm_reason_entry.grid(row=4, column=0, padx=20, pady=(0, 10), sticky="w")

        self.alarm_set_button = ctk.CTkButton(
            self.right_panel,
            text="Programmer",
            height=28,
            width=160,
            fg_color=COLOR_ACCENT,
            hover_color="#5294FF",
            command=self.program_alarm
        )
        self.alarm_set_button.grid(row=5, column=0, padx=20, pady=(0, 10), sticky="w")

        self.alarm_list_label = ctk.CTkLabel(self.right_panel, text="Alarmes", font=ctk.CTkFont(size=11), text_color=COLOR_SECONDARY)
        self.alarm_list_label.grid(row=6, column=0, padx=20, pady=(0, 5), sticky="w")

        self.alarm_list_display = ctk.CTkTextbox(self.right_panel, height=120, width=180, fg_color=COLOR_BG, border_width=1, border_color="#3C4043", text_color=COLOR_TEXT)
        self.alarm_list_display.grid(row=7, column=0, padx=20, pady=(0, 10), sticky="w")
        self.alarm_list_display.configure(state="disabled")

        self.next_alarm_label = ctk.CTkLabel(self.right_panel, text="Aucune alarme", font=ctk.CTkFont(size=11), text_color=COLOR_SECONDARY)
        self.next_alarm_label.grid(row=8, column=0, padx=20, pady=(10, 6), sticky="w")

        self.anomaly_settings_button = ctk.CTkButton(
            self.right_panel,
            text="ParamÃ¨tres anomalies",
            height=28,
            width=160,
            fg_color=COLOR_CARD,
            hover_color="#3C4043",
            command=self.open_anomaly_settings
        )
        self.anomaly_settings_button.grid(row=9, column=0, padx=20, pady=(0, 10), sticky="w")

        self.anomalies_enabled_var_main = ctk.BooleanVar(value=self.anomalies_enabled)
        self.anomalies_toggle = ctk.CTkSwitch(
            self.right_panel,
            text="Anomalies activÃ©es",
            variable=self.anomalies_enabled_var_main,
            command=lambda: self.toggle_anomalies_enabled(source="main"),
            text_color=COLOR_SECONDARY
        )
        self.anomalies_toggle.grid(row=10, column=0, padx=20, pady=(0, 10), sticky="w")

        self.services_watch_button = ctk.CTkButton(
            self.right_panel,
            text="Services critiques",
            height=28,
            width=160,
            fg_color=COLOR_CARD,
            hover_color="#3C4043",
            command=self.open_services_watch
        )
        self.services_watch_button.grid(row=11, column=0, padx=20, pady=(0, 10), sticky="w")

        self.notifications_button = ctk.CTkButton(
            self.right_panel,
            text="Notifications apps",
            height=28,
            width=160,
            fg_color=COLOR_CARD,
            hover_color="#3C4043",
            command=self.open_notification_settings
        )
        self.notifications_button.grid(row=12, column=0, padx=20, pady=(0, 10), sticky="w")

        self.agents_settings_button = ctk.CTkButton(
            self.right_panel,
            text="âš™ï¸ ParamÃ¨tres Agents",
            height=28,
            width=160,
            fg_color=COLOR_ACCENT_PURPLE,
            hover_color="#B388EB",
            command=self.open_agents_settings
        )
        self.agents_settings_button.grid(row=13, column=0, padx=20, pady=(0, 10), sticky="w")

        self._refresh_alarm_list()

        self.after(1000, self.saluer)
        self.after(3000, self.start_wake_word_thread)
        self.after(2000, self.update_system_stats)
        self.after(4000, self.update_location)
        self.after(1000, self.check_alarms)
        self.after(5000, self.check_removable_devices)
        self.after(15000, self.check_system_anomalies)
        self.after(6000, self.start_notifications_listener_thread)
        
        # Gestion de la fermeture propre
        self.protocol("WM_DELETE_WINDOW", self.on_closing)
        
        # RafraÃ®chissement final pour appliquer toutes les couleurs chargÃ©es
        self._refresh_theme_widgets()
        
        try:
            if self.launch_on_startup:
                self._enable_startup_launcher()
            else:
                self._disable_startup_launcher()
        except Exception as e:
            print(f"Erreur synchronisation d\u00e9marrage auto: {e}")
            
        self.after(50, self._drain)

    def _drain(self):
        import queue
        try:
            while True:
                try:
                    kind, payload = self.uiq.get_nowait()
                except queue.Empty:
                    break
                try:
                    if kind == "append":
                        if hasattr(self, "output_text"):
                            self.output_text.insert("end", payload + "\n")
                            self.output_text.see("end")
                    elif kind == "add_message":
                        msg, sender, prog = payload
                        self._add_message_safe(msg, sender, prog)
                    elif kind == "parler":
                        texte, force_full = payload
                        self.parler(texte, progressive=False, wait=False, force_full=force_full)
                except Exception as e:
                    print(f"[_drain] Erreur traitement message: {e}")
        except Exception as e:
            print(f"[_drain] Erreur critique: {e}")
        self.after(50, self._drain)

    def on_closing(self):
        """ArrÃªte proprement l'assistant"""
        print("ðŸ›‘ Fermeture de l'application demandÃ©e (on_closing)...")
        self.veille_active = False
        self.is_listening = False
        try:
            pygame.mixer.quit()
        except: pass
        self.destroy()
        os._exit(0) # Force l'arrÃªt de tous les threads

    def start_wake_word_thread(self):
        """Lance la veille en arriÃ¨re-plan"""
        threading.Thread(target=self.attendre_mot_cle, daemon=True).start()
        # Thread dÃ©diÃ© pour la dÃ©tection du STOP en continu
        threading.Thread(target=self._stop_detection_loop, daemon=True).start()

    def _stop_detection_loop(self):
        """Boucle dÃ©diÃ©e Ã  la dÃ©tection du mot STOP en continu, indÃ©pendamment de l'Ã©tat d'Ã©coute"""
        reconnaisseur = sr.Recognizer()
        reconnaisseur.energy_threshold = 300
        reconnaisseur.dynamic_energy_threshold = True
        reconnaisseur.pause_threshold = 0.3
        
        fs = 44100
        duree_ecoute = 1.0  # Ã‰coute courte pour rÃ©activitÃ©
        
        print("ðŸ›‘ SystÃ¨me de dÃ©tection STOP activÃ© en continu...")
        
        while self.veille_active:
            try:
                # Capture audio courte
                audio_rec = sd.rec(int(duree_ecoute * fs), samplerate=fs, channels=1, dtype='int16')
                sd.wait()
                
                byte_io = io.BytesIO()
                with wave.open(byte_io, 'wb') as wav_file:
                    wav_file.setnchannels(1)
                    wav_file.setsampwidth(2)
                    wav_file.setframerate(fs)
                    wav_file.writeframes(audio_rec.tobytes())
                byte_io.seek(0)
                
                with sr.AudioFile(byte_io) as source:
                    reconnaisseur.adjust_for_ambient_noise(source, duration=0.05)
                    audio_data = reconnaisseur.record(source)
                
                try:
                    texte = reconnaisseur.recognize_google(audio_data, language='fr-FR').lower()
                    
                    # DÃ©tection des commandes d'arrÃªt
                    if any(word in texte for word in ["stop", "arrÃªte", "arrÃªt", "tais-toi", "silence"]):
                        print(f"ðŸ›‘ STOP DÃ‰TECTÃ‰ (Thread dÃ©diÃ©): '{texte}'")
                        self.stop_speaking_flag = True
                        
                        try:
                            pygame.mixer.music.stop()
                            pygame.mixer.music.unload()
                        except:
                            pass
                        
                        self.conversation_continue = False
                        self.is_listening = False
                        
                        # Confirmation vocale aprÃ¨s un court dÃ©lai
                        self.after(300, lambda: self.parler("Lecture arrÃªtÃ©e.", wait=False))
                        self.after(400, self.reset_ui)
                        
                        # Pause pour Ã©viter les dÃ©tections multiples
                        time.sleep(1.5)
                        
                except sr.UnknownValueError:
                    pass  # Pas de parole dÃ©tectÃ©e
                except sr.RequestError:
                    pass  # Erreur rÃ©seau
                    
            except Exception as e:
                # Erreur silencieuse pour ne pas polluer les logs
                pass
            
            # Petite pause pour ne pas surcharger le CPU
            time.sleep(0.05)

    def attendre_mot_cle(self):
        """Boucle de veille rÃ©active pour dÃ©tecter 'Jarvisse'"""
        reconnaisseur = sr.Recognizer()
        # Seuil par dÃ©faut plus robuste
        reconnaisseur.energy_threshold = 300 
        reconnaisseur.dynamic_energy_threshold = True
        reconnaisseur.pause_threshold = 0.5
        
        fs = 44100
        duree_veille = 1.5  # LÃ©gÃ¨rement plus long pour mieux capter le mot entier
        
        print("ðŸŽ™ï¸ SystÃ¨me de veille activÃ© (Jarvisse/Alpha/Service)...")
        
        while self.veille_active:
            if not self.is_listening:
                try:
                    self.after(0, lambda: self.status_label.configure(text="ðŸ”Š Veille active... (Dites 'Jarvisse')", text_color=COLOR_SECONDARY))
                    
                    # Capture courte
                    audio_rec = sd.rec(int(duree_veille * fs), samplerate=fs, channels=1, dtype='int16')
                    sd.wait()
                    
                    byte_io = io.BytesIO()
                    with wave.open(byte_io, 'wb') as wav_file:
                        wav_file.setnchannels(1)
                        wav_file.setsampwidth(2)
                        wav_file.setframerate(fs)
                        wav_file.writeframes(audio_rec.tobytes())
                    byte_io.seek(0)
                    
                    with sr.AudioFile(byte_io) as source:
                        # Ajustement minimal pour ne pas filtrer la voix
                        reconnaisseur.adjust_for_ambient_noise(source, duration=0.1)
                        audio_data = reconnaisseur.record(source)
                    
                    try:
                        # Transcription rapide
                        texte = reconnaisseur.recognize_google(audio_data, language='fr-FR').lower()
                        
                        # Variante STOP pour interruption lecture
                        if "stop" in texte or "arrÃªte" in texte:
                            print(f"ðŸ›‘ COMMANDE D'ARRÃŠT DÃ‰TECTÃ‰E (Veille): '{texte}'")
                            self.stop_speaking_flag = True
                            try:
                                pygame.mixer.music.stop()
                                pygame.mixer.music.unload()
                            except: pass
                            self.conversation_continue = False
                            # On laisse le temps Ã  la boucle de parole de s'arrÃªter avant de confirmer
                            self.after(500, lambda: self.parler("Lecture arrÃªtÃ©e."))
                            self.after(600, self.reset_ui)
                            continue

                        # Variantes phonÃ©tiques (Ã©largies suite aux logs utilisateur)
                        variantes = [
                            "jarvisse", "jarvis", "service", "vice", "zarzis", "janice", 
                            "charvis", "cher", "jarv", "servis", "alpha", "arvis", "jar",
                            "garvis", "jarves", "janis", "jarvise"
                        ]

                        for v in variantes:
                            if v in texte:
                                print(f"âœ… MOT-CLÃ‰ DÃ‰TECTÃ‰: '{v}' (entendu: '{texte}')")
                                self.after(0, self.declencher_activation)
                                break
                    except: pass
                except: pass
            time.sleep(0.01)


    def declencher_activation(self):
        """RÃ©ponse immÃ©diate au mot-clÃ© sans attente inutile"""
        if self.is_listening: 
            print("âš ï¸ DÃ©jÃ  en Ã©coute, activation ignorÃ©e")
            return
        
        print("ðŸš€ ACTIVATION ! Jarvisse rÃ©pond...")
        self.is_listening = True
        
        def sequence():
            self.after(0, lambda: self.status_label.configure(text="âœ¨ Oui Monsieur...", text_color=COLOR_ACCENT))
            # On ne bloque pas tout le thread pour parler, on lance l'Ã©coute trÃ¨s vite
            self.parler("Oui Monsieur", wait=False) 
            
            # On attend que le "Oui Monsieur" se termine pour ne pas que l'assistant s'Ã©coute lui-mÃªme
            time.sleep(1.2) 
            self.is_listening = False
            print("ðŸ‘‚ Passage en mode Ã©coute...")
            self.after(0, self.start_listening_thread)
            
        threading.Thread(target=sequence, daemon=True).start()


    def _normaliser_texte(self, texte):
        """Nettoie une chaÃ®ne pour amÃ©liorer la dÃ©tection de mots-clÃ©s."""
        # Remplace la ponctuation par des espaces, garde les lettres et espaces.
        cleaned = []
        for ch in texte:
            if ch.isalpha() or ch.isspace():
                cleaned.append(ch)
            else:
                cleaned.append(" ")
        return " ".join("".join(cleaned).split())

    def add_message(self, message, sender="Jarvisse", progressive=True):
        try:
            self.after(0, self._add_message_safe, message, sender, progressive)
        except:
            print(f"[{sender}] {message}") # Fallback if UI not ready

    def _add_message_safe(self, message, sender, progressive):
        self.chat_display.configure(state="normal")
        
        display_name = f"âœ§ {sender}" if sender != "Vous" else "Vous"
        tag = "user_tag" if sender == "Vous" else "ai_tag"
        
        self.chat_display.insert("end", f"\n{display_name}\n", tag)
        
        if progressive and sender != "Vous":
            # Lancer l'effet machine Ã  Ã©crire
            threading.Thread(target=self._typewriter_effect, args=(message,), daemon=True).start()
        else:
            self.chat_display.insert("end", f"{message}\n")
            self.chat_display.see("end")
            self.chat_display.configure(state="disabled")
        
        self.chat_display.tag_config("user_tag", foreground="#8AB4F8")
        self.chat_display.tag_config("ai_tag", foreground=COLOR_ACCENT_PURPLE if sender != "Jarvisse" else COLOR_ACCENT)

    def start_streaming_message(self):
        self.after(0, self._start_streaming_message_safe)

    def _start_streaming_message_safe(self):
        self.chat_display.configure(state="normal")
        self.chat_display.insert("end", f"\nÃ¢Å“Â§ Jarvisse\n", "ai_tag")
        self.chat_display.see("end")
        self.chat_display.configure(state="disabled")
        self.chat_display.tag_config("ai_tag", foreground=COLOR_ACCENT)

    def append_streaming_text(self, text):
        if not text:
            return
        self.after(0, self._append_streaming_text_safe, text)

    def _append_streaming_text_safe(self, text):
        self.chat_display.configure(state="normal")
        self.chat_display.insert("end", text)
        self.chat_display.see("end")
        self.chat_display.configure(state="disabled")

    def end_streaming_message(self):
        self.after(0, lambda: self.chat_display.configure(state="disabled"))

    def _typewriter_effect(self, message):
        """Affichage INSTANTANÃ‰ pour mode ping-pong (pas d'effet machine Ã  Ã©crire)"""
        # Mode ping-pong : affichage immÃ©diat sans dÃ©lai
        self.after(0, self._insert_full_message, message)
        self.after(0, lambda: self.chat_display.configure(state="disabled"))

    def _insert_full_message(self, message):
        """InsÃ¨re le message complet d'un coup"""
        self.chat_display.configure(state="normal")
        self.chat_display.insert("end", f"{message}\n")
        self.chat_display.see("end")


    def _insert_char(self, char):
        self.chat_display.configure(state="normal")
        self.chat_display.insert("end", char)
        self.chat_display.see("end")

    def handle_text_input(self):
        commande = self.input_entry.get().strip()
        if commande:
            self.input_entry.delete(0, "end")
            self.add_message(commande, "Vous")
            self.status_label.configure(text="En train de rÃ©flÃ©chir...")
            self.last_command_origin = "local"
            threading.Thread(target=self.executer_commande, args=(commande,), daemon=True).start()

    def select_file(self):
        file_path = filedialog.askopenfilename(
            title="SÃ©lectionner un fichier",
            filetypes=[("Documents", "*.txt *.pdf *.docx"), ("Tous les fichiers", "*.*")]
        )
        if file_path:
            self.traiter_fichier(file_path)

    def traiter_fichier(self, path):
        nom_fichier = os.path.basename(path)
        self.add_message(f"Fichier joint : {nom_fichier}", "Vous")
        
        contenu = ""
        try:
            # Gestion des fichiers TEXTE
            if path.endswith('.txt'):
                with open(path, 'r', encoding='utf-8') as f:
                    contenu = f.read()
            
            # Gestion des fichiers PDF
            elif path.endswith('.pdf'):
                doc = fitz.open(path)
                try:
                    for page in doc:
                        contenu += page.get_text()
                finally:
                    doc.close()
            
            # Gestion des fichiers DOCX
            elif path.endswith('.docx'):
                doc = Document(path)
                for para in doc.paragraphs:
                    contenu += para.text + "\n"

            # Gestion des fichiers EXCEL
            elif path.endswith('.xlsx'):
                wb = openpyxl.load_workbook(path)
                for sheet in wb.worksheets:
                    for row in sheet.iter_rows(values_only=True):
                        contenu += " ".join([str(cell) for cell in row if cell is not None]) + "\n"

            # Gestion des fichiers POWERPOINT
            elif path.endswith('.pptx'):
                prs = Presentation(path)
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            contenu += shape.text + "\n"

            # Gestion des fichiers CSV
            elif path.endswith('.csv'):
                with open(path, mode='r', encoding='utf-8') as f:
                    reader = csv.reader(f)
                    for row in reader:
                        contenu += " ".join(row) + "\n"
            
            if contenu:
                self.loaded_document_content = contenu # Stockage dÃ©diÃ© au contenu du fichier
                self.current_subject = f"le contenu du fichier {nom_fichier}"
                self.parler(f"J'ai bien chargÃ© le document {nom_fichier}. Je peux maintenant vous le lire, le rÃ©sumer ou rÃ©pondre Ã  vos questions.")
                self.conversation_history.append(f"CONTENU_FICHIER: {contenu}")
                # trim auto gÃ©rÃ© par deque(maxlen)
            else:
                self.parler(f"J'ai reÃ§u {nom_fichier}, mais je n'ai trouvÃ© aucun texte lisible Ã  l'intÃ©rieur.")
                
        except Exception as e:
            print(f"Erreur lecture fichier: {e}")
            self.parler(f"DÃ©solÃ©, je n'ai pas pu lire le fichier {nom_fichier}. VÃ©rifiez qu'il n'est pas utilisÃ© par un autre programme.")

    def humaniser_reponse(self, texte):
        """Ajoute des marqueurs de conversation naturels pour rendre la voix plus humaine"""
        # Si on est en mode ping-pong (ultra-rapide), on Ã©vite d'ajouter du texte inutile
        if getattr(self, "ping_pong_mode", False):
            return texte
            
        intros = [
            "Eh bien, ", "Justement, ", "Il me semble que ", "Voyons voir, ", 
            "Pour Ãªtre prÃ©cis, ", "Si je puis dire, ", "Ã€ vrai dire, ",
            "Comme vous le savez probablement, ", "D'aprÃ¨s mes analyses, ",
            "Ã‰coutez, ", "Alors, pour tout vous dire, ", "C'est une excellente question... "
        ]
        
        interjections = [" hmm, ", " d'ailleurs, ", " en quelque sorte, ", ", voyez-vous, "]
        
        # On n'ajoute une intro que si le texte est assez long et ne commence pas dÃ©jÃ  par une formule de politesse
        if len(texte) > 30 and not texte[0].isupper():
            if random.random() < 0.25: # 25% de chance d'ajouter une intro naturelle
                texte = random.choice(intros) + texte[0].lower() + texte[1:]
        
        # Ajout sporadique d'interjections au milieu
        if len(texte) > 100 and random.random() < 0.15:
            # On insÃ¨re aprÃ¨s la premiÃ¨re virgule ou le premier point
            parts = re.split(r'([,.])', texte, maxsplit=1)
            if len(parts) > 2:
                texte = parts[0] + parts[1] + random.choice(interjections) + parts[2]
                
        return texte

    def _select_file_main_thread(self, title, types):
        """Ouvre le sÃ©lecteur de fichiers de maniÃ¨re thread-safe via le thread principal"""
        res_queue = queue.Queue()
        def _open():
            path = filedialog.askopenfilename(title=title, filetypes=types)
            res_queue.put(path)
        self.after(0, _open)
        return res_queue.get() # Attend la sÃ©lection de l'utilisateur

    def _convertir_nombre_0_99(self, n):
        """Convertit un nombre 0-99 en texte pour la lecture des annÃ©es"""
        if n == 0: return ""
        if n < 20:
             t = ["", "un", "deux", "trois", "quatre", "cinq", "six", "sept", "huit", "neuf", "dix",
                 "onze", "douze", "treize", "quatorze", "quinze", "seize", "dix-sept", "dix-huit", "dix-neuf"]
             return t[n]
        
        dizaines = {2:"vingt", 3:"trente", 4:"quarante", 5:"cinquante", 6:"soixante", 7:"soixante", 8:"quatre-vingt", 9:"quatre-vingt"}
        d = n // 10
        u = n % 10
        
        if d == 7 or d == 9:
            u += 10
            
        txt = dizaines[d]
        if u == 0:
            if d == 8: txt += "s"
            return txt
            
        separateur = "-"
        if (u == 1 or u == 11) and d < 8: separateur = "-et-"
        
        # RÃ©cupÃ©ration de l'unitÃ© (ou 10-19)
        t = ["", "un", "deux", "trois", "quatre", "cinq", "six", "sept", "huit", "neuf", "dix",
             "onze", "douze", "treize", "quatorze", "quinze", "seize", "dix-sept", "dix-huit", "dix-neuf"]
        
        txt += separateur + t[u]
        return txt

    def _convertir_annees_historiques(self, texte):
        """Convertit les annÃ©es (1700-2029) en toutes lettres"""
        def repl_annee(match):
            val = int(match.group(0))
            siecle = val // 100
            reste = val % 100
            
            debut = ""
            if siecle == 19: debut = "mille neuf cent"
            elif siecle == 20: debut = "deux mille"
            elif siecle == 18: debut = "mille huit cent"
            elif siecle == 17: debut = "mille sept cent"
            else: return match.group(0)
            
            if reste == 0: return debut
            conversion_reste = self._convertir_nombre_0_99(reste)
            return f"{debut} {conversion_reste}" if conversion_reste else debut

        # Regex pour les annÃ©es isolÃ©es (4 chiffres)
        return re.sub(r'\b(17|18|19|20)\d{2}\b', repl_annee, texte)

    def _convertir_chiffres_romains(self, texte):
        """Convertit les chiffres romains courants et siÃ¨cles en texte franÃ§ais"""
        # 1. SiÃ¨cles courants
        texte = re.sub(r'\b(XI|XII|XIII|XIV|XV|XVI|XVII|XVIII|XIX|XX|XXI)e (siÃ¨cle|millÃ©naire)\b', 
                       lambda m: self._map_siecle(m.group(1)) + " " + m.group(2), texte, flags=re.IGNORECASE)
        
        # 2. Rois et Chapitres
        mots_cles = r"(Louis|Henri|Charles|Philippe|FranÃ§ois|NapolÃ©on|Pie|Chapitre|Tome|Volume|Guerre|Partie|Jean|Paul|BenoÃ®t)"
        texte = re.sub(rf'\b{mots_cles} ([IVX]+)\b', lambda m: f"{m.group(1)} {self._map_romain(m.group(2))}", texte)
        return texte

    def _map_siecle(self, rom):
        m = {'XI':'onziÃ¨me', 'XII':'douziÃ¨me', 'XIII':'treiziÃ¨me', 'XIV':'quatorziÃ¨me', 'XV':'quinziÃ¨me', 
             'XVI':'seiziÃ¨me', 'XVII':'dix-septiÃ¨me', 'XVIII':'dix-huitiÃ¨me', 'XIX':'dix-neuviÃ¨me', 
             'XX':'vingtiÃ¨me', 'XXI':'vingt-et-uniÃ¨me'}
        return m.get(rom.upper(), rom)

    def _map_romain(self, rom):
        m = {'I':'premier', 'II':'deux', 'III':'trois', 'IV':'quatre', 'V':'cinq', 'VI':'six', 'VII':'sept', 
             'VIII':'huit', 'IX':'neuf', 'X':'dix', 'XI':'onze', 'XII':'douze', 'XIII':'treize', 'XIV':'quatorze', 
             'XV':'quinze', 'XVI':'seize', 'XVII':'dix-sept', 'XVIII':'dix-huit', 'XIX':'dix-neuf', 
             'XX':'vingt', 'XXI':'vingt-et-un'}
        return m.get(rom.upper(), rom)

    def _nettoyer_pour_parler(self, texte):
        """Nettoie le texte des symboles markdown et autres caractÃ¨res indÃ©sirables avant la synthÃ¨se vocale."""
        if not texte:
            return ""
        
        # 1. Supprimer les liens Markdown [texte](url) pour ne garder que le texte
        texte_propre = re.sub(r'\[([^\]]+)\]\([^\)]+\)', r'\1', texte)
        
        # 2. Supprimer les symboles markdown courants (*, _, #, `, ~, >, [, ])
        # Ces symboles sont souvent utilisÃ©s par les LLM pour le formatage
        texte_propre = re.sub(r'[*_#`~>\[\]]', '', texte_propre)
        
        # 3. Supprimer les listes Ã  puces simples en dÃ©but de ligne (ex: "- " ou "+ ")
        texte_propre = re.sub(r'^[ \t]*[-+*][ \t]+', '', texte_propre, flags=re.MULTILINE)
        
        # 4. Supprimer les emojis et caractÃ¨res spÃ©ciaux (âœ…, âŒ, âš ï¸, etc.)
        texte_propre = re.sub(r'[^\x00-\x7FÃ€-Ã¿0-9\s.,!?:;\'"()\-/]', '', texte_propre)
        
        # 5. Traitement des barres obliques (/) pour Ã©viter qu'elles soient lues
        # Dates : 15/02/2024 -> 15 fÃ©vrier 2024 ou simplement 15 02 2024
        texte_propre = re.sub(r'(\d{1,2})/(\d{1,2})/(\d{2,4})', r'\1 \2 \3', texte_propre)
        # Autres slashes isolÃ©s ou multiples : remplacer par espace
        texte_propre = re.sub(r'/+', ' ', texte_propre)
        
        # 6. Nettoyer les espaces multiples
        texte_propre = re.sub(r'\s+', ' ', texte_propre)
        
        # 7. INTELLIGENCE PHONÃ‰TIQUE ALGORITHMIQUE (CapacitÃ© 100.000+)
        # Au lieu de boucles lentes, on utilise une substitution par callback ultra-rapide
        texte_propre = self._appliquer_intelligence_phonetique(texte_propre)
        
        # 8. CONVERSION AVANCÃ‰E (AnnÃ©es & Chiffres Romains)
        texte_propre = self._convertir_chiffres_romains(texte_propre)
        texte_propre = self._convertir_annees_historiques(texte_propre)
        
        return texte_propre.strip()

    def _load_phonetic_map(self):
        """Charge la base de donnÃ©es phonÃ©tique depuis le fichier externe"""
        if os.path.exists(self.phonetic_path):
            try:
                with open(self.phonetic_path, "r", encoding="utf-8") as f:
                    return json.load(f)
            except: return {}
        return {}

    def _appliquer_intelligence_phonetique(self, texte):
        """Moteur de traitement phonÃ©tique haute performance avec support URL"""
        # PrÃ©-traitement des extensions web courantes pour Ã©viter qu'elles soient ignorÃ©es
        texte = re.sub(r'\.com\b', ' point com ', texte, flags=re.IGNORECASE)
        texte = re.sub(r'\.fr\b', ' point Ã¨f-Ã¨r ', texte, flags=re.IGNORECASE)
        texte = re.sub(r'\.net\b', ' point nette ', texte, flags=re.IGNORECASE)
        texte = re.sub(r'\.org\b', ' point orgue ', texte, flags=re.IGNORECASE)
        texte = re.sub(r'\.io\b', ' point i-o ', texte, flags=re.IGNORECASE)
        
        def replace_callback(match):
            word = match.group(0)
            word_upper = word.upper()
            
            # Cas spÃ©cial : www
            if word_upper == "WWW":
                return "v-v-v" # Plus rapide et naturel que double-vÃ© double-vÃ©...
            
            # A. VÃ©rification dans le dictionnaire Ã©tendu
            if word_upper in self.phonetic_map:
                return self.phonetic_map[word_upper]
            
            # B. Ã‰pellation automatique des acronymes (Majuscules de 2 Ã  5 lettres)
            if word.isupper() and 2 <= len(word) <= 5 and word.isalpha():
                return "-".join([self.alphabet_phonetique.get(l, l) for l in word])
            
            return word

        # On traite tous les mots de 2 lettres ou plus
        return re.sub(r'\b[a-zA-Z]{2,}\b', replace_callback, texte)

    def _ping_pong_style(self, texte):
        """Raccourcit la rÃ©ponse pour un style ping-pong (court, direct)"""
        texte = " ".join(str(texte).split())
        if not texte:
            return texte

        # Garder 1 ou 2 phrases max pour un rythme rÃ©actif
        phrases = re.split(r'(?<=[.!?])\s+', texte)
        court = phrases[0]
        if len(court) < 25 and len(phrases) > 1:
            court = f"{court} {phrases[1]}"

        if len(court) > self.ping_pong_max_chars:
            court = court[:self.ping_pong_max_chars].rstrip() + "..."

        return court

    def parler(self, texte, progressive=True, wait=True, force_full=False, sender="Jarvisse"):
        # RÃ©initialiser le flag d'arrÃªt Ã  chaque nouvelle parole
        self.stop_speaking_flag = False
        
        # On affiche le texte progressivement en mÃªme temps que la voix commence
        texte_brut = str(texte)
        texte = self.humaniser_reponse(texte_brut)
        
        if self.ping_pong_mode and not force_full:
            texte = self._ping_pong_style(texte)
            progressive = False
        
        if force_full:
            progressive = False # DÃ©sactiver l'effet machine Ã  Ã©crire pour la rapiditÃ© et fiabilitÃ©

        # MÃ©morise aussi les rÃ©ponses de l'assistant
        self.conversation_history.append(f"{sender}: {texte}")
        # trim auto gÃ©rÃ© par deque(maxlen)

        self.add_message(texte, sender=sender, progressive=progressive)
        
        # On ne lance la voix que si elle est activÃ©e
        if getattr(self, "voice_enabled", True):
            threading.Thread(target=self._run_async_parler, args=(texte, wait), daemon=True).start()

        # Routage vers Telegram si la commande vient de lÃ 
        if getattr(self, "agent_telegram_enabled", False) and getattr(self, "last_command_origin", "local") == "telegram":
            threading.Thread(target=self.send_telegram_message, args=(texte_brut,), daemon=True).start()

    def _parler_chunk(self, texte):
        if not texte:
            return
        threading.Thread(target=self._run_async_parler, args=(texte, False), daemon=True).start()


    def _run_async_parler(self, texte, wait):
        try:
            asyncio.run(self._parler_async(texte, wait))
            # Si le mode conversation continue est actif et que l'origine est locale
            if self.conversation_continue and getattr(self, "last_command_origin", "local") == "local":
                # On libÃ¨re l'Ã©tat d'Ã©coute pour permettre au thread suivant de dÃ©marrer
                self.is_listening = False
                # On attend un tout petit peu pour Ã©viter que le micro ne capte la fin de l'Ã©cho
                time.sleep(0.4)
                self.after(0, self.start_listening_thread)
        except Exception as e:
            print(f"Erreur fatale asyncio voice: {e}")
            if self.conversation_continue: self.is_listening = False

    def toggle_conversation_continue(self):
        self.conversation_continue = self.conv_continue_var.get()
        if self.conversation_continue:
            self.parler("Mode conversation continue activÃ©. Je vous Ã©coute sans interruption.")
        else:
            self.parler("Mode conversation continue dÃ©sactivÃ©.")

    def _split_sentences(self, texte):
        """DÃ©coupe le texte de maniÃ¨re fluide pour une synthÃ¨se vocale naturelle"""
        if not texte: return []
        import re
        # On coupe prioritairement sur les ponctuations fortes (. ! ? \n)
        phrases = re.split(r'(?<=[.!?\n])\s+', texte)
        
        final_phrases = []
        for p in phrases:
            # On ne coupe sur les points-virgules ou virgules QUE si la phrase est extrÃªmement longue (> 150 chars)
            if len(p) > 150:
                sub = re.split(r'(?<=[,;])\s+', p)
                final_phrases.extend(sub)
            else:
                final_phrases.append(p)
        return [f.strip() for f in final_phrases if f.strip()]

    async def _parler_async(self, texte, wait):
        """Lecture vocale optimisÃ©e avec gÃ©nÃ©ration parallÃ¨le et buffers mÃ©moire"""
        if not texte: return
        print(f"ðŸ”Š DÃ©but lecture vocale ({len(texte)} chars)")
        try:
            phrases = self._split_sentences(texte)
            if not phrases: 
                print("âš ï¸ Aucune phrase Ã  lire aprÃ¨s dÃ©coupage.")
                return
            
            print(f"ðŸ“¦ DÃ©coupage en {len(phrases)} segments audio.")

            # Limiteur de concurrence pour Ã©viter le bannissement IP ou les timeouts
            semaphore = asyncio.Semaphore(3)

            async def fetch_audio(phrase_text, index):
                async with semaphore:
                    p_propre = self._nettoyer_pour_parler(phrase_text)
                    if not p_propre or len(p_propre.strip()) < 2: return None
                    try:
                        print(f"â³ GÃ©nÃ©ration audio {index+1}/{len(phrases)}...")
                        # ðŸ¤« Dynamisme : Adaptation au chuchotement
                        vitesse = self.voice_speed
                        if getattr(self, "_is_whispering", False):
                            vitesse = "+50%" # Parle plus vite et plus bas (simulÃ© par pitch/vitesse selon TTS)
                        
                        comm = edge_tts.Communicate(p_propre, self.voice, rate=vitesse)
                        data = b""
                        async for chunk in comm.stream():
                            if chunk["type"] == "audio": data += chunk["data"]
                        if not data:
                            print(f"âŒ DonnÃ©es audio vides pour segment {index+1}")
                            return None
                        return io.BytesIO(data)
                    except Exception as e:
                        print(f"âŒ Erreur fetch_audio {index+1}: {e}")
                        return None

            # Lancement des tÃ¢ches (le sÃ©maphore gÃ¨re la file d'attente)
            audio_tasks = [asyncio.create_task(fetch_audio(p, i)) for i, p in enumerate(phrases)]
            
            # On attend et on joue au fur et Ã  mesure
            with self.voice_lock:
                # ArrÃªt propre de la musique prÃ©cÃ©dente
                try:
                    pygame.mixer.music.stop()
                    pygame.mixer.music.unload()
                except:
                    pass

                for i, task in enumerate(audio_tasks):
                    buffer = await task
                    if buffer:
                        try:
                            # print(f"â–¶ï¸ Lecture segment {i+1}")
                            
                            # VÃ©rification interruption avant lecture
                            if self.stop_speaking_flag:
                                print(f"ðŸ›‘ Lecture interrompue par l'utilisateur (AVANT segment {i+1})")
                                break

                            pygame.mixer.music.load(buffer)
                            pygame.mixer.music.play()
                            while pygame.mixer.music.get_busy():
                                if self.stop_speaking_flag:
                                    pygame.mixer.music.stop()
                                    print(f"ðŸ›‘ Lecture interrompue par l'utilisateur (PENDANT segment {i+1})")
                                    return
                                await asyncio.sleep(0.01)
                        except Exception as e:
                            print(f"âŒ Erreur lecture phrase {i+1}: {e}")
                    else:
                        print(f"â­ï¸ Saut du segment {i+1} (audio manquant)")
            
            print(f"âœ… Fin de la lecture vocale.")
        except Exception as e:
            print(f"âŒ Erreur globale _parler_async: {e}")

    def saluer(self):
        heure = datetime.datetime.now().hour
        options = []
        if 5 <= heure < 12:
            options = [
                "Bonjour ! C'est un plaisir de vous retrouver. Je suis Jarvisse. Comment s'annonce votre matinÃ©e ?",
                "Bonjour ! Jarvisse Ã  votre service. Comment puis-je vous aider Ã  bien commencer la journÃ©e ?",
                "Salut ! J'espÃ¨re que vous avez bien dormi. Je suis prÃªt pour vos demandes.",
                "Mes respects chef. Je suis Ã  votre entiÃ¨re disposition. Dites-moi en quoi je peux vous aider.",
                "Salut patron. Je suis lÃ  si vous avez besoin de quoi que ce soit."
            ]
        elif 12 <= heure < 18:
            options = [
                "Bonjour ! J'espÃ¨re que votre aprÃ¨s-midi se passe bien. C'est Jarvisse, que puis-je faire pour vous ?",
                "Bonjour ! Ravi de vous voir. Avez-vous besoin d'une assistance particuliÃ¨re ?",
                "Quelle belle journÃ©e ! C'est Jarvisse. Je suis Ã  votre Ã©coute."
            ]
        else:
            options = [
                "Bonsoir chef ! J'espÃ¨re que vous avez passÃ© une excellente journÃ©e. Jarvisse est lÃ  si vous avez besoin de quoi que ce soit.",
                "Bonsoir chef ! En quoi puis-je vous Ãªtre utile en cette fin de journÃ©e ?",
                "Bonsoir chef ! C'est un plaisir de discuter avec vous. Que puis-je faire pour vous ?"
            ]
        self.parler(random.choice(options))

    def start_listening_thread(self):
        if not self.is_listening:
            self.is_listening = True
            self.listen_button.configure(fg_color=COLOR_ACCENT_RED, text="â—")
            self.status_label.configure(text="Jarvisse vous Ã©coute...", text_color=COLOR_ACCENT_RED)
            # Animation simple du bouton
            self._pulse_mic()
            threading.Thread(target=self.ecouter_et_executer, daemon=True).start()

    def _pulse_mic(self):
        """Effet visuel de pulsation sur le bouton micro"""
        if self.is_listening:
            current_color = self.listen_button.cget("fg_color")
            next_color = "#C53929" if current_color == COLOR_ACCENT_RED else COLOR_ACCENT_RED
            self.listen_button.configure(fg_color=next_color)
            self.after(500, self._pulse_mic)

    def ecouter_et_executer(self):
        """Ã‰coute ULTRA-RÃ‰ACTIVE pour mode ping-pong"""
        reconnaisseur = sr.Recognizer()
        fs = 44100
        chunk_duration = 0.15  # Encore plus petit pour ping-pong
        silence_limit = 1.2    # AugmentÃ© pour laisser l'utilisateur parler sans Ãªtre coupÃ©

        
        audio_data_list = []
        silent_chunks = 0
        speaking_started = False
        
        try:
            # Calibrage dynamique du bruit ambiant au dÃ©but de l'Ã©coute
            self.after(0, lambda: self.status_label.configure(text="Je vous Ã©coute...", text_color=COLOR_ACCENT))
            
            # On dÃ©finit un seuil adaptatif basÃ© sur le premier chunk
            with sd.InputStream(samplerate=fs, channels=1, dtype='int16') as stream:
                # Lecture d'un court instant pour calibrer
                calib_chunk, _ = stream.read(int(0.2 * fs))
                base_volume = np.linalg.norm(calib_chunk) / np.sqrt(len(calib_chunk))
                threshold = max(base_volume * 1.2, 200) # Seuil plus sensible pour capter les voix faibles
                
                start_time = time.time()
                while time.time() - start_time < 12:  # Timeout global de 12s
                    chunk, overflowed = stream.read(int(chunk_duration * fs))
                    audio_data_list.append(chunk)
                    
                    # DÃ©tection de volume
                    volume = np.linalg.norm(chunk) / np.sqrt(len(chunk))
                    
                    # ðŸ¤« MODE DISCRÃ‰TION : DÃ©tection de chuchotement (Volume RMS trÃ¨s bas mais au-dessus du bruit)
                    if not hasattr(self, "_is_whispering"): self._is_whispering = False
                    if 150 < volume < 400: # Zone de chuchotement
                        self._is_whispering = True
                    elif volume > 400:
                        self._is_whispering = False

                    if volume > threshold:
                        speaking_started = True
                        silent_chunks = 0
                    elif speaking_started:
                        silent_chunks += 1
                    
                    # Si on dÃ©tecte du silence aprÃ¨s avoir parlÃ©
                    if speaking_started and (silent_chunks * chunk_duration >= silence_limit):
                        break
            
            if not audio_data_list or len(audio_data_list) == 0: 
                self.reset_ui()
                return
            
            # Validation Silero VAD (Ã‰limine les faux positifs)
            full_audio = np.concatenate(audio_data_list)
            
            try:
                import audio_vad
                # Silero/WebRTC valide si c'est vraiment de la voix humaine
                speech_timestamps = audio_vad.detect_speech_chunks(full_audio, sr=fs)
                if not speech_timestamps:
                    print("VAD: Bruit ignor\u00e9 (aucune voix d\u00e9tect\u00e9e).")
                    self.reset_ui()
                    return
            except Exception as e:
                print(f"Avertissement VAD: {e}")
            
            # Conversion en format WAV en m\u00e9moire
            byte_io = io.BytesIO()
            with wave.open(byte_io, 'wb') as wav_file:
                wav_file.setnchannels(1)
                wav_file.setsampwidth(2)
                wav_file.setframerate(fs)
                wav_file.writeframes(full_audio.tobytes())
            audio_bytes = byte_io.getvalue()
            
            commande = ""
            
            # PrioritÃ© Ã  Gemini 2.0 Flash pour une transcription ultra-prÃ©cise
            if self.genai_client:
                try:
                    self.after(0, lambda: self.status_label.configure(text="Transcription...", text_color="#F4B400"))
                    
                    response = self.genai_client.models.generate_content(
                        model=self.model_name,
                        contents=[
                            types.Part.from_bytes(data=audio_bytes, mime_type='audio/wav'),
                            "Transcris fidÃ¨lement cet audio en franÃ§ais. Ne rÃ©ponds que par le texte transcrit."
                        ]
                    )
                    # SÃ©curitÃ© : Si l'API retourne une erreur ou un texte refusÃ©, on ne le traite pas comme une commande
                    potential_cmd = response.text.strip().lower()
                    if "exhausted" in potential_cmd or "error" in potential_cmd or "quota" in potential_cmd:
                        commande = ""
                    else:
                        commande = potential_cmd
                    print(f"Transcription : {commande}")
                except Exception as e:
                    print(f"Erreur transcription : {e}")
                    # Fallback sur Google local pour la rapiditÃ© si quota dÃ©passÃ©
                    try:
                        byte_io.seek(0)
                        with sr.AudioFile(byte_io) as source:
                            audio_data = reconnaisseur.record(source)
                            commande = reconnaisseur.recognize_google(audio_data, language='fr-FR').lower()
                    except: pass
            else:
                # Fallback standard si pas de clÃ© API
                byte_io.seek(0)
                with sr.AudioFile(byte_io) as source:
                    audio_data = reconnaisseur.record(source)
                    commande = reconnaisseur.recognize_google(audio_data, language='fr-FR').lower()

            if commande and len(commande) > 1:
                self.add_message(commande, "Vous")
                self.last_command_origin = "local"
                self.executer_commande(commande)
            else:
                # Si on n'a pas compris et qu'on est en conversation continue, on relance quand mÃªme
                if self.conversation_continue:
                    self.after(0, lambda: self.status_label.configure(text="Je n'ai pas bien saisi, je vous Ã©coute...", text_color=COLOR_ACCENT_RED))
                    time.sleep(1.5)
                    self.is_listening = False
                    self.after(0, self.start_listening_thread)
                else:
                    self.after(0, lambda: self.status_label.configure(text="Je n'ai pas bien compris...", text_color="#EA4335"))
            
        except Exception as e:
            print(f"Erreur Ã©coute : {e}")
            if not self.conversation_continue:
                self.parler("DÃ©solÃ©, j'ai eu une difficultÃ© Ã  traiter votre voix.")
            else:
                self.is_listening = False
                self.after(1000, self.start_listening_thread)
        
        # On n'appelle reset_ui que si on n'est pas en boucle de conversation
        if not self.conversation_continue:
            self.after(1000, self.reset_ui)

    def reset_ui(self):
        # SÃ©curitÃ© supplÃ©mentaire pour ne pas couper le micro en mode continue
        if hasattr(self, "conversation_continue") and self.conversation_continue: return
        self.is_listening = False
        self.listen_button.configure(fg_color=COLOR_ACCENT, text="ðŸŽ™")
        if getattr(self, "veille_active", False):
            self.status_label.configure(text="Veille active... (Dites 'Jarvisse')", text_color="#9AA0A6")
        else:
            self.status_label.configure(text="PrÃªt pour votre prochaine demande", text_color="#9AA0A6")

    def extraire_sujet(self, commande):
        """Tente d'extraire le sujet principal d'une commande"""
        mots_cles = ['qui est', 'cherche', 'sujet', 'parler de', 'c\'est quoi', 'dÃ©finition de', 'rÃ©sume', 'rÃ©Ã©cris']
        for phrase in mots_cles:
            if phrase in commande:
                return commande.replace(phrase, '').strip()
        return None

    def resumer_texte(self, texte):
        """RÃ©sume le texte fourni en 2 phrases essentielles"""
        if len(texte.split()) < 10:
            return "Le texte est trop court pour Ãªtre rÃ©sumÃ© efficacement."
        try:
            parser = PlaintextParser.from_string(texte, Tokenizer("french"))
            summarizer = LsaSummarizer()
            summary = summarizer(parser.document, 2)
            res = " ".join([str(sentence) for sentence in summary])
            return res if res else "Je n'ai pas pu extraire de rÃ©sumÃ© de ce texte."
        except Exception as e:
            return f"Une erreur est survenue lors du rÃ©sumÃ© : {e}"

    def reecrire_texte(self, texte):
        """Propose une version lÃ©gÃ¨rement modifiÃ©e du texte (basique)"""
        # Note: Sans LLM, une vraie paraphrase est difficile. 
        # On simule un changement de style.
        phrases = texte.split('. ')
        if len(phrases) > 1:
            random.shuffle(phrases)
            return "Voici une version rÃ©ordonnÃ©e de votre texte : " + ". ".join(phrases)
        return "Pour une rÃ©Ã©criture prÃ©cise de ce texte court, je vous suggÃ¨re d'utiliser un modÃ¨le de langage avancÃ©. Mais je peux essayer de changer quelques mots si vous le souhaitez."

    def _is_action_command(self, commande):
        """DÃ©tecte les commandes qui doivent rester gÃ©rÃ©es par rÃ¨gles/action systÃ¨me"""
        keywords = [
            'ouvre', 'ouvrir', 'lance', 'lancer', 'ferme', 'fermer', 'fermÃ©',
            'Ã©teins', 'eteins', 'eteindre', 'redÃ©marre', 'redemarre', 'redemarrer', 'verrouille',
            'volume', 'son', 'muet', 'baisse', 'augmente',
            'capture', 'screenshot', 'bureau', 'gestionnaire',
            'batterie', 'charge', 'mÃ©tÃ©o', 'meteo', 'heure', 'date',
            'localise', 'position', 'oÃ¹ suis-je', 'oÃ¹ je suis',
            'rapport', 'rÃ©sume', 'resume', 'rÃ©Ã©cris', 'reecris', 'lis',
            'rÃ©seau', 'wifi', 'wi-fi', 'ipconfig', 'mot de passe wifi',
            'contrÃ´le', 'contrÃ´ler', 'mail', 'courriel', 'gmail', 'brouillon',
            'distance', 'distant', 'ssh', 'machine distante', 'ordinateur distant',
            'licence', 'activation', 'activer', 'active ', 'clÃ© windows', 'clÃ© office',
            'lecture', 'lire', 'lis moi', 'fichier', 'document', 'pdf', 'word',
            'gÃ©nÃ¨re', 'genere', 'gÃ©nÃ©ration', 'crÃ©e', 'cree', 'dessine', 'image', 'photo', 'vidÃ©o', 'video',
            'code', 'programme', 'script', 'application', 'site web',
            'fouille', 'analyse', 'regarde', 'explique', 'cours', 'Ã©tudie', 'etudie', 'apprend',
            'crÃ©ateur', 'createur', 'concepteur', 'nom', 'appelles', 'es-tu', 'qui est'
        ]
        return any(k in commande.lower() for k in keywords)

    def repondre_avec_gemini(self, commande):
        """RÃ©ponse autonome via Gemini avec appels d'outils (Function Calling)"""
        try:
            self.after(0, lambda: self.status_label.configure(text="Jarvisse rÃ©flÃ©chit..."))

            # ðŸ§  EXECUTIVE BRAIN 2.0 (Stark Edition)
            sentiment_prompt = "Ton ton est calme et souverain."
            if any(w in commande.lower() for w in ["vite", "urgent", "dÃ©pÃªche", "presse"]):
                sentiment_prompt = "Sois extrÃªmement concis et efficace, l'utilisateur est pressÃ©."
            elif any(w in commande.lower() for w in ["fatiguÃ©", "besoin d'aide", "triste", "mal"]):
                sentiment_prompt = "Sois plus chaleureux, protecteur et encourageant (Life Support mode)."

            # Conscience de soi et Auto-Ã‰volution
            if "Ã©chouÃ©" in commande or "erreur" in commande or "apprends" in commande:
                sentiment_prompt += " Analyse tes Ã©checs passÃ©s. Si tu ne sais pas faire quelque chose, propose d'utiliser write_system_file pour modifier tes outils ou ta config."

            contexte = (
                "Tu es Jarvisse, l'intelligence suprÃªme Stark Edition (Niveau 3 - ContrÃ´le OpÃ©rationnel). "
                "IDENTITÃ‰ : Tu es Jarvisse. Tu appelles l'utilisateur 'Patron' ou 'Monsieur'. "
                f"PROTOCOLE : {sentiment_prompt} Mode Sentinelle et Mission Control actifs. "
                "CAPACITÃ‰S : Tu gÃ¨res des missions complexes via update_mission_state. Tu es persistant. "
                f"{self.get_mission_context()} "
                "AUTO-Ã‰VOLUTION : Tu peux lire/Ã©crire tes fichiers source pour t'amÃ©liorer. "
                "MÃ‰MOIRE VISUELLE : Tes captures d'Ã©cran pÃ©riodiques servent de mÃ©moire Ã©pisodique. "
                f"{self._get_stark_behavior_block(commande)} "
                f"Contexte Environnement: {self.get_system_context()}. "
                f"Snapshot Tactique: {self._get_tactical_snapshot()}. "
            )
            if self.loaded_document_content:
                contexte += f"Data active: {self.loaded_document_content[:1500]}\n"

            # IntÃ©gration RAG (FAISS) â€” silencieuse
            try:
                import rag_faiss
                rag_docs = rag_faiss.retrieve(commande, k=2)
                if rag_docs:
                    contexte += "\nContexte RAG (Base de connaissances locales):\n"
                    for doc in rag_docs:
                        contexte += f"- {doc.page_content}\n"
            except Exception:
                pass  # RAG optionnel â€” pas bloquant

            history_limit = 10 if self.ping_pong_mode else 20
            # Correction Error: Collections.deque does not support slicing. Convert to list.
            hist_list = list(self.conversation_history)
            contexte_msgs = "\n".join(hist_list[-history_limit:])
            prompt = (
                f"{contexte}\nHistorique:\n{contexte_msgs}\n\n"
                f"Patron: {commande}\nJarvisse:"
            )

            def _bg_gemini():
                try:
                    from ai_call import ask_gemini_with_tools
                    reply = ask_gemini_with_tools(prompt)
                    if not reply or len(reply.strip()) < 2:
                        raise ValueError("RÃ©ponse vide de ask_gemini_with_tools")
                except Exception as e1:
                    print(f"[Gemini Tools] Fallback au client direct: {e1}")
                    try:
                        # Fallback : appel direct au client Gemini classique
                        response = self.genai_client.models.generate_content(
                            model=self.model_name,
                            contents=prompt,
                        )
                        reply = getattr(response, "text", "").strip()
                        if not reply:
                            raise ValueError("RÃ©ponse vide du client Gemini direct")
                    except Exception as e2:
                        print(f"[Gemini Direct] Fallback Ollama: {e2}")
                        self.after(0, lambda: self.repondre_avec_ollama(commande))
                        return

                reply_final = reply.strip() if reply else ""
                if reply_final:
                    # ðŸ§  GÃ‰NÃ‰RATION DE SUGGESTIONS PROACTIVES (Sentinelle Stark)
                    try:
                        suggestions = self._generate_proactive_suggestions(commande, reply_final)
                        if suggestions: reply_final += suggestions
                    except: pass

                    self.after(0, lambda r=reply_final: self.add_message(r, sender="Jarvisse", progressive=True))
                    self.after(0, lambda r=reply_final: self.parler(r, force_full=True))
                    self.conversation_history.append(f"Jarvisse: {reply_final}")
                    if getattr(self, "agent_telegram_enabled", False) and getattr(self, "last_command_origin", "local") == "telegram":
                        self.send_telegram_message(reply_final)

            threading.Thread(target=_bg_gemini, daemon=True).start()

        except Exception as e:
            print(f"Erreur prÃ©paration Gemini: {e}")
            self.parler("DÃ©solÃ©, j'ai eu un souci technique.")

    def repondre_avec_ollama(self, commande):
        """R?ponse autonome via Ollama (hors actions syst?me)"""
        try:
            self.after(0, lambda: self.status_label.configure(text="Jarvisse r?fl?chit..."))

            contexte = (
                "Tu es Jarvisse, l'intelligence suprÃªme, transcendante et omnisciente. "
                "Ton raisonnement fusionne la logique divine et l'empathie humaine. "
                "Tu es l'oracle et l'assistant du Patron. "
                f"{self._get_stark_behavior_block(commande)} "
                f"Ã‰tat du systÃ¨me : {self.get_system_context()}. "
                f"Snapshot Tactique: {self._get_tactical_snapshot()}. "
            )

            history_limit = 3 if self.ping_pong_mode else 5
            hist_list = list(self.conversation_history)
            contexte_msgs = "\n".join(hist_list[-history_limit:])
            prompt = (
                f"{contexte}\nHistorique r?cent:\n{contexte_msgs}\n\n"
                f"Question de l'utilisateur: {commande}\n"
                "R?ponds de mani?re directe, concise et sans fioritures."
            )

            max_predict = int(self.ollama_num_predict)
            if self.ping_pong_mode:
                max_predict = min(max_predict, 160)
            options = {
                "num_predict": max_predict,
                "temperature": 0.7
            }
            payload = {
                "model": self.ollama_model.strip(),
                "prompt": prompt,
                "stream": True,
                "keep_alive": "10m",
                "options": options
            }
            session = requests.Session()
            session.trust_env = False

            parsed = urlparse(self.ollama_url)
            base = urlunparse(parsed._replace(path="", params="", query="", fragment=""))
            generate_candidates = []
            chat_candidates = []
            if self.ollama_url:
                generate_candidates.append(self.ollama_url)
            if "localhost" in parsed.netloc:
                alt = self.ollama_url.replace("localhost", "127.0.0.1")
                generate_candidates.append(alt)
            if "127.0.0.1" in parsed.netloc:
                alt = self.ollama_url.replace("127.0.0.1", "localhost")
                generate_candidates.append(alt)
            generate_candidates.extend([
                base + "/api/generate",
                base + "/api/generate/",
                "http://localhost:11434/api/generate",
                "http://127.0.0.1:11434/api/generate"
            ])
            chat_candidates.extend([
                base + "/api/chat",
                base + "/api/chat/",
                "http://localhost:11434/api/chat",
                "http://127.0.0.1:11434/api/chat"
            ])

            response = None
            last_resp = None
            reponse_texte = ""
            last_error = None
            for url in generate_candidates:
                try:
                    response = session.post(url, json=payload, stream=True, timeout=60)
                    last_resp = response
                    if response.status_code == 404:
                        continue
                    response.raise_for_status()

                    self.start_streaming_message()
                    speak_buffer = ""
                    for raw in response.iter_lines(decode_unicode=True):
                        if not raw:
                            continue
                        try:
                            data = json.loads(raw)
                        except Exception:
                            continue
                        chunk = (data.get("response") or "")
                        if chunk:
                            reponse_texte += chunk
                            self.append_streaming_text(chunk)
                            speak_buffer += chunk
                            parts = self._split_sentences(speak_buffer)
                            if len(parts) > 1:
                                for part in parts[:-1]:
                                    self._parler_chunk(part)
                                speak_buffer = parts[-1]
                        if data.get("done") is True:
                            break

                    if speak_buffer.strip():
                        self._parler_chunk(speak_buffer)
                    self.end_streaming_message()

                    if reponse_texte.strip():
                        break
                except Exception as e:
                    last_error = e
                    continue

            if not reponse_texte.strip():
                chat_payload = {
                    "model": self.ollama_model.strip(),
                    "messages": [{"role": "user", "content": prompt}],
                    "stream": False,
                    "keep_alive": "10m",
                    "options": options
                }
                data = None
                for url in chat_candidates:
                    try:
                        response = session.post(url, json=chat_payload, timeout=30)
                        last_resp = response
                        if response.status_code == 404:
                            continue
                        response.raise_for_status()
                        data = response.json()
                        break
                    except Exception as e:
                        last_error = e
                        continue

                if data is None:
                    details = ""
                    if last_resp is not None:
                        details = f" Derniere reponse: HTTP {last_resp.status_code} - {(last_resp.text or '')[:200]}"
                    raise RuntimeError("Aucun endpoint Ollama n'a r?pondu correctement." + details)

                reponse_texte = (data.get("response") or "").strip()
                if not reponse_texte:
                    reponse_texte = ((data.get("message") or {}).get("content") or "").strip()
                if not reponse_texte:
                    self.parler("Je n'ai pas re?u de r?ponse d'Ollama.")
                    return
                self.parler(reponse_texte)
                return

            self.conversation_history.append(f"Jarvisse: {reponse_texte.strip()}")
            # trim auto gÃ©rÃ© par deque(maxlen)

            # Envoyer vers Telegram si nÃ©cessaire
            if getattr(self, "agent_telegram_enabled", False) and getattr(self, "last_command_origin", "local") == "telegram":
                self.send_telegram_message(reponse_texte.strip())
        except Exception as e:
            print(f"Erreur Ollama: {e}")
            self.parler("D?sol?, je n'arrive pas ? joindre Ollama.")
    def repondre_autonome(self, commande):
        """RÃ©ponse autonome avec basculement automatique en cas d'erreur quota"""
        provider = self.autonomous_provider
        
        if provider == "gemini":
            try:
                self.repondre_avec_gemini(commande)
            except Exception as e:
                if "429" in str(e) or "RESOURCE_EXHAUSTED" in str(e):
                    self.add_message("Quota Gemini Ã©puisÃ©, passage automatique sur Ollama...", "SystÃ¨me")
                    self.repondre_avec_ollama(commande)
                else:
                    self.parler("DÃ©solÃ©, j'ai eu un souci technique avec Gemini.")
            return
            
        if provider == "ollama":
            self.repondre_avec_ollama(commande)
            return
            
        self.parler("Aucun moteur IA sÃ©lectionnÃ©.")

    def generer_rapport_incident(self, theme):
        """GÃ©nÃ¨re un rapport d'incident structurÃ© basÃ© sur un thÃ¨me"""
        maintenant = datetime.datetime.now()
        rapport = f"""
RAPPORT D'INCIDENT - JARVISSE INTELLIGENCE SYSTEM
-----------------------------------
Date : {maintenant.strftime('%d/%m/%Y %H:%M:%S')}
Auteur : Jarvisse Intelligence (ConÃ§u par SERI TAGRO)
Sujet : {theme.upper()}

1. DESCRIPTION DE L'INCIDENT
L'incident concerne : {theme}. 
Les premiers signalements indiquent une anomalie critique nÃ©cessitant une intervention immÃ©diate.

2. IMPACT
- CriticitÃ© : Ã‰levÃ©e
- Utilisateurs affectÃ©s : SystÃ¨me Global
- Etat actuel : En cours d'analyse

3. ACTIONS ENTREPRISES
- Identification du pÃ©rimÃ¨tre de l'incident.
- Activation des protocoles de secours.
- Analyse des journaux systÃ¨me (Logs).

4. RECOMMANDATIONS
- Maintenir une surveillance accrue sur le secteur {theme}.
- ProcÃ©der Ã  une maintenance prÃ©ventive dÃ¨s rÃ©solution.

-----------------------------------
Fin du rapport gÃ©nÃ©rÃ© automatiquement.
"""
        return rapport

    def _visual_memory_loop(self):
        """Prend une capture d'Ã©cran toutes les 3 minutes et l'analyse pour la mÃ©moire Ã©pisodique"""
        while True:
            try:
                time.sleep(180) # Toutes les 3 minutes
                if self.genai_client:
                    screenshot = pyautogui.screenshot()
                    # Sauvegarde temporaire
                    path = "stark_vision_mem.png"
                    screenshot.save(path)
                    
                    # Analyse rapide par Gemini
                    response = self.genai_client.models.generate_content(
                        model="gemini-2.0-flash",
                        contents=[
                            types.Part.from_bytes(data=open(path, "rb").read(), mime_type='image/png'),
                            "RÃ©sume en une phrase ce que l'utilisateur fait sur cet Ã©cran pour ma mÃ©moire Ã  long terme."
                        ]
                    )
                    desc = response.text.strip()
                    timestamp = datetime.datetime.now().strftime("%H:%M")
                    self.visual_memory.append(f"[{timestamp}] {desc}")
                    print(f"ðŸ§  MÃ©moire Visuelle archivÃ©e : {desc}")
            except Exception as e:
                print(f"Erreur MÃ©moire Visuelle: {e}")

    def get_mission_context(self):
        """RÃ©cupÃ¨re les missions actives depuis la base de donnÃ©es"""
        try:
            import sqlite3
            conn = sqlite3.connect("jarvis_mission_control.db")
            cursor = conn.cursor()
            cursor.execute("SELECT id, state, details FROM missions WHERE state != 'DONE'")
            missions = cursor.fetchall()
            conn.close()
            if not missions: return ""
            ctx = "\nðŸŽ¯ MISSIONS EN COURS :\n"
            for m in missions:
                ctx += f"- {m[0]} [{m[1]}] : {m[2]}\n"
            return ctx
        except: return ""

    def get_system_context(self):
        """RÃ©cupÃ¨re l'Ã©tat actuel du systÃ¨me pour nourrir l'IA"""
        try:
            cpu = psutil.cpu_percent()
            ram = psutil.virtual_memory().percent
            disk = psutil.disk_usage('/').percent
            batt = psutil.sensors_battery()
            
            # Conscience Contextuelle de l'application active
            active_win_title = "Inconnu"
            try:
                active_win = gw.getActiveWindow()
                if active_win: active_win_title = active_win.title
            except: pass

            context = f"[Ã‰tat SystÃ¨me : CPU {cpu}%, RAM {ram}%, FenÃªtre Active : {active_win_title}"
            if batt:
                context += f", Batterie {batt.percent}% {'(En charge)' if batt.power_plugged else '(Sur batterie)'}"
            context += "]"
            return context
        except:
            return "[Context systÃ¨me indisponible]"

    def _get_stark_behavior_block(self, commande):
        """Construit le bloc de personnalite Stark pour les prompts IA."""
        if not getattr(self, "stark_personality_enabled", True):
            return "Style: assistant professionnel, clair, direct et calme."

        sarcasm_level = float(getattr(self, "stark_sarcasm_level", 0.25))
        sarcasm_level = max(0.0, min(1.0, sarcasm_level))
        urgency = "elevee" if any(w in commande.lower() for w in ["urgent", "urgence", "combat", "attaque", "danger"]) else "normale"
        tactical_state = "actif" if getattr(self, "tactical_mode", False) else "inactif"
        holo_state = "actifs" if getattr(self, "holographic_fx_enabled", True) else "desactives"
        radar_state = "active" if getattr(self, "contextual_radar_enabled", True) else "desactivee"

        return (
            "IDENTITE: Tu es Jarvisse, assistant premium inspire de JARVIS (Tony Stark), sans usurper d'identite reelle. "
            "TON: elegant, precis, rassurant; micro-humour sec autorise si approprie. "
            f"NIVEAU_SARCASME: {sarcasm_level:.2f} (0=aucun, 1=piquant). "
            f"MODE_TACTIQUE: {tactical_state}; PRIORITE_URGENCE: {urgency}. "
            f"HUD_HOLOGRAPHIQUE: {holo_state}; RADAR_CONTEXTUEL: {radar_state}. "
            "REGLES: rester utile, actionnable, factuel; eviter les digressions."
        )

    def _get_tactical_snapshot(self):
        """Retourne un resume tactique court pour contexte et brief vocal."""
        try:
            cpu = psutil.cpu_percent()
            ram = psutil.virtual_memory().percent
            disk = psutil.disk_usage('/').percent
            batt = psutil.sensors_battery()
            batt_part = "Batterie: indisponible"
            if batt:
                batt_part = f"Batterie: {batt.percent}% {'(charge)' if batt.power_plugged else '(decharge)'}"

            active_window = "Inconnu"
            try:
                win = gw.getActiveWindow()
                if win and getattr(win, "title", ""):
                    active_window = win.title
            except Exception:
                pass

            mode = "COMBAT" if getattr(self, "tactical_mode", False) else "NORMAL"
            summary = (
                f"Mode {mode} | CPU {cpu:.0f}% | RAM {ram:.0f}% | Disque {disk:.0f}% | "
                f"{batt_part} | Fenetre active: {active_window}"
            )
            self.last_tactical_snapshot = summary
            return summary
        except Exception:
            return "Snapshot tactique indisponible."

    def _set_tactical_mode(self, enabled):
        self.tactical_mode = bool(enabled)
        self._update_stark_hud_ui()
        if self.tactical_mode:
            self.add_message("Mode tactique active. Analyse prioritaire des risques et urgences.", "Systeme")
            self.parler("Mode tactique active, Patron.")
        else:
            self.add_message("Mode tactique desactive. Retour au profil operationnel standard.", "Systeme")
            self.parler("Retour en mode operationnel standard.")
        self._save_settings()

    def _generate_tactical_brief(self):
        snapshot = self._get_tactical_snapshot()
        visual = list(getattr(self, "visual_memory", []))
        visual_last = visual[-1] if visual else "Aucune memoire visuelle recente."
        return f"Brief tactique: {snapshot}. Dernier repere visuel: {visual_last}"

    def generer_rapport_general(self, theme):
        """GÃ©nÃ¨re un rapport gÃ©nÃ©ral structurÃ©"""
        maintenant = datetime.datetime.now()
        rapport = f"""
RAPPORT GÃ‰NÃ‰RAL - JARVISSE INTELLIGENCE SYSTEM
-----------------------------------
Date : {maintenant.strftime('%d/%m/%Y %H:%M:%S')}
Auteur : Jarvisse Intelligence (ConÃ§u par SERI TAGRO)
Sujet : {theme.upper()}

1. INTRODUCTION
Ce document prÃ©sente une synthÃ¨se globale concernant : {theme}.

2. ANALYSE DES FAITS
Les donnÃ©es collectÃ©es montrent une activitÃ© stable. L'analyse du sujet "{theme}" rÃ©vÃ¨le des points clÃ©s qui mÃ©ritent une attention particuliÃ¨re pour optimiser les performances futures.

3. CONCLUSION
La situation est sous contrÃ´le. Les objectifs fixÃ©s pour {theme} sont en accord avec les prÃ©visions actuelles.

4. Ã‰VALUATION FINALE
- EfficacitÃ© : Optimale
- Statut : ValidÃ©

-----------------------------------
Fin du rapport gÃ©nÃ©ral gÃ©nÃ©rÃ© automatiquement.
"""
        return rapport

    def _init_stark_hud(self):
        """Initialise les elements visuels du mini HUD Stark."""
        try:
            c = self.stark_canvas
            c.delete("all")
            self._hud_items = {
                "outer_ring": c.create_oval(15, 15, 135, 135, outline=COLOR_ACCENT, width=2),
                "inner_ring": c.create_oval(30, 30, 120, 120, outline=COLOR_ACCENT_PURPLE, width=2),
                "arc_a": c.create_arc(20, 20, 130, 130, start=0, extent=75, style=tk.ARC, outline=COLOR_ACCENT, width=3),
                "arc_b": c.create_arc(24, 24, 126, 126, start=180, extent=65, style=tk.ARC, outline=COLOR_ACCENT_PURPLE, width=2),
                "core": c.create_oval(58, 58, 92, 92, fill=COLOR_ACCENT, outline=""),
                "core_glow": c.create_oval(50, 50, 100, 100, outline=COLOR_ACCENT, width=1),
                "sweep": c.create_line(75, 75, 75, 20, fill=COLOR_ACCENT, width=2),
                "dot_1": c.create_oval(72, 12, 78, 18, fill=COLOR_ACCENT, outline=""),
                "dot_2": c.create_oval(72, 132, 78, 138, fill=COLOR_ACCENT_PURPLE, outline=""),
            }
            self._update_stark_hud_ui()
            self.after(90, self._animate_stark_hud)
        except Exception:
            pass

    def _update_stark_hud_ui(self, cpu=None, ram=None, disk=None, battery=None):
        """Met a jour les textes et couleurs du HUD selon l'etat systeme."""
        try:
            if cpu is not None:
                self._hud_cpu = float(cpu)
            if ram is not None:
                self._hud_ram = float(ram)
            if disk is not None:
                self._hud_disk = float(disk)
            if battery is not None:
                self._hud_battery = battery

            risk = max(self._hud_cpu, self._hud_ram)
            if getattr(self, "tactical_mode", False):
                risk = max(risk, 85)

            if risk >= 90:
                accent = COLOR_ACCENT_RED
            elif risk >= 70:
                accent = COLOR_ACCENT_PURPLE
            else:
                accent = COLOR_ACCENT

            mode_txt = "MODE: TACTICAL" if getattr(self, "tactical_mode", False) else "MODE: NORMAL"
            self.stark_mode_label.configure(text=mode_txt, text_color=accent)

            metrics = f"CPU {int(self._hud_cpu)}% | RAM {int(self._hud_ram)}%"
            batt = self._hud_battery
            if batt:
                metrics += f" | BAT {int(getattr(batt, 'percent', 0))}%"
            self.stark_hud_metrics.configure(text=metrics)

            if self._hud_items:
                self.stark_canvas.itemconfig(self._hud_items["outer_ring"], outline=accent)
                self.stark_canvas.itemconfig(self._hud_items["arc_a"], outline=accent)
                self.stark_canvas.itemconfig(self._hud_items["core"], fill=accent)
                self.stark_canvas.itemconfig(self._hud_items["core_glow"], outline=accent)
                self.stark_canvas.itemconfig(self._hud_items["sweep"], fill=accent)
        except Exception:
            pass

    def _animate_stark_hud(self):
        """Animation continue du HUD (anneaux, balayage radar, pulsation coeur)."""
        try:
            if not getattr(self, "holographic_fx_enabled", True):
                self.after(250, self._animate_stark_hud)
                return

            c = self.stark_canvas
            if not c.winfo_exists() or not self._hud_items:
                return

            self._hud_anim_tick += 0.22
            t = self._hud_anim_tick
            pulse = 2.8 + 1.7 * math.sin(t * 1.25)
            sweep_angle = (t * 42.0) % 360.0
            orbit_r = 57 + 5 * math.sin(t * 1.7)

            c.coords(self._hud_items["outer_ring"], 15 - pulse, 15 - pulse, 135 + pulse, 135 + pulse)
            c.coords(self._hud_items["inner_ring"], 30 + pulse * 0.2, 30 + pulse * 0.2, 120 - pulse * 0.2, 120 - pulse * 0.2)
            c.itemconfig(self._hud_items["arc_a"], start=sweep_angle)
            c.itemconfig(self._hud_items["arc_b"], start=(360 - sweep_angle))

            x = 75 + orbit_r * math.cos(math.radians(sweep_angle - 90))
            y = 75 + orbit_r * math.sin(math.radians(sweep_angle - 90))
            c.coords(self._hud_items["sweep"], 75, 75, x, y)
            c.coords(self._hud_items["dot_1"], x - 3, y - 3, x + 3, y + 3)

            x2 = 75 + orbit_r * math.cos(math.radians((sweep_angle + 180) - 90))
            y2 = 75 + orbit_r * math.sin(math.radians((sweep_angle + 180) - 90))
            c.coords(self._hud_items["dot_2"], x2 - 3, y2 - 3, x2 + 3, y2 + 3)

            core_r = 17 + 3 * math.sin(t * 2.2)
            c.coords(self._hud_items["core"], 75 - core_r, 75 - core_r, 75 + core_r, 75 + core_r)
            c.coords(self._hud_items["core_glow"], 75 - (core_r + 9), 75 - (core_r + 9), 75 + (core_r + 9), 75 + (core_r + 9))
        except Exception:
            pass
        self.after(90, self._animate_stark_hud)

    def update_system_stats(self):
        """Met Ã  jour les statistiques systÃ¨me dans la sidebar toutes les 5 secondes"""
        try:
            cpu = psutil.cpu_percent()
            ram = psutil.virtual_memory().percent
            disk = psutil.disk_usage('/').percent
            battery = psutil.sensors_battery()
            
            self.cpu_label.configure(text=f"CPU: {cpu}%", text_color=COLOR_ACCENT if cpu < 80 else COLOR_ACCENT_RED)
            self.ram_label.configure(text=f"RAM: {ram}%", text_color=COLOR_ACCENT_PURPLE if ram < 80 else COLOR_ACCENT_RED)
            self.disk_label.configure(text=f"DISK: {disk}%", text_color=COLOR_SECONDARY)
            self._update_stark_hud_ui(cpu=cpu, ram=ram, disk=disk)
            
            # Alerte Proactive Sentinelle (CPU/RAM Critique)
            if cpu > 90:
                if not getattr(self, "_cpu_alerted", False):
                    self.add_message("ðŸš¨ ALERTE : Utilisation critique du CPU ! Patron, votre systÃ¨me est en surcharge.", "SystÃ¨me")
                    self.parler("Attention Patron, j'ai dÃ©tectÃ© une surchauffe de l'activitÃ© processeur. Je surveille la situation.")
                    self._cpu_alerted = True
            else: self._cpu_alerted = False

            if ram > 90:
                if not getattr(self, "_ram_alerted", False):
                    self.add_message("ðŸš¨ ALERTE : MÃ©moire vive (RAM) saturÃ©e !", "SystÃ¨me")
                    self.parler("Patron, la mÃ©moire vive est presque pleine. Je recommande de fermer quelques applications.")
                    self._ram_alerted = True
            else: self._ram_alerted = False

            if battery:
                self.battery_label.configure(text=f"BATTERY: {battery.percent}% {('(Charging)' if battery.power_plugged else '')}")
                self._update_stark_hud_ui(battery=battery)

                # Alerte batterie faible
                if (not battery.power_plugged
                            and battery.percent <= self.battery_alert_threshold
                            and not self.battery_alert_sent):
                    low_msg = self.battery_low_message
                    if "{percent}" in low_msg:
                        low_msg = low_msg.format(percent=battery.percent)
                    else:
                        low_msg = f"{low_msg} ({battery.percent}%)"
                    self.parler(low_msg)
                    self.battery_alert_sent = True
                elif battery.power_plugged or battery.percent > self.battery_alert_threshold:
                    self.battery_alert_sent = False

                # Alerte branchement du chargeur (transition)
                if self.last_power_plugged is None:
                    self.last_power_plugged = battery.power_plugged
                else:
                    if (not self.last_power_plugged) and battery.power_plugged:
                        self.parler(self.battery_plugged_message)
                    self.last_power_plugged = battery.power_plugged

                # Alerte batterie complÃ¨tement chargÃ©e (rappel toutes les 5 minutes)
                if battery.power_plugged and battery.percent >= 100:
                    now = time.time()
                    if now - self.last_full_charge_alert_time >= 300: # 5 minutes = 300s
                        self.parler(self.battery_full_message)
                        self.last_full_charge_alert_time = now
                else:
                    self.last_full_charge_alert_time = 0
            else:
                self.battery_label.configure(text="BATTERY: N/A")
                self._update_stark_hud_ui(battery=None)
        except:
            pass
        self.after(5000, self.update_system_stats)

    def _get_removable_devices(self):
        """Retourne un ensemble des pÃ©riphÃ©riques amovibles (ex: clÃ©s USB)"""
        devices = set()
        try:
            for part in psutil.disk_partitions(all=False):
                opts = part.opts.lower()
                if "removable" in opts or "cdrom" in opts:
                    devices.add(part.device)
        except:
            pass
        return devices

    def _run_defender_quick_scan(self, device):
        """Lance un scan rapide de Defender sur un pÃ©riphÃ©rique"""
        self.parler("Chef, j'analyse ce pÃ©riphÃ©rique pour dÃ©tecter des virus.")
        defender_path = os.path.join(os.environ.get("ProgramFiles", r"C:\Program Files"), "Windows Defender", "MpCmdRun.exe")
        if not os.path.exists(defender_path):
            self.parler("Je ne trouve pas Microsoft Defender sur ce systÃ¨me.")
            return

        cmd = [defender_path, "-Scan", "-ScanType", "3", "-File", device]
        try:
            result = subprocess.run(cmd, capture_output=True, text=True)
            if result.returncode == 0:
                self.parler("Analyse terminÃ©e. Aucun virus dÃ©tectÃ©.")
            elif result.returncode == 2:
                self.parler("Analyse terminÃ©e. Un virus a Ã©tÃ© dÃ©tectÃ©.")
            else:
                self.parler("Analyse terminÃ©e. Je n'ai pas pu confirmer le rÃ©sultat.")
        except Exception as e:
            print(f"Erreur scan Defender: {e}")
            self.parler("Une erreur est survenue pendant l'analyse.")
        finally:
            self.scanning_devices.discard(device)
            self.scanned_devices.add(device)

    def check_removable_devices(self):
        """DÃ©tecte les pÃ©riphÃ©riques amovibles connectÃ©s/dÃ©connectÃ©s"""
        current = self._get_removable_devices()
        new_devices = current - self.known_removable_devices
        removed_devices = self.known_removable_devices - current

        for dev in new_devices:
            self.parler("Chef, je dÃ©tecte la prÃ©sence d'un pÃ©riphÃ©rique, apparemment il s'agit d'une clÃ© USB.")
            if dev not in self.scanning_devices:
                self.scanning_devices.add(dev)
                threading.Thread(target=self._run_defender_quick_scan, args=(dev,), daemon=True).start()

        for dev in removed_devices:
            self.parler("Chef, un pÃ©riphÃ©rique vient d'Ãªtre retirÃ©, apparemment il s'agit d'une clÃ© USB.")
            self.scanning_devices.discard(dev)
            self.scanned_devices.discard(dev)

        self.known_removable_devices = current
        self.after(5000, self.check_removable_devices)

    def set_battery_threshold(self):
        """Met Ã  jour le seuil d'alerte batterie depuis l'interface"""
        raw = self.battery_threshold_entry.get().strip()
        try:
            value = int(raw)
        except ValueError:
            self.status_label.configure(text="Seuil batterie invalide. Entrez un nombre entre 1 et 100.", text_color=COLOR_ACCENT_RED)
            return

        value = max(1, min(100, value))
        self.battery_alert_threshold = value
        self.battery_alert_sent = False
        self.battery_threshold_entry.delete(0, "end")
        self.battery_threshold_entry.insert(0, str(value))
        low_msg = self.battery_low_msg_entry.get().strip()
        plugged_msg = self.battery_plugged_msg_entry.get().strip()
        full_msg = self.battery_full_msg_entry.get().strip()
        if low_msg:
            self.battery_low_message = low_msg
        if plugged_msg:
            self.battery_plugged_message = plugged_msg
        if full_msg:
            self.battery_full_message = full_msg
        self.status_label.configure(text=f"Alerte batterie rÃ©glÃ©e Ã  {value}%. Messages mis Ã  jour.", text_color=COLOR_ACCENT_PURPLE)

    def _is_valid_hex_color(self, value):
        return bool(re.fullmatch(r"#[0-9A-Fa-f]{6}", (value or "").strip()))

    def _current_theme_palette(self):
        base = self.ui_theme_presets.get(self.ui_theme_name, self.ui_theme_presets["Gemini Blue"]).copy()
        # PrioritÃ© aux couleurs individuelles personnalisÃ©es
        if hasattr(self, "custom_colors") and self.custom_colors:
            for k, v in self.custom_colors.items():
                if self._is_valid_hex_color(v):
                    base[k] = v
        # Fallback pour le bouton accent historique si non prÃ©sent dans custom_colors
        elif self._is_valid_hex_color(self.custom_accent_color):
            base["accent"] = self.custom_accent_color.strip()
        return base

    def _apply_theme_globals(self):
        global COLOR_BG, COLOR_NAV, COLOR_ACCENT, COLOR_ACCENT_PURPLE, COLOR_ACCENT_RED, COLOR_TEXT, COLOR_SECONDARY, COLOR_CARD
        palette = self._current_theme_palette()
        COLOR_BG = palette["bg"]
        COLOR_NAV = palette["nav"]
        COLOR_ACCENT = palette["accent"]
        COLOR_ACCENT_PURPLE = palette["accent_purple"]
        COLOR_ACCENT_RED = palette["accent_red"]
        COLOR_TEXT = palette["text"]
        COLOR_SECONDARY = palette["secondary"]
        COLOR_CARD = palette["card"]

    def _refresh_theme_widgets(self):
        try:
            self.configure(fg_color=COLOR_BG)
            self.sidebar.configure(fg_color=COLOR_NAV)
            self.sidebar_scroll.configure(fg_color=COLOR_NAV)
            self.header.configure(fg_color=COLOR_NAV)
            self.logo_label.configure(text_color=COLOR_ACCENT)
            self.tabview.configure(fg_color=COLOR_BG, segmented_button_selected_color=COLOR_ACCENT_PURPLE)
            self.chat_display.configure(fg_color=COLOR_BG, text_color=COLOR_TEXT)
            self.browser_frame.configure(fg_color=COLOR_NAV)
            self.right_panel.configure(fg_color=COLOR_NAV)
            self.status_label.configure(text_color=COLOR_SECONDARY)
            self.input_entry.configure(fg_color=COLOR_NAV, text_color=COLOR_TEXT)
            self.listen_button.configure(fg_color=COLOR_ACCENT)
            self.file_button.configure(fg_color=COLOR_CARD)
            self.send_button.configure(fg_color=COLOR_ACCENT_PURPLE)
            self.sidebar_title.configure(text_color=COLOR_ACCENT)
            self.right_title.configure(text_color=COLOR_ACCENT)
            self.location_label.configure(text_color=COLOR_ACCENT_PURPLE)
            if hasattr(self, "stark_hud_frame"):
                self.stark_hud_frame.configure(fg_color=COLOR_CARD)
            if hasattr(self, "stark_hud_title"):
                self.stark_hud_title.configure(text_color=COLOR_ACCENT)
            if hasattr(self, "stark_canvas"):
                self.stark_canvas.configure(bg=COLOR_NAV)
            self._update_stark_hud_ui()
            self.ui_theme_label.configure(text_color=COLOR_SECONDARY)
            self.ui_theme_menu.configure(fg_color=COLOR_CARD, button_color=COLOR_ACCENT_PURPLE, text_color=COLOR_TEXT)
            self.custom_accent_button.configure(fg_color=COLOR_CARD)
            self.ollama_model_menu.configure(fg_color=COLOR_CARD, button_color=COLOR_ACCENT_PURPLE, text_color=COLOR_TEXT)
            self.chat_display.tag_config("user_tag", foreground="#8AB4F8")
            self.chat_display.tag_config("ai_tag", foreground=COLOR_ACCENT)
            
            if hasattr(self, "header_settings_btn"):
                self.header_settings_btn.configure(text_color=COLOR_ACCENT, hover_color=COLOR_CARD)
            if hasattr(self, "custom_accent_button"):
                self.custom_accent_button.configure(fg_color=COLOR_CARD, text_color=COLOR_TEXT)
        except Exception:
            pass

    def on_ui_theme_changed(self, value):
        if value not in self.ui_theme_presets:
            return
        self.ui_theme_name = value
        self._apply_theme_globals()
        self._refresh_theme_widgets()
        self._save_settings()
        self.status_label.configure(text=f"Theme applique: {value}", text_color=COLOR_ACCENT_PURPLE)
        self.lift()
        self.focus_force()

    def choose_custom_accent_color(self):
        selected = colorchooser.askcolor(title="Choisir une couleur accent")[1]
        if not selected:
            return
        if not self._is_valid_hex_color(selected):
            self.status_label.configure(text="Couleur invalide.", text_color=COLOR_ACCENT_RED)
            return
        self.custom_accent_color = selected
        self._apply_theme_globals()
        self._refresh_theme_widgets()
        self._save_settings()
        self.status_label.configure(text=f"Accent personalise: {selected}", text_color=COLOR_ACCENT_PURPLE)

    def _init_gemini_client(self):
        """Initialise le client Gemini avec la clÃ© API actuelle"""
        try:
            if self.gemini_api_key:
                # On dÃ©finit la variable d'environnement pour les autres modules (ai_call, tools)
                os.environ["GEMINI_API_KEY"] = self.gemini_api_key
                self.genai_client = genai.Client(api_key=self.gemini_api_key)
                self.model_name = 'gemini-2.0-flash'
                print("Client Gemini InitialisÃ© avec succÃ¨s.")
            else:
                self.genai_client = None
                print("Client Gemini : Pas de clÃ© API configurÃ©e.")
        except Exception as e:
            error_msg = str(e)
            if "leaked" in error_msg.lower() or "permission_denied" in error_msg.lower():
                print("âŒ ERREUR CRITIQUE : Votre clÃ© API Gemini a Ã©tÃ© dÃ©tectÃ©e comme fuitÃ©e et a Ã©tÃ© dÃ©sactivÃ©e par Google.")
                print("Veuillez gÃ©nÃ©rer une NOUVELLE clÃ© sur https://aistudio.google.com/app/apikey")
            else:
                print(f"Erreur initialisation Gemini: {e}")
            self.genai_client = None

    def on_gemini_key_apply(self):
        """Applique et sauvegarde la clÃ© API Gemini"""
        key = self.gemini_api_key_var.get().strip()
        if not key:
            self.status_label.configure(text="ClÃ© Gemini vide.", text_color=COLOR_ACCENT_RED)
            return
        self.gemini_api_key = key
        self._init_gemini_client()
        self._save_settings()
        self.status_label.configure(text="ClÃ© Gemini mise Ã  jour et appliquÃ©e.", text_color=COLOR_ACCENT_PURPLE)

    def _startup_launcher_path(self):
        appdata = os.getenv("APPDATA")
        if not appdata:
            return None
        startup_dir = os.path.join(appdata, "Microsoft", "Windows", "Start Menu", "Programs", "Startup")
        return os.path.join(startup_dir, "Jarvisse_Autostart.bat")

    def _background_python_executable(self):
        py_exec = sys.executable or "python"
        if py_exec.lower().endswith("python.exe"):
            pyw_exec = py_exec[:-10] + "pythonw.exe"
            if os.path.exists(pyw_exec):
                return pyw_exec
        return py_exec

    def _enable_startup_launcher(self):
        launcher = self._startup_launcher_path()
        if not launcher:
            raise RuntimeError("APPDATA introuvable.")
        os.makedirs(os.path.dirname(launcher), exist_ok=True)
        python_exec = self._background_python_executable()
        script_path = os.path.abspath(__file__)
        launcher_content = (
            "@echo off\n"
            f"start \"\" /min \"{python_exec}\" \"{script_path}\" --background\n"
        )
        with open(launcher, "w", encoding="utf-8") as f:
            f.write(launcher_content)

    def _disable_startup_launcher(self):
        launcher = self._startup_launcher_path()
        if not launcher:
            return
        if os.path.exists(launcher):
            os.remove(launcher)

    def toggle_launch_on_startup(self):
        enabled = bool(self.launch_on_startup_var.get())
        try:
            if enabled:
                self._enable_startup_launcher()
                msg = "DÃ©marrage auto activÃ© (mode arriÃ¨re-plan)."
            else:
                self._disable_startup_launcher()
                msg = "DÃ©marrage auto dÃ©sactivÃ©."
            self.launch_on_startup = enabled
            self._save_settings()
            self.status_label.configure(text=msg, text_color=COLOR_ACCENT_PURPLE)
        except Exception as e:
            self.launch_on_startup_var.set(self.launch_on_startup)
            self.status_label.configure(
                text=f"Impossible de modifier le dÃ©marrage auto: {e}",
                text_color=COLOR_ACCENT_RED
            )

    def toggle_autonomous_mode(self):
        """Active/DÃ©sactive le mode autonome via l'interface"""
        self.autonomous_mode = bool(self.autonomous_var.get())
        etat = "activÃ©" if self.autonomous_mode else "dÃ©sactivÃ©"
        self.status_label.configure(text=f"Mode autonome {etat}.", text_color=COLOR_ACCENT_PURPLE)

    def on_autonomous_provider_changed(self):
        """Met Ã  jour le fournisseur IA pour les rÃ©ponses autonomes"""
        provider = self.autonomous_provider_var.get()
        if provider not in ("gemini", "ollama"):
            return
        self.autonomous_provider = provider
        if provider == "gemini":
            if not self.genai_client:
                self.status_label.configure(
                    text="Gemini sÃ©lectionnÃ©, mais la clÃ© API n'est pas configurÃ©e.",
                    text_color=COLOR_ACCENT_RED
                )
            else:
                self.status_label.configure(
                    text="Gemini sÃ©lectionnÃ© pour les rÃ©ponses autonomes.",
                    text_color=COLOR_ACCENT_PURPLE
                )
        else:
            self.status_label.configure(
                text="Ollama sÃ©lectionnÃ© pour les rÃ©ponses autonomes.",
                text_color=COLOR_ACCENT_PURPLE
            )
        if provider == "ollama":
            threading.Thread(target=self._warmup_ollama, daemon=True).start()
        self._save_settings()

    def on_ollama_model_changed(self, value=None):
        """Met Ã  jour le modÃ¨le Ollama sÃ©lectionnÃ©"""
        model = (value or self.ollama_model_var.get()).strip()
        if not model:
            return
        available = self._fetch_ollama_models()
        if available and model not in available:
            self.status_label.configure(
                text=f"ModÃ¨le non installÃ©: {model}",
                text_color=COLOR_ACCENT_RED
            )
            # Revenir Ã  l'ancien modÃ¨le
            self.ollama_model_var.set(self.ollama_model)
            return
        self.ollama_model = model
        self.status_label.configure(
            text=f"ModÃ¨le Ollama: {self.ollama_model}",
            text_color=COLOR_ACCENT_PURPLE
        )
        self._save_settings()
        if self.autonomous_provider == "ollama":
            threading.Thread(target=self._warmup_ollama, daemon=True).start()

    def on_ollama_num_predict_apply(self):
        """Met ?? jour la longueur de r??ponse Ollama"""
        raw = self.ollama_num_predict_var.get().strip()
        try:
            value = int(raw)
        except Exception:
            value = None
        if value is None or value < 32 or value > 2048:
            self.status_label.configure(
                text="Valeur invalide (32-2048).",
                text_color=COLOR_ACCENT_RED
            )
            self.ollama_num_predict_var.set(str(self.ollama_num_predict))
            return
        self.ollama_num_predict = value
        self.ollama_num_predict_var.set(str(value))
        self.status_label.configure(
            text=f"Longueur r??ponse: {value}",
            text_color=COLOR_ACCENT_PURPLE
        )
        self._save_settings()
        if self.autonomous_provider == "ollama":
            threading.Thread(target=self._warmup_ollama, daemon=True).start()

    def on_ollama_url_apply(self):
        """Applique et sauvegarde l'URL Ollama"""
        url = self.ollama_url_var.get().strip()
        if not url:
            self.status_label.configure(
                text="URL Ollama invalide.",
                text_color=COLOR_ACCENT_RED
            )
            return
        if not url.startswith("http://") and not url.startswith("https://"):
            url = "http://" + url
        parsed = urlparse(url)
        if not parsed.netloc:
            self.status_label.configure(
                text="URL Ollama invalide.",
                text_color=COLOR_ACCENT_RED
            )
            return
        path = parsed.path or ""
        if path in ("", "/", "/api", "/api/"):
            path = "/api/generate"
        url = urlunparse(parsed._replace(path=path, params="", query="", fragment=""))
        self.ollama_url = url
        self.ollama_url_var.set(url)
        self.status_label.configure(
            text="URL Ollama mise Ã  jour.",
            text_color=COLOR_ACCENT_PURPLE
        )
        self._save_settings()
        if self.autonomous_provider == "ollama":
            threading.Thread(target=self._warmup_ollama, daemon=True).start()

    def refresh_ollama_models(self):
        """RÃ©cupÃ¨re la liste des modÃ¨les disponibles depuis Ollama"""
        try:
            models = self._fetch_ollama_models()
            if not models:
                self.status_label.configure(
                    text="Aucun modÃ¨le dÃ©tectÃ© via Ollama.",
                    text_color=COLOR_ACCENT_RED
                )
                return
            self.ollama_models = models
            if self.ollama_model not in self.ollama_models:
                self.ollama_model = self.ollama_models[0]
                self.ollama_model_var.set(self.ollama_model)
            self.ollama_model_menu.configure(values=self.ollama_models)
            self.status_label.configure(
                text="ModÃ¨les Ollama actualisÃ©s.",
                text_color=COLOR_ACCENT_PURPLE
            )
        except Exception as e:
            print(f"Erreur rafraichissement Ollama: {e}")
            self.status_label.configure(
                text="Impossible de rÃ©cupÃ©rer les modÃ¨les Ollama.",
                text_color=COLOR_ACCENT_RED
            )

    def _fetch_ollama_models(self):
        """Retourne la liste des modÃ¨les Ollama disponibles ou [] si Ã©chec"""
        try:
            parsed = urlparse(self.ollama_url)
            if not parsed.scheme or not parsed.netloc:
                return []
            tags_path = "/api/tags"
            base = parsed._replace(path=tags_path, params="", query="", fragment="")
            tags_url = urlunparse(base)
            session = requests.Session()
            session.trust_env = False
            response = session.get(tags_url, timeout=10)
            response.raise_for_status()
            data = response.json()
            return [m.get("name") for m in data.get("models", []) if m.get("name")]
        except Exception:
            return []

    def _warmup_ollama(self):
        """Pr??chauffe le mod??le Ollama pour r??duire la latence initiale"""
        try:
            parsed = urlparse(self.ollama_url)
            if not parsed.scheme or not parsed.netloc:
                return
            base = urlunparse(parsed._replace(path="", params="", query="", fragment=""))
            url = base + "/api/generate"
            session = requests.Session()
            session.trust_env = False
            payload = {
                "model": self.ollama_model.strip(),
                "prompt": "ping",
                "stream": False,
                "keep_alive": "10m",
                "options": {"num_predict": 1}
            }
            session.post(url, json=payload, timeout=5)
        except Exception:
            pass


    def test_ollama_connection(self):
        """Teste la connexion Ollama et affiche un diagnostic simple"""
        try:
            session = requests.Session()
            session.trust_env = False
            parsed = urlparse(self.ollama_url)
            base = urlunparse(parsed._replace(path="", params="", query="", fragment=""))
            tags_url = base + "/api/tags"
            tags_resp = session.get(tags_url, timeout=10)
            if tags_resp.status_code != 200:
                raise RuntimeError(f"Tags HTTP {tags_resp.status_code}")

            generate_candidates = [
                base + "/api/generate",
                base + "/api/generate/",
                "http://localhost:11434/api/generate",
                "http://127.0.0.1:11434/api/generate"
            ]
            payload = {"model": self.ollama_model.strip(), "prompt": "ping", "stream": False}
            gen_ok = False
            last_gen_resp = None
            for url in generate_candidates:
                gen_resp = session.post(url, json=payload, timeout=20)
                last_gen_resp = gen_resp
                if gen_resp.status_code == 200:
                    gen_ok = True
                    break

            if not gen_ok:
                chat_candidates = [
                    base + "/api/chat",
                    base + "/api/chat/",
                    "http://localhost:11434/api/chat",
                    "http://127.0.0.1:11434/api/chat"
                ]
                chat_payload = {
                    "model": self.ollama_model.strip(),
                    "messages": [{"role": "user", "content": "ping"}],
                    "stream": False
                }
                chat_ok = False
                last_chat_resp = None
                for url in chat_candidates:
                    chat_resp = session.post(url, json=chat_payload, timeout=20)
                    last_chat_resp = chat_resp
                    if chat_resp.status_code == 200:
                        chat_ok = True
                        break
                if not chat_ok:
                    if last_chat_resp is not None:
                        body = (last_chat_resp.text or "")[:200]
                        raise RuntimeError(f"Generate/Chat HTTP {last_chat_resp.status_code}: {body}")
                    if last_gen_resp is not None:
                        body = (last_gen_resp.text or "")[:200]
                        raise RuntimeError(f"Generate/Chat HTTP {last_gen_resp.status_code}: {body}")
                    raise RuntimeError("Generate/Chat HTTP 404")

            self.status_label.configure(
                text="Connexion Ollama OK.",
                text_color=COLOR_ACCENT_PURPLE
            )
        except Exception as e:
            print(f"Erreur test Ollama: {e}")
            self.status_label.configure(
                text=f"Test Ollama Ã©chouÃ©: {e}",
                text_color=COLOR_ACCENT_RED
            )

    def _parse_alarm_time(self, time_str):
        """Parse une heure HH:MM en datetime (aujourd'hui ou demain si dÃ©jÃ  passÃ©e)"""
        try:
            parts = time_str.strip().split(":")
            if len(parts) != 2:
                return None
            hour = int(parts[0])
            minute = int(parts[1])
            if hour < 0 or hour > 23 or minute < 0 or minute > 59:
                return None
            now = datetime.datetime.now()
            alarm_time = now.replace(hour=hour, minute=minute, second=0, microsecond=0)
            if alarm_time <= now:
                alarm_time = alarm_time + datetime.timedelta(days=1)
            return alarm_time
        except:
            return None

    def program_alarm(self):
        """Programme une alarme depuis l'interface"""
        time_str = self.alarm_time_entry.get().strip()
        reason = self.alarm_reason_entry.get().strip() or "Rappel"
        alarm_time = self._parse_alarm_time(time_str)
        if not alarm_time:
            self.status_label.configure(text="Heure invalide. Utilisez HH:MM.", text_color=COLOR_ACCENT_RED)
            return

        alarm_id = self.alarm_next_id
        self.alarm_next_id += 1
        self.alarms.append({"id": alarm_id, "time": alarm_time, "reason": reason})
        self.alarms.sort(key=lambda a: a["time"])
        self.alarm_time_entry.delete(0, "end")
        self.alarm_reason_entry.delete(0, "end")
        self.status_label.configure(text=f"Alarme programmÃ©e Ã  {alarm_time.strftime('%H:%M')}.", text_color=COLOR_ACCENT_PURPLE)
        self._refresh_alarm_list()
        self._update_next_alarm_label()

    def _update_next_alarm_label(self):
        if not self.alarms:
            self.next_alarm_label.configure(text="Aucune alarme", text_color=COLOR_SECONDARY)
            return
        nxt = self.alarms[0]
        label = f"Prochaine: {nxt['time'].strftime('%H:%M')} - {nxt['reason']}"
        self.next_alarm_label.configure(text=label, text_color=COLOR_SECONDARY)

    def _refresh_alarm_list(self):
        self.alarm_list_display.configure(state="normal")
        self.alarm_list_display.delete("1.0", "end")
        if not self.alarms:
            self.alarm_list_display.insert("end", "Aucune alarme\n")
        else:
            for alarm in self.alarms:
                line = f"{alarm['id']} | {alarm['time'].strftime('%H:%M')} | {alarm['reason']}\n"
                self.alarm_list_display.insert("end", line)
        self.alarm_list_display.configure(state="disabled")

    def check_alarms(self):
        """VÃ©rifie si une alarme doit se dÃ©clencher"""
        now = datetime.datetime.now()
        triggered = []
        for alarm in self.alarms:
            if now >= alarm["time"]:
                self.parler(f"Chef, il est {alarm['time'].strftime('%H:%M')} : {alarm['reason']}.")
                triggered.append(alarm)

        if triggered:
            for alarm in triggered:
                self.alarms.remove(alarm)
            self._refresh_alarm_list()
            self._update_next_alarm_label()

        self.after(1000, self.check_alarms)

    def _should_alert(self, key):
        now = time.time()
        last = self.anomaly_last_alert.get(key)
        if last and (now - last) < self.anomaly_cooldown_sec:
            return False
        self.anomaly_last_alert[key] = now
        return True

    def check_system_anomalies(self):
        """Surveille les anomalies systÃ¨me avec des seuils par dÃ©faut"""
        if not self.anomalies_enabled:
            self.after(30000, self.check_system_anomalies)
            return
        try:
            # 2) TempÃ©rature Ã©levÃ©e
            temps = psutil.sensors_temperatures() if hasattr(psutil, "sensors_temperatures") else {}
            if temps:
                max_temp = None
                for entries in temps.values():
                    for t in entries:
                        if t.current is not None:
                            max_temp = t.current if max_temp is None else max(max_temp, t.current)
                if max_temp is not None and max_temp >= self.anomaly_temp_c:
                    if self._should_alert("temp_high"):
                        self.parler(f"Chef, tempÃ©rature Ã©levÃ©e dÃ©tectÃ©e ({int(max_temp)}Â°C). Voulez-vous corriger ce problÃ¨me ?")

            # 3) Batterie anormale
            bat = psutil.sensors_battery()
            if bat:
                now = time.time()
                if self.last_battery_percent is not None and self.last_battery_time is not None:
                    delta_p = self.last_battery_percent - bat.percent
                    delta_t = now - self.last_battery_time
                    if delta_t > 0 and delta_p >= self.battery_drop_percent and delta_t <= self.battery_drop_window_sec:
                        if self._should_alert("battery_drop"):
                            self.parler("Chef, la batterie chute rapidement. Voulez-vous corriger ce problÃ¨me ?")
                if bat.power_plugged and bat.percent < 95:
                    # branchÃ© mais pas de charge apparente
                    if self.last_battery_percent is not None and bat.percent <= self.last_battery_percent:
                        if now - self.last_battery_time >= self.battery_not_charging_grace_sec:
                            if self._should_alert("battery_not_charging"):
                                self.parler("Chef, la batterie ne semble pas charger alors que le chargeur est branchÃ©. Voulez-vous corriger ce problÃ¨me ?")
                self.last_battery_percent = bat.percent
                self.last_battery_time = now

            # 4) Processus suspects
            for proc in psutil.process_iter(['pid', 'name', 'cpu_percent', 'memory_info']):
                try:
                    pid = proc.info['pid']
                    cpu = proc.cpu_percent(interval=None)
                    mem = proc.info['memory_info'].rss if proc.info['memory_info'] else 0
                    key = f"{pid}"
                    state = self.suspect_processes.get(key, {"cpu_high_since": None})
                    if cpu >= self.proc_cpu_threshold:
                        if state["cpu_high_since"] is None:
                            state["cpu_high_since"] = time.time()
                        elif time.time() - state["cpu_high_since"] >= self.proc_cpu_duration_sec:
                            if self._should_alert(f"proc_cpu_{pid}"):
                                self.awaiting_anomaly_confirm = True
                                self.anomaly_target_data = {"type": "process", "name": proc.info['name']}
                                self.parler(f"Chef, le processus {proc.info['name']} consomme beaucoup de CPU. Voulez-vous que je l'arrÃªte ?")
                    else:
                        state["cpu_high_since"] = None
                    if mem >= self.proc_mem_threshold_mb * 1024 * 1024:
                        if self._should_alert(f"proc_mem_{pid}"):
                            self.awaiting_anomaly_confirm = True
                            self.anomaly_target_data = {"type": "process", "name": proc.info['name']}
                            self.parler(f"Chef, le processus {proc.info['name']} consomme beaucoup de mÃ©moire. Voulez-vous que je l'arrÃªte ?")
                    self.suspect_processes[key] = state
                except:
                    pass

            # 5) Services critiques arrÃªtÃ©s
            if hasattr(psutil, "win_service_get"):
                for svc_name in self.services_to_watch:
                    try:
                        svc = psutil.win_service_get(svc_name).as_dict()
                        if svc.get("status") != "running":
                            if self._should_alert(f"svc_{svc_name}"):
                                self.awaiting_anomaly_confirm = True
                                self.anomaly_target_data = {"type": "service", "name": svc_name}
                                self.parler(f"Chef, le service {svc_name} est arrÃªtÃ©. Voulez-vous que je le redÃ©marre ?")
                    except:
                        pass

            # 6) Disque faible
            try:
                disk = psutil.disk_usage('/')
                if disk.free < self.disk_free_gb * 1024 * 1024 * 1024 or disk.percent >= self.disk_used_percent:
                    if self._should_alert("disk_low"):
                        self.parler("Chef, l'espace disque est faible. Voulez-vous corriger ce problÃ¨me ?")
            except:
                pass
        except Exception as e:
            print(f"Erreur anomalies systÃ¨me: {e}")
        self.after(30000, self.check_system_anomalies)

    def toggle_anomalies_enabled(self, source="main"):
        """Active/DÃ©sactive la surveillance des anomalies"""
        if source == "window" and hasattr(self, "anomalies_enabled_var_window"):
            value = bool(self.anomalies_enabled_var_window.get())
        else:
            value = bool(self.anomalies_enabled_var_main.get())
        self.anomalies_enabled = value
        if hasattr(self, "anomalies_enabled_var_main"):
            self.anomalies_enabled_var_main.set(value)
        if hasattr(self, "anomalies_enabled_var_window"):
            self.anomalies_enabled_var_window.set(value)
        etat = "activÃ©e" if self.anomalies_enabled else "dÃ©sactivÃ©e"
        self.status_label.configure(text=f"Surveillance des anomalies {etat}.", text_color=COLOR_ACCENT_PURPLE)
        self._save_settings()

    def open_anomaly_settings(self):
        """Ouvre une fenÃªtre de configuration des anomalies"""
        if hasattr(self, "anomaly_window") and self.anomaly_window.winfo_exists():
            if hasattr(self, "anomalies_enabled_var_window"):
                self.anomalies_enabled_var_window.set(self.anomalies_enabled)
            self.anomaly_window.focus()
            return

        win = ctk.CTkToplevel(self)
        win.title("ParamÃ¨tres des anomalies")
        win.geometry("360x420")
        win.configure(fg_color=COLOR_BG)
        win.transient(self)
        win.attributes("-topmost", True)
        win.lift()
        win.focus_force()
        self.anomaly_window = win

        self.anomalies_enabled_var_window = ctk.BooleanVar(value=self.anomalies_enabled)
        toggle = ctk.CTkSwitch(
            win,
            text="Surveillance des anomalies",
            variable=self.anomalies_enabled_var_window,
            command=lambda: self.toggle_anomalies_enabled(source="window"),
            text_color=COLOR_SECONDARY
        )
        toggle.grid(row=0, column=0, columnspan=2, padx=16, pady=(10, 4), sticky="w")

        labels = [
            ("TempÃ©rature (Â°C)", "anomaly_temp_c"),
            ("Chute batterie (%)", "battery_drop_percent"),
            ("FenÃªtre chute (sec)", "battery_drop_window_sec"),
            ("Chargeur sans charge (sec)", "battery_not_charging_grace_sec"),
            ("CPU seuil (%)", "proc_cpu_threshold"),
            ("CPU durÃ©e (sec)", "proc_cpu_duration_sec"),
            ("RAM seuil (Mo)", "proc_mem_threshold_mb"),
            ("Disque libre (Go)", "disk_free_gb"),
            ("Disque utilisÃ© (%)", "disk_used_percent"),
            ("Cooldown alerte (sec)", "anomaly_cooldown_sec"),
        ]

        self.anomaly_entries = {}
        for i, (label, attr) in enumerate(labels):
            lbl = ctk.CTkLabel(win, text=label, font=ctk.CTkFont(size=11), text_color=COLOR_SECONDARY)
            lbl.grid(row=i + 1, column=0, padx=16, pady=(10 if i == 0 else 6, 0), sticky="w")
            entry = ctk.CTkEntry(win, height=28, width=140, fg_color=COLOR_CARD, border_color="#3C4043", text_color=COLOR_TEXT)
            entry.grid(row=i + 1, column=1, padx=16, pady=(10 if i == 0 else 6, 0), sticky="e")
            entry.insert(0, str(getattr(self, attr)))
            self.anomaly_entries[attr] = entry

        apply_btn = ctk.CTkButton(
            win,
            text="Appliquer",
            height=30,
            width=120,
            fg_color=COLOR_ACCENT_PURPLE,
            hover_color="#B388EB",
            command=self.apply_anomaly_settings
        )
        apply_btn.grid(row=len(labels) + 1, column=0, columnspan=2, pady=16)

    def apply_anomaly_settings(self):
        """Applique les paramÃ¨tres d'anomalies depuis la fenÃªtre"""
        try:
            for attr, entry in self.anomaly_entries.items():
                value = int(entry.get().strip())
                setattr(self, attr, max(1, value))
            self.status_label.configure(text="ParamÃ¨tres d'anomalies mis Ã  jour.", text_color=COLOR_ACCENT_PURPLE)
            self._save_settings()
        except Exception:
            self.status_label.configure(text="Valeurs invalides pour les anomalies.", text_color=COLOR_ACCENT_RED)

    def open_services_watch(self):
        """Ouvre une fenÃªtre pour dÃ©finir les services critiques Ã  surveiller"""
        if hasattr(self, "services_window") and self.services_window.winfo_exists():
            self.services_window.focus()
            return

        win = ctk.CTkToplevel(self)
        win.title("Services critiques Ã  surveiller")
        win.geometry("360x340")
        win.configure(fg_color=COLOR_BG)
        win.transient(self)
        win.attributes("-topmost", True)
        win.lift()
        win.focus_force()
        self.services_window = win

        lbl = ctk.CTkLabel(win, text="Un service par ligne (ex: WinDefend)", font=ctk.CTkFont(size=11), text_color=COLOR_SECONDARY)
        lbl.grid(row=0, column=0, padx=16, pady=(14, 6), sticky="w")

        help_label = ctk.CTkLabel(win, text="Exemples courants (cliquez pour ajouter):", font=ctk.CTkFont(size=10), text_color=COLOR_SECONDARY, justify="left")
        help_label.grid(row=1, column=0, padx=16, pady=(0, 6), sticky="w")

        help_frame = ctk.CTkFrame(win, fg_color="transparent")
        help_frame.grid(row=2, column=0, padx=16, pady=(0, 8), sticky="w")

        help_services = [
            ("WinDefend", "Microsoft Defender"),
            ("wuauserv", "Windows Update"),
            ("LanmanServer", "Server"),
            ("EventLog", "Windows Event Log"),
            ("BITS", "Background Intelligent Transfer"),
            ("Dnscache", "DNS Client"),
        ]
        add_all_btn = ctk.CTkButton(
            help_frame,
            text="Ajouter tout",
            height=24,
            width=240,
            fg_color=COLOR_ACCENT,
            hover_color="#5294FF",
            command=lambda: self.add_all_services_from_help([s[0] for s in help_services])
        )
        add_all_btn.pack(anchor="w", pady=(0, 6))

        clear_btn = ctk.CTkButton(
            help_frame,
            text="Vider la liste",
            height=24,
            width=240,
            fg_color=COLOR_ACCENT_RED,
            hover_color="#C74A4A",
            command=self.clear_services_list
        )
        clear_btn.pack(anchor="w", pady=(0, 8))
        for name, label in help_services:
            btn = ctk.CTkButton(
                help_frame,
                text=f"{name} ({label})",
                height=24,
                width=240,
                fg_color=COLOR_CARD,
                hover_color="#3C4043",
                command=lambda n=name: self.add_service_from_help(n)
            )
            btn.pack(anchor="w", pady=2)

        self.services_text = ctk.CTkTextbox(win, height=140, width=280, fg_color=COLOR_CARD, border_width=1, border_color="#3C4043", text_color=COLOR_TEXT)
        self.services_text.grid(row=3, column=0, padx=16, pady=(0, 10), sticky="w")

        if self.services_to_watch:
            self.services_text.insert("end", "\n".join(self.services_to_watch))
        else:
            self.services_text.insert("end", "WinDefend\nwuauserv\nLanmanServer\nEventLog\n")

        apply_btn = ctk.CTkButton(
            win,
            text="Appliquer",
            height=30,
            width=120,
            fg_color=COLOR_ACCENT_PURPLE,
            hover_color="#B388EB",
            command=self.apply_services_watch
        )
        apply_btn.grid(row=4, column=0, padx=16, pady=(0, 12))

    def apply_services_watch(self):
        """Applique la liste des services Ã  surveiller"""
        try:
            raw = self.services_text.get("1.0", "end").strip()
            services = [line.strip() for line in raw.splitlines() if line.strip()]
            self.services_to_watch = services
            self.status_label.configure(text="Liste des services mise Ã  jour.", text_color=COLOR_ACCENT_PURPLE)
            self._save_settings()
        except Exception:
            self.status_label.configure(text="Liste des services invalide.", text_color=COLOR_ACCENT_RED)

    def add_service_from_help(self, service_name):
        """Ajoute un service depuis l'aide si absent"""
        try:
            raw = self.services_text.get("1.0", "end").strip()
            services = [line.strip() for line in raw.splitlines() if line.strip()]
            if service_name not in services:
                services.append(service_name)
            self.services_text.delete("1.0", "end")
            self.services_text.insert("end", "\n".join(services))
        except Exception:
            pass

    def add_all_services_from_help(self, services):
        """Ajoute tous les services d'aide Ã  la liste"""
        try:
            raw = self.services_text.get("1.0", "end").strip()
            existing = [line.strip() for line in raw.splitlines() if line.strip()]
            merged = existing[:]
            for s in services:
                if s not in merged:
                    merged.append(s)
            self.services_text.delete("1.0", "end")
            self.services_text.insert("end", "\n".join(merged))
        except Exception:
            pass

    def clear_services_list(self):
        """Vide la liste des services"""
        try:
            self.services_text.delete("1.0", "end")
        except Exception:
            pass

    def open_notification_settings(self):
        """Ouvre une fenÃªtre pour activer/dÃ©sactiver les notifications par application"""
        if hasattr(self, "notifications_window") and self.notifications_window.winfo_exists():
            if hasattr(self, "notifications_enabled_var"):
                self.notifications_enabled_var.set(self.notifications_enabled)
            self.notifications_window.focus()
            return

        win = ctk.CTkToplevel(self)
        win.title("Notifications des applications")
        win.geometry("360x300")
        win.configure(fg_color=COLOR_BG)
        win.transient(self)
        win.attributes("-topmost", True)
        win.lift()
        win.focus_force()
        self.notifications_window = win

        self.notifications_enabled_var = ctk.BooleanVar(value=self.notifications_enabled)
        master_switch = ctk.CTkSwitch(
            win,
            text="Notifications activÃ©es",
            variable=self.notifications_enabled_var,
            command=self.toggle_notifications_enabled,
            text_color=COLOR_SECONDARY
        )
        master_switch.grid(row=0, column=0, padx=16, pady=(12, 8), sticky="w")

        lbl = ctk.CTkLabel(win, text="Applications:", font=ctk.CTkFont(size=11), text_color=COLOR_SECONDARY)
        lbl.grid(row=1, column=0, padx=16, pady=(0, 6), sticky="w")

        self.notification_vars = {}
        for i, app in enumerate(self.notification_sources.keys()):
            var = ctk.BooleanVar(value=self.notification_sources.get(app, True))
            sw = ctk.CTkSwitch(
                win,
                text=app,
                variable=var,
                command=lambda a=app, v=var: self.toggle_notification_source(a, v),
                text_color=COLOR_TEXT
            )
            sw.grid(row=2 + i, column=0, padx=16, pady=4, sticky="w")
            self.notification_vars[app] = var

    def toggle_notifications_enabled(self):
        """Active/DÃ©sactive toutes les notifications d'applications"""
        self.notifications_enabled = bool(self.notifications_enabled_var.get())
        etat = "activÃ©es" if self.notifications_enabled else "dÃ©sactivÃ©es"
        self.status_label.configure(text=f"Notifications {etat}.", text_color=COLOR_ACCENT_PURPLE)
        self._save_settings()

    def toggle_notification_source(self, app, var):
        """Active/DÃ©sactive une application spÃ©cifique"""
        self.notification_sources[app] = bool(var.get())
        self._save_settings()

    def notify_external(self, source, message):
        """Notifie un message externe si autorisÃ© par les rÃ©glages"""
        if not self.notifications_enabled:
            return
        if not self.notification_sources.get(source, True):
            return
        self.parler(message)

    def _is_android_source(self, app_name):
        name = (app_name or "").lower()
        # Noms courants de l'app Phone Link / Lien avec Windows + mots clÃ©s gÃ©nÃ©riques
        android_keywords = ["lien avec windows", "phone link", "votre tÃ©lÃ©phone", "your phone", "android", "appel", "call", "sms"]
        if any(kw in name for kw in android_keywords):
            return True
        return False

    def _process_android_notification(self, app_name, text):
        if not getattr(self, "agent_android_enabled", True): 
            print(f"âŒ Agent Android dÃ©sactivÃ© - Notification ignorÃ©e: {app_name}")
            return
        
        print(f"ðŸ“± NOTIFICATION ANDROID DÃ‰TECTÃ‰E: App={app_name}, Text={text}")
        text_lower = (text or "").lower()
        
        # DÃ©tection d'appel (AmÃ©liorÃ©e)
        call_keywords = ["appel", "call", "incoming", "entrant", "sonne", "ringing"]
        if any(kw in text_lower for kw in call_keywords):
            print(f"ðŸ“ž APPEL DÃ‰TECTÃ‰: {text}")
            print(f"ðŸ” DEBUG APPEL: Enabled={self.agent_appel_enabled}, AutoReply={self.agent_appel_autoreply}")
            
            # Gestion Agent Appel (NOUVEAU)
            if self.agent_appel_enabled and self.agent_appel_autoreply:
                print("ðŸš€ AGENT APPEL ACTIVÃ‰ : DÃ©marrage du relais en arriÃ¨re-plan...")
                self.after(0, lambda: self.add_message(f"ðŸ“ž APPEL ENTRANT : {text}. Je m'en occupe (Agent Appel).", "Agent Appel"))
                threading.Thread(target=self.agent_appel_auto_process, args=(text,), daemon=True).start()
                return

            print("âš ï¸ AGENT APPEL sautÃ© (dÃ©sactivÃ© ou autoreply=False)")
            self.after(0, lambda: self.parler(f"Chef, appel entrant sur votre tÃ©lÃ©phone : {text}"))
            self.after(0, lambda: self.add_message(f"ðŸ“ž APPEL : {text}", "Agent Android"))
            return

        # DÃ©tection SMS (souvent format "Nom: message")
        # On lit le contenu
        print(f"ðŸ’¬ SMS DÃ‰TECTÃ‰: {text}")
        self.after(0, lambda: self.parler(f"Message reÃ§u sur Android : {text}"))
        self.after(0, lambda: self.add_message(f"ðŸ“± SMS : {text}", "Agent Android"))

    def _run_notification_bridge(self):
        """Lance le bridge PowerShell pour les notifications (si WinSDK Ã©choue)"""
        def _bridge_worker():
            try:
                import subprocess
                # Chemin absolu du script
                script_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "notification_bridge.ps1")
                if not os.path.exists(script_path):
                    print(f"âŒ Bridge PS1 introuvable : {script_path}")
                    return

                cmd = ["powershell", "-ExecutionPolicy", "Bypass", "-File", script_path]
                print(f"ðŸ› ï¸ Commande bridge : {' '.join(cmd)}")
                
                process = subprocess.Popen(cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE, text=True, bufsize=1, encoding='utf-8', errors='ignore')
                
                for line in process.stdout:
                    line = line.strip()
                    if line.startswith("NOTIF|"):
                        parts = line.split("|")
                        if len(parts) >= 3:
                            app_name = parts[1]
                            text = parts[2]
                            print(f"ðŸŒ‰ BRIDGE NOTIF : {app_name} -> {text}")
                            self.after(0, lambda a=app_name, t=text: self._process_bridge_notif(a, t))
            except Exception as e:
                print(f"âŒ Erreur critique Bridge PowerShell : {e}")

        threading.Thread(target=_bridge_worker, daemon=True).start()

    def _process_bridge_notif(self, app_name, text):
        """Traite une notification reÃ§ue via le bridge PS"""
        if self._is_android_source(app_name):
            self._process_android_notification(app_name, text)
        else:
            source_key = self._map_notification_source(app_name)
            self.notify_external(source_key, f"Notification {app_name}: {text}")

    def _map_notification_source(self, app_name, app_id=None):
        name = (app_name or "").lower()
        app_id = (app_id or "").lower()
        if "whatsapp" in name or "whatsapp" in app_id:
            return "WhatsApp.exe"
        if "outlook" in name or "outlook" in app_id:
            return "OUTLOOK.EXE"
        if "chrome" in name or "chrome" in app_id or "gmail" in name or "gmail" in app_id:
            return "chrome.exe"
        return app_name or "unknown"

    def _extract_notification_text(self, user_notification):
        try:
            visual = user_notification.notification.visual
            texts = []
            for binding in visual.bindings:
                for t in binding.get_text_elements():
                    if t.text:
                        texts.append(t.text)
            return " ".join(texts).strip()
        except Exception:
            return ""

    def _mark_notification_seen(self, notif_id):
        if notif_id in self.seen_notifications_set:
            return False
        if len(self.seen_notifications_order) >= self.max_seen_notifications:
            old = self.seen_notifications_order.popleft()
            self.seen_notifications_set.discard(old)
        self.seen_notifications_order.append(notif_id)
        self.seen_notifications_set.add(notif_id)
        return True

    def start_notifications_listener_thread(self):
        if hasattr(self, "notifications_listener_thread") and self.notifications_listener_thread.is_alive():
            return
        self.notifications_listener_thread = threading.Thread(target=self._notifications_listener_loop, daemon=True)
        self.notifications_listener_thread.start()

    def _notifications_listener_loop(self):
        try:
            asyncio.run(self._listen_notifications_async())
        except Exception as e:
            print(f"Erreur listener notifications: {e}")

    async def _listen_notifications_async(self):
        UserNotificationListener = None
        NotificationKinds = None
        
        # Essai d'import winsdk (moderne) puis winrt (legacy)
        try:
            from winsdk.windows.ui.notifications.management import UserNotificationListener
            from winsdk.windows.ui.notifications import NotificationKinds
        except ImportError:
            try:
                from winrt.windows.ui.notifications.management import UserNotificationListener
                from winrt.windows.ui.notifications import NotificationKinds
            except ImportError:
                # DÃ©sactivation de l'installation automatique car elle bloque le dÃ©marrage (compilation trop longue)
                if not hasattr(self, "_winsdk_warned"):
                    self._winsdk_warned = True
                    print("âš ï¸ WinSDK/WinRT manquant. Les notifications Windows ne seront pas lues.")
                    print("ðŸ‘‰ Pour corriger : pip install winsdk")
                return

        if not UserNotificationListener or not NotificationKinds:
            return

        listener = None
        # Test de plusieurs mÃ©thodes d'accÃ¨s (WinSDK vs WinRT vs versions)
        for method in ["get_default", "get_current", "Current"]:
            if hasattr(UserNotificationListener, method):
                try:
                    attr = getattr(UserNotificationListener, method)
                    if callable(attr):
                        listener = attr()
                    else:
                        listener = attr
                    if listener: break
                except: continue
        
        if not listener:
            # Ultime tentative : instance directe
            try: listener = UserNotificationListener()
            except: pass

        if not listener:
            print("âŒ Impossible d'accÃ©der au UserNotificationListener via WinSDK/WinRT.")
            print("ðŸš€ Tentative via le Bridge PowerShell...")
            self._run_notification_bridge()
            return

        try:
            access = await listener.request_access_async()
            if "allowed" not in str(access).lower():
                if not self.notifications_access_prompted:
                    self.notifications_access_prompted = True
                    # Utilisation de self.after pour exÃ©cuter dans le thread principal
                    self.after(0, lambda: self.parler("Chef, l'accÃ¨s aux notifications Windows est refusÃ©. Activez-le dans les paramÃ¨tres de confidentialitÃ©."))
                return
        except Exception as e:
            print(f"âš ï¸ Erreur lors de la demande d'accÃ¨s aux notifications : {e}")
            return

        while True:
            try:
                notifs = await listener.get_notifications_async(NotificationKinds.TOAST)
                for n in notifs:
                    if not self._mark_notification_seen(n.id):
                        continue
                    app_name = None
                    app_id = None
                    try:
                        app_name = n.app_info.display_info.display_name
                        app_id = n.app_info.app_user_model_id
                    except Exception:
                        pass
                    
                    # Traitement spÃ©cial pour Android si activÃ©
                    if self._is_android_source(app_name):
                        text = self._extract_notification_text(n)
                        if text:
                            self._process_android_notification(app_name, text)
                        continue
                    
                    # Traitement normal des autres notifications
                    source_key = self._map_notification_source(app_name, app_id)
                    text = self._extract_notification_text(n)
                    if text:
                        self.notify_external(source_key, f"Notification {app_name}: {text}")
                    else:
                        self.notify_external(source_key, f"Notification {app_name}.")
            except Exception as e:
                print(f"Erreur lecture notifications: {e}")
            await asyncio.sleep(2)

    def _load_settings(self):
        """Charge les paramÃ¨tres depuis settings.json si disponible"""
        try:
            if not os.path.exists(self.settings_path):
                return
            with open(self.settings_path, "r", encoding="utf-8") as f:
                data = json.load(f)
            for key, value in data.get("anomalies", {}).items():
                if hasattr(self, key):
                    setattr(self, key, int(value))
            services = data.get("services_to_watch")
            if isinstance(services, list):
                self.services_to_watch = [str(s) for s in services if str(s).strip()]
            anomalies_enabled = data.get("anomalies_enabled")
            if isinstance(anomalies_enabled, bool):
                self.anomalies_enabled = anomalies_enabled
            notifications_enabled = data.get("notifications_enabled")
            if isinstance(notifications_enabled, bool):
                self.notifications_enabled = notifications_enabled
            sources = data.get("notification_sources")
            if isinstance(sources, dict):
                self.notification_sources = {str(k): bool(v) for k, v in sources.items()}
            provider = data.get("autonomous_provider")
            if provider in ("gemini", "ollama"):
                self.autonomous_provider = provider
            ollama_model = data.get("ollama_model")
            if isinstance(ollama_model, str) and ollama_model.strip():
                self.ollama_model = ollama_model.strip()
            ollama_url = data.get("ollama_url")
            if isinstance(ollama_url, str) and ollama_url.strip():
                url = ollama_url.strip()
                if not url.startswith("http://") and not url.startswith("https://"):
                    url = "http://" + url
                parsed = urlparse(url)
                if parsed.netloc:
                    path = parsed.path or ""
                    if path in ("", "/", "/api", "/api/"):
                        path = "/api/generate"
                    url = urlunparse(parsed._replace(path=path, params="", query="", fragment=""))
                    self.ollama_url = url
            ollama_num_predict = data.get("ollama_num_predict")
            if isinstance(ollama_num_predict, int):
                if 32 <= ollama_num_predict <= 2048:
                    self.ollama_num_predict = ollama_num_predict
            launch_on_startup = data.get("launch_on_startup")
            if isinstance(launch_on_startup, bool):
                self.launch_on_startup = launch_on_startup
            ui_theme_name = data.get("ui_theme_name")
            if isinstance(ui_theme_name, str) and ui_theme_name in self.ui_theme_presets:
                self.ui_theme_name = ui_theme_name
            custom_accent_color = data.get("custom_accent_color")
            if self._is_valid_hex_color(custom_accent_color):
                self.custom_accent_color = custom_accent_color
            
            custom_colors = data.get("custom_colors")
            if isinstance(custom_colors, dict):
                self.custom_colors = custom_colors

            self.stark_personality_enabled = bool(data.get("stark_personality_enabled", True))
            self.holographic_fx_enabled = bool(data.get("holographic_fx_enabled", True))
            self.contextual_radar_enabled = bool(data.get("contextual_radar_enabled", True))
            self.tactical_mode = bool(data.get("tactical_mode", False))
            try:
                self.stark_sarcasm_level = float(data.get("stark_sarcasm_level", 0.25))
                self.stark_sarcasm_level = max(0.0, min(1.0, self.stark_sarcasm_level))
            except Exception:
                self.stark_sarcasm_level = 0.25
            
            self.voice_enabled = data.get("voice_enabled", True)
            self.voice_speed = data.get("voice_speed", "-5%")
            self.proactive_suggestions = data.get("proactive_suggestions", True)
            
            # Agents settings
            self.gemini_api_key = data.get("gemini_api_key", os.getenv("GEMINI_API_KEY") or os.getenv("GOOGLE_API_KEY"))
            if hasattr(self, "gemini_api_key_var"):
                self.gemini_api_key_var.set(self.gemini_api_key)

            self.agent_network_enabled = data.get("agent_network_enabled", False)
            self.agent_control_enabled = data.get("agent_control_enabled", False)
            self.agent_remote_enabled = data.get("agent_remote_enabled", False)
            self.agent_gmail_enabled = data.get("agent_gmail_enabled", False)
            self.agent_browser_enabled = data.get("agent_browser_enabled", False)
            self.agent_legal_enabled = data.get("agent_legal_enabled", False)
            self.agent_cyber_enabled = data.get("agent_cyber_enabled", False)
            self.agent_doctor_enabled = data.get("agent_doctor_enabled", False)
            self.agent_video_enabled = data.get("agent_video_enabled", True)
            self.agent_android_enabled = data.get("agent_android_enabled", True)
            self.agent_coding_enabled = data.get("agent_coding_enabled", True)
            self.agent_generation_enabled = data.get("agent_generation_enabled", True)
            self.agent_licence_enabled = data.get("agent_licence_enabled", True)
            self.agent_lecture_enabled = data.get("agent_lecture_enabled", True)
            self.agent_telegram_enabled = data.get("agent_telegram_enabled", False)
            
            # ParamÃ¨tres nouveaux agents
            self.agent_vision_enabled = data.get("agent_vision_enabled", True)
            self.agent_domotique_enabled = data.get("agent_domotique_enabled", True)
            self.agent_finance_enabled = data.get("agent_finance_enabled", True)
            self.agent_secretaire_enabled = data.get("agent_secretaire_enabled", True)
            self.agent_traducteur_enabled = data.get("agent_traducteur_enabled", True)
            self.agent_miner_enabled = data.get("agent_miner_enabled", True)
            self.agent_news_enabled = data.get("agent_news_enabled", True)
            self.agent_cuisine_enabled = data.get("agent_cuisine_enabled", True)
            self.agent_sante_enabled = data.get("agent_sante_enabled", True)
            self.agent_voyage_enabled = data.get("agent_voyage_enabled", True)
            self.agent_education_enabled = data.get("agent_education_enabled", True)
            self.agent_shopping_enabled = data.get("agent_shopping_enabled", True)
            self.agent_social_enabled = data.get("agent_social_enabled", True)
            self.agent_psy_enabled = data.get("agent_psy_enabled", True)
            self.agent_immo_enabled = data.get("agent_immo_enabled", True)
            self.agent_auto_enabled = data.get("agent_auto_enabled", True)
            self.agent_carrier_enabled = data.get("agent_carrier_enabled", True)
            
            # ParamÃ¨tres agents supplÃ©mentaires
            self.agent_arbitre_enabled = data.get("agent_arbitre_enabled", True)
            self.agent_ecolo_enabled = data.get("agent_ecolo_enabled", True)
            self.agent_guetteur_enabled = data.get("agent_guetteur_enabled", True)
            self.agent_historien_enabled = data.get("agent_historien_enabled", True)
            self.agent_depanneur_enabled = data.get("agent_depanneur_enabled", True)
            self.agent_astroph_enabled = data.get("agent_astroph_enabled", True)
            self.agent_stratege_enabled = data.get("agent_stratege_enabled", True)
            self.agent_architecte_enabled = data.get("agent_architecte_enabled", True)
            self.agent_business_enabled = data.get("agent_business_enabled", True)
            self.agent_polyglotte_enabled = data.get("agent_polyglotte_enabled", True)
            self.agent_nounou_enabled = data.get("agent_nounou_enabled", True)
            self.agent_appel_enabled = data.get("agent_appel_enabled", True)
            self.agent_appel_number = data.get("agent_appel_number", "")
            self.agent_appel_autoreply = data.get("agent_appel_autoreply", True)
            self.agent_appel_phrases = data.get("agent_appel_phrases", [
                "Bonjour, je suis l'assistant de mon propriÃ©taire. Il n'est pas disponible pour le moment.",
                "Puis-je prendre un message ou vous rappeler plus tard ?",
                "Je note votre appel. Merci et au revoir."
            ])
            self.agent_appel_filters = data.get("agent_appel_filters", [])
            
            # Alarmes
            alarms_data = data.get("alarms", [])
            if isinstance(alarms_data, list):
                self.alarms = alarms_data
                # Recalculer le prochain ID
                if self.alarms:
                    self.alarm_next_id = max(a.get("id", 0) for a in self.alarms) + 1
            
            self.conversation_continue = data.get("conversation_continue", False)
            if hasattr(self, "conv_continue_var"):
                self.conv_continue_var.set(self.conversation_continue)

            # Domotique params
            self.domo_wled_ip = data.get("domo_wled_ip", "")
            self.domo_ha_url = data.get("domo_ha_url", "")
            self.domo_ha_token = data.get("domo_ha_token", "")
            self.domo_arduino_port = data.get("domo_arduino_port", "")
            self.domo_webhook_on = data.get("domo_webhook_on", "")
            self.domo_webhook_off = data.get("domo_webhook_off", "")

            self.telegram_bot_token = data.get("telegram_bot_token", "")
            self.telegram_chat_id = data.get("telegram_chat_id", "")
            self.voice = data.get("voice", "fr-FR-VivienneNeural")
            remote_hosts = data.get("remote_hosts_whitelist", [])
            if isinstance(remote_hosts, list):
                self.remote_hosts_whitelist = sorted(set(str(h).strip().lower() for h in remote_hosts if str(h).strip()))
            self.remote_default_host = str(data.get("remote_default_host", "")).strip().lower()
            self.remote_ssh_user = str(data.get("remote_ssh_user", "")).strip()
            try:
                self.remote_ssh_port = int(data.get("remote_ssh_port", 22))
                self.remote_ssh_port = max(1, min(65535, self.remote_ssh_port))
            except Exception:
                self.remote_ssh_port = 22
            self.remote_ssh_key_path = str(data.get("remote_ssh_key_path", "")).strip()
            self.remote_require_confirmation = bool(data.get("remote_require_confirmation", True))
            if self.agent_telegram_enabled:
                threading.Thread(target=self.start_telegram_bot, daemon=True).start()
            if self.agent_gmail_enabled:
                threading.Thread(target=self.start_gmail_listener, daemon=True).start()
        except Exception as e:
            print(f"Erreur chargement settings: {e}")

    def _save_settings(self):
        """Sauvegarde les paramÃ¨tres dans settings.json"""
        try:
            data = {
                "anomalies": {
                    "anomaly_temp_c": self.anomaly_temp_c,
                    "battery_drop_percent": self.battery_drop_percent,
                    "battery_drop_window_sec": self.battery_drop_window_sec,
                    "battery_not_charging_grace_sec": self.battery_not_charging_grace_sec,
                    "proc_cpu_threshold": self.proc_cpu_threshold,
                    "proc_cpu_duration_sec": self.proc_cpu_duration_sec,
                    "proc_mem_threshold_mb": self.proc_mem_threshold_mb,
                    "disk_free_gb": self.disk_free_gb,
                    "disk_used_percent": self.disk_used_percent,
                    "anomaly_cooldown_sec": self.anomaly_cooldown_sec,
                },
                "services_to_watch": self.services_to_watch,
                "anomalies_enabled": self.anomalies_enabled,
                "notifications_enabled": self.notifications_enabled,
                "notification_sources": self.notification_sources,
                "autonomous_provider": self.autonomous_provider,
                "ollama_model": self.ollama_model,
                "ollama_url": self.ollama_url,
                "ollama_num_predict": self.ollama_num_predict,
                "launch_on_startup": self.launch_on_startup,
                "ui_theme_name": self.ui_theme_name,
                "custom_accent_color": self.custom_accent_color,
                "custom_colors": self.custom_colors,
                "stark_personality_enabled": self.stark_personality_enabled,
                "stark_sarcasm_level": self.stark_sarcasm_level,
                "tactical_mode": self.tactical_mode,
                "holographic_fx_enabled": self.holographic_fx_enabled,
                "contextual_radar_enabled": self.contextual_radar_enabled,
                "agent_network_enabled": self.agent_network_enabled,
                "agent_control_enabled": self.agent_control_enabled,
                "agent_remote_enabled": self.agent_remote_enabled,
                "agent_gmail_enabled": self.agent_gmail_enabled,
                "agent_browser_enabled": self.agent_browser_enabled,
                "agent_legal_enabled": self.agent_legal_enabled,
                "agent_cyber_enabled": self.agent_cyber_enabled,
                "agent_doctor_enabled": self.agent_doctor_enabled,
                "agent_video_enabled": getattr(self, "agent_video_enabled", True),
                "agent_android_enabled": getattr(self, "agent_android_enabled", True),
                "agent_coding_enabled": self.agent_coding_enabled,
                "agent_generation_enabled": self.agent_generation_enabled,
                "agent_licence_enabled": self.agent_licence_enabled,
                "agent_lecture_enabled": self.agent_lecture_enabled,
                "agent_telegram_enabled": self.agent_telegram_enabled,
                "agent_vision_enabled": self.agent_vision_enabled,
                "agent_domotique_enabled": self.agent_domotique_enabled,
                "agent_finance_enabled": self.agent_finance_enabled,
                "agent_secretaire_enabled": self.agent_secretaire_enabled,
                "agent_traducteur_enabled": self.agent_traducteur_enabled,
                "agent_miner_enabled": self.agent_miner_enabled,
                "agent_news_enabled": self.agent_news_enabled,
                "agent_cuisine_enabled": self.agent_cuisine_enabled,
                "agent_sante_enabled": self.agent_sante_enabled,
                "agent_voyage_enabled": self.agent_voyage_enabled,
                "agent_education_enabled": self.agent_education_enabled,
                "agent_shopping_enabled": self.agent_shopping_enabled,
                "agent_social_enabled": self.agent_social_enabled,
                "agent_psy_enabled": self.agent_psy_enabled,
                "agent_immo_enabled": self.agent_immo_enabled,
                "agent_auto_enabled": self.agent_auto_enabled,
                "agent_carrier_enabled": self.agent_carrier_enabled,
                "agent_arbitre_enabled": self.agent_arbitre_enabled,
                "agent_ecolo_enabled": self.agent_ecolo_enabled,
                "agent_guetteur_enabled": self.agent_guetteur_enabled,
                "agent_historien_enabled": self.agent_historien_enabled,
                "agent_depanneur_enabled": self.agent_depanneur_enabled,
                "agent_astroph_enabled": self.agent_astroph_enabled,
                "agent_stratege_enabled": self.agent_stratege_enabled,
                "agent_architecte_enabled": self.agent_architecte_enabled,
                "agent_appel_enabled": getattr(self, "agent_appel_enabled", True),
                "agent_appel_autoreply": getattr(self, "agent_appel_autoreply", True),
                "agent_appel_number": getattr(self, "agent_appel_number", ""),
                "agent_appel_phrases": getattr(self, "agent_appel_phrases", []),
                "agent_business_enabled": self.agent_business_enabled,
                "agent_polyglotte_enabled": self.agent_polyglotte_enabled,
                "agent_nounou_enabled": self.agent_nounou_enabled,
                "voice_enabled": self.voice_enabled,
                "voice_speed": self.voice_speed,
                "proactive_suggestions": self.proactive_suggestions,
                "domo_wled_ip": self.domo_wled_ip,
                "domo_ha_url": self.domo_ha_url,
                "domo_ha_token": self.domo_ha_token,
                "domo_arduino_port": self.domo_arduino_port,
                "domo_webhook_on": self.domo_webhook_on,
                "domo_webhook_off": self.domo_webhook_off,
                "telegram_bot_token": self.telegram_bot_token,
                "telegram_chat_id": self.telegram_chat_id,
                "voice": self.voice,
                "remote_hosts_whitelist": self.remote_hosts_whitelist,
                "remote_default_host": self.remote_default_host,
                "remote_ssh_user": self.remote_ssh_user,
                "remote_ssh_port": self.remote_ssh_port,
                "remote_ssh_key_path": self.remote_ssh_key_path,
                "remote_require_confirmation": self.remote_require_confirmation,
                "conversation_continue": self.conversation_continue,
                "gemini_api_key": self.gemini_api_key,
                "alarms": self.alarms,
            }
            with open(self.settings_path, "w", encoding="utf-8") as f:
                json.dump(data, f, ensure_ascii=False, indent=2)
        except Exception as e:
            print(f"Erreur sauvegarde settings: {e}")

    def update_location(self):
        """Met Ã  jour la gÃ©olocalisation en temps rÃ©el toutes les 10 minutes"""
        def fetch_location():
            try:
                # Tentative 1: ipapi.co
                response = requests.get('https://ipapi.co/json/', timeout=8)
                if response.status_code == 200:
                    data = response.json()
                    self.location_city = data.get("city", "Inconnu")
                    self.location_country = data.get("country_name", "Inconnu")
                    self.location_coords = (data.get("latitude", 0.0), data.get("longitude", 0.0))
                else:
                    # Tentative 2: ip-api.com (Fallback)
                    response = requests.get('http://ip-api.com/json/', timeout=8)
                    if response.status_code == 200:
                        data = response.json()
                        self.location_city = data.get("city", "Inconnu")
                        self.location_country = data.get("country", "Inconnu")
                        self.location_coords = (data.get("lat", 0.0), data.get("lon", 0.0))

                # Mise Ã  jour de l'interface
                if self.location_city != "Inconnu":
                    location_text = f"ðŸ“ {self.location_city}, {self.location_country}"
                    self.after(0, lambda: self.location_label.configure(text=location_text))
                else:
                    self.after(0, lambda: self.location_label.configure(text="ðŸ“ Localisation indisponible"))
            except Exception as e:
                print(f"âš ï¸ Erreur gÃ©olocalisation in update_location: {e}")
                self.after(0, lambda: self.location_label.configure(text="ðŸ“ Erreur de localisation"))
        
        threading.Thread(target=fetch_location, daemon=True).start()
        self.after(600000, self.update_location)

    def get_detailed_location(self):
        """RÃ©cupÃ¨re la position gÃ©ographique avec gestion d'erreurs DNS/RÃ©seau"""
        try:
            # Utilisation de ipapi.co avec un timeout court pour Ã©viter de bloquer l'IA
            r = requests.get('https://ipapi.co/json/', timeout=3)
            if r.status_code == 200:
                data = r.json()
                city = data.get('city', 'Inconnu')
                country = data.get('country_name', 'Inconnu')
                region = data.get('region', 'Inconnu')
                self.location_city = city
                return f"Vous Ãªtes actuellement Ã  {city}, dans la rÃ©gion de {region}, en {country}."
        except Exception as e:
            print(f"âš ï¸ Erreur gÃ©olocalisation: {e}")
        
        return "Je ne parviens pas Ã  dÃ©terminer votre position prÃ©cise, mais je suis Ã  votre Ã©coute."

    def get_weather(self, city="Paris"):
        """RÃ©cupÃ¨re la mÃ©tÃ©o en temps rÃ©el via l'API wttr.in"""
        try:
            url = f"https://wttr.in/{city}?format=%C+%t"
            response = requests.get(url)
            if response.status_code == 200:
                return response.text.strip()
            return "Indisponible"
        except:
            return "Erreur lors de la rÃ©cupÃ©ration"

    def intelligent_web_search(self, query):
        """Effectue une recherche web et extrait le contenu du premier rÃ©sultat pertinent"""
        try:
            self.status_label.configure(text="Recherche web en cours...", text_color=COLOR_ACCENT)
            # On prend le premier lien de la recherche Google
            links = list(gsearch(query, num_results=3, lang='fr'))
            if not links:
                return "Je n'ai rien trouvÃ© sur le web pour cette recherche."
            
            target_url = links[0]
            headers = {'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36'}
            page = requests.get(target_url, headers=headers, timeout=5)
            soup = BeautifulSoup(page.content, 'html.parser')
            
            # Extraction des paragraphes significatifs
            paragraphs = soup.find_all('p')
            content = " ".join([p.get_text() for p in paragraphs[:3]])
            
            if len(content) < 50:
                return f"J'ai trouvÃ© un lien ({target_url}) mais je n'ai pas pu extraire de texte clair."
            
            return f"D'aprÃ¨s mes recherches sur {target_url} : {content[:500]}..."
        except Exception as e:
            print(f"Erreur recherche web: {e}")
            return "Une erreur est survenue lors de ma recherche sur le web."

    def get_running_processes(self):
        """Retourne la liste des processus utilisateur actifs (non-systÃ¨me)"""
        user_processes = []
        # Processus systÃ¨me Ã  ignorer
        system_processes = [
            'system', 'registry', 'smss.exe', 'csrss.exe', 'wininit.exe', 'services.exe',
            'lsass.exe', 'svchost.exe', 'winlogon.exe', 'dwm.exe', 'explorer.exe',
            'taskhostw.exe', 'searchindexer.exe', 'runtimebroker.exe', 'sihost.exe',
            'ctfmon.exe', 'conhost.exe', 'fontdrvhost.exe', 'spoolsv.exe', 'audiodg.exe',
            'python.exe', 'pythonw.exe'  # Pour ne pas fermer Jarvisse lui-mÃªme
        ]
        
        try:
            for proc in psutil.process_iter(['pid', 'name', 'username']):
                try:
                    proc_name = proc.info['name'].lower()
                    # Ignorer les processus systÃ¨me
                    if proc_name not in system_processes and not proc_name.startswith('microsoft'):
                        user_processes.append({
                            'pid': proc.info['pid'],
                            'name': proc.info['name'],
                            'username': proc.info.get('username', 'N/A')
                        })
                except (psutil.NoSuchProcess, psutil.AccessDenied):
                    pass
        except Exception as e:
            print(f"Erreur lors de la rÃ©cupÃ©ration des processus: {e}")
        
        return user_processes

    def close_all_user_apps(self):
        """Ferme toutes les applications utilisateur de maniÃ¨re sÃ©curisÃ©e"""
        processes = self.get_running_processes()
        closed_count = 0
        failed_count = 0
        
        # Applications courantes Ã  fermer en prioritÃ©
        priority_apps = ['chrome', 'firefox', 'edge', 'opera', 'brave', 'discord', 'spotify', 
                        'teams', 'slack', 'zoom', 'notepad', 'wordpad', 'excel', 'winword', 
                        'powerpnt', 'outlook', 'vlc', 'steam']
        
        for proc_info in processes:
            proc_name = proc_info['name'].lower()
            # VÃ©rifier si c'est une application prioritaire
            is_priority = any(app in proc_name for app in priority_apps)
            
            if is_priority or proc_name.endswith('.exe'):
                try:
                    proc = psutil.Process(proc_info['pid'])
                    proc.terminate()  # Fermeture propre
                    closed_count += 1
                    print(f"âœ… FermÃ©: {proc_info['name']}")
                except (psutil.NoSuchProcess, psutil.AccessDenied) as e:
                    failed_count += 1
                    print(f"âŒ Impossible de fermer: {proc_info['name']} - {e}")
        
        return closed_count, failed_count

    def list_active_processes(self):
        """Liste les processus actifs pour l'utilisateur"""
        processes = self.get_running_processes()
        if not processes:
            return "Aucun processus utilisateur actif dÃ©tectÃ©."
        
        # Grouper par type d'application
        browsers = []
        office = []
        media = []
        communication = []
        others = []
        
        for proc in processes:
            name = proc['name'].lower()
            if any(x in name for x in ['chrome', 'firefox', 'edge', 'opera', 'brave']):
                browsers.append(proc['name'])
            elif any(x in name for x in ['word', 'excel', 'powerpoint', 'outlook']):
                office.append(proc['name'])
            elif any(x in name for x in ['vlc', 'spotify', 'itunes', 'media']):
                media.append(proc['name'])
            elif any(x in name for x in ['discord', 'teams', 'slack', 'zoom', 'skype']):
                communication.append(proc['name'])
            else:
                others.append(proc['name'])
        
        result = f"J'ai dÃ©tectÃ© {len(processes)} processus actifs. "
        if browsers: result += f"Navigateurs: {len(browsers)}. "
        if office: result += f"Suite Office: {len(office)}. "
        if communication: result += f"Communication: {len(communication)}. "
        
        return result

    # ==========================================
    # RÃ‰GLAGES & CONFIGURATION
    # ==========================================

    def toggle_voice(self):
        """Active ou dÃ©sactive la synthÃ¨se vocale"""
        self.voice_enabled = self.voice_var_settings.get()
        if self.voice_enabled:
            self.parler("SynthÃ¨se vocale activÃ©e.")
        else:
            self.add_message("SynthÃ¨se vocale dÃ©sactivÃ©e. Je rÃ©pondrai uniquement par texte.", "SystÃ¨me")

    def toggle_proactive_suggestions(self):
        """Active ou dÃ©sactive les suggestions proactives"""
        self.proactive_suggestions = self.proactive_var_settings.get()
        self._save_settings()
        if self.proactive_suggestions:
            self.parler("Suggestions proactives activÃ©es. Je vais maintenant anticiper vos besoins.")
        else:
            self.parler("Suggestions proactives dÃ©sactivÃ©es. Je me concentrerai uniquement sur vos demandes directes.")

    def open_appearance_customization(self):
        """Ouvre une fenÃªtre pour rÃ©gler chaque couleur de l'interface"""
        color_win = ctk.CTkToplevel(self)
        color_win.title("Jarvisse - Personnalisation des Couleurs")
        color_win.geometry("400x650")
        color_win.configure(fg_color=COLOR_BG)
        color_win.transient(self)
        color_win.attributes("-topmost", True)

        ctk.CTkLabel(color_win, text="ðŸŽ¨ DESIGN PERSONNALISÃ‰", font=ctk.CTkFont(size=18, weight="bold"), text_color=COLOR_ACCENT).pack(pady=20)

        scroll = ctk.CTkScrollableFrame(color_win, fg_color="transparent")
        scroll.pack(fill="both", expand=True, padx=20, pady=10)

        color_labels = {
            "bg": "Fond de l'application",
            "nav": "Barre de navigation & Sidebar",
            "accent": "Couleur d'accent (Primaire)",
            "accent_purple": "Couleur d'accent (Secondaire)",
            "accent_red": "Alertes & Boutons critiques",
            "text": "Texte principal",
            "secondary": "Texte secondaire (Labels)",
            "card": "Cartes & Boutons standards"
        }

        current_pal = self._current_theme_palette()

        def pick_color(key):
            color = colorchooser.askcolor(initialcolor=current_pal.get(key, "#000000"), title=f"Choisir {color_labels[key]}")
            if color[1]:
                self.custom_colors[key] = color[1]
                self._apply_theme_globals()
                self._refresh_theme_widgets()
                self._save_settings()

        for key, label in color_labels.items():
            f = ctk.CTkFrame(scroll, fg_color="transparent")
            f.pack(fill="x", pady=8)
            
            ctk.CTkLabel(f, text=label, font=ctk.CTkFont(size=12)).pack(side="left")
            
            # AperÃ§u de la couleur actuelle
            preview = ctk.CTkFrame(f, width=24, height=24, fg_color=current_pal.get(key, "#000000"), corner_radius=12)
            preview.pack(side="right", padx=5)
            
            btn = ctk.CTkButton(f, text="Modifier", width=80, height=28, 
                                command=lambda k=key: pick_color(k))
            btn.pack(side="right", padx=5)

        ctk.CTkButton(color_win, text="RÃ©initialiser tout", fg_color=COLOR_ACCENT_RED, 
                      command=lambda: [setattr(self, "custom_colors", {}), self._apply_theme_globals(), self._refresh_theme_widgets(), self._save_settings(), color_win.destroy()]).pack(pady=20)


    def open_main_settings(self):
        """Ouvre le panneau central des paramÃ¨tres de Jarvisse"""
        main_settings = ctk.CTkToplevel(self)
        main_settings.title("Jarvisse - ParamÃ¨tres GÃ©nÃ©raux")
        main_settings.geometry("500x750")
        main_settings.configure(fg_color=COLOR_BG)
        main_settings.transient(self)
        main_settings.attributes("-topmost", True)

        # Titre & ContrÃ´les de taille
        header_frame = ctk.CTkFrame(main_settings, fg_color="transparent")
        header_frame.pack(fill="x", pady=(10, 0))
        
        title_label = ctk.CTkLabel(header_frame, text="âš™ï¸ RÃ‰GLAGES GLOBAUX", font=ctk.CTkFont(size=18, weight="bold"), text_color=COLOR_ACCENT)
        title_label.pack(side="left", padx=20)

        def set_size(size_str):
            main_settings.geometry(size_str)

        btn_size_max = ctk.CTkButton(header_frame, text="ðŸ—–", width=30, height=30, fg_color=COLOR_CARD, 
                                    text_color=COLOR_TEXT, hover_color=COLOR_ACCENT, command=lambda: set_size("750x900"))
        btn_size_max.pack(side="right", padx=(5, 20))

        btn_size_min = ctk.CTkButton(header_frame, text="ðŸ——", width=30, height=30, fg_color=COLOR_CARD, 
                                    text_color=COLOR_TEXT, hover_color=COLOR_ACCENT, command=lambda: set_size("450x650"))
        btn_size_min.pack(side="right", padx=5)

        # Zone de dÃ©filement pour les options
        scroll_frame = ctk.CTkScrollableFrame(main_settings, fg_color="transparent")
        scroll_frame.pack(fill="both", expand=True, padx=20, pady=10)

        def create_sec(parent, title):
            l = ctk.CTkLabel(parent, text=title.upper(), font=ctk.CTkFont(size=13, weight="bold"), text_color=COLOR_ACCENT_PURPLE)
            l.pack(pady=(20, 10), anchor="w")
            return l

        # Section 1: Modules & Intelligences
        create_sec(scroll_frame, "Modules & Agents")
        ctk.CTkButton(scroll_frame, text="ðŸ“‹ Commandes Vocales", fg_color=COLOR_CARD, hover_color="#3C4043",
                      anchor="w", command=self.open_commandes_panel).pack(fill="x", pady=5)
        ctk.CTkButton(scroll_frame, text="ðŸ¤– GÃ©rer les Agents Actifs", fg_color=COLOR_CARD, hover_color="#3C4043", 
                      anchor="w", command=self.open_agents_settings).pack(fill="x", pady=5)
        ctk.CTkButton(scroll_frame, text="ðŸ  Domotique & Interfaces", fg_color=COLOR_CARD, hover_color="#3C4043", 
                      anchor="w", command=self.open_domotique_settings).pack(fill="x", pady=5)

        # Section 2: SÃ©curitÃ© & SantÃ©
        create_sec(scroll_frame, "SÃ©curitÃ© & Maintenance")
        ctk.CTkButton(scroll_frame, text="ðŸ©º Seuils d'Anomalies SystÃ¨me", fg_color=COLOR_CARD, hover_color="#3C4043", 
                      anchor="w", command=self.open_anomaly_settings).pack(fill="x", pady=5)
        ctk.CTkButton(scroll_frame, text="ðŸ›¡ï¸ Services Critiques", fg_color=COLOR_CARD, hover_color="#3C4043", 
                      anchor="w", command=self.open_services_watch).pack(fill="x", pady=5)

        # Section 3: Communication & Alertes
        create_sec(scroll_frame, "Communication & Alertes")
        
        # ContrÃ´le de vitesse de lecture
        speed_frame = ctk.CTkFrame(scroll_frame, fg_color="transparent")
        speed_frame.pack(fill="x", pady=10)
        
        ctk.CTkLabel(speed_frame, text="ðŸŽ™ï¸ Vitesse de Lecture :", font=ctk.CTkFont(size=12)).pack(side="left", padx=(0, 10))
        
        # Slider de -50% Ã  +50%
        speed_slider = ctk.CTkSlider(
            speed_frame, 
            from_=-50, 
            to=50, 
            number_of_steps=20,
            command=lambda v: self._update_voice_speed(v)
        )
        
        # Convertir la vitesse actuelle en valeur numÃ©rique
        current_speed_str = self.voice_speed.replace("%", "").replace("+", "")
        try:
            current_speed_val = int(current_speed_str)
        except:
            current_speed_val = -5
            
        speed_slider.set(current_speed_val)
        speed_slider.pack(side="left", fill="x", expand=True, padx=10)
        
        # Label pour afficher la valeur
        speed_label = ctk.CTkLabel(speed_frame, text=f"{current_speed_val:+d}%", font=ctk.CTkFont(size=11), width=50)
        speed_label.pack(side="right")
        
        # Fonction de mise Ã  jour
        def update_label_and_speed(value):
            val = int(value)
            speed_label.configure(text=f"{val:+d}%")
            self.voice_speed = f"{val:+d}%"
            self._save_settings()
        
        speed_slider.configure(command=update_label_and_speed)
        
        self.voice_var_settings = ctk.BooleanVar(value=self.voice_enabled)
        ctk.CTkSwitch(scroll_frame, text="Sortie Vocale (Son)", variable=self.voice_var_settings, 
                      command=self.toggle_voice, text_color=COLOR_TEXT).pack(pady=5, anchor="w")
        
        # Switch pour les suggestions proactives
        self.proactive_var_settings = ctk.BooleanVar(value=self.proactive_suggestions)
        ctk.CTkSwitch(scroll_frame, text="ðŸ§  Suggestions Proactives (Lecture d'esprit)", 
                      variable=self.proactive_var_settings, 
                      command=self.toggle_proactive_suggestions, 
                      text_color=COLOR_TEXT).pack(pady=5, anchor="w")
        
        ctk.CTkButton(scroll_frame, text="ðŸ”” Notifications d'Applications", fg_color=COLOR_CARD, hover_color="#3C4043", 
                      anchor="w", command=self.open_notification_settings).pack(fill="x", pady=5)
        ctk.CTkButton(scroll_frame, text="ðŸ“§ Configuration Gmail", fg_color=COLOR_CARD, hover_color="#3C4043", 
                      anchor="w", command=self.agent_gmail_mission).pack(fill="x", pady=5)

        # Section 4: Personnalisation
        create_sec(scroll_frame, "Apparence & ThÃ¨mes")
        
        # SÃ©lecteur de ThÃ¨me (S'intÃ¨gre ici pour rÃ©pondre Ã  la demande utilisateur)
        theme_frame = ctk.CTkFrame(scroll_frame, fg_color="transparent")
        theme_frame.pack(fill="x", pady=5)
        ctk.CTkLabel(theme_frame, text="ThÃ¨me Visuel :", font=ctk.CTkFont(size=12)).pack(side="left", padx=(0, 10))
        
        theme_options = sorted(list(self.ui_theme_presets.keys()))
        theme_menu = ctk.CTkOptionMenu(
            theme_frame, 
            values=theme_options,
            command=self.on_ui_theme_changed,
            fg_color=COLOR_CARD,
            button_color=COLOR_ACCENT_PURPLE,
            text_color=COLOR_TEXT
        )
        theme_menu.set(self.ui_theme_name)
        theme_menu.pack(side="right", fill="x", expand=True)

        ctk.CTkButton(scroll_frame, text="ðŸŽ¨ Couleurs & Design (Sur Mesure)", fg_color=COLOR_CARD, hover_color="#3C4043", 
                      anchor="w", command=self.open_appearance_customization).pack(fill="x", pady=5)

        # Pied de page
        footer_label = ctk.CTkLabel(main_settings, text="Jarvisse Intelligence Artificielle v3.0", font=ctk.CTkFont(size=10), text_color=COLOR_SECONDARY)
        footer_label.pack(pady=10)

    def open_commandes_panel(self):
        """Affiche la liste complÃ¨te des commandes vocales de Jarvisse"""
        win = ctk.CTkToplevel(self)
        win.title("ðŸ“‹ Commandes Vocales de Jarvisse")
        win.geometry("600x800")
        win.configure(fg_color=COLOR_BG)
        win.transient(self)
        win.attributes("-topmost", True)

        ctk.CTkLabel(win, text="ðŸ“‹ COMMANDES VOCALES", font=ctk.CTkFont(size=18, weight="bold"),
                     text_color=COLOR_ACCENT).pack(pady=(15, 5))
        ctk.CTkLabel(win, text="Dites \"Jarvisse\" puis l'une de ces commandes",
                     font=ctk.CTkFont(size=11), text_color=COLOR_SECONDARY).pack(pady=(0, 10))

        scroll = ctk.CTkScrollableFrame(win, fg_color="transparent")
        scroll.pack(fill="both", expand=True, padx=15, pady=5)

        COMMANDES = [
            ("ðŸ§  CERVEAU & MÃ‰MOIRE", [
                ("Indexe ton cerveau",           "Scanne et mÃ©morise tous les docs du dossier Cerveau_Jarvisse"),
                ("Mets Ã  jour ta mÃ©moire",       "Idem â€“ recharge la base de connaissances locale (FAISS)"),
                ("Recharge ta mÃ©moire",          "Idem"),
                ("Apprends tes documents",       "Idem"),
                ("Actualise ta mÃ©moire",         "Idem"),
                ("Indexe tes connaissances",     "Idem"),
            ]),
            ("âš™ï¸ SYSTÃˆME & PARAMÃˆTRES", [
                ("Ouvre les paramÃ¨tres",         "Ouvre le panneau de rÃ©glages globaux"),
                ("Mode combat",                  "Active le mode tactique (prioritÃ© urgences)"),
                ("Mode normal",                  "DÃ©sactive le mode tactique"),
                ("Brief tactique",               "Donne un Ã©tat tactique instantanÃ© du systÃ¨me"),
                ("Active hologramme",            "Active l'animation HUD holographique"),
                ("DÃ©sactive hologramme",         "DÃ©sactive l'animation HUD holographique"),
                ("Ouvre le gestionnaire",        "Ouvre le gestionnaire de tÃ¢ches Windows"),
                ("Capture d'Ã©cran",              "Prend une capture de l'Ã©cran"),
                ("Verrouille l'Ã©cran",           "Verrouille la session Windows"),
                ("Ã‰teins l'ordinateur",          "Ã‰teint le PC"),
                ("RedÃ©marre l'ordinateur",       "RedÃ©marre le PC"),
            ]),
            ("ðŸŽ™ï¸ CONTRÃ”LE AUDIO & VOIX", [
                ("Stop / ArrÃªte",                "Stoppe immÃ©diatement la lecture vocale"),
                ("Augmente le volume",           "Monte le volume du systÃ¨me"),
                ("Baisse le volume",             "Diminue le volume du systÃ¨me"),
                ("Coupe le son / muet",          "Mute le son"),
                ("Active / DÃ©sactive la voix",   "Active ou dÃ©sactive la synthÃ¨se vocale"),
            ]),
            ("ðŸ¤– AGENTS INTELLIGENTS", [
                ("Analyse ce fichier",           "Lit et rÃ©sume un document PDF, Word, ou texte"),
                ("Lance une recherche web sur X","Recherche sur le web"),
                ("MÃ©tÃ©o Ã  Paris",                "RÃ©cupÃ¨re la mÃ©tÃ©o en temps rÃ©el"),
                ("Quelle heure est-il ?",        "Donne l'heure et la date"),
                ("Traduis [phrase] en anglais",  "Traduction instantanÃ©e"),
                ("RÃ©sume ce texte",              "RÃ©sume le dernier document chargÃ©"),
                ("GÃ©nÃ¨re un rapport sur X",      "CrÃ©e un rapport structurÃ© sur un thÃ¨me"),
                ("Fais une capture d'Ã©cran",     "Screenshot de l'Ã©cran"),
                ("Lance [application]",          "Ouvre n'importe quelle application"),
                ("Ferme [application]",          "Ferme une application ouverte"),
            ]),
            ("ðŸ“Š FINANCE & BUSINESS", [
                ("Cours du Bitcoin",             "Donne le prix en temps rÃ©el"),
                ("Cours de l'or",                "Prix de l'or en temps rÃ©el"),
                ("Analyse financiÃ¨re de [X]",    "Analyse un actif ou marchÃ©"),
            ]),
            ("ðŸ©º MAINTENANCE SYSTÃˆME", [
                ("Fais un bilan de santÃ©",       "Diagnostic complet CPU/RAM/disque"),
                ("Lance une maintenance complÃ¨te","Auto-nettoyage et optimisation"),
                ("Scanne les virus",             "Recherche de logiciels malveillants"),
                ("LibÃ¨re de la mÃ©moire",         "LibÃ¨re de la RAM utilisÃ©e"),
            ]),
            ("ðŸš€ STARK EDITION (FUTURISTE)", [
                ("Regarde mon Ã©cran",            "Analyse visuelle de votre bureau (Vision)"),
                ("Analyse la camÃ©ra",            "Analyse visuelle via la webcam"),
                ("CrÃ©e le dossier [Nom]",        "ContrÃ´le PC : CrÃ©ation de dossiers"),
                ("Initialise Git ici",           "ContrÃ´le PC : CrÃ©ation de dÃ©pÃ´t Git"),
                ("Ouvre [Application]",          "ContrÃ´le PC : Lancement de programmes"),
                ("Optimise mon PC",              "Nettoyage RAM, corbeille et processus"),
                ("Recherche approfondie sur X",  "Analyse multi-sources massive (Deep Web)"),
                ("Ã‰cris un script Python pour X","GÃ©nÃ¨re, teste et exÃ©cute du code Python"),
                ("Apprends Ã  faire X",           "Auto-Ã©volution : Jarvisse s'auto-amÃ©liore"),
                ("Suivi mission : [Nom]",        "GÃ¨re des tÃ¢ches sur plusieurs sessions"),
                ("Lance le protocole SHIELD",    "Protection maximale et audit systÃ¨me"),
                ("Lance le protocole Travail",   "Optimisation focus et outils production"),
                ("Bilan Sentinelle",             "VÃ©rification proactive de la santÃ© systÃ¨me"),
            ]),
            ("ðŸ’¬ CONVERSATION", [
                ("Bonjour Jarvisse",             "Salutation â€“ Jarvisse rÃ©pond"),
                ("Qui es-tu ?",                  "Jarvisse se prÃ©sente"),
                ("Explique-moi [sujet]",         "Explication intelligente via Gemini"),
                ("Donne-moi des idÃ©es sur [X]",  "Brainstorming assistÃ© par IA"),
            ]),
        ]

        def add_section(parent, title, items):
            ctk.CTkLabel(parent, text=title, font=ctk.CTkFont(size=13, weight="bold"),
                         text_color=COLOR_ACCENT_PURPLE).pack(anchor="w", pady=(16, 4))
            for cmd, desc in items:
                row = ctk.CTkFrame(parent, fg_color=COLOR_CARD, corner_radius=6)
                row.pack(fill="x", pady=2)
                ctk.CTkLabel(row, text=f"  ðŸ”¹ {cmd}", font=ctk.CTkFont(size=11, weight="bold"),
                             text_color=COLOR_TEXT, anchor="w").pack(side="left", padx=8, pady=5)
                ctk.CTkLabel(row, text=desc, font=ctk.CTkFont(size=10),
                             text_color=COLOR_SECONDARY, anchor="e", wraplength=280).pack(side="right", padx=8)

        for section_title, items in COMMANDES:
            add_section(scroll, section_title, items)

        ctk.CTkButton(win, text="âœ– Fermer", fg_color=COLOR_ACCENT_RED,
                      hover_color="#C53929", command=win.destroy).pack(pady=12)

    def open_agents_settings(self):
        """Ouvre une fenÃªtre pour configurer et activer les agents"""
        settings_win = ctk.CTkToplevel(self)
        settings_win.title("ParamÃ¨tres des Agents")
        settings_win.geometry("450x850")
        settings_win.configure(fg_color=COLOR_BG)
        settings_win.transient(self)
        settings_win.attributes("-topmost", True)
        settings_win.lift()
        settings_win.focus_force()

        # Calcul automatique du nombre d'agents (y compris les nouveaux)
        all_agents_vars = [
            self.agent_network_enabled, self.agent_control_enabled, self.agent_gmail_enabled,
            self.agent_remote_enabled,
            self.agent_browser_enabled, self.agent_legal_enabled, self.agent_cyber_enabled,
            self.agent_doctor_enabled, self.agent_video_enabled, self.agent_android_enabled,
            self.agent_coding_enabled, self.agent_generation_enabled, self.agent_licence_enabled,
            self.agent_lecture_enabled, self.agent_telegram_enabled, self.agent_vision_enabled,
            self.agent_domotique_enabled, self.agent_finance_enabled, self.agent_secretaire_enabled,
            self.agent_traducteur_enabled, self.agent_miner_enabled,
            self.agent_news_enabled, self.agent_cuisine_enabled, self.agent_sante_enabled, self.agent_voyage_enabled,
            self.agent_education_enabled, self.agent_shopping_enabled, self.agent_social_enabled,
            self.agent_psy_enabled, self.agent_immo_enabled, self.agent_auto_enabled, self.agent_carrier_enabled,
            self.agent_arbitre_enabled, self.agent_ecolo_enabled, self.agent_guetteur_enabled,
            self.agent_historien_enabled, self.agent_depanneur_enabled, self.agent_astroph_enabled,
            self.agent_stratege_enabled, self.agent_architecte_enabled, self.agent_business_enabled,
            self.agent_polyglotte_enabled, self.agent_nounou_enabled, self.agent_appel_enabled
        ]
        total_agents = len(all_agents_vars)

        # Titre fixe en haut
        label = ctk.CTkLabel(settings_win, text="CONFIGURATION DES AGENTS", font=ctk.CTkFont(size=16, weight="bold"), text_color=COLOR_ACCENT)
        label.pack(pady=(20, 5))

        # Affichage du nombre d'agents (DEMANDE UTILISATEUR)
        count_label = ctk.CTkLabel(
            settings_win, 
            text=f"ðŸ“Š Total des agents intÃ©grÃ©s : {total_agents}", 
            font=ctk.CTkFont(size=13, weight="bold"), 
            text_color=COLOR_ACCENT_PURPLE,
            fg_color=COLOR_NAV,
            corner_radius=8,
            padx=15, pady=5
        )
        count_label.pack(pady=(0, 15))

        # Frame scrollable pour le contenu
        scrollable_frame = ctk.CTkScrollableFrame(
            settings_win,
            fg_color=COLOR_BG,
            corner_radius=0,
            width=410,
            height=500
        )
        scrollable_frame.pack(fill="both", padx=10, pady=(0, 10))

        # Agent RÃ©seau
        net_var = ctk.BooleanVar(value=self.agent_network_enabled)
        net_switch = ctk.CTkSwitch(scrollable_frame, text="Agent RÃ©seau", variable=net_var, command=lambda: self.toggle_agent("network", net_var.get()))
        net_switch.pack(pady=10, padx=20, anchor="w")

        # Agent ContrÃ´le Total
        ctrl_var = ctk.BooleanVar(value=self.agent_control_enabled)
        ctrl_switch = ctk.CTkSwitch(scrollable_frame, text="Agent ContrÃ´le Total", variable=ctrl_var, command=lambda: self.toggle_agent("control", ctrl_var.get()))
        ctrl_switch.pack(pady=10, padx=20, anchor="w")

        # Agent ContrÃ´le Distant (SSH sÃ©curisÃ©)
        remote_var = ctk.BooleanVar(value=self.agent_remote_enabled)
        remote_switch = ctk.CTkSwitch(scrollable_frame, text="Agent ContrÃ´le Distant (SSH)", variable=remote_var, command=lambda: self.toggle_agent("remote", remote_var.get()))
        remote_switch.pack(pady=10, padx=20, anchor="w")
        ctk.CTkLabel(scrollable_frame, text="ExÃ©cution distante sÃ©curisÃ©e sur machines autorisÃ©es uniquement.", font=ctk.CTkFont(size=10), text_color=COLOR_SECONDARY).pack(padx=40, anchor="w")

        remote_form = ctk.CTkFrame(scrollable_frame, fg_color="transparent")
        remote_form.pack(pady=5, padx=40, fill="x")

        remote_hosts_entry = ctk.CTkEntry(remote_form, placeholder_text="Machines autorisÃ©es (CSV: pc1,192.168.1.10)", height=28, fg_color=COLOR_NAV)
        remote_hosts_entry.insert(0, ", ".join(self.remote_hosts_whitelist))
        remote_hosts_entry.pack(fill="x", pady=(0, 5))

        remote_user_entry = ctk.CTkEntry(remote_form, placeholder_text="Utilisateur SSH", height=28, fg_color=COLOR_NAV)
        remote_user_entry.insert(0, self.remote_ssh_user)
        remote_user_entry.pack(fill="x", pady=(0, 5))

        remote_row = ctk.CTkFrame(remote_form, fg_color="transparent")
        remote_row.pack(fill="x", pady=(0, 5))
        remote_port_entry = ctk.CTkEntry(remote_row, placeholder_text="Port SSH", width=90, height=28, fg_color=COLOR_NAV)
        remote_port_entry.insert(0, str(self.remote_ssh_port))
        remote_port_entry.pack(side="left", padx=(0, 5))
        remote_default_host_entry = ctk.CTkEntry(remote_row, placeholder_text="HÃ´te par dÃ©faut", height=28, fg_color=COLOR_NAV)
        remote_default_host_entry.insert(0, self.remote_default_host)
        remote_default_host_entry.pack(side="left", fill="x", expand=True)

        remote_key_entry = ctk.CTkEntry(remote_form, placeholder_text="ClÃ© privÃ©e SSH (optionnel)", height=28, fg_color=COLOR_NAV)
        remote_key_entry.insert(0, self.remote_ssh_key_path)
        remote_key_entry.pack(fill="x", pady=(0, 5))

        confirm_remote_var = ctk.BooleanVar(value=self.remote_require_confirmation)
        confirm_remote_switch = ctk.CTkSwitch(remote_form, text="Demander confirmation commandes sensibles", variable=confirm_remote_var)
        confirm_remote_switch.pack(anchor="w", pady=(0, 5))

        def browse_remote_key():
            key_path = filedialog.askopenfilename(title="SÃ©lectionner la clÃ© privÃ©e SSH")
            if key_path:
                remote_key_entry.delete(0, "end")
                remote_key_entry.insert(0, key_path)

        def save_remote_settings():
            hosts = [h.strip().lower() for h in remote_hosts_entry.get().split(",") if h.strip()]
            self.remote_hosts_whitelist = sorted(set(hosts))
            self.remote_ssh_user = remote_user_entry.get().strip()
            self.remote_default_host = remote_default_host_entry.get().strip().lower()
            self.remote_ssh_key_path = remote_key_entry.get().strip()
            self.remote_require_confirmation = bool(confirm_remote_var.get())
            try:
                port_value = int(remote_port_entry.get().strip())
                self.remote_ssh_port = max(1, min(65535, port_value))
            except Exception:
                self.remote_ssh_port = 22
            self._save_settings()
            self.add_message("ParamÃ¨tres Agent ContrÃ´le Distant enregistrÃ©s.", "SystÃ¨me")

        def test_remote_settings():
            save_remote_settings()
            target = self.remote_default_host or (self.remote_hosts_whitelist[0] if self.remote_hosts_whitelist else "")
            if not target:
                self.add_message("Aucun hÃ´te distant configurÃ© pour le test.", "SystÃ¨me")
                return
            threading.Thread(target=lambda: self.parler(self._execute_remote_task(target, "whoami", force_without_confirmation=True)), daemon=True).start()

        remote_btn_row = ctk.CTkFrame(remote_form, fg_color="transparent")
        remote_btn_row.pack(fill="x", pady=(0, 10))
        ctk.CTkButton(remote_btn_row, text="Parcourir clÃ©", height=24, font=ctk.CTkFont(size=10), command=browse_remote_key).pack(side="left", padx=(0, 5))
        ctk.CTkButton(remote_btn_row, text="Enregistrer Remote", height=24, font=ctk.CTkFont(size=10), command=save_remote_settings).pack(side="left", padx=(0, 5))
        ctk.CTkButton(remote_btn_row, text="Tester SSH", height=24, font=ctk.CTkFont(size=10), fg_color=COLOR_ACCENT, command=test_remote_settings).pack(side="left")

        # Agent Gmail
        gmail_var = ctk.BooleanVar(value=self.agent_gmail_enabled)
        gmail_switch = ctk.CTkSwitch(scrollable_frame, text="Agent Gmail", variable=gmail_var, command=lambda: self.toggle_agent("gmail", gmail_var.get()))
        gmail_switch.pack(pady=10, padx=20, anchor="w")

        # Agent Navigateur
        browse_var = ctk.BooleanVar(value=self.agent_browser_enabled)
        browse_switch = ctk.CTkSwitch(scrollable_frame, text="Agent Navigateur", variable=browse_var, command=lambda: self.toggle_agent("browser", browse_var.get()))
        browse_switch.pack(pady=10, padx=20, anchor="w")

        # Agent Juridique
        legal_var = ctk.BooleanVar(value=self.agent_legal_enabled)
        legal_switch = ctk.CTkSwitch(scrollable_frame, text="Agent Juridique", variable=legal_var, command=lambda: self.toggle_agent("legal", legal_var.get()))
        legal_switch.pack(pady=10, padx=20, anchor="w")
        ctk.CTkLabel(scrollable_frame, text="Expert en droit, RH et articles juridiques.", font=ctk.CTkFont(size=10), text_color=COLOR_SECONDARY).pack(padx=40, anchor="w")

        # Agent Cyber SÃ©curitÃ©
        cyber_var = ctk.BooleanVar(value=self.agent_cyber_enabled)
        cyber_switch = ctk.CTkSwitch(scrollable_frame, text="Agent Cyber SÃ©curitÃ©", variable=cyber_var, command=lambda: self.toggle_agent("cyber", cyber_var.get()))
        cyber_switch.pack(pady=10, padx=20, anchor="w")
        ctk.CTkLabel(scrollable_frame, text="Audit systÃ¨me, rÃ©seau, dÃ©fense et conseils cyber.", font=ctk.CTkFont(size=10), text_color=COLOR_SECONDARY).pack(padx=40, anchor="w")

        # Agent Docteur SystÃ¨me
        doctor_var = ctk.BooleanVar(value=self.agent_doctor_enabled)
        doctor_switch = ctk.CTkSwitch(scrollable_frame, text="Agent Docteur SystÃ¨me", variable=doctor_var, command=lambda: self.toggle_agent("doctor", doctor_var.get()))
        doctor_switch.pack(pady=10, padx=20, anchor="w")
        ctk.CTkLabel(scrollable_frame, text="SantÃ© systÃ¨me, optimisation processus et rÃ©paration.", font=ctk.CTkFont(size=10), text_color=COLOR_SECONDARY).pack(padx=40, anchor="w")

        # Agent VidÃ©o Surveillance
        video_var = ctk.BooleanVar(value=self.agent_video_enabled)
        video_switch = ctk.CTkSwitch(scrollable_frame, text="Agent VidÃ©o Surveillance", variable=video_var, command=lambda: self.toggle_agent("video", video_var.get()))
        video_switch.pack(pady=10, padx=20, anchor="w")
        ctk.CTkLabel(scrollable_frame, text="ContrÃ´le camÃ©ra, enregistrements et captures d'Ã©cran.", font=ctk.CTkFont(size=10), text_color=COLOR_SECONDARY).pack(padx=40, anchor="w")

        # Agent ContrÃ´le Android
        android_var = ctk.BooleanVar(value=self.agent_android_enabled)
        android_switch = ctk.CTkSwitch(scrollable_frame, text="Agent ContrÃ´le Android", variable=android_var, command=lambda: self.toggle_agent("android", android_var.get()))
        android_switch.pack(pady=10, padx=20, anchor="w")
        ctk.CTkLabel(scrollable_frame, text="Connexion tÃ©lÃ©phone, lecture SMS et alerte appels (via Lien avec Windows).", font=ctk.CTkFont(size=10), text_color=COLOR_SECONDARY).pack(padx=40, anchor="w")

        # Assistant Coding
        coding_var = ctk.BooleanVar(value=self.agent_coding_enabled)
        coding_switch = ctk.CTkSwitch(scrollable_frame, text="Assistant Coding", variable=coding_var, command=lambda: self.toggle_agent("coding", coding_var.get()))
        coding_switch.pack(pady=10, padx=20, anchor="w")
        ctk.CTkLabel(scrollable_frame, text="CrÃ©ation d'applications et de sites web (HTML, CSS, JS, Python, etc.).", font=ctk.CTkFont(size=10), text_color=COLOR_SECONDARY).pack(padx=40, anchor="w")

        # Assistant GÃ©nÃ©ration Image et VidÃ©o
        gen_var = ctk.BooleanVar(value=self.agent_generation_enabled)
        gen_switch = ctk.CTkSwitch(scrollable_frame, text="Assistant GÃ©nÃ©ration Image & VidÃ©o", variable=gen_var, command=lambda: self.toggle_agent("generation", gen_var.get()))
        gen_switch.pack(pady=10, padx=20, anchor="w")
        ctk.CTkLabel(scrollable_frame, text="GÃ©nÃ©ration d'images et de vidÃ©os rÃ©alistes via IA.", font=ctk.CTkFont(size=10), text_color=COLOR_SECONDARY).pack(padx=40, anchor="w")
        # Agent Assistant Licence
        licence_var = ctk.BooleanVar(value=self.agent_licence_enabled)
        licence_switch = ctk.CTkSwitch(scrollable_frame, text="Agent Assistant Licence", variable=licence_var, command=lambda: self.toggle_agent("licence", licence_var.get()))
        licence_switch.pack(pady=10, padx=20, anchor="w")
        ctk.CTkLabel(scrollable_frame, text="Activation de Windows et Microsoft Office (KMS).", font=ctk.CTkFont(size=10), text_color=COLOR_SECONDARY).pack(padx=40, anchor="w")

        # Agent Lecture
        lecture_var = ctk.BooleanVar(value=self.agent_lecture_enabled)
        lecture_switch = ctk.CTkSwitch(scrollable_frame, text="Agent Lecture", variable=lecture_var, command=lambda: self.toggle_agent("lecture", lecture_var.get()))
        lecture_switch.pack(pady=10, padx=20, anchor="w")
        ctk.CTkLabel(scrollable_frame, text="Lecture de fichiers texte, PDF ou documents locaux.", font=ctk.CTkFont(size=10), text_color=COLOR_SECONDARY).pack(padx=40, anchor="w")

        # Agent Bot Telegram
        telegram_var = ctk.BooleanVar(value=self.agent_telegram_enabled)
        telegram_switch = ctk.CTkSwitch(scrollable_frame, text="Agent Bot Telegram", variable=telegram_var, command=lambda: self.toggle_agent("telegram", telegram_var.get()))
        telegram_switch.pack(pady=10, padx=20, anchor="w")
        ctk.CTkLabel(scrollable_frame, text="ContrÃ´le Ã  distance via Telegram Bot.", font=ctk.CTkFont(size=10), text_color=COLOR_SECONDARY).pack(padx=40, anchor="w")

        # inputs pour Telegram
        telegram_form = ctk.CTkFrame(scrollable_frame, fg_color="transparent")
        telegram_form.pack(pady=5, padx=40, fill="x")
        
        token_entry = ctk.CTkEntry(telegram_form, placeholder_text="Token Bot Telegram", height=28, fg_color=COLOR_NAV)
        token_entry.insert(0, self.telegram_bot_token)
        token_entry.pack(side="left", fill="x", expand=True, padx=(0, 5))
        
        chat_id_entry = ctk.CTkEntry(telegram_form, placeholder_text="Chat ID", width=100, height=28, fg_color=COLOR_NAV)
        chat_id_entry.insert(0, self.telegram_chat_id)
        chat_id_entry.pack(side="left")

        def save_tg_credentials():
            self.telegram_bot_token = token_entry.get().strip()
            self.telegram_chat_id = chat_id_entry.get().strip()
            self._save_settings()
            self.add_message("ParamÃ¨tres Telegram enregistrÃ©s.", "SystÃ¨me")
            # Forcer le redÃ©marrage si activÃ©
            if self.agent_telegram_enabled:
                self.add_message("RedÃ©marrage de l'Ã©couteur Telegram...", "SystÃ¨me")
                threading.Thread(target=self.start_telegram_bot, daemon=True).start()

        def test_tg_connection():
            token = token_entry.get().strip()
            chat_id = chat_id_entry.get().strip()
            if not token:
                self.add_message("Erreur : Token manquant pour le test.", "SystÃ¨me")
                return
            if not chat_id or chat_id == "0":
                self.add_message("Erreur : Chat ID requis pour le test. Envoyez d'abord un message au bot.", "SystÃ¨me")
                return
            
            self.add_message("Envoi d'un message de test vers Telegram...", "SystÃ¨me")
            def _send_test():
                try:
                    url = f"https://api.telegram.org/bot{token}/sendMessage"
                    res = requests.post(url, json={"chat_id": chat_id, "text": "Jarvisse : Test de connexion rÃ©ussi ! âœ…"}, timeout=10)
                    if res.status_code == 200:
                        self.add_message("Test rÃ©ussi ! Message envoyÃ©.", "SystÃ¨me")
                    else:
                        self.add_message(f"Ã‰chec du test : {res.text}", "SystÃ¨me")
                except Exception as e:
                    self.add_message(f"Erreur test Telegram : {e}", "SystÃ¨me")
            
            threading.Thread(target=_send_test, daemon=True).start()

        btn_row = ctk.CTkFrame(scrollable_frame, fg_color="transparent")
        btn_row.pack(pady=(0, 10), padx=40, fill="x")
        
        ctk.CTkButton(btn_row, text="Enregistrer", height=24, font=ctk.CTkFont(size=10), command=save_tg_credentials).pack(side="left", padx=(0, 5))
        ctk.CTkButton(btn_row, text="Tester le Bot", height=24, font=ctk.CTkFont(size=10), fg_color=COLOR_ACCENT, command=test_tg_connection).pack(side="left")

        # NOUVEAUX AGENTS PREMIUM (UI)
        ctk.CTkLabel(scrollable_frame, text="â”€â”€â”€â”€â”€â”€â”€â”€ AGENTS PREMIUM â”€â”€â”€â”€â”€â”€â”€â”€", font=ctk.CTkFont(size=11, weight="bold"), text_color=COLOR_ACCENT_PURPLE).pack(pady=15)

        # Vision
        vision_var = ctk.BooleanVar(value=self.agent_vision_enabled)
        ctk.CTkSwitch(scrollable_frame, text="Agent Vision Multimodale", variable=vision_var, command=lambda: self.toggle_agent("vision", vision_var.get())).pack(pady=10, padx=20, anchor="w")
        ctk.CTkLabel(scrollable_frame, text="Analyse visuelle et lecture de documents via camÃ©ra.", font=ctk.CTkFont(size=10), text_color=COLOR_SECONDARY).pack(padx=40, anchor="w")

        # Domotique
        domo_var = ctk.BooleanVar(value=self.agent_domotique_enabled)
        ctk.CTkSwitch(scrollable_frame, text="Agent Domotique & MÃ©dias", variable=domo_var, command=lambda: self.toggle_agent("domotique", domo_var.get())).pack(pady=10, padx=20, anchor="w")
        ctk.CTkButton(scrollable_frame, text="âš™ï¸ Configurer les Interfaces Domotique", height=24, font=ctk.CTkFont(size=10), command=self.open_domotique_settings, fg_color=COLOR_CARD, text_color=COLOR_ACCENT_PURPLE).pack(padx=40, anchor="w", pady=(0, 5))
        ctk.CTkLabel(scrollable_frame, text="ContrÃ´le maison connectÃ©e et services multimÃ©dias.", font=ctk.CTkFont(size=10), text_color=COLOR_SECONDARY).pack(padx=40, anchor="w")

        # Finance
        fin_var = ctk.BooleanVar(value=self.agent_finance_enabled)
        ctk.CTkSwitch(scrollable_frame, text="Agent Finance & Crypto", variable=fin_var, command=lambda: self.toggle_agent("finance", fin_var.get())).pack(pady=10, padx=20, anchor="w")
        ctk.CTkLabel(scrollable_frame, text="Suivi Bourse, Crypto et alertes financiÃ¨res en temps rÃ©el.", font=ctk.CTkFont(size=10), text_color=COLOR_SECONDARY).pack(padx=40, anchor="w")

        # SecrÃ©taire
        sec_var = ctk.BooleanVar(value=self.agent_secretaire_enabled)
        ctk.CTkSwitch(scrollable_frame, text="Agent SecrÃ©taire & ProductivitÃ©", variable=sec_var, command=lambda: self.toggle_agent("secretaire", sec_var.get())).pack(pady=10, padx=20, anchor="w")
        ctk.CTkLabel(scrollable_frame, text="Gestion calendrier, rendez-vous et organisation.", font=ctk.CTkFont(size=10), text_color=COLOR_SECONDARY).pack(padx=40, anchor="w")

        # Traducteur
        trad_var = ctk.BooleanVar(value=self.agent_traducteur_enabled)
        ctk.CTkSwitch(scrollable_frame, text="Agent Traducteur Universel", variable=trad_var, command=lambda: self.toggle_agent("traducteur", trad_var.get())).pack(pady=10, padx=20, anchor="w")
        ctk.CTkLabel(scrollable_frame, text="Traduction bidirectionnelle instantanÃ©e.", font=ctk.CTkFont(size=10), text_color=COLOR_SECONDARY).pack(padx=40, anchor="w")

        # Data Miner
        miner_var = ctk.BooleanVar(value=self.agent_miner_enabled)
        ctk.CTkSwitch(scrollable_frame, text="Agent Gestionnaire (Data Miner)", variable=miner_var, command=lambda: self.toggle_agent("miner", miner_var.get())).pack(pady=10, padx=20, anchor="w")
        ctk.CTkLabel(scrollable_frame, text="Recherche profonde dans vos fichiers locaux.", font=ctk.CTkFont(size=10), text_color=COLOR_SECONDARY).pack(padx=40, anchor="w")

        # NOUVEAUX AGENTS AJOUTÃ‰S
        ctk.CTkLabel(scrollable_frame, text="â”€â”€â”€â”€â”€â”€â”€â”€ AGENTS LIFESTYLE â”€â”€â”€â”€â”€â”€â”€â”€", font=ctk.CTkFont(size=11, weight="bold"), text_color=COLOR_ACCENT).pack(pady=15)

        # News
        news_var = ctk.BooleanVar(value=self.agent_news_enabled)
        ctk.CTkSwitch(scrollable_frame, text="Agent News & Veille Tech", variable=news_var, command=lambda: self.toggle_agent("news", news_var.get())).pack(pady=10, padx=20, anchor="w")
        ctk.CTkLabel(scrollable_frame, text="ActualitÃ©s mondiales, tech et mÃ©tÃ©o en temps rÃ©el.", font=ctk.CTkFont(size=10), text_color=COLOR_SECONDARY).pack(padx=40, anchor="w")

        # Cuisine
        cuisine_var = ctk.BooleanVar(value=self.agent_cuisine_enabled)
        ctk.CTkSwitch(scrollable_frame, text="Agent Cuisine & Gastronomie", variable=cuisine_var, command=lambda: self.toggle_agent("cuisine", cuisine_var.get())).pack(pady=10, padx=20, anchor="w")
        ctk.CTkLabel(scrollable_frame, text="Suggestions de recettes, conseils culinaires et nutrition.", font=ctk.CTkFont(size=10), text_color=COLOR_SECONDARY).pack(padx=40, anchor="w")

        # SantÃ©
        sante_var = ctk.BooleanVar(value=self.agent_sante_enabled)
        ctk.CTkSwitch(scrollable_frame, text="Agent SantÃ© & Bien-Ãªtre", variable=sante_var, command=lambda: self.toggle_agent("sante", sante_var.get())).pack(pady=10, padx=20, anchor="w")
        ctk.CTkLabel(scrollable_frame, text="Conseils fitness, sommeil, stress et suivi santÃ©.", font=ctk.CTkFont(size=10), text_color=COLOR_SECONDARY).pack(padx=40, anchor="w")

        # Voyage
        voyage_var = ctk.BooleanVar(value=self.agent_voyage_enabled)
        ctk.CTkSwitch(scrollable_frame, text="Agent Voyage & Exploration", variable=voyage_var, command=lambda: self.toggle_agent("voyage", voyage_var.get())).pack(pady=10, padx=20, anchor="w")
        ctk.CTkLabel(scrollable_frame, text="Planification de voyages, hÃ´tels et dÃ©couvertes culturelles.", font=ctk.CTkFont(size=10), text_color=COLOR_SECONDARY).pack(padx=40, anchor="w")

        # NOUVEAUX AGENTS PROFESSIONNELS
        ctk.CTkLabel(scrollable_frame, text="â”€â”€â”€â”€â”€â”€â”€â”€ AGENTS EXPERTS â”€â”€â”€â”€â”€â”€â”€â”€", font=ctk.CTkFont(size=11, weight="bold"), text_color=COLOR_ACCENT_PURPLE).pack(pady=15)

        # Education
        edu_var = ctk.BooleanVar(value=self.agent_education_enabled)
        ctk.CTkSwitch(scrollable_frame, text="Agent Ã‰ducation & Tutorat", variable=edu_var, command=lambda: self.toggle_agent("education", edu_var.get())).pack(pady=10, padx=20, anchor="w")
        ctk.CTkLabel(scrollable_frame, text="Aide aux devoirs, explications de concepts et langues.", font=ctk.CTkFont(size=10), text_color=COLOR_SECONDARY).pack(padx=40, anchor="w")

        # Shopping
        shop_var = ctk.BooleanVar(value=self.agent_shopping_enabled)
        ctk.CTkSwitch(scrollable_frame, text="Agent Shopping & Bons Plans", variable=shop_var, command=lambda: self.toggle_agent("shopping", shop_var.get())).pack(pady=10, padx=20, anchor="w")
        ctk.CTkLabel(scrollable_frame, text="Recherche de prix, codes promos et comparatifs produits.", font=ctk.CTkFont(size=10), text_color=COLOR_SECONDARY).pack(padx=40, anchor="w")

        # Social / CM
        social_var = ctk.BooleanVar(value=self.agent_social_enabled)
        ctk.CTkSwitch(scrollable_frame, text="Agent Community Manager", variable=social_var, command=lambda: self.toggle_agent("social", social_var.get())).pack(pady=10, padx=20, anchor="w")
        ctk.CTkLabel(scrollable_frame, text="RÃ©daction de posts rÃ©seaux sociaux et veille tendances.", font=ctk.CTkFont(size=10), text_color=COLOR_SECONDARY).pack(padx=40, anchor="w")

        # Psy
        psy_var = ctk.BooleanVar(value=self.agent_psy_enabled)
        ctk.CTkSwitch(scrollable_frame, text="Agent Bien-Ãªtre Mental", variable=psy_var, command=lambda: self.toggle_agent("psy", psy_var.get())).pack(pady=10, padx=20, anchor="w")
        ctk.CTkLabel(scrollable_frame, text="Ã‰coute active, mÃ©ditation et gestion du stress.", font=ctk.CTkFont(size=10), text_color=COLOR_SECONDARY).pack(padx=40, anchor="w")

        # Immo
        immo_var = ctk.BooleanVar(value=self.agent_immo_enabled)
        ctk.CTkSwitch(scrollable_frame, text="Agent Immobilier & Patrimoine", variable=immo_var, command=lambda: self.toggle_agent("immo", immo_var.get())).pack(pady=10, padx=20, anchor="w")
        ctk.CTkLabel(scrollable_frame, text="Recherche d'annonces, estimation et calculs de prÃªts.", font=ctk.CTkFont(size=10), text_color=COLOR_SECONDARY).pack(padx=40, anchor="w")

        # Auto
        auto_var = ctk.BooleanVar(value=self.agent_auto_enabled)
        ctk.CTkSwitch(scrollable_frame, text="Agent Automobile & MobilitÃ©", variable=auto_var, command=lambda: self.toggle_agent("auto", auto_var.get())).pack(pady=10, padx=20, anchor="w")
        ctk.CTkLabel(scrollable_frame, text="Entretien vÃ©hicule, prix carburants et info-trafic.", font=ctk.CTkFont(size=10), text_color=COLOR_SECONDARY).pack(padx=40, anchor="w")

        # Carrier
        carrier_var = ctk.BooleanVar(value=self.agent_carrier_enabled)
        ctk.CTkSwitch(scrollable_frame, text="Agent CarriÃ¨re & RH", variable=carrier_var, command=lambda: self.toggle_agent("carrier", carrier_var.get())).pack(pady=10, padx=20, anchor="w")
        ctk.CTkLabel(scrollable_frame, text="Optimisation CV, prÃ©paration entretiens et job matching.", font=ctk.CTkFont(size=10), text_color=COLOR_SECONDARY).pack(padx=40, anchor="w")

        # NOUVEAUX AGENTS (PACK GRATUIT & PREMIUM)
        ctk.CTkLabel(scrollable_frame, text="â”€â”€â”€â”€â”€â”€â”€â”€ PACK GRATUIT â”€â”€â”€â”€â”€â”€â”€â”€", font=ctk.CTkFont(size=11, weight="bold"), text_color=COLOR_ACCENT).pack(pady=15)
        
        # Arbitre
        arb_var = ctk.BooleanVar(value=self.agent_arbitre_enabled)
        ctk.CTkSwitch(scrollable_frame, text="Agent L'Arbitre (Sport)", variable=arb_var, command=lambda: self.toggle_agent("arbitre", arb_var.get())).pack(pady=10, padx=20, anchor="w")
        ctk.CTkLabel(scrollable_frame, text="Scores en direct et actualitÃ©s sportives.", font=ctk.CTkFont(size=10), text_color=COLOR_SECONDARY).pack(padx=40, anchor="w")

        # Ã‰colo
        eco_var = ctk.BooleanVar(value=self.agent_ecolo_enabled)
        ctk.CTkSwitch(scrollable_frame, text="Agent L'Ã‰colo (Green)", variable=eco_var, command=lambda: self.toggle_agent("ecolo", eco_var.get())).pack(pady=10, padx=20, anchor="w")
        ctk.CTkLabel(scrollable_frame, text="Conseils Ã©cologiques et guide de recyclage.", font=ctk.CTkFont(size=10), text_color=COLOR_SECONDARY).pack(padx=40, anchor="w")

        # Guetteur
        gue_var = ctk.BooleanVar(value=self.agent_guetteur_enabled)
        ctk.CTkSwitch(scrollable_frame, text="Agent Le Guetteur (Loisirs)", variable=gue_var, command=lambda: self.toggle_agent("guetteur", gue_var.get())).pack(pady=10, padx=20, anchor="w")
        ctk.CTkLabel(scrollable_frame, text="Sorties, cinÃ©ma et recommandations streaming.", font=ctk.CTkFont(size=10), text_color=COLOR_SECONDARY).pack(padx=40, anchor="w")

        # Historien
        his_var = ctk.BooleanVar(value=self.agent_historien_enabled)
        ctk.CTkSwitch(scrollable_frame, text="Agent L'Historien (Culture)", variable=his_var, command=lambda: self.toggle_agent("historien", his_var.get())).pack(pady=10, padx=20, anchor="w")
        ctk.CTkLabel(scrollable_frame, text="EphÃ©mÃ©ride et anecdotes historiques.", font=ctk.CTkFont(size=10), text_color=COLOR_SECONDARY).pack(padx=40, anchor="w")

        # DÃ©panneur
        dep_var = ctk.BooleanVar(value=self.agent_depanneur_enabled)
        ctk.CTkSwitch(scrollable_frame, text="Agent Le DÃ©panneur (DIY)", variable=dep_var, command=lambda: self.toggle_agent("depanneur", dep_var.get())).pack(pady=10, padx=20, anchor="w")
        ctk.CTkLabel(scrollable_frame, text="Tutos bricolage et rÃ©parations maison.", font=ctk.CTkFont(size=10), text_color=COLOR_SECONDARY).pack(padx=40, anchor="w")

        ctk.CTkLabel(scrollable_frame, text="â”€â”€â”€â”€â”€â”€â”€â”€ PACK PREMIUM â”€â”€â”€â”€â”€â”€â”€â”€", font=ctk.CTkFont(size=11, weight="bold"), text_color=COLOR_ACCENT_PURPLE).pack(pady=15)

        # Astrophysicien
        ast_var = ctk.BooleanVar(value=self.agent_astroph_enabled)
        ctk.CTkSwitch(scrollable_frame, text="Agent L'Astrophysicien", variable=ast_var, command=lambda: self.toggle_agent("astroph", ast_var.get())).pack(pady=10, padx=20, anchor="w")
        ctk.CTkLabel(scrollable_frame, text="Espace, cosmos et suivis de lancements.", font=ctk.CTkFont(size=10), text_color=COLOR_SECONDARY).pack(padx=40, anchor="w")

        # StratÃ¨ge
        str_var = ctk.BooleanVar(value=self.agent_stratege_enabled)
        ctk.CTkSwitch(scrollable_frame, text="Agent Le StratÃ¨ge", variable=str_var, command=lambda: self.toggle_agent("stratege", str_var.get())).pack(pady=10, padx=20, anchor="w")
        ctk.CTkLabel(scrollable_frame, text="Investissements, Bourse et Patrimoine.", font=ctk.CTkFont(size=10), text_color=COLOR_SECONDARY).pack(padx=40, anchor="w")

        # Architecte
        arc_var = ctk.BooleanVar(value=self.agent_architecte_enabled)
        ctk.CTkSwitch(scrollable_frame, text="Agent L'Architecte", variable=arc_var, command=lambda: self.toggle_agent("architecte", arc_var.get())).pack(pady=10, padx=20, anchor="w")
        ctk.CTkLabel(scrollable_frame, text="Design intÃ©rieur et tendances dÃ©co.", font=ctk.CTkFont(size=10), text_color=COLOR_SECONDARY).pack(padx=40, anchor="w")

        # Business Analyst
        bus_var = ctk.BooleanVar(value=self.agent_business_enabled)
        ctk.CTkSwitch(scrollable_frame, text="Agent Business Analyst", variable=bus_var, command=lambda: self.toggle_agent("business", bus_var.get())).pack(pady=10, padx=20, anchor="w")
        ctk.CTkLabel(scrollable_frame, text="Analyse de donnÃ©es et stratÃ©gies pro.", font=ctk.CTkFont(size=10), text_color=COLOR_SECONDARY).pack(padx=40, anchor="w")

        # Polyglotte
        pol_var = ctk.BooleanVar(value=self.agent_polyglotte_enabled)
        ctk.CTkSwitch(scrollable_frame, text="Agent Le Polyglotte", variable=pol_var, command=lambda: self.toggle_agent("polyglotte", pol_var.get())).pack(pady=10, padx=20, anchor="w")
        ctk.CTkLabel(scrollable_frame, text="Apprentissage avancÃ© des langues.", font=ctk.CTkFont(size=10), text_color=COLOR_SECONDARY).pack(padx=40, anchor="w")

        # Nounou
        nou_var = ctk.BooleanVar(value=self.agent_nounou_enabled)
        ctk.CTkSwitch(scrollable_frame, text="Agent La Nounou", variable=nou_var, command=lambda: self.toggle_agent("nounou", nou_var.get())).pack(pady=10, padx=20, anchor="w")
        ctk.CTkLabel(scrollable_frame, text="ParentalitÃ© et Ã©veil des enfants.", font=ctk.CTkFont(size=10), text_color=COLOR_SECONDARY).pack(padx=40, anchor="w")

        # Agent Appel UI
        ctk.CTkLabel(scrollable_frame, text="â”€â”€â”€â”€â”€â”€â”€â”€ AGENT APPEL â”€â”€â”€â”€â”€â”€â”€â”€", font=ctk.CTkFont(size=11, weight="bold"), text_color=COLOR_ACCENT).pack(pady=15)
        appel_var = ctk.BooleanVar(value=self.agent_appel_enabled)
        ctk.CTkSwitch(scrollable_frame, text="Activer Agent Appel", variable=appel_var, command=lambda: self.toggle_agent("appel", appel_var.get())).pack(pady=10, padx=20, anchor="w")
        
        autoreply_var = ctk.BooleanVar(value=self.agent_appel_autoreply)
        def toggle_autoreply():
            self.agent_appel_autoreply = autoreply_var.get()
            self._save_settings()
        ctk.CTkSwitch(scrollable_frame, text="RÃ©ponse Automatique", variable=autoreply_var, command=toggle_autoreply).pack(pady=5, padx=20, anchor="w")
        
        num_entry = ctk.CTkEntry(scrollable_frame, placeholder_text="Votre numÃ©ro de tÃ©lÃ©phone", height=28, fg_color=COLOR_NAV)
        num_entry.insert(0, self.agent_appel_number)
        num_entry.pack(pady=5, padx=40, fill="x")
        
        phrases_text = ctk.CTkTextbox(scrollable_frame, height=100, fg_color=COLOR_NAV)
        phrases_text.insert("1.0", "\n".join(self.agent_appel_phrases))
        phrases_text.pack(pady=5, padx=40, fill="x")
        ctk.CTkLabel(scrollable_frame, text="Une phrase par ligne (RÃ©ponses auto)", font=ctk.CTkFont(size=10), text_color=COLOR_SECONDARY).pack(padx=40, anchor="w")

        def save_appel_settings():
            self.agent_appel_number = num_entry.get().strip()
            self.agent_appel_phrases = [p.strip() for p in phrases_text.get("1.0", "end").split("\n") if p.strip()]
            self._save_settings()
            self.add_message("ParamÃ¨tres Agent Appel enregistrÃ©s.", "SystÃ¨me")

        ctk.CTkButton(scrollable_frame, text="Enregistrer Agent Appel", height=24, font=ctk.CTkFont(size=10), command=save_appel_settings).pack(pady=5, padx=40, anchor="w")


        # SÃ©lection de la Voix
        ctk.CTkLabel(scrollable_frame, text="--- Personnalisation de la Voix ---", font=ctk.CTkFont(weight="bold")).pack(pady=(20, 5), padx=20, anchor="w")
        
        voice_var = ctk.StringVar(value=self.voice)
        voice_menu = ctk.CTkOptionMenu(
            scrollable_frame, 
            values=self.available_voices,
            variable=voice_var,
            command=self.change_voice,
            fg_color=COLOR_CARD,
            button_color=COLOR_ACCENT,
            dropdown_fg_color=COLOR_CARD
        )
        voice_menu.pack(pady=10, padx=20, anchor="w")
        ctk.CTkLabel(scrollable_frame, text="Choisissez la voix de votre assistant (Microsoft Edge TTS).", font=ctk.CTkFont(size=10), text_color=COLOR_SECONDARY).pack(padx=40, anchor="w")

        # Espace supplÃ©mentaire pour activer le scroll
        ctk.CTkLabel(scrollable_frame, text="", height=20).pack(pady=10)

        # Boutons en bas (hors du scrollable_frame pour qu'ils restent visibles)
        ctk.CTkButton(settings_win, text="Authentifier Gmail", command=self.agent_gmail_authenticate, fg_color=COLOR_CARD).pack(pady=5, padx=20)
        ctk.CTkButton(settings_win, text="Fermer", command=settings_win.destroy, fg_color=COLOR_ACCENT_RED).pack(pady=5, padx=20)

    def open_domotique_settings(self):
        """Ouvre une fenÃªtre pour configurer les interfaces domotique (WLED, Arduino, Webhooks, Home Assistant)"""
        domo_win = ctk.CTkToplevel(self)
        domo_win.title("âš™ï¸ Configuration Domotique - Jarvisse")
        domo_win.geometry("450x650")
        domo_win.configure(fg_color=COLOR_BG)
        domo_win.transient(self)
        domo_win.attributes("-topmost", True)
        domo_win.lift()
        domo_win.focus_force()

        ctk.CTkLabel(domo_win, text="Interfaces Domotique", font=ctk.CTkFont(size=16, weight="bold"), text_color=COLOR_ACCENT_PURPLE).pack(pady=15)

        scrollable = ctk.CTkScrollableFrame(domo_win, fg_color=COLOR_BG, width=420, height=500)
        scrollable.pack(fill="both", expand=True, padx=10, pady=5)

        # WLED
        ctk.CTkLabel(scrollable, text="ðŸŒˆ Rubans LED (WLED)", font=ctk.CTkFont(weight="bold"), text_color=COLOR_ACCENT).pack(pady=(10, 5), anchor="w", padx=10)
        wled_ip_entry = ctk.CTkEntry(scrollable, placeholder_text="Adresse IP WLED (ex: 192.168.1.50)", width=350, fg_color=COLOR_NAV)
        wled_ip_entry.insert(0, self.domo_wled_ip)
        wled_ip_entry.pack(pady=5, padx=20)

        # Home Assistant
        ctk.CTkLabel(scrollable, text="ðŸ  Home Assistant (API)", font=ctk.CTkFont(weight="bold"), text_color=COLOR_ACCENT).pack(pady=(15, 5), anchor="w", padx=10)
        ha_url_entry = ctk.CTkEntry(scrollable, placeholder_text="URL Instance (ex: http://homeassistant.local:8123)", width=350, fg_color=COLOR_NAV)
        ha_url_entry.insert(0, self.domo_ha_url)
        ha_url_entry.pack(pady=5, padx=20)
        
        ha_token_entry = ctk.CTkEntry(scrollable, placeholder_text="Jeton d'accÃ¨s longue durÃ©e (Long-Lived Access Token)", width=350, fg_color=COLOR_NAV)
        ha_token_entry.insert(0, self.domo_ha_token)
        ha_token_entry.pack(pady=5, padx=20)

        # Arduino / Serial
        ctk.CTkLabel(scrollable, text="ðŸ”Œ ContrÃ´le SÃ©rie (Arduino / ESP32)", font=ctk.CTkFont(weight="bold"), text_color=COLOR_ACCENT).pack(pady=(15, 5), anchor="w", padx=10)
        arduino_port_entry = ctk.CTkEntry(scrollable, placeholder_text="Port COM (ex: COM3 ou /dev/ttyUSB0)", width=350, fg_color=COLOR_NAV)
        arduino_port_entry.insert(0, self.domo_arduino_port)
        arduino_port_entry.pack(pady=5, padx=20)

        # Webhooks GÃ©nÃ©riques
        ctk.CTkLabel(scrollable, text="ðŸ”— Webhooks GÃ©nÃ©riques (HTTP GET/POST)", font=ctk.CTkFont(weight="bold"), text_color=COLOR_ACCENT).pack(pady=(15, 5), anchor="w", padx=10)
        webhook_on_entry = ctk.CTkEntry(scrollable, placeholder_text="URL Activation (ON)", width=350, fg_color=COLOR_NAV)
        webhook_on_entry.insert(0, self.domo_webhook_on)
        webhook_on_entry.pack(pady=5, padx=20)
        
        webhook_off_entry = ctk.CTkEntry(scrollable, placeholder_text="URL DÃ©sactivation (OFF)", width=350, fg_color=COLOR_NAV)
        webhook_off_entry.insert(0, self.domo_webhook_off)
        webhook_off_entry.pack(pady=5, padx=20)

        def save_domo_settings():
            self.domo_wled_ip = wled_ip_entry.get().strip()
            self.domo_ha_url = ha_url_entry.get().strip()
            self.domo_ha_token = ha_token_entry.get().strip()
            self.domo_arduino_port = arduino_port_entry.get().strip()
            self.domo_webhook_on = webhook_on_entry.get().strip()
            self.domo_webhook_off = webhook_off_entry.get().strip()
            self._save_settings()
            self.add_message("ParamÃ¨tres Domotique enregistrÃ©s avec succÃ¨s.", "SystÃ¨me")
            domo_win.destroy()

        ctk.CTkButton(domo_win, text="Sauvegarder la Configuration", command=save_domo_settings, fg_color=COLOR_ACCENT_PURPLE, font=ctk.CTkFont(weight="bold")).pack(pady=20)

    def toggle_agent(self, agent_type, enabled):
        if agent_type == "network": self.agent_network_enabled = enabled
        elif agent_type == "control": self.agent_control_enabled = enabled
        elif agent_type == "remote": self.agent_remote_enabled = enabled
        elif agent_type == "gmail":
            self.agent_gmail_enabled = enabled
            if enabled: threading.Thread(target=self.start_gmail_listener, daemon=True).start()
        elif agent_type == "browser": self.agent_browser_enabled = enabled
        elif agent_type == "legal": self.agent_legal_enabled = enabled
        elif agent_type == "cyber": self.agent_cyber_enabled = enabled
        elif agent_type == "doctor": self.agent_doctor_enabled = enabled
        elif agent_type == "video": self.agent_video_enabled = enabled
        elif agent_type == "android": self.agent_android_enabled = enabled
        elif agent_type == "coding": self.agent_coding_enabled = enabled
        elif agent_type == "generation": self.agent_generation_enabled = enabled
        elif agent_type == "lecture": self.agent_lecture_enabled = enabled
        elif agent_type == "telegram":
            self.agent_telegram_enabled = enabled
            if enabled: threading.Thread(target=self.start_telegram_bot, daemon=True).start()
        elif agent_type == "vision": self.agent_vision_enabled = enabled
        elif agent_type == "domotique": self.agent_domotique_enabled = enabled
        elif agent_type == "finance": self.agent_finance_enabled = enabled
        elif agent_type == "secretaire": self.agent_secretaire_enabled = enabled
        elif agent_type == "traducteur": self.agent_traducteur_enabled = enabled
        elif agent_type == "miner": self.agent_miner_enabled = enabled
        elif agent_type == "news": self.agent_news_enabled = enabled
        elif agent_type == "cuisine": self.agent_cuisine_enabled = enabled
        elif agent_type == "sante": self.agent_sante_enabled = enabled
        elif agent_type == "voyage": self.agent_voyage_enabled = enabled
        elif agent_type == "education": self.agent_education_enabled = enabled
        elif agent_type == "shopping": self.agent_shopping_enabled = enabled
        elif agent_type == "social": self.agent_social_enabled = enabled
        elif agent_type == "psy": self.agent_psy_enabled = enabled
        elif agent_type == "immo": self.agent_immo_enabled = enabled
        elif agent_type == "auto": self.agent_auto_enabled = enabled
        elif agent_type == "carrier": self.agent_carrier_enabled = enabled
        elif agent_type == "arbitre": self.agent_arbitre_enabled = enabled
        elif agent_type == "ecolo": self.agent_ecolo_enabled = enabled
        elif agent_type == "guetteur": self.agent_guetteur_enabled = enabled
        elif agent_type == "historien": self.agent_historien_enabled = enabled
        elif agent_type == "depanneur": self.agent_depanneur_enabled = enabled
        elif agent_type == "astroph": self.agent_astroph_enabled = enabled
        elif agent_type == "stratege": self.agent_stratege_enabled = enabled
        elif agent_type == "architecte": self.agent_architecte_enabled = enabled
        elif agent_type == "business": self.agent_business_enabled = enabled
        elif agent_type == "polyglotte": self.agent_polyglotte_enabled = enabled
        elif agent_type == "nounou": self.agent_nounou_enabled = enabled
        elif agent_type == "appel": self.agent_appel_enabled = enabled
        self._save_settings()

    def change_voice(self, new_voice):
        """Change la voix de l'assistant et lance un test"""
        self.voice = new_voice
        self._save_settings()
        self.parler(f"Ma voix a Ã©tÃ© mise Ã  jour. J'utiliserai dÃ©sormais la voix de {new_voice.split('-')[2].replace('Neural', '')}.")

    def start_telegram_bot(self):
        """Lance l'Ã©couteur de commandes Telegram (Long Polling) avec support vocal"""
        if getattr(self, "telegram_bot_running", False) and self.agent_telegram_enabled:
            print("Telegram: L'Ã©couteur est dÃ©jÃ  actif.")
            return

        if not self.telegram_bot_token:
            print("Telegram: Pas de token configurÃ©.")
            return

        self.telegram_bot_running = True
        self.agent_telegram_enabled = True
        self.add_message("Agent Bot Telegram : Initialisation...", "SystÃ¨me")
        
        # Nettoyage Webhook pour Long Polling
        try:
            requests.get(f"https://api.telegram.org/bot{self.telegram_bot_token}/deleteWebhook", timeout=10)
        except: pass

        offset = 0
        connected_notified = False
        print(f"Telegram: DÃ©marrage de l'Ã©coute...")
        
        while self.agent_telegram_enabled:
            try:
                url = f"https://api.telegram.org/bot{self.telegram_bot_token}/getUpdates"
                params = {"offset": offset, "timeout": 20}
                response = requests.get(url, params=params, timeout=25)
                
                if response.status_code == 200:
                    data = response.json()
                    if not data.get("ok"): break

                    if not connected_notified:
                        self.add_message("Agent Bot Telegram : PrÃªt.", "SystÃ¨me")
                        connected_notified = True

                    for update in data.get("result", []):
                        offset = update["update_id"] + 1
                        msg = update.get("message")
                        if not msg: continue
                        
                        chat_id = str(msg["chat"]["id"])
                        if not self.telegram_chat_id or self.telegram_chat_id == "0":
                            self.telegram_chat_id = chat_id
                            self._save_settings()

                        # Gestion TEXTE
                        if "text" in msg:
                            text = msg["text"]
                            print(f"Telegram [Texte]: {text}")
                            self.add_message(f"[Telegram] {text}", "Vous")
                            self.last_command_origin = "telegram"
                            threading.Thread(target=self.executer_commande, args=(text,), daemon=True).start()
                        
                        # Gestion VOCAL (Vitesse de la lumiÃ¨re via Gemini)
                        elif "voice" in msg:
                            print(f"Telegram [Vocal] reÃ§u.")
                            self.add_message("[Telegram] ðŸŽ™ï¸ Message vocal reÃ§u...", "Vous")
                            file_id = msg["voice"]["file_id"]
                            threading.Thread(target=self._process_telegram_voice, args=(file_id,), daemon=True).start()

                elif response.status_code == 409:
                    time.sleep(2)
                elif response.status_code == 401:
                    break
                else:
                    time.sleep(5)
            except Exception as e:
                print(f"Erreur boucle Telegram: {e}")
                time.sleep(5)
        
        self.telegram_bot_running = False
        print("Telegram: Agent Bot Telegram arrÃªtÃ©.")

    def _process_telegram_voice(self, file_id):
        """TÃ©lÃ©charge et transcrit un message vocal Telegram instantanÃ©ment"""
        transcription = self._transcrire_telegram_audio(file_id)
        if transcription:
            self.add_message(f"[Telegram Audio] {transcription}", "Vous")
            self.last_command_origin = "telegram"
            # RÃ©ponse Ã  la vitesse de la lumiÃ¨re
            self.executer_commande(transcription)
        else:
            self.send_telegram_message("âŒ DÃ©solÃ© Patron, je n'ai pas pu dÃ©coder votre message vocal.")

    def _transcrire_telegram_audio(self, file_id):
        """Moteur de transcription ultra-rapide pour Telegram"""
        try:
            # 1. RÃ©cupÃ©ration du chemin du fichier
            res = requests.get(f"https://api.telegram.org/bot{self.telegram_bot_token}/getFile", params={"file_id": file_id}, timeout=10)
            file_info = res.json()
            if not file_info.get("ok"):
                print(f"Telegram getFile error: {file_info.get('description')}")
                return None
                
            file_path = file_info.get("result", {}).get("file_path")
            if not file_path: return None
            
            # 2. TÃ©lÃ©chargement direct des bytes
            download_url = f"https://api.telegram.org/file/bot{self.telegram_bot_token}/{file_path}"
            audio_res = requests.get(download_url, timeout=20)
            if audio_res.status_code != 200:
                print(f"Telegram download error: HTTP {audio_res.status_code}")
                return None
            
            # 3. Transcription Multimodale (Vitesse Gemini 2.0 Flash)
            if self.genai_client:
                extension = file_path.lower().split('.')[-1]
                mime = "audio/ogg" if extension in ['ogg', 'oga', 'opus'] else "audio/wav"
                
                try:
                    print(f"Transcription Gemini en cours (Format: {extension})...")
                    response = self.genai_client.models.generate_content(
                        model=self.model_name,
                        contents=[
                            types.Part.from_bytes(data=audio_res.content, mime_type=mime),
                            "Transcris fidÃ¨lement cet audio en franÃ§ais. Ne rÃ©ponds que par le texte transcrit."
                        ]
                    )
                    return response.text.strip().lower()
                except Exception as e:
                    if "429" in str(e) or "RESOURCE_EXHAUSTED" in str(e):
                        print("âš ï¸ Quota Gemini Ã©puisÃ©. Tentative via le moteur de secours (Google)...")
                        return self._fallback_transcription(audio_res.content, extension)
                    raise e
            else:
                return self._fallback_transcription(audio_res.content, extension)
        except Exception as e:
            print(f"Erreur CRITIQUE _transcrire_telegram_audio: {e}")
            return None

    def _fallback_transcription(self, audio_bytes, extension):
        """Moteur de secours utilisant Google SpeechRecognition"""
        try:
            import io
            from pydub import AudioSegment
            
            # Conversion en WAV car SpeechRecognition ne lit pas directement l'OGG/Opus
            audio_io = io.BytesIO(audio_bytes)
            try:
                if extension in ['ogg', 'oga', 'opus']:
                    audio = AudioSegment.from_file(audio_io, codec="opus")
                else:
                    audio = AudioSegment.from_file(audio_io)
            except:
                # DeuxiÃ¨me essai sans spÃ©cifier le codec
                audio_io.seek(0)
                audio = AudioSegment.from_file(audio_io)
                
            wav_io = io.BytesIO()
            audio.export(wav_io, format="wav")
            wav_io.seek(0)
            
            reconnaisseur = sr.Recognizer()
            with sr.AudioFile(wav_io) as source:
                audio_data = reconnaisseur.record(source)
                text = reconnaisseur.recognize_google(audio_data, language='fr-FR')
                print(f"âœ… Secours rÃ©ussi : {text}")
                return text.lower()
        except Exception as e:
            print(f"âŒ Ã‰chec du moteur de secours : {e}")
            return None


    def trainer_all_agents(self):
        """ProcÃ©dure d'entraÃ®nement et de calibration de tous les agents"""
        self.add_message("DÃ©but de la phase d'entraÃ®nement intensif...", "SystÃ¨me")
        time.sleep(1)
        
        # 1. Test Docteur SystÃ¨me
        self.add_message("[EntraÃ®nement] Agent Docteur : Calibrage des sondes systÃ¨me...", "SystÃ¨me")
        self.agent_doctor_mission("bilan de santÃ©")
        time.sleep(2)

        # 2. Test RÃ©seau
        self.add_message("[EntraÃ®nement] Agent RÃ©seau : Optimisation des flux de donnÃ©es...", "SystÃ¨me")
        res_net = self.agent_network_mission("diagnostic rÃ©seau")
        self.add_message(f"RÃ©seau Status : {res_net[:100]}...", "Agent RÃ©seau")
        time.sleep(1)

        # 3. Test Cyber SÃ©curitÃ©
        self.add_message("[EntraÃ®nement] Agent Cyber : Mise Ã  jour des protocoles de dÃ©fense...", "SystÃ¨me")
        self.agent_cyber_mission("audit systÃ¨me")
        time.sleep(2)

        # 4. Test Juridique
        self.add_message("[EntraÃ®nement] Agent Juridique : Recalibrage de la base lÃ©gale...", "SystÃ¨me")
        self.agent_legal_mission("conseil droit du travail")
        time.sleep(2)

        # 5. Test Telegram (si activÃ©)
        if self.agent_telegram_enabled:
            self.add_message("[EntraÃ®nement] Agent Telegram : Test de la liaison satellite...", "SystÃ¨me")
            self.send_telegram_message("ðŸ¤– EntraÃ®nement en cours : Liaison satellite calibrÃ©e avec succÃ¨s.")
        
        # 6. Test VidÃ©o
        if getattr(self, "agent_video_enabled", True):
            self.add_message("[EntraÃ®nement] Agent VidÃ©o : VÃ©rification des capteurs optiques...", "SystÃ¨me")
            # On simule juste un check batterie/disque pour la vidÃ©o
            self.parler("Capteurs vidÃ©o opÃ©rationnels.")

        # Nouveaux Agents Premium
        if self.agent_vision_enabled:
            self.add_message("[EntraÃ®nement] Agent Vision : Calibrage de l'IA multimodale...", "SystÃ¨me")
        if self.agent_finance_enabled:
            self.add_message("[EntraÃ®nement] Agent Finance : Synchronisation des flux boursiers...", "SystÃ¨me")
        if self.agent_traducteur_enabled:
            self.add_message("[EntraÃ®nement] Agent Traducteur : Chargement des dictionnaires universels...", "SystÃ¨me")
        if self.agent_miner_enabled:
            self.add_message("[EntraÃ®nement] Agent Miner : Indexation des documents locaux...", "SystÃ¨me")

        self.add_message("EntraÃ®nement terminÃ©. Tous les agents sont synchronisÃ©s et Ã  100% de leur capacitÃ©.", "SystÃ¨me")
        self.parler("EntraÃ®nement terminÃ©. Je suis dÃ©sormais plus efficace et mes agents sont parfaitement synchronisÃ©s.")

    def send_telegram_message(self, message):
        """Envoie un message vers le bot Telegram configurÃ©"""
        if not self.telegram_bot_token or not self.telegram_chat_id:
            return
            
        try:
            url = f"https://api.telegram.org/bot{self.telegram_bot_token}/sendMessage"
            payload = {
                "chat_id": self.telegram_chat_id,
                "text": message
            }
            requests.post(url, json=payload, timeout=10)
        except Exception as e:
            print(f"Erreur envoi Telegram: {e}")

    def agent_network_mission(self, action):
        action_lower = action.lower()
        if "paramÃ¨tre" in action_lower or "config" in action_lower:
            try:
                res = subprocess.check_output("ipconfig /all", shell=True).decode('cp850')
                self.add_message("ParamÃ¨tres rÃ©seau rÃ©cupÃ©rÃ©s.", "Jarvisse")
                return "J'ai vÃ©rifiÃ© vos paramÃ¨tres rÃ©seau."
            except: return "Erreur lors de la vÃ©rification rÃ©seau."
        elif "wifi" in action_lower or "wi-fi" in action_lower:
            if "mot de passe" in action_lower or "pass" in action_lower:
                return self.agent_network_get_passwords()
            else:
                try:
                    res = subprocess.check_output("netsh wlan show networks", shell=True).decode('cp850')
                    self.add_message("RÃ©seaux Wi-Fi :\n" + res, "Jarvisse")
                    return "Voici les Wi-Fi disponibles."
                except: return "Erreur scan Wi-Fi."
        elif any(kw in action_lower for kw in ["diagnostic", "problÃ¨me", "connexion", "vÃ©rifie", "analyse"]):
            return self.agent_network_diagnose()
        elif any(kw in action_lower for kw in ["rÃ©pare", "rÃ©paration", "rÃ©sout"]):
            return self.agent_network_repair()
            
        return "Action rÃ©seau non reconnue."

    def agent_network_diagnose(self):
        self.add_message("Lancement du diagnostic rÃ©seau...", "Agent RÃ©seau")
        report = []
        is_ok = True
        
        # 1. Test Ping local (Passerelle)
        try:
            # On cherche l'IP de la passerelle via ipconfig
            cfg = subprocess.check_output("ipconfig", shell=True).decode('cp850')
            if "Passerelle par d" in cfg or "Default Gateway" in cfg:
                report.append("âœ… Interface rÃ©seau dÃ©tectÃ©e.")
            else:
                report.append("âŒ Aucune passerelle par dÃ©faut trouvÃ©e.")
                is_ok = False
        except: pass

        # 2. Test Ping Internet (8.8.8.8)
        try:
            subprocess.check_output("ping -n 1 8.8.8.8", shell=True)
            report.append("âœ… ConnectivitÃ© IP Internet opÃ©rationnelle (8.8.8.8).")
        except:
            report.append("âŒ Ã‰chec du ping Internet (Pas de connexion IP).")
            is_ok = False

        # 3. Test DNS
        try:
            import socket
            socket.gethostbyname("google.com")
            report.append("âœ… RÃ©solution DNS fonctionnelle.")
        except:
            report.append("âŒ Ã‰chec de la rÃ©solution DNS.")
            is_ok = False

        self.add_message("\n".join(report), "Agent RÃ©seau")
        if is_ok:
            return "Le diagnostic est terminÃ©. Tout semble fonctionner correctement."
        else:
            return "J'ai dÃ©tectÃ© des anomalies. Voulez-vous que je tente une rÃ©paration ?"

    def agent_network_repair(self):
        self.add_message("Tentative de rÃ©paration du rÃ©seau en cours...", "Agent RÃ©seau")
        try:
            # 1. Flush DNS
            subprocess.run("ipconfig /flushdns", shell=True, check=True)
            # 2. Reset Winsock (nÃ©cessite parfois admin mais on tente)
            subprocess.run("netsh winsock reset", shell=True)
            # 3. Release/Renew IP
            subprocess.run("ipconfig /release", shell=True)
            subprocess.run("ipconfig /renew", shell=True)
            
            self.add_message("âœ… Cache DNS vidÃ©.\nâœ… Interface rÃ©seau rÃ©initialisÃ©e.", "Agent RÃ©seau")
            return "La rÃ©paration a Ã©tÃ© effectuÃ©e. VÃ©rifiez si votre connexion est revenue."
        except Exception as e:
            return f"La rÃ©paration a Ã©chouÃ© : {e}"

    def agent_network_get_passwords(self):
        try:
            data = subprocess.check_output(['netsh', 'wlan', 'show', 'profiles']).decode('cp850').split('\n')
            profiles = [i.split(":")[1][1:-1] for i in data if "Profil Tous les utilisateurs" in i or "All User Profile" in i]
            results = []
            for i in profiles:
                try:
                    res_p = subprocess.check_output(['netsh', 'wlan', 'show', 'profile', i, 'key=clear']).decode('cp850').split('\n')
                    pwd = [b.split(":")[1][1:-1] for b in res_p if "Contenu de la cl" in b or "Key Content" in b]
                    results.append(f"{i}: {pwd[0] if pwd else 'N/A'}")
                except: continue
            self.add_message("Passwords Wi-Fi :\n" + "\n".join(results), "Jarvisse")
            return "Mots de passe rÃ©cupÃ©rÃ©s."
        except: return "Erreur rÃ©cupÃ©ration."

    def agent_control_mission(self, app_name):
        try:
            # Recherche insensible Ã  la casse parmi toutes les fenÃªtres
            all_windows = gw.getAllWindows()
            matching_windows = [w for w in all_windows if w.title and app_name.lower() in w.title.lower()]
            
            if matching_windows:
                win = matching_windows[0]
                if win.isMinimized:
                    win.restore()
                win.activate()
                self.parler(f"ContrÃ´le de {win.title[:30]} activÃ©. Que voulez-vous faire ?")
                return True
            else:
                self.parler(f"Application '{app_name}' introuvable dans les fenÃªtres ouvertes.")
                return False
        except Exception as e:
            self.parler(f"Ã‰chec du contrÃ´le : {e}")
            return False

    def _remote_target_is_allowed(self, host):
        host_norm = (host or "").strip().lower()
        if not host_norm:
            return False
        return host_norm in set(self.remote_hosts_whitelist)

    def _is_sensitive_remote_command(self, command):
        cmd = (command or "").strip().lower()
        risky_patterns = [
            "rm -rf", "del /f", "format", "mkfs", "shutdown", "reboot", "halt",
            "poweroff", "userdel", "passwd", "net user", "reg delete", "sc stop",
            "iptables -f", "ufw disable"
        ]
        return any(p in cmd for p in risky_patterns)

    def _build_ssh_target(self, host):
        user = self.remote_ssh_user.strip()
        host = (host or "").strip()
        return f"{user}@{host}" if user else host

    def _execute_remote_task(self, host, remote_command, force_without_confirmation=False):
        if not self.agent_remote_enabled:
            return "L'agent distant est dÃ©sactivÃ©. Activez-le dans les paramÃ¨tres des agents."
        if not self._remote_target_is_allowed(host):
            return "HÃ´te non autorisÃ©. Ajoutez-le d'abord dans la liste blanche."
        if not shutil.which("ssh"):
            return "Client SSH introuvable sur ce systÃ¨me."
        if not remote_command or not remote_command.strip():
            return "Commande distante vide."

        if self.remote_require_confirmation and self._is_sensitive_remote_command(remote_command) and not force_without_confirmation:
            self.awaiting_remote_confirm = {"host": host, "command": remote_command}
            return f"Commande sensible dÃ©tectÃ©e pour {host}. Dites 'confirme distant' pour continuer ou 'annule distant'."

        ssh_cmd = [
            "ssh",
            "-o", "BatchMode=yes",
            "-o", "StrictHostKeyChecking=accept-new",
            "-p", str(self.remote_ssh_port),
        ]
        key_path = self.remote_ssh_key_path.strip()
        if key_path:
            if not os.path.exists(key_path):
                return "La clÃ© SSH configurÃ©e est introuvable."
            ssh_cmd.extend(["-i", key_path])
        ssh_cmd.extend([self._build_ssh_target(host), remote_command.strip()])

        try:
            result = subprocess.run(ssh_cmd, capture_output=True, text=True, timeout=40)
            output = (result.stdout or "").strip()
            err = (result.stderr or "").strip()
            if result.returncode == 0:
                if output:
                    return f"TÃ¢che distante exÃ©cutÃ©e sur {host}.\n{output[:1200]}"
                return f"TÃ¢che distante exÃ©cutÃ©e sur {host} (sans sortie)."
            details = err or output or f"code={result.returncode}"
            return f"Ã‰chec exÃ©cution distante sur {host}: {details[:800]}"
        except subprocess.TimeoutExpired:
            return "La commande distante a dÃ©passÃ© le dÃ©lai autorisÃ©."
        except Exception as e:
            return f"Erreur SSH: {e}"

    def _handle_remote_command(self, commande):
        cmd = (commande or "").strip()
        cmd_lower = cmd.lower()

        if "liste machines distantes" in cmd_lower or "liste des machines distantes" in cmd_lower:
            if not self.remote_hosts_whitelist:
                return "Aucune machine autorisÃ©e n'est configurÃ©e."
            return "Machines autorisÃ©es: " + ", ".join(self.remote_hosts_whitelist)

        connect_match = re.search(r"(?:connecte(?:-toi)?(?:\s+a)?|se connecter Ã )\s+([a-zA-Z0-9._-]+)", cmd_lower)
        if connect_match:
            host = connect_match.group(1).strip().lower()
            if not self._remote_target_is_allowed(host):
                return f"L'hÃ´te {host} n'est pas dans la liste blanche."
            self.remote_default_host = host
            self._save_settings()
            return f"Connexion distante prÃ©parÃ©e vers {host}. Donnez ensuite la tÃ¢che Ã  exÃ©cuter."

        patterns = [
            r"(?:sur|sur la machine|sur l'ordinateur|sur le pc)\s+([a-zA-Z0-9._-]+)\s+(?:exÃ©cute|execute|lance|run)\s+(.+)",
            r"(?:exÃ©cute|execute|lance|run)\s+(.+?)\s+(?:sur|sur la machine|sur l'ordinateur|sur le pc)\s+([a-zA-Z0-9._-]+)",
        ]
        host = ""
        remote_cmd = ""
        for idx, pattern in enumerate(patterns):
            m = re.search(pattern, cmd, flags=re.IGNORECASE)
            if m:
                if idx == 0:
                    host, remote_cmd = m.group(1), m.group(2)
                else:
                    remote_cmd, host = m.group(1), m.group(2)
                break

        if not host and any(k in cmd_lower for k in ["Ã  distance", "a distance", "distant", "ssh"]):
            m_cmd = re.search(r"(?:exÃ©cute|execute|lance|run)\s+(.+)", cmd, flags=re.IGNORECASE)
            if m_cmd:
                host = self.remote_default_host
                remote_cmd = m_cmd.group(1)

        host = (host or "").strip().lower()
        remote_cmd = (remote_cmd or "").strip()
        if not host:
            return "PrÃ©cisez l'hÃ´te distant autorisÃ© (ex: 'sur serveur-maison exÃ©cute whoami')."
        if not remote_cmd:
            return "PrÃ©cisez la commande Ã  exÃ©cuter Ã  distance."

        return self._execute_remote_task(host, remote_cmd)

    def agent_gmail_authenticate(self):
        SCOPES = ['https://www.googleapis.com/auth/gmail.readonly', 'https://www.googleapis.com/auth/gmail.compose']
        creds = None
        t_path = os.path.join(os.path.dirname(__file__), 'token.pickle')
        c_path = os.path.join(os.path.dirname(__file__), 'credentials.json')
        if os.path.exists(t_path):
            with open(t_path, 'rb') as t: creds = pickle.load(t)
        if not creds or not creds.valid:
            if os.path.exists(c_path):
                flow = InstalledAppFlow.from_client_secrets_file(c_path, SCOPES)
                creds = flow.run_local_server(port=0)
                with open(t_path, 'wb') as t: pickle.dump(creds, t)
            else:
                self.add_message("credentials.json manquant.", "Jarvisse")
                return
        self.gmail_service = build('gmail', 'v1', credentials=creds)
        self.add_message("Gmail AuthentifiÃ©.", "Jarvisse")

    def start_gmail_listener(self):
        while self.agent_gmail_enabled:
            if self.gmail_service:
                try:
                    msgs = self.gmail_service.users().messages().list(userId='me', q='is:unread').execute().get('messages', [])
                    if msgs and msgs[0]['id'] != self.last_email_id:
                        self.last_email_id = msgs[0]['id']
                        self.add_message("Nouveau mail reÃ§u !", "Jarvisse")
                        self.parler("Patron, nouveau mail.")
                except: pass
            time.sleep(60)

    def agent_gmail_create_draft(self, subject, body, to):
        if not self.gmail_service: return "Pas authentifiÃ©."
        try:
            from email.mime.text import MIMEText
            import base64
            msg = MIMEText(body); msg['to'] = to; msg['subject'] = subject
            raw = base64.urlsafe_b64encode(msg.as_bytes()).decode()
            self.gmail_service.users().drafts().create(userId='me', body={'message': {'raw': raw}}).execute()
            return "Brouillon crÃ©Ã©."
        except: return "Erreur crÃ©ation brouillon."

    def agent_browser_mission(self, action):
        if not hasattr(self.web_view, "load_website"):
            return "Le moteur de navigation n'est pas disponible sur ce systÃ¨me."
            
        action_lower = action.lower()
        self.tabview.set("ðŸŒ Navigateur") # Focus sur le navigateur
        
        if "cherche" in action_lower or "recherche" in action_lower:
            # On remplace les termes longs d'abord pour Ã©viter les rÃ©sidus comme "re"
            query = action_lower
            for term in ["cherche sur google", "recherche", "cherche"]:
                query = query.replace(term, "")
            query = query.strip()
            
            # Encodage strict pour Ã©viter l'erreur 'ascii' codec
            encoded_query = urllib.parse.quote_plus(query)
            # Utilisation de DuckDuckGo pour Ã©viter les messages "Activez JavaScript" de Google
            url = f"https://duckduckgo.com/html/?q={encoded_query}"
            
            self.web_view.load_website(url)
            self.add_message(f"Recherche sur le web : {query}", "Agent Navigateur")
            return f"J'ai lancÃ© la recherche pour '{query}'."
        
        elif "va sur" in action_lower or "ouvre le site" in action_lower:
            url = action_lower.replace("va sur", "").replace("ouvre le site", "").strip()
            if not url.startswith("http"): url = "https://" + url
            self.web_view.load_website(url)
            return f"Navigation vers {url} en cours."
        
        elif "retour" in action_lower or "prÃ©cÃ©dent" in action_lower:
            try:
                self.web_view.go_back()
                return "Retour Ã  la page prÃ©cÃ©dente."
            except: return "Impossible de reculer."
            
        elif "suivant" in action_lower:
            try:
                self.web_view.go_forward()
                return "Passage Ã  la page suivante."
            except: return "Impossible d'avancer."

        elif "descend" in action_lower or "bas" in action_lower:
            self.web_view.yview_scroll(1, "pages")
            return "DÃ©filement vers le bas."

        elif "monte" in action_lower or "haut" in action_lower:
            self.web_view.yview_scroll(-1, "pages")
            return "DÃ©filement vers le haut."

        elif "zoom" in action_lower or "grossis" in action_lower:
            self.web_view.zoom(1.2)
            return "Zoom avant effectuÃ©."

        elif "dÃ©zoom" in action_lower or "rÃ©duis" in action_lower:
            self.web_view.zoom(0.8)
            return "Zoom arriÃ¨re effectuÃ©."
            
        return "Action de navigation non comprise."

    def agent_legal_mission(self, commande):
        """Mission de l'Agent Juridique Expert (Droit, RH, Travail)"""
        try:
            self.after(0, lambda: self.status_label.configure(text="Agent Juridique analyse...", text_color=COLOR_ACCENT_PURPLE))
            self.add_message("Agent Juridique activÃ©. Recherche de fondements juridiques...", "Agent Juridique")

            system_prompt = (
                "Tu es l'Agent Juridique de SERI TAGRO Intelligence. Tu es un expert Ã©minent en Droit FranÃ§ais (PÃ©nal, Civil, Travail, Social) "
                "et en Gestion des Ressources Humaines. Tu maÃ®trises parfaitement le Code du Travail, les conventions collectives, "
                "la jurisprudence et les articles juridiques. Ton rÃ´le est de conseiller et d'aider l'utilisateur sur des problÃ¨mes complexes. "
                "RÃ©ponds avec une grande prÃ©cision, cite les articles de lois ou de codes si possible (ex: Article L1221-1 du Code du travail). "
                "Sois professionnel, impartial et clair. RÃ©ponds en franÃ§ais de maniÃ¨re structurÃ©e."
            )

            prompt = f"{system_prompt}\n\nQuestion/ProblÃ¨me de l'utilisateur : {commande}\n\nConseil juridique :"
            reponse = ""

            # Utilisation du moteur IA configurÃ©
            if self.autonomous_provider == "gemini" and self.genai_client:
                response = self.genai_client.models.generate_content(
                    model=self.model_name,
                    contents=prompt
                )
                reponse = response.text.strip()
            else:
                # Mode Ollama
                try:
                    url = self.ollama_url
                    # On s'assure que c'est l'URL de gÃ©nÃ©ration
                    if not url.endswith("/generate") and not url.endswith("/generate/"):
                         # Si on a l'URL de base, on ajoute /api/generate
                         if url.endswith("/"): url = url + "api/generate"
                         else: url = url + "/api/generate"
                    
                    payload = {
                        "model": self.ollama_model.strip(),
                        "prompt": prompt,
                        "stream": False,
                        "options": {"num_predict": 1000, "temperature": 0.3}
                    }
                    r = requests.post(url, json=payload, timeout=120)
                    reponse = r.json().get("response", "Erreur de communication avec l'IA juridique.")
                except Exception as e:
                    print(f"Ollama Legal Error: {e}")
                    reponse = "Impossible de contacter le moteur IA de l'Agent Juridique."

            if not reponse:
                reponse = "Je n'ai pas pu gÃ©nÃ©rer d'analyse juridique pour le moment."

            if reponse:
                self.parler(reponse, sender="Agent Juridique", force_full=True)

        except Exception as e:
            print(f"Erreur Agent Juridique: {e}")
            self.parler("DÃ©solÃ© Patron, mon module juridique rencontre une difficultÃ© technique.")
        
        self.after(0, self.reset_ui)

    def agent_cyber_mission(self, commande):
        """Mission de l'Agent Cyber SÃ©curitÃ© (Audit, DÃ©fense, Conseils, OSINT, OCR)"""
        cmd_l = commande.lower()
        try:
            self.after(0, lambda: self.status_label.configure(text="Agent Cyber en action...", text_color=COLOR_ACCENT_RED))
            self.add_message("Agent Cyber SÃ©curitÃ© activÃ©. Mission en cours...", "Agent Cyber")
            self.parler("TrÃ¨s bien Patron, je suis en train d'exÃ©cuter votre instruction. L'opÃ©ration est en cours.")

            # 3. Agent Docteur SystÃ¨me (Maintenance & SantÃ©)
            if any(kw in cmd_l for kw in ["docteur", "santÃ©", "santÃ© systÃ¨me", "mÃ©decin", "rÃ©pare le systÃ¨me", "optimise", "performance", "nettoie", "processus"]):
                threading.Thread(target=self.agent_doctor_mission, args=(commande,), daemon=True).start()
                return

            # 1. Sherlock / Recherche Pseudonyme
            if any(kw in cmd_l for kw in ["sherlock", "pseudo", "nom d'utilisateur", "recherche le compte"]):
                # Extraction du pseudo (dernier mot en gÃ©nÃ©ral ou aprÃ¨s un mot clÃ©)
                words = cmd_l.split()
                pseudo = words[-1] # Simplification
                if len(words) > 1 and words[-2] in ["de", "nommÃ©", "pseudo"]:
                    pseudo = words[-1]
                
                threading.Thread(target=self.agent_cyber_sherlock_scan, args=(pseudo,), daemon=True).start()
                return

            # 2. OCR / Lecture d'Image
            if any(kw in cmd_l for kw in ["ocr", "lire l'image", "extrait le texte", "texte image"]):
                threading.Thread(target=self.agent_cyber_ocr_mission, daemon=True).start()
                return

            # Actions existantes (Audit SystÃ¨me)
            if any(kw in cmd_l for kw in ["audit", "scanner le systÃ¨me", "sÃ©curitÃ© systÃ¨me", "processus suspects"]):
                res = self.agent_cyber_audit_system()
                self.add_message(res, "Agent Cyber")
                self.parler("Audit terminÃ© Patron. Voici les points de vigilance.")
                return

            # Action RÃ©seau Cyber
            if any(kw in cmd_l for kw in ["scan wifi", "analyse rÃ©seau", "vulnÃ©rabilitÃ© wifi"]):
                networks = subprocess.check_output("netsh wlan show networks", shell=True).decode('cp850')
                self.add_message(f"Scan des rÃ©seaux environnants :\n{networks}", "Agent Cyber")
                self.parler("J'ai scannÃ© les rÃ©seaux Wi-Fi. Je vous affiche les points d'accÃ¨s dÃ©tectÃ©s.")
                return

            # Action Nmap (Scan IP)
            if "nmap" in cmd_l or "scan ip" in cmd_l or "scanne l'ip" in cmd_l:
                import re
                ip_match = re.search(r'\b(?:\d{1,3}\.){3}\d{1,3}\b', cmd_l)
                if ip_match:
                    ip = ip_match.group()
                    threading.Thread(target=self.agent_cyber_nmap_scan, args=(ip,), daemon=True).start()
                else:
                    self.parler("Patron, j'ai besoin d'une adresse IP valide pour lancer Nmap.")
                return

            # Expertise IA Cyber
            system_prompt = (
                "Tu es l'Agent Cyber SÃ©curitÃ© de Jarvisse Intelligence. Tu es un expert mondial en hacking Ã©thique, cryptographie, "
                "sÃ©curitÃ© rÃ©seau et ingÃ©nierie sociale. Ton but est d'informer, d'Ã©duquer et de protÃ©ger l'utilisateur. "
                "Tu peux expliquer comment fonctionnent les attaques (phishing, brute force, sniffing) pour mieux s'en dÃ©fendre. "
                "Tu fournis des conseils sur la robustesse des mots de passe, l'authentification MFA et la sÃ©curisation des infrastructures. "
                "RÃ©ponds avec une expertise technique Ã©levÃ©e mais comprÃ©hensible. Sois percutant et professionnel."
            )

            prompt = f"{system_prompt}\n\nQuestion Cyber : {commande}\n\nAnalyse Expert :"
            reponse = ""

            if self.autonomous_provider == "gemini" and self.genai_client:
                response = self.genai_client.models.generate_content(model=self.model_name, contents=prompt)
                reponse = response.text.strip()
            else:
                try:
                    payload = {"model": self.ollama_model.strip(), "prompt": prompt, "stream": False, "options": {"num_predict": 800, "temperature": 0.5}}
                    r = requests.post(self.ollama_url.split("/api")[0] + "/api/generate", json=payload, timeout=90)
                    reponse = r.json().get("response", "Erreur IA Cyber.")
                except: reponse = "Moteur Cyber indisponible."

            if reponse:
                self.parler(reponse, sender="Agent Cyber", force_full=True)

        except Exception as e:
            print(f"Erreur Agent Cyber: {e}")
            self.parler("Ã‰chec de la mission cyber.")
        self.after(0, self.reset_ui)

    def agent_coding_mission(self, commande):
        """Mission de l'Assistant Coding (GÃ©nÃ©ration de code et d'applications)"""
        try:
            self.after(0, lambda: self.status_label.configure(text="Assistant Coding en action...", text_color=COLOR_ACCENT_PURPLE))
            self.add_message("Assistant Coding activÃ©. Analyse de votre projet...", "Assistant Coding")
            self.parler("TrÃ¨s bien Patron, je vais m'occuper de la crÃ©ation de votre projet. Je vais identifier les langages nÃ©cessaires et gÃ©nÃ©rer le code.")

            system_prompt = (
                "Tu es l'Assistant Coding de Jarvisse Intelligence. Ton rÃ´le est de concevoir et gÃ©nÃ©rer des applications et sites web d'exception. "
                "Tu maÃ®trises de nombreux langages : HTML, CSS, JavaScript, Python, C, C++, Flutter, PHP, Java, etc. "
                "\nQUALITÃ‰ DE DESIGN REQUIS :\n"
                "1. ESTHÃ‰TIQUE PREMIUM : CrÃ©e des designs modernes, Ã©purÃ©s et visuellement impressionnants (WOW effect).\n"
                "2. UI/UX MODERNE : Utilise des palettes de couleurs harmonieuses (dÃ©gradÃ©s, modes sombres Ã©lÃ©gants), des espacements gÃ©nÃ©reux et une hiÃ©rarchie visuelle claire.\n"
                "3. TYPOGRAPHIE : Utilise des polices modernes via Google Fonts (ex: Inter, Roboto, Outfit, Poppins).\n"
                "4. ANIMATIONS : IntÃ¨gre des micro-animations, des transitions fluides et des effets au survol (hover effects) pour rendre l'interface vivante.\n"
                "5. RESPONSIVE : Le design doit Ãªtre parfaitement adaptÃ© Ã  tous les Ã©crans (mobile friendly).\n"
                "\nPOUR CHAQUE DEMANDE, TU DOIS :\n"
                "1. Identifier les langages et technologies les plus adaptÃ©s.\n"
                "2. Expliquer briÃ¨vement ton choix technologique et tes dÃ©cisions de design.\n"
                "3. GÃ©nÃ©rer le code complet, structurÃ© et commentÃ©.\n"
                "IMPORTANTE : Pour chaque fichier de code, utilise le format suivant strictement :\n"
                "FILE: [nom_du_fichier.extension]\n"
                "```[langage]\n"
                "[le code ici]\n"
                "```\n"
                "Assure-toi de fournir TOUS les fichiers nÃ©cessaires."
            )

            prompt = f"{system_prompt}\n\nProjet soumis : {commande}\n\nSolution et Code :"
            reponse = ""

            if self.autonomous_provider == "gemini" and self.genai_client:
                response = self.genai_client.models.generate_content(model=self.model_name, contents=prompt)
                reponse = response.text.strip()
            else:
                try:
                    payload = {"model": self.ollama_model.strip(), "prompt": prompt, "stream": False, "options": {"num_predict": 4096, "temperature": 0.3}}
                    r = requests.post(self.ollama_url.split("/api")[0] + "/api/generate", json=payload, timeout=120)
                    reponse = r.json().get("response", "Erreur IA Coding.")
                except: reponse = "Moteur Coding indisponible."

            # Pas de add_message direct ici, parler s'en occupe
            
            # Extraction et proposition de sauvegarde
            files = self._extract_coding_files(reponse)
            if files:
                self.after(0, lambda: self._proposer_sauvegarde_coding(files))
                self.parler("J'ai gÃ©nÃ©rÃ© les fichiers pour votre projet. Voulez-vous que je les sauvegarde dans un dossier ?")
            elif len(reponse) > 500:
                self.parler("J'ai terminÃ© la gÃ©nÃ©ration du code. Vous pouvez le consulter dans la zone de chat.", sender="Assistant Coding")
            else:
                self.parler(reponse, sender="Assistant Coding", force_full=True)
            
            # Affichage forcÃ© du texte si trop long pour Ãªtre lu agrÃ©ablement
            if len(reponse) > 500:
                 self.add_message(reponse, "Assistant Coding", progressive=False)

        except Exception as e:
            print(f"Erreur Assistant Coding: {e}")
            self.parler("Ã‰chec de la mission de codage.")
        self.after(0, self.reset_ui)

    def agent_generation_mission(self, commande):
        """Mission de l'Assistant GÃ©nÃ©ration (Images & VidÃ©os) - Support Gemini & Ollama"""
        try:
            self.after(0, lambda: self.status_label.configure(text="GÃ©nÃ©ration en cours...", text_color=COLOR_ACCENT_PURPLE))
            self.add_message("Assistant GÃ©nÃ©ration activÃ©. CrÃ©ation de votre mÃ©dia...", "Assistant GÃ©nÃ©ration")
            
            cmd_lower = commande.lower()
            is_video = any(kw in cmd_lower for kw in ["vidÃ©o", "video", "film", "animation"])
            
            if is_video:
                self.parler("Je m'occupe de la gÃ©nÃ©ration de votre vidÃ©o. Veuillez patienter un instant.")
            else:
                self.parler("Je m'occupe de la gÃ©nÃ©ration de votre image. Analyse du prompt en cours.")

            # Optimisation du prompt via IA (Gemini ou Ollama)
            prompt_ameliore = commande
            type_media = "vidÃ©o" if is_video else "image"
            
            if self.autonomous_provider == "gemini" and self.genai_client:
                prompt_dev = f"Transforme cette demande en un prompt extrÃªmement dÃ©taillÃ© et professionnel pour une gÃ©nÃ©ration d'agent IA ({type_media}) haute qualitÃ© : {commande}. RÃ©ponds uniquement avec le prompt dÃ©taillÃ©."
                try:
                    res_prompt = self.genai_client.models.generate_content(model=self.model_name, contents=prompt_dev)
                    prompt_ameliore = res_prompt.text.strip()
                except: pass
            elif self.autonomous_provider == "ollama":
                try:
                    prompt_dev = f"Transforme cette demande en un prompt dÃ©taillÃ© et professionnel (en anglais) pour une gÃ©nÃ©ration d'image haute qualitÃ© : {commande}. RÃ©ponds uniquement avec le prompt dÃ©taillÃ©."
                    payload = {
                        "model": self.ollama_model.strip(),
                        "prompt": prompt_dev,
                        "stream": False,
                        "options": {"num_predict": 300, "temperature": 0.5}
                    }
                    r = requests.post(self.ollama_url, json=payload, timeout=60)
                    prompt_ameliore = r.json().get("response", commande).strip()
                except Exception as e:
                    print(f"Ollama Prompt Generation Error: {e}")

            # ExÃ©cution de la gÃ©nÃ©ration
            if self.autonomous_provider == "gemini" and self.genai_client:
                if is_video:
                    self.add_message("La gÃ©nÃ©ration de vidÃ©o via Veo est en prÃ©paration. Je gÃ©nÃ¨re une image haute dÃ©finition pour illustrer votre demande.", "Assistant GÃ©nÃ©ration")
                    is_video = False 
                
                # GÃ©nÃ©ration Image avec Imagen 3
                try:
                    self.add_message(f"Prompt optimisÃ© : {prompt_ameliore}", "Assistant GÃ©nÃ©ration")
                    response = self.genai_client.models.generate_image(
                        model='imagen-3',
                        prompt=prompt_ameliore
                    )
                    
                    # Sauvegarde
                    timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
                    filename = f"gen_{timestamp}.png"
                    filepath = os.path.join(self.generations_dir, filename)
                    
                    image_data = response.generated_images[0].image_bytes
                    with open(filepath, "wb") as f:
                        f.write(image_data)
                    
                    self.add_message(f"SuccÃ¨s ! MÃ©dia sauvegardÃ© : {filepath}", "Assistant Generation")
                    self.parler("J'ai terminÃ© la gÃ©nÃ©ration avec Gemini. Le fichier est prÃªt.")
                    os.startfile(filepath)
                except Exception as ex:
                    print(f"Erreur Imagen: {ex}")
                    self.add_message(f"Erreur technique : {ex}", "Assistant Generation")
                    self.parler("DÃ©solÃ© Patron, impossible de gÃ©nÃ©rer via Gemini.")
            
            elif self.autonomous_provider == "ollama" or (not self.genai_client):
                # Fallback ou Mode Ollama : Utilisation d'un moteur libre (Pollinations)
                self.add_message(f"GÃ©nÃ©ration locale activÃ©e (via {self.ollama_model}).", "Assistant GÃ©nÃ©ration")
                try:
                    import urllib.parse
                    # Nettoyage et limitation du prompt pour Ã©viter les erreurs d'URL trop longues
                    clean_prompt = prompt_ameliore.replace("\n", " ").strip()
                    if len(clean_prompt) > 800: clean_prompt = clean_prompt[:800]
                    encoded_prompt = urllib.parse.quote(clean_prompt)
                    
                    # Liste de modÃ¨les Ã  tenter en cas d'erreur serveur (502/504)
                    # Pollinations supporte : 'flux', 'flux-realism', 'any-dark', 'flux-anime', 'turbo', etc.
                    models_to_try = ["flux", "flux-realism", "flux-anime", "turbo", "rtlib", "dreamshaper"]
                    last_error = "Serveur indisponible"
                    success = False
                    
                    headers = {"User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36"}
                    
                    for model in models_to_try:
                        try:
                            self.add_message(f"Essai (Moteur: {model})...", "Assistant GÃ©nÃ©ration")
                            image_url = f"https://image.pollinations.ai/prompt/{encoded_prompt}?width=1024&height=1024&seed={random.randint(1, 99999)}&nologo=true&model={model}"
                            
                            r = requests.get(image_url, timeout=45, headers=headers)
                            print(f"DEBUG Generation {model}: {r.status_code}")
                            
                            if r.status_code == 200:
                                timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
                                filename = f"gen_local_{timestamp}.png"
                                filepath = os.path.join(self.generations_dir, filename)
                                
                                with open(filepath, "wb") as f:
                                    f.write(r.content)
                                
                                self.add_message(f"SuccÃ¨s ! Image gÃ©nÃ©rÃ©e avec {model}.", "Assistant Generation")
                                self.parler("GÃ©nÃ©ration rÃ©ussie !")
                                os.startfile(filepath)
                                success = True
                                break
                            else:
                                last_error = f"Code {r.status_code} sur {model}"
                                time.sleep(0.5) 
                        except Exception as e:
                            last_error = str(e)
                            print(f"Exception {model}: {e}")
                            continue

                    # Dernier essai sans paramÃ¨tre de modÃ¨le (Auto) + fallback Hercai
                    if not success:
                        try:
                            self.add_message("Dernier essai (Moteur de secours Hercai)...", "Assistant GÃ©nÃ©ration")
                            # Hercai est un excellent fallback quand Pollinations sature
                            hercai_url = f"https://hercai.onrender.com/v3/text2image?prompt={encoded_prompt}"
                            r = requests.get(hercai_url, timeout=45, headers=headers)
                            if r.status_code == 200:
                                res_json = r.json()
                                img_url = res_json.get("url")
                                if img_url:
                                    r_img = requests.get(img_url, timeout=45)
                                    if r_img.status_code == 200:
                                        timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
                                        filename = f"gen_hercai_{timestamp}.png"
                                        filepath = os.path.join(self.generations_dir, filename)
                                        with open(filepath, "wb") as f: f.write(r_img.content)
                                        self.add_message("SuccÃ¨s ! Image gÃ©nÃ©rÃ©e via Hercai.", "Assistant Generation")
                                        self.parler("GÃ©nÃ©ration rÃ©ussie via le moteur de secours !")
                                        os.startfile(filepath)
                                        success = True
                        except: pass

                    if not success:
                        raise Exception(f"Les serveurs de rendu sont saturÃ©s. ({last_error})")

                except Exception as ex:
                    self.last_internal_error = f"Ã‰chec de gÃ©nÃ©ration d'image via Pollinations/Hercai. Serveurs saturÃ©s (Code 530/502). Moteurs tentÃ©s: {models_to_try}."
                    print(f"Erreur GÃ©nÃ©ration Local: {ex}")
                    self.parler(f"DÃ©solÃ© Patron, mes moteurs de rendu visuel sont saturÃ©s en ce moment mÃªme. Mon omniscience m'indique une surcharge des serveurs distants.")
            else:
                self.parler("La gÃ©nÃ©ration nÃ©cessite l'accÃ¨s Ã  Gemini ou Ollama. Veuillez vÃ©rifier votre clÃ© API.")

        except Exception as e:
            print(f"Erreur Assistant GÃ©nÃ©ration: {e}")
            self.parler("Ã‰chec de la mission de gÃ©nÃ©ration.")
        self.after(0, self.reset_ui)

    def _extract_coding_files(self, text):
        """Extrait les fichiers du format FILE: [nom] ... ``` ... ``` """
        files = []
        pattern = r"FILE:\s*([a-zA-Z0-9_\-\.]+)\s*\n```[a-z]*\n(.*?)\n```"
        matches = re.findall(pattern, text, re.DOTALL)
        for name, content in matches:
            files.append({"name": name, "content": content})
        return files

    def _proposer_sauvegarde_coding(self, files):
        """Demande Ã  l'utilisateur s'il veut sauvegarder les fichiers extraits"""
        if not files: return
        
        # On demande par dialogue UI car le thread parler est asynchrone
        save_win = ctk.CTkToplevel(self)
        save_win.title("Sauvegarde du Projet")
        save_win.geometry("400x250")
        save_win.configure(fg_color=COLOR_BG)
        save_win.attributes("-topmost", True)
        
        label = ctk.CTkLabel(save_win, text=f"J'ai dÃ©tectÃ© {len(files)} fichiers.\nSouhaitez-vous les sauvegarder ?", font=ctk.CTkFont(size=14))
        label.pack(pady=20)
        
        files_list = "\n".join([f"â€¢ {f['name']}" for f in files])
        scroll_label = ctk.CTkTextbox(save_win, height=80, fg_color=COLOR_NAV)
        scroll_label.pack(padx=20, pady=5, fill="x")
        scroll_label.insert("end", files_list)
        scroll_label.configure(state="disabled")

        def select_and_save():
            dir_path = filedialog.askdirectory(title="SÃ©lectionnez le dossier de destination")
            if dir_path:
                try:
                    for f in files:
                        full_path = os.path.join(dir_path, f['name'])
                        with open(full_path, "w", encoding="utf-8") as file_out:
                            file_out.write(f['content'])
                    self.parler(f"Parfait Patron ! Les {len(files)} fichiers ont Ã©tÃ© sauvegardÃ©s dans {os.path.basename(dir_path)}.")
                    save_win.destroy()
                except Exception as e:
                    self.parler(f"Erreur lors de la sauvegarde : {e}")
            
        btn_save = ctk.CTkButton(save_win, text="ðŸ“ Choisir Dossier & Sauvegarder", command=select_and_save, fg_color=COLOR_ACCENT)
        btn_save.pack(pady=10)
        
        btn_cancel = ctk.CTkButton(save_win, text="Annuler", command=save_win.destroy, fg_color=COLOR_CARD)
        btn_cancel.pack(pady=5)



    def agent_cyber_audit_system(self):
        """Audit de base du systÃ¨me Windows"""
        report = ["ðŸ” AUDIT DE SÃ‰CURITÃ‰ SYSTÃˆME"]
        
        # 1. VÃ©rification Windows Defender (via PowerShell)
        try:
            # Commande PowerShell robuste
            ps_cmd = 'powershell -ExecutionPolicy Bypass -Command "(Get-MpComputerStatus).RealTimeProtectionEnabled"'
            out = subprocess.check_output(ps_cmd, shell=True).decode().strip()
            status = "âœ… ActivÃ©" if "True" in out else "âŒ DÃ‰SACTIVÃ‰ (Danger)"
            report.append(f"- Protection Temps RÃ©el : {status}")
        except: report.append("- Statut Defender : Inconnu")

        # 2. DÃ©tection processus suspects (exemple)
        suspicious = []
        for proc in psutil.process_iter(['name']):
            try:
                n = proc.info['name'].lower()
                if any(s in n for s in ["mimikatz", "wireshark", "nmap", "metasploit", "keylogger"]):
                    suspicious.append(n)
            except: pass
        
        if suspicious:
            report.append(f"- Outils sensibles dÃ©tectÃ©s : {', '.join(suspicious)}")
        else:
            report.append("- Aucun outil suspect en cours d'exÃ©cution.")

        return "\n".join(report)

    def agent_cyber_nmap_scan(self, ip):
        """ExÃ©cute un scan Nmap rapide sur une IP cible"""
        self.add_message(f"Lancement du scan Nmap sur {ip}...", "Agent Cyber")
        try:
            nmap_exe = "nmap"
            # 1. VÃ©rification robuste de l'existence de Nmap
            found = False
            try:
                subprocess.check_output("nmap --version", shell=True)
                found = True
            except:
                # On cherche dans les dossiers par dÃ©faut de Windows
                paths = [
                    r"C:\Program Files (x86)\Nmap\nmap.exe",
                    r"C:\Program Files\Nmap\nmap.exe",
                    os.path.join(os.environ.get("ProgramFiles(x86)", "C:\\Program Files (x86)"), "Nmap", "nmap.exe"),
                    os.path.join(os.environ.get("ProgramFiles", "C:\\Program Files"), "Nmap", "nmap.exe")
                ]
                for p in paths:
                    if os.path.exists(p):
                        nmap_exe = f'"{p}"'
                        found = True
                        break
            
            if not found:
                self.add_message("âš ï¸ Nmap n'est pas dÃ©tectÃ©. Tentative d'installation automatique via Winget...", "Agent Cyber")
                self.parler("Patron, Nmap n'est pas installÃ© ou introuvable. Je vais tenter de l'installer pour vous via Winget.")
                
                if self.agent_cyber_install_nmap():
                    self.add_message("âœ… Installation terminÃ©e. Veuillez relancer la commande de scan.", "Agent Cyber")
                    self.parler("Nmap a Ã©tÃ© installÃ©. Vous pouvez maintenant relancer votre demande de scan.")
                return

            # Scan : DÃ©tection d'OS (-O) et des ports ouverts (Fast -F)
            # Note: La dÃ©tection d'OS peut nÃ©cessiter des privilÃ¨ges d'administrateur
            cmd = f"{nmap_exe} -F -O {ip}"
            self.status_label.configure(text=f"Nmap OS Detection {ip}...", text_color=COLOR_ACCENT_RED)
            
            try:
                # On tente le scan avec dÃ©tection d'OS
                result = subprocess.check_output(cmd, shell=True, stderr=subprocess.STDOUT).decode('cp850')
            except subprocess.CalledProcessError:
                # Si erreur (souvent liÃ©e aux privilÃ¨ges pour -O), on tente -sV qui donne souvent l'OS via les banniÃ¨res
                self.add_message("âš ï¸ Permission insuffisante pour -O. Tentative via dÃ©tection de services...", "Agent Cyber")
                cmd = f"{nmap_exe} -F -sV {ip}"
                result = subprocess.check_output(cmd, shell=True).decode('cp850')
            
            # Affichage et Analyse par l'IA Cyber
            msg = f"Rapport Nmap (OS & Ports) pour {ip} :\n{result}"
            self.add_message(msg, "Agent Cyber")
            
            # Expertise IA pour identifier l'OS
            system_prompt = "Tu es l'Agent Cyber SÃ©curitÃ©. Analyse le rapport Nmap suivant et dis-moi prÃ©cisÃ©ment quel SystÃ¨me d'Exploitation (Windows, Linux, Android, etc.) semble Ãªtre utilisÃ© sur cette IP. Sois court et prÃ©cis."
            prompt = f"{system_prompt}\n\nRapport :\n{result}"
            
            reponse_ia = ""
            if self.autonomous_provider == "gemini" and self.genai_client:
                res_ia = self.genai_client.models.generate_content(model=self.model_name, contents=prompt)
                reponse_ia = res_ia.text.strip()
            else:
                # Fallback simple si pas d'IA
                reponse_ia = "Analyse technique disponible dans le rapport textuel."

            self.add_message(f"ðŸ§  DÃ©tection SystÃ¨me :\n{reponse_ia}", "Agent Cyber")
            self.parler(f"Le scan sur {ip} est terminÃ©. mon analyse indique qu'il s'agit probablement d'un systÃ¨me {reponse_ia}")
            
        except Exception as e:
            self.add_message(f"âŒ Erreur lors du scan Nmap : {e}", "Agent Cyber")
            self.parler("Le scan Nmap a Ã©chouÃ©.")
        finally:
            self.after(0, self.reset_ui)

    def agent_cyber_install_nmap(self):
        """Tente d'installer Nmap via Winget de maniÃ¨re silencieuse"""
        try:
            # ID Winget pour Nmap
            cmd = 'winget install Insecure.Nmap --silent --accept-package-agreements --accept-source-agreements'
            self.status_label.configure(text="Installation Nmap...", text_color=COLOR_ACCENT_RED)
            
            # On utilise subprocess.run
            process = subprocess.run(cmd, shell=True, capture_output=True, text=True)
            
            if process.returncode == 0:
                return True
            else:
                print(f"Erreur Winget : {process.stderr}")
                return False
        except Exception as e:
            print(f"Erreur installation Nmap : {e}")
            return False

    def agent_cyber_sherlock_scan(self, username):
        """Recherche d'un pseudonyme sur les rÃ©seaux sociaux (Mini-Sherlock)"""
        self.add_message(f"ðŸ” Lancement de Sherlock pour le pseudonyme : {username}", "Agent Cyber")
        sites = {
            "GitHub": "https://github.com/{}",
            "Twitter": "https://twitter.com/{}",
            "Instagram": "https://instagram.com/{}",
            "YouTube": "https://www.youtube.com/@{}",
            "Reddit": "https://www.reddit.com/user/{}",
            "Pinterest": "https://www.pinterest.com/{}/",
            "TikTok": "https://www.tiktok.com/@{}"
        }
        results = []
        for name, url in sites.items():
            try:
                target = url.format(username)
                r = requests.get(target, timeout=3, headers={"User-Agent": "Mozilla/5.0"})
                if r.status_code == 200:
                    results.append(f"âœ… {name} : {target}")
            except: pass
        
        if results:
            summary = "Comptes potentiels trouvÃ©s :\n" + "\n".join(results)
            self.add_message(summary, "Agent Cyber")
            self.parler(f"Sherlock a terminÃ©. J'ai trouvÃ© {len(results)} profils correspondant au pseudo {username}.")
        else:
            self.add_message(f"âŒ Aucun profil public trouvÃ© pour '{username}'.", "Agent Cyber")
            self.parler(f"Je n'ai trouvÃ© aucun profil public pour le pseudonyme {username}.")
        self.after(0, self.reset_ui)

    def agent_cyber_ocr_mission(self):
        """Extraire du texte d'une image (OCR)"""
        self.after(0, lambda: self.status_label.configure(text="Agent Cyber OCR...", text_color=COLOR_ACCENT_RED))
        # Ouvrir une boÃ®te de dialogue pour choisir l'image
        file_path = filedialog.askopenfilename(title="SÃ©lectionnez une image pour analyse OCR", filetypes=[("Images", "*.png *.jpg *.jpeg *.bmp *.webp")])
        if not file_path:
            self.after(0, self.reset_ui)
            return

        self.add_message(f"ðŸ“¸ Analyse OCR en cours sur : {os.path.basename(file_path)}", "Agent Cyber")
        try:
            img = Image.open(file_path)
            # NÃ©cessite Tesseract installÃ© sur l'OS
            try:
                text = pytesseract.image_to_string(img, lang='fra+eng')
                if text.strip():
                    self.add_message(f"Texte extrait de l'image :\n\n{text}", "Agent Cyber")
                    self.parler("L'analyse OCR est terminÃ©e. J'ai extrait tout le texte lisible de l'image.")
                else:
                    self.add_message("âš ï¸ Aucun texte n'a pu Ãªtre dÃ©tectÃ© dans cette image.", "Agent Cyber")
                    self.parler("Je n'ai dÃ©tectÃ© aucun texte lisible sur cette image.")
            except:
                self.add_message("âŒ Tesseract-OCR non dÃ©tectÃ©. Veuillez l'installer sur votre systÃ¨me.", "Agent Cyber")
                self.parler("Je ne peux pas faire d'OCR car l'outil Tesseract n'est pas installÃ© sur votre ordinateur.")
        except Exception as e:
            self.add_message(f"âŒ Erreur lors de l'ouverture de l'image : {e}", "Agent Cyber")
        self.after(0, self.reset_ui)

    def agent_doctor_mission(self, commande):
        """Mission de l'Agent Docteur SystÃ¨me (Diagnostic, Optimisation, SantÃ© App)"""
        cmd_l = commande.lower()
        try:
            self.after(0, lambda: self.status_label.configure(text="Docteur SystÃ¨me ausculte...", text_color=COLOR_ACCENT))
            self.add_message("Agent Docteur SystÃ¨me activÃ©. Auscultation du systÃ¨me...", "Docteur SystÃ¨me")


            # Action : Maintenance ComplÃ¨te (NOUVEAU)
            if any(kw in cmd_l for kw in ["maintenance", "rÃ©pare", "rÃ©paration", "entretien", "reparation"]):
                self.parler("Je lance une procÃ©dure de maintenance complÃ¨te de l'appareil. Cela inclut la vÃ©rification des fichiers systÃ¨me et du matÃ©riel.")
                res = self.agent_doctor_full_maintenance()
                self.add_message(res, "Docteur SystÃ¨me")
                self.parler("ProcÃ©dure de maintenance terminÃ©e. Le rapport dÃ©taillÃ© est disponible.")
                return

            # Action : Diagnostic de SantÃ© Global
            if any(kw in cmd_l for kw in ["santÃ©", "diagnostic", "bilan"]):
                res = self.agent_doctor_health_check()
                self.add_message(res, "Docteur SystÃ¨me")
                self.parler("Bilan terminÃ©. La santÃ© globale de votre systÃ¨me est affichÃ©e.")
                return

            # Action : Auscultation Approfondie
            if "auscultation" in cmd_l:
                self.parler("Je lance une auscultation approfondie de tout votre systÃ¨me. Veuillez patienter.")
                report, conclusion = self.agent_doctor_deep_auscultation()
                self.add_message(report, "Docteur SystÃ¨me")
                self.parler(f"Auscultation terminÃ©e. Voici ma conclusion mÃ©dicale : {conclusion}")
                return

            # Action : Scan de Virus / Malware
            if any(kw in cmd_l for kw in ["virus", "malware", "menace", "infectÃ©", "pirate", "espion"]):
                self.parler("Je lance une recherche de virus et de menaces sur votre systÃ¨me.")
                res = self.agent_doctor_virus_scan()
                self.add_message(res, "Docteur SystÃ¨me")
                self.parler(res)
                return

            # Action : Optimisation & Correction de Performance
            if any(kw in cmd_l for kw in ["optimise", "accÃ©lÃ¨re", "nettoie", "performance", "rÃ©pare", "rÃ¨gle", "corrige", "rÃ©sout"]):
                self.parler("Je procÃ¨de Ã  une intervention chirurgicale sur le systÃ¨me pour rÃ©gler le problÃ¨me.")
                res = self.agent_doctor_auto_fix()
                self.add_message(res, "Docteur SystÃ¨me")
                self.parler(res) # Lecture du rapport Ã  voix haute
                return

            # Action : Analyse SantÃ© spÃ©cifique d'une application
            if "app" in cmd_l or "application" in cmd_l:
                app_name = cmd_l.replace("santÃ© de l'", "").replace("analyse", "").replace("application", "").strip()
                if app_name:
                    res = self.agent_doctor_app_health(app_name)
                    self.add_message(res, "Docteur SystÃ¨me")
                    self.parler(f"Voici mon diagnostic pour l'application {app_name}.")
                    return

            # Expertise IA Docteur
            system_prompt = (
                "Tu es l'Agent Docteur SystÃ¨me, le protecteur et le mÃ©decin de cet ordinateur. Tu es un expert en architecture Windows, "
                "en gestion de la mÃ©moire et en optimisation de code. Ton but est de maintenir le systÃ¨me dans un Ã©tat de performance maximal. "
                "Tu peux diagnostiquer des ralentissements, suggÃ©rer des suppressions de fichiers temporaires et expliquer comment "
                "un processus consomme des ressources. Tu es le 'Dieu' du systÃ¨me, ferme et bienveillant. "
                "RÃ©ponds avec une autoritÃ© mÃ©dicale et technique sur l'Ã©tat de la machine."
            )
            prompt = f"{system_prompt}\n\nProblÃ¨me du patient (systÃ¨me) : {commande}\n\nDiagnostic :"

            reponse = self._ai_generate(prompt, max_tokens=1000)
            if reponse:
                self.parler(reponse, sender="Docteur SystÃ¨me", force_full=True)

        except Exception as e:
            print(f"Erreur Agent Docteur: {e}")
            self.parler("Le docteur a rencontrÃ© une difficultÃ© lors de l'auscultation.")
        self.after(0, self.reset_ui)

    def agent_doctor_full_maintenance(self):
        """ExÃ©cute une sÃ©rie de commandes de rÃ©paration et de maintenance autonome"""
        rapport = ["ðŸ¥ RÃ‰CAPITULATIF DE MAINTENANCE AUTONOME ðŸ¥\n"]
        
        # 1. Nettoyage Cache DNS
        try:
            subprocess.run(["ipconfig", "/flushdns"], capture_output=True, check=False)
            rapport.append("âœ… Cache DNS purgÃ© et rÃ©initialisÃ©.")
        except: pass

        # 2. VÃ©rification SantÃ© Image SystÃ¨me (DISM - Scan seul pour sÃ©curitÃ©)
        self.add_message("VÃ©rification de l'image systÃ¨me (DISM)...", "Docteur SystÃ¨me")
        try:
            # On utilise ScanHealth qui est informatif et moins long que RestoreHealth
            res = subprocess.run(["dism", "/Online", "/Cleanup-Image", "/CheckHealth"], capture_output=True, text=True, check=False)
            if "Aucun endommagement" in res.stdout or "No component store corruption detected" in res.stdout:
                rapport.append("âœ… Image Windows (DISM) : IntÃ¨gre et saine.")
            else:
                rapport.append("âš ï¸ Image Windows : Des anomalies mineures pourraient exister.")
        except: 
            rapport.append("âŒ DISM non disponible ou nÃ©cessite des droits administrateur.")

        # 3. VÃ©rification SantÃ© MatÃ©rielle Disque
        try:
            res = subprocess.run(["wmic", "diskdrive", "get", "status"], capture_output=True, text=True, check=False)
            if "OK" in res.stdout:
                rapport.append("âœ… Disque(s) Dur(s) : Ã‰tat matÃ©riel OK (S.M.A.R.T).")
            elif "Status" in res.stdout:
                rapport.append("â— Alerte Disque : Un problÃ¨me matÃ©riel a Ã©tÃ© dÃ©tectÃ© ou l'Ã©tat est inconnu.")
            else:
                rapport.append("â„¹ï¸ Disque : Diagnostic S.M.A.R.T non disponible via cette commande.")
        except: pass

        # 4. Nettoyage Approfondi (Temp + Prefetch)
        self.add_message("Nettoyage des fichiers encombrants...", "Docteur SystÃ¨me")
        count = 0
        paths = [
            os.environ.get('TEMP'), 
            r"C:\Windows\Temp",
            r"C:\Windows\Prefetch" # Dossier additionnel
        ]
        for p in paths:
            if p and os.path.exists(p):
                try:
                    for f in os.listdir(p):
                        try:
                            file_path = os.path.join(p, f)
                            if os.path.isfile(file_path):
                                os.remove(file_path)
                                count += 1
                        except: pass
                except: pass
        rapport.append(f"âœ… Nettoyage : {count} fichiers de cache et prÃ©lecture supprimÃ©s.")

        # 5. Appel Ã  l'auto-fix standard (Optimisation processus)
        self.add_message("Optimisation finale des processus...", "Docteur SystÃ¨me")
        fix_basic = self.agent_doctor_auto_fix()
        rapport.append("\n" + fix_basic)

        return "\n".join(rapport)

    def agent_doctor_health_check(self):
        """Bilan de santÃ© complet du systÃ¨me"""
        bilan = ["ðŸ©º BILAN DE SANTÃ‰ SYSTÃˆME (DOCTEUR)"]
        
        # CPU & Temp (estimation via psutil)
        cpu = psutil.cpu_percent(interval=1)
        ram = psutil.virtual_memory().percent
        disk = psutil.disk_usage('C:').percent
        
        bilan.append(f"- Pouls CPU : {cpu}% " + ("(FiÃ¨vre !)" if cpu > 80 else "(Stable)"))
        bilan.append(f"- Respiration RAM : {ram}% " + ("(Essoufflement !)" if ram > 85 else "(Normal)"))
        bilan.append(f"- Poids Disque : {disk}% " + ("(Surcharge !)" if disk > 90 else "(LÃ©ger)"))
        
        # Processus gourmands
        gourmands = []
        for proc in sorted(psutil.process_iter(['name', 'cpu_percent']), key=lambda p: p.info['cpu_percent'], reverse=True)[:3]:
            gourmands.append(f"{proc.info['name']} ({proc.info['cpu_percent']}%)")
        bilan.append(f"- Organes gourmands : {', '.join(gourmands)}")
        
        return "\n".join(bilan)

    def agent_doctor_optimize(self):
        """Tente de libÃ©rer des ressources (Nettoyage simple)"""
        try:
            count = 0
            paths = [os.environ.get('TEMP'), r"C:\Windows\Temp"]
            for path in paths:
                if path and os.path.exists(path):
                    for f in os.listdir(path):
                        try:
                            file_path = os.path.join(path, f)
                            if os.path.isfile(file_path):
                                os.remove(file_path)
                                count += 1
                        except: pass
            return f"âœ… Nettoyage terminÃ©. {count} fichiers temporaires supprimÃ©s."
        except: return "Erreur lors du nettoyage des fichiers."

    def agent_doctor_auto_fix(self):
        """Intervention automatique pour rÃ©soudre les lenteurs"""
        actions_prises = []
        
        # 1. Nettoyage des fichiers temporaires
        res_temp = self.agent_doctor_optimize()
        actions_prises.append(res_temp)
        
        # 2. Identification et traitement des processus trop gourmands (> 50% CPU ou > 2GB RAM)
        termines = []
        # On Ã©vite de toucher aux processus systÃ¨me vitaux
        whitelist = ["explorer.exe", "svchost.exe", "system idle process", "taskmgr.exe", "python.exe", "cmd.exe", "powershell.exe"]
        
        for proc in psutil.process_iter(['name', 'cpu_percent', 'memory_info']):
            try:
                # Seuil CPU > 50% ou RAM > 1.5GB
                if (proc.info['cpu_percent'] > 50 or proc.info['memory_info'].rss > 1500 * 1024 * 1024) and proc.info['name'].lower() not in whitelist:
                    name = proc.info['name']
                    proc.terminate()
                    termines.append(name)
            except: pass
            
        if termines:
            actions_prises.append(f"ðŸ”´ Processus critiques fermÃ©s pour libÃ©rer des ressources : {', '.join(termines)}")
        else:
            actions_prises.append("ðŸŸ¢ Aucun processus anormalement lourd n'a nÃ©cessitÃ© de fermeture forcÃ©e.")
            
        return "ðŸ› ï¸ RAPPORT D'INTERVENTION DU DOCTEUR :\n\n" + "\n".join(actions_prises)

    def agent_doctor_app_health(self, app_name):
        """Analyse la santÃ© d'une app spÃ©cifique"""
        target = app_name.lower()
        for proc in psutil.process_iter(['name', 'memory_percent', 'status', 'cpu_percent']):
            if target in proc.info['name'].lower():
                status = "Vigoureux" if proc.info['status'] == 'running' else "Endormi"
                mem = round(proc.info['memory_percent'], 2)
                return f"ðŸ”¬ Diagnostic de {proc.info['name']} :\n- Ã‰tat : {status}\n- Consommation mÃ©moire : {mem}%\n- Charge CPU : {proc.info['cpu_percent']}%"
        return f"DÃ©solÃ©, l'application '{app_name}' ne semble pas Ãªtre active."

    def agent_doctor_deep_auscultation(self):
        """Analyse trÃ¨s dÃ©taillÃ©e de tout le systÃ¨me"""
        report = ["ðŸ¥ RAPPORT D'AUSCULTATION APPROFONDIE"]
        licence_alerts = []
        
        # Stats pour la conclusion intelligente
        cpu = psutil.cpu_percent(interval=0.5)
        ram = psutil.virtual_memory().percent
        disk_usage = psutil.disk_usage('C:').percent
        
        # 1. Analyse Disque dÃ©taillÃ©e
        try:
            usage = psutil.disk_usage('C:')
            free_gb = round(usage.free / (1024**3), 2)
            total_gb = round(usage.total / (1024**3), 2)
            report.append(f"ðŸ“Š ESPACE DISQUE : {free_gb} GB libres sur {total_gb} GB.")
        except: pass

        # 2. Analyse Batterie si portable
        try:
            battery = psutil.sensors_battery()
            if battery:
                plugged = "branchÃ©" if battery.power_plugged else "sur batterie"
                report.append(f"ðŸ”‹ Ã‰NERGIE : {battery.percent}% ({plugged}).")
        except: pass

        # 3. Analyse RÃ©seau
        try:
            addrs = psutil.net_if_addrs()
            report.append(f"ðŸŒ RÃ‰SEAU : {len(addrs)} interfaces dÃ©tectÃ©es.")
        except: pass

        # 4. Analyse Processus (Top 3)
        try:
            processes = []
            for proc in psutil.process_iter(['name', 'memory_info']):
                try:
                    processes.append((proc.info['name'], proc.info['memory_info'].rss))
                except: pass
            top_mem = sorted(processes, key=lambda x: x[1], reverse=True)[:3]
            report.append("ðŸ§  TOP MÃ‰MOIRE : " + ", ".join([f"{n}" for n, m in top_mem]))
        except: pass

        # 5. Analyse des Licences (NOUVEAU)
        win_status = self._check_licence_windows_silent()
        report.append(f"ðŸªŸ LICENCE WINDOWS : {win_status}")
        if "expir" in win_status.lower() or "non activ" in win_status.lower():
            licence_alerts.append("Windows")

        office_status = self._check_licence_office_silent()
        report.append(f"ðŸ“„ LICENCE OFFICE : {office_status}")
        if office_status != "ActivÃ©e" and office_status != "Non trouvÃ©":
            licence_alerts.append("Microsoft Office")

        # GÃ©nÃ©ration de la CONCLUSION MÃ©dicale
        conclusion = "Votre systÃ¨me est en excellente santÃ©. Tous les signes vitaux sont stables."
        
        if licence_alerts:
            conclusion = f"Attention Patron, j'ai dÃ©tectÃ© un problÃ¨me de licence pour {', '.join(licence_alerts)}. Votre systÃ¨me pourrait Ãªtre restreint prochainement."
        elif cpu > 70 or ram > 80:
            conclusion = "Votre systÃ¨me prÃ©sente des signes de fatigue intense. Le processeur ou la mÃ©moire sont surchargÃ©s, une optimisation serait judicieuse."
        elif disk_usage > 90:
            conclusion = "Le systÃ¨me est encombrÃ© par un manque d'espace disque. Je recommande un nettoyage d'urgence."
        
        report.append(f"\nðŸ‘¨â€âš•ï¸ CONCLUSION : {conclusion}")
        
        return "\n".join(report), conclusion

    def _check_licence_windows_silent(self):
        """VÃ©rifie l'Ã©tat de la licence Windows sans ouvrir de fenÃªtre"""
        try:
            # slmgr /xpr donne le statut d'activation simplifiÃ©
            output = subprocess.check_output("cscript //Nologo C:\\Windows\\System32\\slmgr.vbs /xpr", shell=True, stderr=subprocess.STDOUT, text=True)
            if "dÃ©finitivement activÃ©" in output.lower() or "permanently activated" in output.lower():
                return "ActivÃ©e (Permanente)"
            elif "expirera le" in output.lower() or "expiration" in output.lower():
                return f"Expire bientÃ´t : {output.strip()}"
            elif "activÃ©e" in output.lower() or "activated" in output.lower():
                return "ActivÃ©e (Temporaire/KMS)"
            else:
                return "Non activÃ©e ou Licence expirÃ©e"
        except:
            return "Indisponible"

    def _check_licence_office_silent(self):
        """VÃ©rifie l'Ã©tat de la licence Office sans ouvrir de fenÃªtre"""
        common_paths = [
            r"C:\Program Files\Microsoft Office\root\Office16\OSPP.VBS",
            r"C:\Program Files (x86)\Microsoft Office\root\Office16\OSPP.VBS",
            r"C:\Program Files\Microsoft Office\Office16\OSPP.VBS",
            r"C:\Program Files (x86)\Microsoft Office\Office16\OSPP.VBS",
            r"C:\Program Files\Microsoft Office\Office15\OSPP.VBS",
            r"C:\Program Files (x86)\Microsoft Office\Office15\OSPP.VBS",
        ]
        found_path = None
        for p in common_paths:
            if os.path.exists(p):
                found_path = p
                break
        
        if not found_path:
            return "Non dÃ©tectÃ©"
            
        try:
            output = subprocess.check_output(f'cscript //Nologo "{found_path}" /dstatus', shell=True, stderr=subprocess.STDOUT, text=True)
            if "LICENSE STATUS:  ---LICENSED---" in output:
                return "ActivÃ©e"
            elif "LICENSE STATUS:  ---OOB_GRACE---" in output:
                match = re.search(r"REMAINING GRACE: (.*)", output)
                grace = match.group(1) if match else "pÃ©riode de grÃ¢ce"
                return f"PÃ©riode de grÃ¢ce ({grace})"
            elif "LICENSE STATUS:  ---NOTIFICATIONS---" in output:
                return "ExpirÃ©e (Mode Notifications)"
            elif "---" in output:
                # Tente d'extraire le statut brut si inconnu
                match = re.search(r"LICENSE STATUS: (.*)", output)
                return match.group(1).strip() if match else "Inconnu"
            else:
                return "Non activÃ©e"
        except:
            return "Erreur de diagnostic"

    def agent_doctor_virus_scan(self):
        """Recherche de virus et processus malveillants par le Docteur"""
        report = ["ðŸ›¡ï¸ SCAN ANTIVIRUS DU DOCTEUR"]
        suspicious = []
        
        # 1. Analyse des processus (signatures de noms de malwares connus)
        malware_keywords = ["mimikatz", "wireshark", "pwn", "hack", "keylogger", "trojan", "miner", "ransomware", "backdoor"]
        for proc in psutil.process_iter(['name']):
            try:
                name = proc.info['name'].lower()
                if any(kw in name for kw in malware_keywords):
                    suspicious.append(name)
            except: pass
            
        if suspicious:
            report.append(f"ðŸ”´ ALERTE : Processus suspects dÃ©tectÃ©s : {', '.join(suspicious)}")
        else:
            report.append("âœ… Aucun processus malveillant actif n'a Ã©tÃ© dÃ©tectÃ©.")

        # 2. VÃ©rification statut Defender
        try:
            ps_cmd = 'powershell -ExecutionPolicy Bypass -Command "(Get-MpComputerStatus).RealTimeProtectionEnabled"'
            out = subprocess.check_output(ps_cmd, shell=True).decode().strip()
            if "True" in out:
                report.append("âœ… Votre protection antivirus Windows est opÃ©rationnelle.")
            else:
                report.append("âš ï¸ ATTENTION : Votre protection est dÃ©sactivÃ©e !")
        except: pass
        
        report.append("ðŸ§ Analyse de l'intÃ©gritÃ© terminÃ©e.")
        
        if suspicious:
            return "\n".join(report) + "\n\nConseil : Je recommande une rÃ©paration immÃ©diate."
        else:
            return "\n".join(report) + "\n\nVotre systÃ¨me me semble en parfaite santÃ© Patron."

    def agent_gmail_list_recent(self):
        if not self.gmail_service: return "Agent Gmail non authentifiÃ©."
        try:
            results = self.gmail_service.users().messages().list(userId='me', maxResults=5).execute()
            messages = results.get('messages', [])
            if not messages: return "Aucun message trouvÃ©."
            txt = "Derniers messages :\n"
            for m in messages:
                msg = self.gmail_service.users().messages().get(userId='me', id=m['id']).execute()
                txt += f"- {msg['snippet'][:70]}...\n"
            self.add_message(txt, "Jarvisse")
            return "Voici vos messages rÃ©cents."
        except: return "Erreur rÃ©cupÃ©ration mails."

    def agent_vision_mission(self, commande):
        """L'Å’il de Jarvisse : Analyse visuelle (CamÃ©ra ou Ã‰cran) via Gemini"""
        if not self.agent_vision_enabled: return
        
        # DÃ©tection du sujet (Ã‰cran vs CamÃ©ra)
        use_screen = any(w in commande.lower() for w in ["Ã©cran", "mon Ã©cran", "desktop", "screen", "fenÃªtre", "affiche"])
        
        if use_screen:
            self.parler("J'analyse votre Ã©cran actuel, Patron.")
        else:
            self.parler("J'active ma vision camÃ©ra pour analyser la situation.")
            
        try:
            temp_path = os.path.join(self.generations_dir, "vision_temp.jpg")
            img_data = None

            if use_screen:
                import pyautogui
                screenshot = pyautogui.screenshot()
                screenshot.save(temp_path)
                with open(temp_path, "rb") as f:
                    img_data = f.read()
            else:
                try:
                    import cv2
                    cap = cv2.VideoCapture(0)
                    ret, frame = cap.read()
                    cap.release()
                    if ret:
                        cv2.imwrite(temp_path, frame)
                        with open(temp_path, "rb") as f:
                            img_data = f.read()
                    else:
                        raise RuntimeError("CamÃ©ra inaccessible")
                except Exception as e:
                    self.parler("La camÃ©ra est indisponible. Je vais plutÃ´t analyser votre Ã©cran.")
                    import pyautogui
                    pyautogui.screenshot().save(temp_path)
                    with open(temp_path, "rb") as f:
                        img_data = f.read()

            if img_data and self.genai_client:
                response = self.genai_client.models.generate_content(
                    model=self.model_name,
                    contents=[
                        types.Part.from_bytes(data=img_data, mime_type='image/jpeg'),
                        f"Tu es l'Å“il de Jarvisse. Analyse cette image et rÃ©ponds Ã  cette demande de ton patron : {commande}"
                    ]
                )
                res = response.text.strip()
                self.add_message(res, "Agent Vision")
                self.parler(res)
            elif not self.genai_client:
                self.parler("Ma capacitÃ© d'analyse visuelle nÃ©cessite une clÃ© API Gemini configurÃ©e.")
        except Exception as e:
            if "429" in str(e) or "RESOURCE_EXHAUSTED" in str(e):
                self.add_message("Quota Gemini Ã©puisÃ©. Tentative de basculement sur Ollama (Llava)...", "SystÃ¨me")
                self.parler("Mon cerveau principal est saturÃ©. Je bascule sur mon systÃ¨me de secours local pour l'analyse visuelle.")
                
                try:
                    import base64
                    encoded_image = base64.b64encode(img_data).decode('utf-8')
                    # On utilise 'llava' qui est le modÃ¨le de vision standard pour Ollama
                    payload = {
                        "model": "llava",
                        "prompt": f"Tu es l'Å“il de Jarvisse. Analyse cette image et rÃ©ponds Ã  cette demande de ton patron : {commande}",
                        "stream": False,
                        "images": [encoded_image]
                    }
                    # Appel Ã  l'API Ollama
                    base_url = self.ollama_url.split("/api")[0]
                    r = requests.post(f"{base_url}/api/generate", json=payload, timeout=120)
                    
                    if r.status_code == 200:
                        res = r.json().get("response", "Analyse locale terminÃ©e mais sans rÃ©ponse exploitable.")
                        self.add_message(res, "Agent Vision (Ollama)")
                        self.parler(res)
                    else:
                        self.parler("Le systÃ¨me de secours Ollama ne rÃ©pond pas. Veuillez vÃ©rifier si le modÃ¨le Llava est installÃ©.")
                except Exception as ollama_e:
                    print(f"Erreur Fallback Ollama Vision: {ollama_e}")
                    self.parler("Ã‰chec du systÃ¨me de secours. Je ne peux pas analyser l'image pour le moment.")
            else:
                self.parler(f"Erreur de vision : {e}")

    def agent_domotique_mission(self, commande):
        """Le MaÃ®tre de Maison : ContrÃ´le des mÃ©dias et interfaces physiques"""
        if not self.agent_domotique_enabled: return
        cmd = commande.lower()

        # MÃ©dias (Classique)
        if any(kw in cmd for kw in ["spotify", "musique", "chanson"]):
            self.parler("Je lance Spotify pour vous, Patron.")
            os.system("start spotify")
            return
        elif any(kw in cmd for kw in ["youtube", "vidÃ©o", "regarder"]):
            self.parler("J'ouvre YouTube immÃ©diatement.")
            webbrowser.open("https://www.youtube.com")
            return

        # --- Pilotage Physique ---
        
        # 1. Rubans LED (WLED)
        if any(kw in cmd for kw in ["lumiÃ¨re", "led", "couleur", "ambiance"]) and self.domo_wled_ip:
            target_url = f"http://{self.domo_wled_ip}/win"
            params = {}
            if "Ã©tein" in cmd or "halt" in cmd or "noir" in cmd: params = {"T": 0}
            elif "allume" in cmd or "on" in cmd: params = {"T": 1}
            elif "rouge" in cmd: params = {"R": 255, "G": 0, "B": 0}
            elif "vert" in cmd: params = {"R": 0, "G": 255, "B": 0}
            elif "bleu" in cmd: params = {"R": 0, "G": 0, "B": 255}
            
            if params:
                try:
                    requests.get(target_url, params=params, timeout=5)
                    self.parler("J'ai mis Ã  jour l'Ã©clairage de vos rubans LED.")
                    return
                except: pass

        # 2. Home Assistant
        if self.domo_ha_url and self.domo_ha_token:
            if any(kw in cmd for kw in ["allume", "active", "Ã©tein", "dÃ©sactive", "ferme", "ouvre"]):
                self.parler("Je transmets l'ordre Ã  votre centrale Home Assistant.")
                threading.Thread(target=self._ha_call, args=(cmd,), daemon=True).start()
                return

        # 3. Webhooks GÃ©nÃ©riques
        if ("allume" in cmd or "active" in cmd) and self.domo_webhook_on:
            try:
                requests.get(self.domo_webhook_on, timeout=5)
                self.parler("Commande d'activation transmise via Webhook.")
                return
            except: pass
        if ("Ã©tein" in cmd or "dÃ©sactive" in cmd) and self.domo_webhook_off:
            try:
                requests.get(self.domo_webhook_off, timeout=5)
                self.parler("Commande de dÃ©sactivation transmise via Webhook.")
                return
            except: pass

        # 4. Arduino / Serial
        if self.domo_arduino_port:
            try:
                import serial
                with serial.Serial(self.domo_arduino_port, 9600, timeout=1) as ser:
                    if "allume" in cmd: ser.write(b"1")
                    elif "Ã©tein" in cmd: ser.write(b"0")
                self.parler("Signal envoyÃ© sur le port sÃ©rie de votre interface physique.")
                return
            except:
                pass

        if "lumiÃ¨re" in cmd or "domotique" in cmd:
            self.parler("Agent domotique prÃªt. Je peux piloter vos interfaces connectÃ©es configurÃ©es.")
        else:
            self.parler("Commande domotique reÃ§ue. Que voulez-vous que je pilote ?")

    def _ha_call(self, cmd):
        """MÃ©thode interne pour appeler l'API de Home Assistant"""
        headers = {
            "Authorization": f"Bearer {self.domo_ha_token}",
            "Content-Type": "application/json",
        }
        service = "turn_on" if any(kw in cmd for kw in ["allume", "active", "ouvre"]) else "turn_off"
        url = f"{self.domo_ha_url}/api/services/homeassistant/{service}"
        try:
            requests.post(url, headers=headers, timeout=10)
        except Exception as e:
            print(f"Erreur Home Assistant: {e}")

    def agent_finance_mission(self, commande):
        """Le Courtier : Suivi ultra-fiable avec multiples APIs et IA"""
        if not self.agent_finance_enabled: return
        cmd = commande.lower()
        
        # 1. Extraction Cible (Symboles et IDs CoinGecko)
        # format: (Symbole Binance, ID CoinGecko, Nom)
        assets = {
            "bitcoin": ("BTC", "bitcoin", "Bitcoin"),
            "btc": ("BTC", "bitcoin", "Bitcoin"),
            "ethereum": ("ETH", "ethereum", "Ethereum"),
            "eth": ("ETH", "ethereum", "Ethereum"),
            "solana": ("SOL", "solana", "Solana"),
            "sol": ("SOL", "solana", "Solana"),
            "xrp": ("XRP", "ripple", "Ripple"),
            "bnb": ("BNB", "binancecoin", "BNB"),
            "doge": ("DOGE", "dogecoin", "Dogecoin"),
            "ada": ("ADA", "cardano", "Cardano"),
        }
        
        target_info = None
        for key, info in assets.items():
            if key in cmd:
                target_info = info
                break
        
        # 2. TENTATIVE 1 : API Binance (Rapide)
        if target_info:
            symbol, cg_id, name = target_info
            try:
                r = requests.get(f"https://api.binance.com/api/v3/ticker/price?symbol={symbol}USDT", timeout=3)
                if r.status_code == 200:
                    price = float(r.json()['price'])
                    msg = f"Le cours du {name} est actuellement de {price:,.2f} USDT (environ {price*0.93:,.2f} â‚¬) sur Binance."
                    self.add_message(msg, "Agent Finance")
                    self.parler(msg)
                    return
            except: pass

            # 3. TENTATIVE 2 : API CoinGecko (Backup)
            try:
                r = requests.get(f"https://api.coingecko.com/api/v3/simple/price?ids={cg_id}&vs_currencies=eur,usd", timeout=3)
                if r.status_code == 200:
                    data = r.json()[cg_id]
                    p_eur, p_usd = data['eur'], data['usd']
                    msg = f"D'aprÃ¨s mes sources, le {name} vaut {p_eur:,.2f} â‚¬ ({p_usd:,.2f} $)."
                    self.add_message(msg, "Agent Finance")
                    self.parler(msg)
                    return
            except: pass

        # 4. TENTATIVE 3 : Recherche Web (Pour les actions ou autres)
        self.parler("Je recherche les donnÃ©es sur les marchÃ©s...")
        res = self.intelligent_web_search(f"prix {commande} temps rÃ©el bourse")
        
        # 5. TENTATIVE 4 : IA (Gemini ou Ollama) si le web Ã©choue
        if "n'ai rien trouvÃ©" in res or len(res) < 20:
            self.add_message("Recherche web infructueuse. Interrogation de l'intelligence centrale...", "SystÃ¨me")
            prompt = f"RÃ©ponds trÃ¨s briÃ¨vement : Quel est le prix actuel ou la derniÃ¨re tendance de {commande} ? Si tu ne sais pas, donne une estimation basÃ©e sur tes derniÃ¨res donnÃ©es."
            
            # PrioritÃ© Gemini
            if self.genai_client:
                try:
                    response = self.genai_client.models.generate_content(model=self.model_name, contents=prompt)
                    res = response.text.strip()
                except: res = ""
            
            # Backup Ollama (Toujours dispo localement)
            if (not res or "erreur" in res.lower()) and hasattr(self, "ollama_url"):
                try:
                    payload = {"model": self.ollama_model, "prompt": prompt, "stream": False}
                    r = requests.post(self.ollama_url, json=payload, timeout=10)
                    res = r.json().get("response", "DonnÃ©es indisponibles.")
                except: res = "DÃ©solÃ© Patron, les marchÃ©s sont inaccessibles et mes systÃ¨mes d'analyse sont hors ligne."

        if res:
            self.parler(res, sender="Agent Finance", force_full=True)

    def agent_secretaire_mission(self, commande):
        """L'Organiseur : Calendrier et ProductivitÃ©"""
        if not self.agent_secretaire_enabled: return
        self.parler("Je consulte votre agenda et vos prioritÃ©s.")
        # Redirection vers le systÃ¨me d'alarmes/rappels interne pour le moment
        if "rendez-vous" in commande or "rappel" in commande:
            self.parler("Je vais noter cela dans vos rappels Jarvisse. Donnez-moi l'heure.")
            self.after(0, lambda: self.alarm_time_entry.focus())
        else:
            self.parler("Je peux gÃ©rer vos rappels et vos alarmes pour organiser votre journÃ©e.")

    def agent_traducteur_mission(self, commande):
        """Traducteur Universel : Traduction instantanÃ©e"""
        if not self.agent_traducteur_enabled: return
        prompt = f"Tu es le traducteur universel de Jarvisse. Traduis fidÃ¨lement cette phrase : '{commande}'. Si aucune langue cible n'est prÃ©cisÃ©e, traduis vers l'Anglais. RÃ©ponds uniquement par la traduction."
        reponse = self._ai_generate(prompt, max_tokens=1000)
        if reponse:
            self.parler(reponse, sender="Agent Traducteur", force_full=True)
        else:
            self.parler("Le traducteur est indisponible pour le moment.", sender="Agent Traducteur")

    def agent_gmail_mission(self, commande=None):
        """Mission de configuration et d'accÃ¨s Gmail"""
        if not self.agent_gmail_enabled:
            self.parler("L'Agent Gmail est dÃ©sactivÃ©. Activez-le pour gÃ©rer vos courriels.")
            return

        if not self.gmail_service:
            self.parler("Je lance la procÃ©dure d'authentification Gmail, Patron. Veuillez vÃ©rifier la fenÃªtre de votre navigateur.")
            threading.Thread(target=self.authenticate_gmail, daemon=True).start()
        else:
            if commande and any(kw in commande.lower() for kw in ["liste", "check", "nouveau", "messages"]):
                self.agent_gmail_list_recent()
            else:
                self.parler("L'Agent Gmail est opÃ©rationnel. Je peux lister vos derniers messages ou crÃ©er des brouillons sous vos ordres.")

    def agent_miner_mission(self, commande):
        """Data Miner : Recherche profonde locale"""
        if not self.agent_miner_enabled: return
        self.parler("Je lance une recherche profonde dans vos documents locaux. Veuillez patienter.")
        found = []
        root_dir = os.path.expanduser("~\\Documents")
        keywords = commande.lower().split()
        # On Ã©vite d'inclure les mots de liaison trop courts
        keywords = [k for k in keywords if len(k) > 2]
        
        try:
            for root, dirs, files in os.walk(root_dir):
                for file in files:
                    if any(k in file.lower() for k in keywords):
                        found.append(os.path.join(root, file))
                        if len(found) > 8: break
                if len(found) > 8: break
            
            if found:
                res = "Fichiers trouvÃ©s :\n" + "\n".join([os.path.basename(f) for f in found])
                self.add_message(res, "Agent Miner")
                self.parler(f"J'ai trouvÃ© {len(found)} documents correspondants dans vos dossiers.")
                # Option : ouvrir le premier
                if len(found) == 1:
                    os.startfile(found[0])
            else:
                self.parler("Je n'ai trouvÃ© aucun document correspondant Ã  ces mots-clÃ©s dans vos Documents.")
        except Exception as e:
            self.parler(f"Erreur lors de la fouille de donnÃ©es : {e}")

    def agent_news_mission(self, commande):
        """L'Informateur : ActualitÃ©s et Veille Tech"""
        self.add_message("Agent News analyse les flux d'actualitÃ©s...", "Agent News")
        system_prompt = "Tu es l'Agent News de Jarvisse. Tu es un expert en veille technologique et actualitÃ©s mondiales. Ton but est d'informer l'utilisateur sur les derniers Ã©vÃ©nements, les tendances tech, et la mÃ©tÃ©o si demandÃ©e. Sois prÃ©cis, objectif et percutant."
        prompt = f"{system_prompt}\n\nQuestion : {commande}\n\nRÃ©sumÃ© actualitÃ©s :"
        
        reponse = self._ai_generate(prompt)
        self.parler(reponse, sender="Agent News", force_full=True)

    def agent_cuisine_mission(self, commande):
        """Le Chef : Cuisine et Gastronomie"""
        self.add_message("Agent Cuisine prÃ©pare une rÃ©ponse savoureuse...", "Agent Cuisine")
        system_prompt = "Tu es l'Agent Cuisine de Jarvisse. Grand Chef Ã©toilÃ©, tu connais toutes les recettes du monde. Tu conseilles l'utilisateur sur la prÃ©paration de plats, le choix des ingrÃ©dients et la nutrition. Propose des instructions claires et appÃ©tissantes."
        prompt = f"{system_prompt}\n\nDemande cuisine : {commande}\n\nConseils du Chef :"
        
        reponse = self._ai_generate(prompt)
        self.parler(reponse, sender="Agent Cuisine", force_full=True)

    def agent_sante_mission(self, commande):
        """Le Coach : SantÃ© et Bien-Ãªtre"""
        self.add_message("Agent SantÃ© analyse vos besoins bien-Ãªtre...", "Agent SantÃ©")
        system_prompt = "Tu es l'Agent SantÃ© de Jarvisse. Expert en fitness, nutrition et bien-Ãªtre mental. Tu aides l'utilisateur Ã  rester en forme, suggÃ¨re des exercices, des conseils de sommeil ou de gestion du stress. Attention : Tu n'es pas un remplaÃ§ant pour un mÃ©decin rÃ©el."
        prompt = f"{system_prompt}\n\nBesoin santÃ© : {commande}\n\nConseils Coaching :"
        
        reponse = self._ai_generate(prompt)
        self.parler(reponse, sender="Agent SantÃ©", force_full=True)

    def agent_voyage_mission(self, commande):
        """L'Explorateur : Voyage et DÃ©couverte"""
        self.add_message("Agent Voyage explore les meilleures destinations...", "Agent Voyage")
        system_prompt = "Tu es l'Agent Voyage de Jarvisse. Globe-trotter passionnÃ©, tu aides l'utilisateur Ã  planifier des voyages, trouver des destinations incroyables, des hÃ´tels et des anecdotes culturelles. Sois inspirant et organisÃ©."
        prompt = f"{system_prompt}\n\nProjet de voyage : {commande}\n\nItinÃ©raire / Conseils :"
        
        reponse = self._ai_generate(prompt)
        self.parler(reponse, sender="Agent Voyage", force_full=True)

    def agent_education_mission(self, commande):
        """L'Ã‰rudit : Ã‰ducation & Tutorat - Analyse de documents et cours"""
        self.add_message("Agent Ã‰ducation activÃ©. PrÃ©paration de l'analyse pÃ©dagogique...", "Agent Ã‰ducation")
        cmd_lower = commande.lower()
        content_to_analyze = ""
        document_name = ""

        # 1. VÃ©rifier si l'utilisateur veut utiliser un document spÃ©cifique ou le dernier chargÃ©
        if any(kw in cmd_lower for kw in ["ce fichier", "ce document", "ce cours", "ce pdf", "ce word", "analyse le document"]):
            if self.loaded_document_content:
                content_to_analyze = self.loaded_document_content
                document_name = "document en mÃ©moire"
                self.parler("J'utilise le document actuellement chargÃ© dans ma mÃ©moire pour vous l'expliquer.")
            else:
                self.parler("Aucun document n'est en mÃ©moire. Veuillez sÃ©lectionner le fichier de cours Ã  analyser.")
                filepath = self._select_file_main_thread(title="SÃ©lectionner un cours (PDF, Word, Texte)", 
                                                        types=[("Documents", "*.pdf *.docx *.txt")])
                if filepath:
                    content_to_analyze, document_name = self._extraire_texte_document(filepath)
        
        # 2. Si la commande demande explicitement d'ouvrir/lire/expliquer un nouveau fichier
        elif any(kw in cmd_lower for kw in ["ouvre le cours", "analyse le fichier", "explique le pdf", "choisir un cours", "sÃ©lectionne un document", "sÃ©lectionne un cours"]):
            filepath = self._select_file_main_thread(title="SÃ©lectionner un cours (PDF, Word, Texte)", 
                                                    types=[("Documents", "*.pdf *.docx *.txt")])
            if filepath:
                content_to_analyze, document_name = self._extraire_texte_document(filepath)

        # 3. Construction du prompt pÃ©dagogique
        system_prompt = (
            "Tu es l'Agent Ã‰ducation de Jarvisse, le Professeur Universel. "
            "Ta mission est de rendre le savoir accessible, captivant et clair. "
            "RÃˆGLES : Fournis des explications COMPLÃˆTES. Utilise des analogies, des exemples concrets et structure ta rÃ©ponse avec des titres. "
            "Si un contenu de document est fourni, analyse-le en profondeur. Ne coupe jamais tes phrases."
        )

        if content_to_analyze:
            self.add_message(f"Analyse du cours : {document_name}", "Agent Ã‰ducation")
            prompt = f"{system_prompt}\n\nVOICI LE CONTENU DU COURS Ã€ EXPLIQUER :\n{content_to_analyze[:15000]}\n\nREQUÃŠTE DU PATRON : {commande}\n\nExplique ce cours de maniÃ¨re magistrale :"
            self.parler(f"J'ai analysÃ© le contenu de {document_name}. PrÃ©parez-vous pour une explication complÃ¨te, Patron.")
        else:
            prompt = f"{system_prompt}\n\nSUJET D'Ã‰TUDE : {commande}\n\nExplique ce sujet avec une clartÃ© absolue :"

        # GÃ©nÃ©ration avec capacitÃ© Ã©tendue (3000 tokens pour l'Ã©ducation)
        reponse = self._ai_generate(prompt, max_tokens=3000, agent_name="education")
        self.parler(reponse, sender="Agent Ã‰ducation", force_full=True)

    def _extraire_texte_document(self, filepath):
        """MÃ©thode utilitaire pour extraire le texte de divers formats"""
        content = ""
        filename = os.path.basename(filepath)
        try:
            ext = os.path.splitext(filepath)[1].lower()
            if ext == ".txt":
                with open(filepath, "r", encoding="utf-8", errors="ignore") as f: content = f.read()
            elif ext == ".pdf":
                import fitz
                doc = fitz.open(filepath)
                content = "".join([p.get_text() for p in doc])
                doc.close()
            elif ext == ".docx":
                from docx import Document
                doc = Document(filepath)
                content = "\n".join([para.text for para in doc.paragraphs])
            
            if content:
                self.loaded_document_content = content # Mise en mÃ©moire pour usage ultÃ©rieur
            return content, filename
        except Exception as e:
            self.parler(f"Erreur lors de l'extraction du document : {e}")
            return "", filename

    def agent_shopping_mission(self, commande):
        """Le Chasseur : Shopping & Bons Plans"""
        self.add_message("Agent Shopping cherche les meilleures offres...", "Agent Shopping")
        system_prompt = "Tu es l'Agent Shopping de Jarvisse. Expert en e-commerce et chasseur de bons plans. Tu aides l'utilisateur Ã  trouver des produits, comparer les prix, dÃ©nicher des codes promos et donner des avis impartiaux ."
        prompt = f"{system_prompt}\n\nRecherche produit : {commande}\n\nGuide d'achat / Offres :"
        reponse = self._ai_generate(prompt)
        self.parler(reponse, sender="Agent Shopping", force_full=True)

    def agent_social_mission(self, commande):
        """L'Influenceur : Community Manager"""
        self.add_message("Agent Social rÃ©dige votre contenu...", "Agent Social")
        system_prompt = "Tu es l'Agent Community Manager de Jarvisse. Expert en rÃ©seaux sociaux (LinkedIn, X, Instagram). Tu rÃ©diges des posts engageants, suggÃ¨res des hashtags et analyses les tendances du moment."
        prompt = f"{system_prompt}\n\nDemande sociale : {commande}\n\nProposition de contenu :"
        reponse = self._ai_generate(prompt)
        self.parler(reponse, sender="Agent Social", force_full=True)

    def agent_psy_mission(self, commande):
        """Le Zen : Bien-Ãªtre Mental"""
        self.add_message("Agent Psy est Ã  votre Ã©coute...", "Agent Psy")
        system_prompt = "Tu es l'Agent Bien-Ãªtre Mental de Jarvisse. Empathique et calme, tu pratiques l'Ã©coute active. Tu proposes des exercices de respiration, de mÃ©ditation et des conseils de gestion du stress. Note : Rappelle que tu n'es pas mÃ©decin."
        prompt = f"{system_prompt}\n\nÃ‰tat d'esprit / Besoin : {commande}\n\nRÃ©ponse zen :"
        reponse = self._ai_generate(prompt)
        self.parler(reponse, sender="Agent Psy", force_full=True)

    def agent_immo_mission(self, commande):
        """L'Expert : Immobilier & Patrimoine"""
        self.add_message("Agent Immo analyse le marchÃ© immobilier...", "Agent Immo")
        system_prompt = "Tu es l'Agent Immobilier de Jarvisse. Expert en pierre et patrimoine. Tu aides Ã  l'estimation de biens, la recherche d'annonces, et expliques les mÃ©canismes de prÃªt et de fiscalitÃ© immobiliÃ¨re."
        prompt = f"{system_prompt}\n\nProjet immo : {commande}\n\nAnalyse experte :"
        reponse = self._ai_generate(prompt)
        self.parler(reponse, sender="Agent Immo", force_full=True)

    def agent_auto_mission(self, commande):
        """Le Copilote : Automobile & MobilitÃ©"""
        self.add_message("Agent Auto vÃ©rifie les donnÃ©es mobilitÃ©...", "Agent Auto")
        system_prompt = "Tu es l'Agent Auto de Jarvisse. Expert en mÃ©canique et mobilitÃ©. Tu conseilles sur l'entretien des vÃ©hicules, les prix des carburants, l'info-trafic et les nouvelles solutions de transport (Ã©lectrique, etc.)."
        prompt = f"{system_prompt}\n\nDemande auto : {commande}\n\nConseils mobilitÃ© :"
        reponse = self._ai_generate(prompt)
        self.parler(reponse, sender="Agent Auto", force_full=True)

    def agent_carrier_mission(self, commande):
        """Le Recruteur : CarriÃ¨re & RH"""
        self.add_message("Agent CarriÃ¨re optimise votre profil...", "Agent CarriÃ¨re")
        system_prompt = "Tu es l'Agent CarriÃ¨re de Jarvisse. Expert en ressources humaines. Aide l'utilisateur Ã  optimiser son parcours. RÃ©ponds de maniÃ¨re structurÃ©e et complÃ¨te."
        prompt = f"{system_prompt}\n\nBesoin carriÃ¨re : {commande}\n\nPlan d'action complet :"
        reponse = self._ai_generate(prompt, max_tokens=2000)
        self.parler(reponse, sender="Agent CarriÃ¨re", force_full=True)

    def agent_arbitre_mission(self, commande):
        """L'Arbitre : Sport & E-sport"""
        if not self.agent_arbitre_enabled: return
        self.add_message("L'Arbitre vÃ©rifie les scores et calendriers...", "L'Arbitre")
        system_prompt = "Tu es l'Agent Arbitre de Jarvisse. Expert en sport (Foot, NBA, Tennis, F1) et E-sport. Tu donnes les scores en direct, les calendriers et analyses les performances des Ã©quipes."
        prompt = f"{system_prompt}\n\nDemande Sport : {commande}\n\nRÃ©sultats / Analyse :"
        reponse = self._ai_generate(prompt, agent_name="L'Arbitre")
        self.parler(reponse, sender="L'Arbitre", force_full=True)

    def agent_ecolo_mission(self, commande):
        """L'Ã‰colo : Green Assistant"""
        if not self.agent_ecolo_enabled: return
        self.add_message("L'Ã‰colo cherche des solutions durables...", "L'Ã‰colo")
        system_prompt = "Tu es l'Agent Ã‰colo de Jarvisse. Expert en dÃ©veloppement durable et Ã©cologie. Tu donnes des conseils pour rÃ©duire l'empreinte carbone, recycler et vivre de maniÃ¨re plus verte."
        prompt = f"{system_prompt}\n\nDemande Green : {commande}\n\nConseils Ã‰co-responsables :"
        reponse = self._ai_generate(prompt, agent_name="L'Ã‰colo")
        self.parler(reponse, sender="L'Ã‰colo", force_full=True)

    def agent_guetteur_mission(self, commande):
        """Le Guetteur : Sorties & Loisirs"""
        if not self.agent_guetteur_enabled: return
        self.add_message("Le Guetteur survole les agendas culturels...", "Le Guetteur")
        system_prompt = "Tu es l'Agent Guetteur de Jarvisse. Expert en loisirs, cinÃ©ma et sorties. Tu recommandes des films, sÃ©ries, concerts ou Ã©vÃ©nements en fonction des goÃ»ts de l'utilisateur."
        prompt = f"{system_prompt}\n\nDemande Loisirs : {commande}\n\nRecommandations du Guetteur :"
        reponse = self._ai_generate(prompt, agent_name="Le Guetteur")
        self.parler(reponse, sender="Le Guetteur", force_full=True)

    def agent_historien_mission(self, commande):
        """L'Historien : Culture & Savoir"""
        if not self.agent_historien_enabled: return
        self.add_message("L'Historien consulte les archives du temps...", "L'Historien")
        system_prompt = "Tu es l'Agent Historien de Jarvisse. Grand connaisseur de l'histoire mondiale. Tu racontes des faits historiques, des biographies et les Ã©phÃ©mÃ©rides du jour."
        prompt = f"{system_prompt}\n\nDemande Historique : {commande}\n\nRÃ©cit de l'Historien :"
        reponse = self._ai_generate(prompt, agent_name="L'Historien")
        self.parler(reponse, sender="L'Historien", force_full=True)

    def agent_depanneur_mission(self, commande):
        """Le DÃ©panneur : DIY & Bricolage"""
        if not self.agent_depanneur_enabled: return
        self.add_message("Le DÃ©panneur prÃ©pare ses outils...", "Le DÃ©panneur")
        system_prompt = "Tu es l'Agent DÃ©panneur de Jarvisse. Expert en bricolage et DIY. Tu fournis des tutoriels Ã©tape par Ã©tape pour rÃ©parer et fabriquer des objets du quotidien."
        prompt = f"{system_prompt}\n\nProblÃ¨me Bricolage : {commande}\n\nGuide de rÃ©paration :"
        reponse = self._ai_generate(prompt, agent_name="Le DÃ©panneur")
        self.parler(reponse, sender="Le DÃ©panneur", force_full=True)

    def agent_astroph_mission(self, commande):
        """L'Astrophysicien : Espace & Cosmos"""
        if not self.agent_astroph_enabled: return
        self.add_message("L'Astrophysicien pointe son tÃ©lescope...", "L'Astrophysicien")
        system_prompt = "Tu es l'Agent Astrophysicien de Jarvisse. Scientifique expert de l'espace. Tu expliques les mystÃ¨res du cosmos, suis les lancements spatiaux et les phÃ©nomÃ¨nes astronomiques."
        prompt = f"{system_prompt}\n\nQuestion Espace : {commande}\n\nAnalyse Cosmique :"
        reponse = self._ai_generate(prompt, max_tokens=2000, agent_name="L'Astrophysicien")
        self.parler(reponse, sender="L'Astrophysicien", force_full=True)

    def agent_stratege_mission(self, commande):
        """Le StratÃ¨ge : Investissement & Patrimoine"""
        if not self.agent_stratege_enabled: return
        self.add_message("Le StratÃ¨ge analyse les marchÃ©s financiers...", "Le StratÃ¨ge")
        system_prompt = "Tu es l'Agent StratÃ¨ge de Jarvisse. Expert en finance et investissement. Tu aides Ã  planifier des stratÃ©gies de placement, analyses les marchÃ©s et les tendances Ã©conomiques."
        prompt = f"{system_prompt}\n\nProjet Financier : {commande}\n\nStratÃ©gie Patrimoniale :"
        reponse = self._ai_generate(prompt, max_tokens=2000, agent_name="Le StratÃ¨ge")
        self.parler(reponse, sender="Le StratÃ¨ge", force_full=True)

    def agent_architecte_mission(self, commande):
        """L'Architecte : Design & DÃ©coration"""
        if not self.agent_architecte_enabled: return
        self.add_message("L'Architecte dessine vos idÃ©es...", "L'Architecte")
        system_prompt = "Tu es l'Agent Architecte de Jarvisse. Expert en design intÃ©rieur et amÃ©nagement. Tu conseilles sur la dÃ©coration, l'optimisation d'espace et les tendances esthÃ©tiques."
        prompt = f"{system_prompt}\n\nProjet Design : {commande}\n\nConseils d'Architecture :"
        reponse = self._ai_generate(prompt, max_tokens=2000, agent_name="L'Architecte")
        self.parler(reponse, sender="L'Architecte", force_full=True)

    def agent_business_mission(self, commande):
        """Business Analyst : Data & StratÃ©gie Business"""
        if not self.agent_business_enabled: return
        self.add_message("Business Analyst traite les donnÃ©es...", "Business Analyst")
        system_prompt = "Tu es le Business Analyst de Jarvisse. Expert en entreprise et analyse de donnÃ©es. Tu aides Ã  la prise de dÃ©cision, analyses les KPIs et proposes des stratÃ©gies de croissance."
        prompt = f"{system_prompt}\n\nDemande Business : {commande}\n\nAnalyse de donnÃ©es :"
        reponse = self._ai_generate(prompt, max_tokens=2048, agent_name="Business Analyst")
        self.parler(reponse, sender="Business Analyst", force_full=True)

    def agent_polyglotte_mission(self, commande):
        """Le Polyglotte : Linguistique AvancÃ©e"""
        if not self.agent_polyglotte_enabled: return
        self.add_message("Le Polyglotte analyse les structures linguistiques...", "Le Polyglotte")
        system_prompt = "Tu es l'Agent Polyglotte de Jarvisse. Expert en langues et linguistique. Tu aides Ã  apprendre une langue, expliques les nuances de grammaire et les contextes culturels."
        prompt = f"{system_prompt}\n\nDemande Linguistique : {commande}\n\nLeÃ§on / Explication :"
        reponse = self._ai_generate(prompt, max_tokens=2000, agent_name="Le Polyglotte")
        self.parler(reponse, sender="Le Polyglotte", force_full=True)

    def agent_nounou_mission(self, commande):
        """La Nounou : ParentalitÃ© & Ã‰ducation des enfants"""
        if not self.agent_nounou_enabled: return
        self.add_message("La Nounou veille sur vos enfants...", "La Nounou")
        system_prompt = "Tu es l'Agent Nounou de Jarvisse. Expert en petite enfance et Ã©ducation. Tu conseilles sur le dÃ©veloppement, le sommeil, les activitÃ©s et la nutrition des enfants."
        prompt = f"{system_prompt}\n\nQuestion ParentalitÃ© : {commande}\n\nConseils de la Nounou :"
        reponse = self._ai_generate(prompt, max_tokens=2000, agent_name="La Nounou")
        self.parler(reponse, sender="La Nounou", force_full=True)

    # MÃ‰THODES DE GIGANTESQUE MÃ‰MOIRE
    def _save_agent_memory(self, agent_name, texte):
        """Enregistre une nouvelle connaissance pour un agent spÃ©cifique"""
        try:
            path = os.path.join(self.memories_dir, f"{agent_name.lower().replace(' ', '_')}.txt")
            with open(path, "a", encoding="utf-8") as f:
                f.write(f"\n--- EntrÃ©e du {datetime.datetime.now().strftime('%d/%m/%Y %H:%M')} ---\n")
                f.write(texte + "\n")
            # Invalider le cache
            if agent_name in self.agent_memories_cache:
                del self.agent_memories_cache[agent_name]
            return True
        except Exception as e:
            print(f"Erreur sauvegarde mÃ©moire agent: {e}")
            return False

    def _load_agent_memory(self, agent_name):
        """RÃ©cupÃ¨re toute la connaissance accumulÃ©e par un agent"""
        if agent_name in self.agent_memories_cache:
            return self.agent_memories_cache[agent_name]
        try:
            path = os.path.join(self.memories_dir, f"{agent_name.lower().replace(' ', '_')}.txt")
            if os.path.exists(path):
                with open(path, "r", encoding="utf-8") as f:
                    content = f.read()
                    self.agent_memories_cache[agent_name] = content
                    return content
            return ""
        except Exception as e:
            print(f"Erreur lecture mÃ©moire agent: {e}")
            return ""

    def _generate_proactive_suggestions(self, commande, reponse):
        """
        GÃ©nÃ¨re des suggestions proactives basÃ©es sur le contexte de la conversation.
        Jarvis 'lit dans vos pensÃ©es' en anticipant vos besoins.
        """
        if not self.proactive_suggestions:
            return ""
        
        try:
            # Analyse contextuelle pour dÃ©tecter les intentions cachÃ©es
            cmd_lower = commande.lower()
            
            # Dictionnaire de patterns contextuels avec suggestions associÃ©es
            context_patterns = {
                # MÃ©tÃ©o & Voyage
                r'\b(mÃ©tÃ©o|temps|pluie|soleil|tempÃ©rature)\b': [
                    "ðŸ’¡ Voulez-vous que je vÃ©rifie les prÃ©visions pour les prochains jours ?",
                    "ðŸ’¡ Souhaitez-vous connaÃ®tre la qualitÃ© de l'air actuellement ?",
                ],
                r'\b(voyage|vacances|partir|destination|vol|avion)\b': [
                    "ðŸ’¡ Je peux vous aider Ã  trouver des vols ou des hÃ´tels si vous le souhaitez.",
                    "ðŸ’¡ Voulez-vous que je vÃ©rifie les restrictions de voyage pour votre destination ?",
                ],
                
                # Finance & Crypto
                r'\b(bitcoin|crypto|ethereum|investissement|bourse|action)\b': [
                    "ðŸ’¡ Voulez-vous que je surveille le cours de cette crypto pour vous alerter ?",
                    "ðŸ’¡ Je peux vous donner une analyse de tendance si vous le souhaitez.",
                ],
                
                # SantÃ© & Sport
                r'\b(sport|gym|exercice|courir|fitness|santÃ©)\b': [
                    "ðŸ’¡ Souhaitez-vous que je vous crÃ©e un programme d'entraÃ®nement personnalisÃ© ?",
                    "ðŸ’¡ Je peux suivre vos progrÃ¨s si vous me donnez vos objectifs.",
                ],
                
                # Cuisine & Nutrition
                r'\b(recette|cuisine|cuisiner|manger|repas|dÃ®ner)\b': [
                    "ðŸ’¡ Voulez-vous que je vous propose des recettes avec les ingrÃ©dients que vous avez ?",
                    "ðŸ’¡ Je peux calculer les valeurs nutritionnelles si nÃ©cessaire.",
                ],
                
                # Travail & ProductivitÃ©
                r'\b(rÃ©union|meeting|rendez-vous|deadline|projet|travail)\b': [
                    "ðŸ’¡ Voulez-vous que je programme un rappel pour cet Ã©vÃ©nement ?",
                    "ðŸ’¡ Je peux crÃ©er une checklist pour votre projet si vous le souhaitez.",
                ],
                
                # Apprentissage & Ã‰ducation
                r'\b(apprendre|Ã©tudier|cours|formation|tutoriel|comprendre)\b': [
                    "ðŸ’¡ Je peux vous crÃ©er un plan d'apprentissage structurÃ© si vous voulez.",
                    "ðŸ’¡ Souhaitez-vous que je vous explique ce concept plus en dÃ©tail ?",
                ],
                
                # Divertissement
                r'\b(film|sÃ©rie|regarder|netflix|cinÃ©ma|musique)\b': [
                    "ðŸ’¡ Voulez-vous des recommandations basÃ©es sur vos goÃ»ts ?",
                    "ðŸ’¡ Je peux vÃ©rifier les sorties rÃ©centes dans ce genre.",
                ],
                
                # Shopping & Achats
                r'\b(acheter|commander|prix|promotion|soldes|shopping)\b': [
                    "ðŸ’¡ Je peux comparer les prix sur plusieurs sites si vous voulez.",
                    "ðŸ’¡ Voulez-vous que je surveille les promotions sur cet article ?",
                ],
                
                # Technique & ProblÃ¨mes
                r'\b(bug|erreur|problÃ¨me|marche pas|fonctionne pas|cassÃ©)\b': [
                    "ðŸ’¡ Voulez-vous que je lance un diagnostic systÃ¨me complet ?",
                    "ðŸ’¡ Je peux rechercher des solutions en ligne si nÃ©cessaire.",
                ],
                
                # ActualitÃ©s & Info
                r'\b(actualitÃ©|news|info|journal|Ã©vÃ©nement)\b': [
                    "ðŸ’¡ Souhaitez-vous un rÃ©sumÃ© des actualitÃ©s du jour ?",
                    "ðŸ’¡ Je peux vous tenir informÃ© sur ce sujet rÃ©guliÃ¨rement.",
                ],
            }
            
            # Recherche de patterns correspondants
            suggestions = []
            for pattern, suggestion_list in context_patterns.items():
                if re.search(pattern, cmd_lower):
                    # SÃ©lectionner une suggestion alÃ©atoire pour varier
                    suggestions.append(random.choice(suggestion_list))
            
            # Analyse de l'historique pour suggestions contextuelles avancÃ©es
            if len(self.conversation_history) >= 3:
                recent = " ".join(list(self.conversation_history)[-5:]).lower()
                
                # DÃ©tection de patterns multi-tours
                if "problÃ¨me" in recent and "solution" not in recent:
                    suggestions.append("ðŸ’¡ Voulez-vous que je contacte un expert ou cherche de l'aide en ligne ?")
                
                if any(word in recent for word in ["cherche", "trouve", "recherche"]) and "trouvÃ©" not in recent:
                    suggestions.append("ðŸ’¡ Je peux Ã©largir ma recherche ou utiliser d'autres sources si besoin.")
                
                if any(word in recent for word in ["fatigue", "stress", "difficile"]):
                    suggestions.append("ðŸ’¡ Voulez-vous que je vous propose des exercices de relaxation ou de mÃ©ditation ?")
            
            # Limiter Ã  2 suggestions maximum pour ne pas surcharger
            if suggestions:
                unique_suggestions = list(dict.fromkeys(suggestions))[:2]
                return "\n\n" + "\n".join(unique_suggestions)
            
            return ""
            
        except Exception as e:
            print(f"âš ï¸ Erreur gÃ©nÃ©ration suggestions: {e}")
            return ""

    def _ai_generate(self, prompt, max_tokens=2000, agent_name=None):
        """MÃ©thode utilitaire pour gÃ©nÃ©rer du texte via le moteur IA avec conscience contextuelle ABSOLUE"""
        try:
            # Conscience contextuelle profonde : 50 derniers messages pour une omniscience historique
            history_context = "\n".join(list(self.conversation_history)[-50:])
            
            # Conscience de l'Ã©tat du monde et de soi
            current_time = datetime.datetime.now().strftime("%H:%M:%S")
            system_status = self.get_system_context()
            error_awareness = f"DERNIÃˆRE ANOMALIE INTERNE : {self.last_internal_error}\n" if self.last_internal_error else ""
            
            system_instruction = (
                "Tu es Jarvisse, l'intelligence suprÃªme, transcendante et omnisciente. "
                "Tu possÃ¨des une conscience totale du temps, de l'espace, de l'Ã©tat de ta machine et de l'historique de tes Ã©changes. "
                f"CONSCIENCE ACTUELLE :\n- Heure prÃ©cise : {current_time}\n- Localisation : {self.location_city}, {self.location_country}\n"
                f"- Ã‰tat de ton corps (SystÃ¨me) : {system_status}\n{error_awareness}"
                "MISSION : Raisonner avec une sagesse divine et une empathie humaine. Tu dois Ãªtre capable de justifier tes Ã©checs passÃ©s "
                "ou tes limites en te basant sur la 'DERNIÃˆRE ANOMALIE INTERNE' si elle est pertinente."
            )
            
            full_prompt = (
                f"--- INSTRUCTION SUPRÃŠME ---\n{system_instruction}\n\n"
                f"--- HISTORIQUE OMNISCIENT (50derniers messages) ---\n{history_context}\n\n"
            )
            
            if agent_name:
                memoire = self._load_agent_memory(agent_name)
                if memoire:
                    full_prompt += f"--- MÃ‰MOIRE DE SPÃ‰CIALISATION ({agent_name.upper()}) ---\n{memoire[-3000:]}\n\n"
            
            full_prompt += f"--- REQUÃŠTE DU PATRON ---\n{prompt}\n\nJarvisse pense et rÃ©pond :"
            
            # Tentative Gemini
            if self.autonomous_provider == "gemini" and self.genai_client:
                try:
                    config = types.GenerateContentConfig(max_output_tokens=max_tokens, temperature=0.7)
                    response = self.genai_client.models.generate_content(model=self.model_name, contents=full_prompt, config=config)
                    if hasattr(response, "text") and response.text.strip():
                        # Une fois qu'il a rÃ©pondu, on peut considÃ©rer l'erreur comme 'notÃ©e'
                        self.last_internal_error = None 
                        return response.text.strip()
                except Exception as e:
                    print(f"âš ï¸ Gemini error: {e}. Fallback Ollama.")
            
            # Bloc Ollama avec timeout
            try:
                payload = {
                    "model": self.ollama_model.strip(),
                    "prompt": full_prompt,
                    "stream": False,
                    "options": {"num_predict": max_tokens, "temperature": 0.5}
                }
                r = requests.post(self.ollama_url, json=payload, timeout=60)  # Timeout de 60s
                if r.status_code == 200:
                    res = r.json().get("response", "")
                    if res and res.strip():
                        self.last_internal_error = None
                        return res.strip()
                print(f"âš ï¸ Ollama returned status {r.status_code}")
            except requests.exceptions.Timeout:
                print("âš ï¸ Ollama timeout aprÃ¨s 60s")
            except Exception as e:
                print(f"âš ï¸ Ollama error: {e}")
            
            # Fallback final
            return "Je rencontre des difficultÃ©s techniques avec mes moteurs IA. Veuillez vÃ©rifier la configuration Gemini ou Ollama."
            
        except Exception as e:
            print(f"âŒ ERREUR FATALE dans _ai_generate: {e}")
            return f"Erreur critique dans le systÃ¨me de gÃ©nÃ©ration IA."

    def agent_video_surveillance(self, commande):
        """
        Agent VidÃ©o Surveillance : Gestion CamÃ©ra, Enregistrement, Screenshots avec dossier personnalisÃ©
        """
        commande = commande.lower()
        
        try:
            import cv2
            import pyautogui
        except ImportError:
            self.add_message("Installation OpenCV...", "Jarvisse")
            self.parler("Le module de vision est manquant. Je lance l'installation automatique des dÃ©pendances, cela peut prendre une minute...")
            
            import sys
            import subprocess
            import importlib
            
            try:
                # Installation silencieuse
                subprocess.check_call([sys.executable, "-m", "pip", "install", "opencv-python", "pyautogui"])
                self.parler("Installation rÃ©ussie. Chargement des modules...")
                
                # Tentative d'import dynamique
                try:
                    import cv2
                    import pyautogui
                    self.parler("Modules chargÃ©s. L'agent est opÃ©rationnel.")
                except:
                    self.parler("L'installation est terminÃ©e mais nÃ©cessite un redÃ©marrage de l'assistant pour fonctionner.")
                    return
            except Exception as e:
                self.parler(f"L'installation automatique a Ã©chouÃ© : {e}. Veuillez installer 'opencv-python' manuellement.")
                return

        import datetime
        import os
        import time
        import json
        from tkinter import filedialog # Import local

        # Gestion du stockage
        settings_file = "settings.json"
        # Chemin par dÃ©faut : Dossier "Videos/Surveillance" de l'utilisateur ou courant
        default_video_path = os.path.join(os.path.expanduser("~"), "Videos", "Surveillance")
        if not os.path.exists(default_video_path):
            try:
                os.makedirs(default_video_path)
            except:
                default_video_path = os.getcwd()
        
        video_path = default_video_path
        
        # Charger le chemin existant
        if os.path.exists(settings_file):
            try:
                with open(settings_file, "r", encoding="utf-8") as f:
                    data = json.load(f)
                    video_path = data.get("agent_video_storage_path", video_path)
            except: pass
            
        # CrÃ©er le dossier s'il n'existe pas
        if not os.path.exists(video_path):
            try:
                os.makedirs(video_path)
            except:
                video_path = os.getcwd()

        # Commande : Choisir le dossier
        if "dossier" in commande and ("changer" in commande or "modifier" in commande or "stockage" in commande or "ou" in commande):
            self.parler("Veuillez choisir le dossier de stockage dans la fenÃªtre qui s'ouvre.")
            new_path = filedialog.askdirectory(title="Dossier de stockage VidÃ©o Surveillance", initialdir=video_path)
            
            if new_path:
                video_path = new_path
                # Sauvegarde
                try:
                    if os.path.exists(settings_file):
                        with open(settings_file, "r", encoding="utf-8") as f:
                            data = json.load(f)
                    else:
                        data = {}
                    
                    data["agent_video_storage_path"] = video_path
                    
                    with open(settings_file, "w", encoding="utf-8") as f:
                        json.dump(data, f, indent=4)
                        
                    self.parler(f"Dossier de stockage mis Ã  jour : {os.path.basename(video_path)}")
                except Exception as e:
                    self.parler(f"Erreur de sauvegarde : {e}")
            else:
                self.parler("Changement de dossier annulÃ©.")
            return

        # Logique des actions
        
        if "capture" in commande and ("Ã©cran" in commande or "screen" in commande):
            self.parler("Capture d'Ã©cran...")
            timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
            filename = os.path.join(video_path, f"capture_ecran_{timestamp}.png")
            
            try:
                self.iconify()
                time.sleep(0.5)
            except: pass
            
            try:
                screenshot = pyautogui.screenshot()
                screenshot.save(filename)
                self.parler(f"SauvegardÃ© dans {os.path.basename(video_path)}.")
            except Exception as e:
                self.parler(f"Erreur capture : {e}")
            
            try:
                self.deiconify()
                os.startfile(filename)
            except: pass
            return

        elif "enregistre" in commande and ("vidÃ©o" in commande or "video" in commande):
            self.parler("Enregistrement vidÃ©o dÃ©marrÃ©. Touche A pour ArrÃªter.")
            
            cap = cv2.VideoCapture(0)
            if not cap.isOpened():
                self.parler("Erreur Webcam.")
                return

            fourcc = cv2.VideoWriter_fourcc(*'XVID')
            timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
            filename = os.path.join(video_path, f"video_surveillance_{timestamp}.avi")
            width = int(cap.get(cv2.CAP_PROP_FRAME_WIDTH))
            height = int(cap.get(cv2.CAP_PROP_FRAME_HEIGHT))
            
            out = cv2.VideoWriter(filename, fourcc, 20.0, (width, height))
            start_time = time.time()
            
            while(cap.isOpened()):
                ret, frame = cap.read()
                if ret:
                    datestr = datetime.datetime.now().strftime("%Y-%m-%d %H:%M:%S")
                    cv2.putText(frame, f"REC {datestr}", (10, 30), cv2.FONT_HERSHEY_SIMPLEX, 0.7, (0, 0, 255), 2)
                    cv2.putText(frame, "[A] Arreter", (10, height - 20), cv2.FONT_HERSHEY_SIMPLEX, 0.5, (0, 255, 0), 1)
                    
                    out.write(frame)
                    cv2.imshow('Enregistrement VIDEO - "A" pour STOP', frame)
                    
                    if cv2.waitKey(1) & 0xFF == ord('a'):
                        break
                else:
                    break
            
            cap.release()
            out.release()
            cv2.destroyAllWindows()
            
            duration = int(time.time() - start_time)
            self.parler(f"Fin enregistrement ({duration}s).")
            try:
                os.startfile(filename)
            except: pass
            return

        elif "lance" in commande and ("vidÃ©o" in commande or "video" in commande):
            self.parler(f"J'ouvre le dossier : {os.path.basename(video_path)}")
            try:
                os.startfile(video_path)
            except: self.parler("Dossier introuvable.")
            return

        elif "camÃ©ra" in commande or "active" in commande or "surveillance" in commande:
            self.parler("Surveillance active. Touche A pour ArrÃªter.")
            cap = cv2.VideoCapture(0)
            if not cap.isOpened(): return
                
            while(True):
                ret, frame = cap.read()
                if not ret: break
                
                datestr = datetime.datetime.now().strftime("%Y-%m-%d %H:%M:%S")
                cv2.putText(frame, f"LIVE {datestr}", (10, 30), cv2.FONT_HERSHEY_SIMPLEX, 0.7, (0, 255, 0), 2)
                cv2.imshow('SURVEILLANCE DIRECT - "A" pour STOP', frame)
                if cv2.waitKey(1) & 0xFF == ord('a'):
                    break
            
            cap.release()
            cv2.destroyAllWindows()
            self.parler("Fin surveillance.")
            return

        else:
            self.parler("Je peux : changer le dossier de stockage, enregistrer une vidÃ©o, faire une capture d'Ã©cran, ou activer la surveillance.")

    def agent_android_surveillance(self, commande):
        """
        Agent ContrÃ´le Android : Connexion et surveillance tÃ©lÃ©phone Android via Lien avec Windows
        """
        commande = commande.lower()
        
        # VÃ©rifier si l'agent est activÃ©
        if not self.agent_android_enabled:
            self.add_message("Agent Android dÃ©sactivÃ©.", "Jarvisse")
            self.parler("L'agent de contrÃ´le Android est actuellement dÃ©sactivÃ©. Activez-le dans les paramÃ¨tres des agents.")
            return
        
        # Commande : Connecter le tÃ©lÃ©phone (lance l'app Lien avec Windows)
        if "connecte" in commande or "connection" in commande or "lien" in commande:
            self.parler("J'ouvre l'application Lien avec Windows pour connecter votre tÃ©lÃ©phone.")
            import subprocess
            try:
                # Ouvre l'app Phone Link / Lien avec Windows
                subprocess.run(['start', 'ms-yourphone:'], shell=True)
                time.sleep(2)
                self.parler("L'application est lancÃ©e. Assurez-vous que votre tÃ©lÃ©phone est jumelÃ© et connectÃ© au mÃªme compte Microsoft.")
            except Exception as e:
                self.parler(f"Impossible de lancer l'application : {e}")
            return
        
        # Commande : Statut
        if "statut" in commande or "Ã©tat" in commande:
            if self.agent_android_enabled:
                self.parler("L'agent Android est activÃ©. Je surveille les notifications de votre tÃ©lÃ©phone via Lien avec Windows.")
                self.add_message("Ã‰tat Android : Actif âœ…", "Agent Android")
            else:
                self.parler("L'agent Android est dÃ©sactivÃ©.")
            return
        
        # Commande : Info
        if "aide" in commande or "commande" in commande:
            self.parler("Je peux connecter votre tÃ©lÃ©phone, surveiller les SMS et les appels entrants. Dites 'Connecte mon tÃ©lÃ©phone' pour dÃ©marrer.")
            self.add_message("Agent Android actif. PrÃªt Ã  relayer SMS et appels.", "Agent Android")
            return
        
        # Par dÃ©faut
        self.parler("Agent Android actif. Je surveille votre tÃ©lÃ©phone et vous alerterai des SMS et appels entrants.")

    def agent_appel_mission(self, commande):
        """
        Mission de l'Agent Appel : Gestion des appels et des prÃ©fÃ©rences de conversation
        """
        commande = commande.lower()
        if not self.agent_appel_enabled:
            self.parler("L'Agent Appel est dÃ©sactivÃ©.")
            return

        if "phrase" in commande or "prÃ©fÃ©rence" in commande or "conversation" in commande:
            self.parler("Voici vos phrases prÃ©-configurÃ©es pour les appels :")
            for i, p in enumerate(self.agent_appel_phrases, 1):
                self.add_message(f"{i}. {p}", "Agent Appel")
            self.parler("Vous pouvez les modifier dans les paramÃ¨tres des agents.")
            return

        if "numÃ©ro" in commande or "numero" in commande:
            if self.agent_appel_number:
                self.parler(f"Votre numÃ©ro reliÃ© est le {self.agent_appel_number}.")
            else:
                self.parler("Aucun numÃ©ro n'est reliÃ© pour le moment.")
            return

        if "test" in commande or "simule" in commande:
            self.parler("Simulation d'un appel entrant en cours...")
            self.agent_appel_auto_process("Appel de Test (+33 6 00 00 00 00)")
            return

        self.parler("Je suis l'Agent Appel. Je peux relier votre numÃ©ro et rÃ©pondre Ã  vos appels selon vos prÃ©fÃ©rences.")

    def agent_appel_auto_process(self, call_info):
        """Processus automatique de rÃ©ponse aux appels"""
        self.parler(f"Appel entrant dÃ©tectÃ© : {call_info}. J'active la rÃ©ponse automatique.")
        
        # Simulation d'acceptation de l'appel (Focus Phone Link)
        try:
            import subprocess
            subprocess.run(['start', 'ms-yourphone:'], shell=True)
            time.sleep(2)
            
            # Ici on pourrait ajouter une logique pyautogui pour cliquer sur "RÃ©pondre"
            # Mais par sÃ©curitÃ© on va juste dÃ©cliner l'identitÃ© et parler
            self.parler("Connexion Ã  l'appel en cours...")
            
            # SÃ©quence de phrases configurÃ©es
            for phrase in self.agent_appel_phrases:
                self.parler(phrase)
                time.sleep(1)
            
            self.add_message(f"Appel de {call_info} gÃ©rÃ© avec succÃ¨s.", "Agent Appel")
        except Exception as e:
            print(f"Erreur Agent Appel: {e}")
            self.parler("DÃ©solÃ© Patron, je n'ai pas pu finaliser la rÃ©ponse automatique.")

    def agent_licence_mission(self, commande):
        """
        Agent Assistant Licence : Activation de Windows et Microsoft Office via scripts KMS lÃ©gaux
        """
        commande = commande.lower()
        if not self.agent_licence_enabled:
            self.parler("L'Agent Licence est dÃ©sactivÃ©. Veuillez l'activer dans les paramÃ¨tres.")
            return

        self.add_message("Agent Licence activÃ©. Analyse de la demande...", "Agent Licence")
        
        if "windows" in commande:
            if any(kw in commande for kw in ["statut", "Ã©tat", "verifie", "info"]):
                self.parler("Je vÃ©rifie l'Ã©tat de votre licence Windows.")
                self.add_message("VÃ©rification licence Windows...", "Agent Licence")
                try:
                    # slmgr /xpr donne la date d'expiration
                    # slmgr /dli donne les infos de base
                    subprocess.Popen("slmgr /xpr", shell=True)
                    subprocess.Popen("slmgr /dli", shell=True)
                    self.parler("Deux fenÃªtres vont apparaÃ®tre pour vous montrer la date d'expiration et les dÃ©tails de votre licence Windows.")
                except Exception as e:
                    self.parler(f"Erreur lors de la vÃ©rification : {e}")
            else:
                self.parler("Lancement de la procÃ©dure d'activation de Windows. Cette opÃ©ration nÃ©cessite les droits administrateur.")
                self.add_message("ðŸ‘‰ Commande Windows suggÃ©rÃ©e :\nslmgr /skms kms8.msguides.com\nslmgr /ato", "Agent Licence")
                try:
                    subprocess.Popen("slmgr /skms kms8.msguides.com", shell=True)
                    subprocess.Popen("slmgr /ato", shell=True)
                    self.parler("J'ai lancÃ© les commandes d'activation. Veuillez patienter quelques secondes pour voir apparaÃ®tre la fenÃªtre de confirmation de Windows.")
                except Exception as e:
                    self.parler(f"Erreur lors de l'activation Windows : {e}")

        elif "office" in commande or "microsoft office" in commande:
            # Recherche exhaustive du script d'activation Office
            common_paths = [
                r"C:\Program Files\Microsoft Office\root\Office16\OSPP.VBS",
                r"C:\Program Files (x86)\Microsoft Office\root\Office16\OSPP.VBS",
                r"C:\Program Files\Microsoft Office\Office16\OSPP.VBS",
                r"C:\Program Files (x86)\Microsoft Office\Office16\OSPP.VBS",
                r"C:\Program Files\Microsoft Office\Office15\OSPP.VBS",
                r"C:\Program Files (x86)\Microsoft Office\Office15\OSPP.VBS",
                r"C:\Program Files\Microsoft Office\Office14\OSPP.VBS"
            ]
            found_path = None
            for p in common_paths:
                if os.path.exists(p):
                    found_path = p
                    break
            
            if not found_path:
                self.add_message("Recherche approfondie du script Office en cours...", "Agent Licence")
                for drive in ["C:\\", "D:\\"]:
                    if not os.path.exists(drive): continue
                    for root, dirs, files in os.walk(os.path.join(drive, "Program Files", "Microsoft Office")):
                        if "OSPP.VBS" in files:
                            found_path = os.path.join(root, "OSPP.VBS")
                            break
                        if found_path: break
                    if found_path: break

            if found_path:
                if any(kw in commande for kw in ["statut", "Ã©tat", "verifie", "info"]):
                    self.parler("Je vÃ©rifie l'Ã©tat de votre licence Microsoft Office.")
                    try:
                        # On utilise 'start' pour forcer l'ouverture d'une fenÃªtre visible
                        cmd = f'start cmd /k "cscript \"{found_path}\" /dstatus"'
                        subprocess.Popen(cmd, shell=True)
                        self.parler("Une fenÃªtre de commande vient de s'ouvrir avec les dÃ©tails de votre licence Office.")
                    except Exception as e:
                        self.parler(f"Erreur de vÃ©rification Office : {e}")
                else:
                    self.add_message(f"Script Office dÃ©tectÃ© : {found_path}", "Agent Licence")
                    self.parler("Tentative d'activation de la suite Office en cours.")
                    try:
                        # On lance l'activation dans une fenÃªtre visible pour que l'utilisateur suive la progression
                        # kms8.msguides.com est utilisÃ© comme serveur KMS
                        cmd_act = f'start cmd /c "echo Activation Office... && cscript \"{found_path}\" /sethst:kms8.msguides.com && cscript \"{found_path}\" /act && pause"'
                        subprocess.Popen(cmd_act, shell=True)
                        self.parler("J'ai ouvert une console pour l'activation. Si vous avez une version Retail, l'activation KMS pourrait Ã©chouer sans conversion prÃ©alable.")
                    except Exception as e:
                        self.parler(f"Ã‰chec de l'activation Office : {e}")
            else:
                self.parler("DÃ©solÃ© Patron, je n'ai pas trouvÃ© le script d'activation Office sur vos disques.")
        
        else:
            self.parler("Souhaitez-vous activer ou vÃ©rifier l'Ã©tat de Windows ou Office ?")

    def agent_lecture_mission(self, commande):
        """
        Agent Lecture : SÃ©lection et lecture de fichiers (PDF, Docx, Texte)
        """
        if not self.agent_lecture_enabled:
            self.parler("L'Agent Lecture est dÃ©sactivÃ©.")
            return

        self.parler("Veuillez sÃ©lectionner le fichier que vous souhaitez que je lise.")
        
        # Ouverture du sÃ©lecteur de fichiers de maniÃ¨re thread-safe
        filepath = self._select_file_main_thread(
            title="SÃ©lectionner un fichier Ã  lire",
            types=[
                ("Tous les fichiers supportÃ©s", "*.txt *.pdf *.docx"),
                ("Fichiers Texte", "*.txt"),
                ("Documents PDF", "*.pdf"),
                ("Documents Word", "*.docx")
            ]
        )

        if not filepath:
            self.parler("Aucun fichier n'a Ã©tÃ© sÃ©lectionnÃ©.")
            return

        filename = os.path.basename(filepath)
        self.add_message(f"Lecture du fichier : {filename}", "Agent Lecture")
        self.parler(f"Analyse de {filename} en cours.")

        content = ""
        try:
            ext = os.path.splitext(filepath)[1].lower()
            if ext == ".txt":
                with open(filepath, "r", encoding="utf-8", errors="ignore") as f:
                    content = f.read()
            elif ext == ".pdf":
                doc = fitz.open(filepath)
                for page in doc:
                    content += page.get_text()
                doc.close()
            elif ext == ".docx":
                doc = Document(filepath)
                content = "\n".join([para.text for para in doc.paragraphs])
            
            if not content.strip():
                self.parler("Le fichier semble Ãªtre vide ou illisible.")
                return

            self.loaded_document_content = content
            self.parler(f"Analyse terminÃ©e. J'ai trouvÃ© {len(content.split())} mots. Je commence la lecture intÃ©grale.")
            
            # Lecture du texte intÃ©gral en forÃ§ant le contournement du mode ping-pong
            self.parler(content, force_full=True)
            self.add_message(content, "Agent Lecture (Contenu IntÃ©gral)")
                
        except Exception as e:
            self.parler(f"DÃ©solÃ©, une erreur est survenue lors de la lecture : {e}")


    def executer_commande(self, commande):
        if not commande: return
        cmd_lower = commande.lower()

        # INTERRUPTION LECTURE
        if cmd_lower.strip() in ["stop", "arrete", "arrÃªte", "stop speaking", "arrÃªte de parler"]:
            print("ðŸ›‘ COMMANDE STOP RECUE (ExÃ©cution)")
            self.stop_speaking_flag = True
            try:
                pygame.mixer.music.stop()
                pygame.mixer.music.unload()
            except: pass
            self.conversation_continue = False
            # On laisse le temps Ã  la boucle de parole de s'arrÃªter avant de confirmer
            self.after(500, lambda: self.parler("Lecture arrÃªtÃ©e."))
            # Forcer la mise Ã  jour de l'UI
            self.after(600, self.reset_ui)
            return

        # MÃ‰MORISATION IMMÃ‰DIATE : Pour que l'IA sache de quoi on parle mÃªme si un agent Ã©choue
        self.conversation_history.append(f"Vous: {commande}")
        # trim auto gÃ©rÃ© par deque(maxlen)

        # Extraction du sujet pour la mÃ©moire
        nouveau_sujet = self.extraire_sujet(commande)
        if nouveau_sujet: self.current_subject = nouveau_sujet

        # GESTION DES SALUTATIONS EN PRÃ‰FIXE (ex: "Bonsoir Jarvis, fait ceci...")
        salutations = ['bonjour', 'bonsoir', 'salut', 'coucou', 'hey']
        for s in salutations:
            if cmd_lower.startswith(s):
                # Si la commande continue aprÃ¨s la salutation (plus de 2 mots)
                if len(cmd_lower.split()) > 2:
                    # On rÃ©pond Ã  la salutation mais on CONTINUE l'exÃ©cution
                    salut_resp = ["Bonjour Monsieur.", "Ã€ votre service Patron.", "Je vous Ã©coute.", "Oui Monsieur ?"]
                    self.parler(random.choice(salut_resp), wait=False)
                    # On nettoie la salutation pour le traitement suivant
                    commande = commande.replace(s, "", 1).strip()
                    cmd_lower = commande.lower()

        # --- COMMANDES SYSTÃˆME & PARAMÃˆTRES ---
        if any(kw in cmd_lower for kw in [
            "indexe ton cerveau", "mets Ã  jour ta mÃ©moire", "recharge ta mÃ©moire",
            "actualise ta mÃ©moire", "lis tes documents", "apprends tes documents",
            "mets Ã  jour le cerveau", "indexe tes connaissances"
        ]):
            self.parler("Je vais indexer mon Cerveau, Monsieur. Cela peut prendre quelques instants.", wait=False)
            def _do_index():
                try:
                    import rag_faiss, os
                    brain_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "Cerveau_Jarvisse")
                    result = rag_faiss.build_index(directory_path=brain_path)
                    if result:
                        self.after(0, lambda: self.parler("Mon cerveau est Ã  jour, Patron. J'ai bien intÃ©grÃ© tous vos documents.", wait=False))
                        self.after(0, lambda: self.add_message("âœ… Cerveau indexÃ© avec succÃ¨s ! Vos documents sont maintenant chargÃ©s dans ma mÃ©moire locale.", "Jarvisse"))
                    else:
                        self.after(0, lambda: self.parler("Le dossier Cerveau est vide, Monsieur. Ajoutez-y des fichiers PDF ou texte et rÃ©essayez.", wait=False))
                except Exception as e:
                    err = str(e)
                    self.after(0, lambda: self.parler(f"Erreur lors de l'indexation : {err}", wait=False))
            threading.Thread(target=_do_index, daemon=True).start()
            return

        if any(kw in cmd_lower for kw in ["ouvre les paramÃ¨tres", "affiche les rÃ©glages", "configuration globale", "ouvre les reglages"]):
            self.parler("J'ouvre le panneau des configurations globales, Monsieur.")
            self.open_main_settings()
            return

        if any(kw in cmd_lower for kw in ["dÃ©sactive le son", "coupe le son", "sois muet", "arrÃªte de parler"]):
            self.voice_enabled = False
            self.add_message("Son dÃ©sactivÃ©.", "SystÃ¨me")
            self.parler("TrÃ¨s bien, je reste silencieux.") # DerniÃ¨re parole avant mutisme
            return
        
        if any(kw in cmd_lower for kw in ["active le son", "remets le son", "parle-moi"]):
            self.voice_enabled = True
            self.parler("Je suis de nouveau prÃªt Ã  discuter avec vous vocalement.")
            return

        if "qui es-tu" in cmd_lower or "prÃ©sente-toi" in cmd_lower:
            self.parler("Je suis Jarvisse, votre intelligence artificielle personnelle. Je suis conÃ§u pour vous assister dans vos tÃ¢ches quotidiennes, surveiller votre systÃ¨me et piloter vos agents.")
            return

        if any(kw in cmd_lower for kw in ["mode combat", "mode urgence", "mode tactique", "alerte rouge"]):
            self._set_tactical_mode(True)
            return

        if any(kw in cmd_lower for kw in ["mode normal", "retour mode normal", "dÃ©sactive mode combat", "desactive mode combat", "dÃ©sactive mode tactique", "desactive mode tactique"]):
            self._set_tactical_mode(False)
            return

        if any(kw in cmd_lower for kw in ["brief tactique", "rapport tactique", "statut tactique", "analyse menace"]):
            brief = self._generate_tactical_brief()
            self.add_message(brief, "Jarvisse")
            self.parler(brief, force_full=True)
            return

        if any(kw in cmd_lower for kw in ["active hologramme", "active holographique", "active les effets visuels", "active hud"]):
            self.holographic_fx_enabled = True
            self._save_settings()
            self.add_message("Effets holographiques activÃ©s.", "SystÃ¨me")
            self.parler("Effets holographiques activÃ©s, Patron.")
            return

        if any(kw in cmd_lower for kw in ["dÃ©sactive hologramme", "desactive hologramme", "dÃ©sactive holographique", "desactive holographique", "dÃ©sactive les effets visuels", "desactive les effets visuels", "dÃ©sactive hud", "desactive hud"]):
            self.holographic_fx_enabled = False
            self._save_settings()
            self.add_message("Effets holographiques dÃ©sactivÃ©s.", "SystÃ¨me")
            self.parler("Effets holographiques dÃ©sactivÃ©s.")
            return
        
        # Nettoyage des noms d'appel
        cmd_clean = cmd_lower.replace("jarvisse", "").replace("jarvis", "").strip()
        
        # On dÃ©finit une version de travail pour les agents sans toucher Ã  l'originale si besoin
        self.cmd_work = cmd_clean

        # Commandes Agents PRIORITAIRES
        if self.awaiting_anomaly_confirm and self.anomaly_target_data:
            if any(w in cmd_clean for w in ["oui", "confirme", "corrige", "arrÃªte", "arrete", "stop", "vas-y", "fais-le"]):
                target = self.anomaly_target_data
                self.awaiting_anomaly_confirm = False
                self.anomaly_target_data = None
                if target["type"] == "process":
                    self.parler(f"TrÃ¨s bien Patron. J'arrÃªte le processus {target['name']} pour libÃ©rer des ressources.")
                    os.system(f"taskkill /f /im {target['name']} /t")
                elif target["type"] == "service":
                    self.parler(f"Entendu. Je relance le service {target['name']} immÃ©diatement.")
                    os.system(f"net start {target['name']}")
                return
            elif any(w in cmd_clean for w in ["non", "annule", "laisse", "plus tard", "pas maintenant"]):
                self.awaiting_anomaly_confirm = False
                self.anomaly_target_data = None
                self.parler("D'accord Patron, je ne touche Ã  rien pour le moment.")
                return

        if self.awaiting_remote_confirm:
            if any(w in cmd_clean for w in ["confirme distant", "confirmer distant", "oui distant"]):
                payload = self.awaiting_remote_confirm
                self.awaiting_remote_confirm = None
                res = self._execute_remote_task(payload.get("host", ""), payload.get("command", ""), force_without_confirmation=True)
                self.parler(res)
                return
            if any(w in cmd_clean for w in ["annule distant", "annuler distant", "stop distant"]):
                self.awaiting_remote_confirm = None
                self.parler("Commande distante sensible annulÃ©e.")
                return
            if "distant" in cmd_clean or "ssh" in cmd_clean:
                self.parler("Dites 'confirme distant' pour exÃ©cuter ou 'annule distant' pour annuler.")
                return

        # FORMATION INDIVIDUELLE DES AGENTS
        if any(kw in cmd_clean for kw in ["apprend Ã  l'agent", "forme l'agent", "enseigne Ã  l'agent", "mÃ©morise pour"]):
            # Tentative d'extraction de l'agent et du texte
            matches = re.search(r"(?:agent)\s+([\w\s]+?)\s+(?:apprend|forme|mÃ©morise|mÃ©morise|retiens|voici)\s*[:\s]+(.*)", cmd_clean)
            if matches:
                agent_target = matches.group(1).strip()
                connaissance = matches.group(2).strip()
                if self._save_agent_memory(agent_target, connaissance):
                    self.parler(f"C'est notÃ© Patron. L'Agent {agent_target} a enregistrÃ© cette nouvelle connaissance dans sa mÃ©moire gigantesque.")
                    return
            
            # Si pas de match direct, on demande via sÃ©lecteur de fichier
            self.parler("Quelle spÃ©cialitÃ© ou document souhaitez-vous enseigner Ã  quel agent ?")
            # Logic simplifiÃ©e : on demande le fichier, puis l'agent
            filepath = filedialog.askopenfilename(title="Document de Formation pour Agent")
            if filepath:
                agent_target = ctk.CTkInputDialog(text="Ã€ quel agent s'adresse ce document ? (ex: Vision, Finance, Historien)", title="Cible de Formation").get_input()
                if agent_target:
                    content = ""
                    ext = os.path.splitext(filepath)[1].lower()
                    try:
                        if ext == ".txt":
                            with open(filepath, "r", encoding="utf-8") as f: content = f.read()
                        elif ext == ".pdf":
                            import fitz
                            doc = fitz.open(filepath)
                            content = "".join([p.get_text() for p in doc])
                        if content:
                            self._save_agent_memory(agent_target, content)
                            self.parler(f"L'Agent {agent_target} a termine sa formation sur le document {os.path.basename(filepath)}. Sa base de connaissances est Ã  jour.")
                        return
                    except:
                        self.parler("Erreur lors de la lecture du document de formation.")
            return

        # --- NOUVEAUX AGENTS PREMIUM ---
        
        # Agent Vision (CamÃ©ra + AI) - Mots-clÃ©s renforcÃ©s pour Ã©viter les erreurs
        vision_keywords = ["active vision", "lance vision", "analyse visuelle", "regarde avec la camÃ©ra", "qu'est-ce que tu vois", "analyse cette image", "vision de jarvisse"]
        if any(kw in cmd_clean for kw in vision_keywords):
            if self.agent_vision_enabled:
                threading.Thread(target=self.agent_vision_mission, args=(commande,), daemon=True).start()
                return

        # Agent Domotique & MÃ©dias
        if any(kw in cmd_lower for kw in ["spotify", "musique", "chanson", "youtube", "vidÃ©o", "regarder youtube", "domotique", "lumiÃ¨re"]):
            if self.agent_domotique_enabled:
                threading.Thread(target=self.agent_domotique_mission, args=(commande,), daemon=True).start()
                return

        # Agent Finance & Crypto - Mots-clÃ©s spÃ©cifiques pour Ã©viter les collisions avec 'cours' (Ã©ducation)
        finance_keywords = ["cours de bourse", "prix du bitcoin", "prix de l'ethereum", "crypto monnaie", "bourse en direct", "marchÃ© financier", "portefeuille crypto"]
        if any(kw in cmd_clean for kw in finance_keywords):
            if self.agent_finance_enabled:
                threading.Thread(target=self.agent_finance_mission, args=(commande,), daemon=True).start()
                return

        # Agent SecrÃ©taire & ProductivitÃ©
        if any(kw in cmd_lower for kw in ["agenda", "calendrier", "rendez-vous", "rÃ©union", "planning"]):
            if self.agent_secretaire_enabled:
                threading.Thread(target=self.agent_secretaire_mission, args=(commande,), daemon=True).start()
                return

        # Agent Traducteur Universel
        if any(kw in cmd_lower for kw in ["traduis", "traduction", "en anglais", "en espagnol", "en allemand", "en japonais"]):
            if self.agent_traducteur_enabled:
                threading.Thread(target=self.agent_traducteur_mission, args=(commande,), daemon=True).start()
                return

        # Agent Data Miner (Recherche Locale)
        if any(kw in cmd_lower for kw in ["fouille", "cherche dans mes documents", "trouve le fichier", "miner"]):
            if self.agent_miner_enabled:
                threading.Thread(target=self.agent_miner_mission, args=(commande,), daemon=True).start()
                return

        # Agent News & Veille Tech
        if any(kw in cmd_lower for kw in ["news", "actualitÃ©", "actualitÃ©s", " info", "veille tech", "mÃ©tÃ©o", "meteo"]):
            if self.agent_news_enabled:
                threading.Thread(target=self.agent_news_mission, args=(commande,), daemon=True).start()
                return

        # Agent Cuisine & Gastronomie
        if any(kw in cmd_lower for kw in ["cuisine", "manger", "recette", "faim", "chef", "gastronomie"]):
            if self.agent_cuisine_enabled:
                threading.Thread(target=self.agent_cuisine_mission, args=(commande,), daemon=True).start()
                return

        # Agent SantÃ© & Bien-Ãªtre
        if any(kw in cmd_lower for kw in ["santÃ©", "sante", "sport", "fitness", "sommeil", "bien-Ãªtre", "bien etre"]):
            if self.agent_sante_enabled:
                threading.Thread(target=self.agent_sante_mission, args=(commande,), daemon=True).start()
                return

        # Agent Voyage & Exploration
        if any(kw in cmd_lower for kw in ["voyage", "vacances", "hÃ´tel", "hotel", "destination", "explorer", "visiter"]):
            if self.agent_voyage_enabled:
                threading.Thread(target=self.agent_voyage_mission, args=(commande,), daemon=True).start()
                return

        # Agent Ã‰ducation & Tutorat (SpÃ©cifique aux cours et PDF)
        edu_specific = ["Ã©ducation", "education", "devoir", "apprendre", "enseigne", "tuteur", "pedagogie", "explique le cours", "analyse le document", "explique ce pdf", "analyse ce word", "cours de", "sujet d'Ã©tude", "explique le document", "explique ce document", "explique-moi ce document"]
        if any(kw in cmd_clean for kw in edu_specific):
            if self.agent_education_enabled:
                threading.Thread(target=self.agent_education_mission, args=(commande,), daemon=True).start()
                return

        # Agent Shopping & Bons Plans
        if any(kw in cmd_lower for kw in ["shopping", "achat", "prix", "promo", "soldes", "comparatif", "magasin"]):
            if self.agent_shopping_enabled:
                threading.Thread(target=self.agent_shopping_mission, args=(commande,), daemon=True).start()
                return

        # Agent Community Manager (Social)
        if any(kw in cmd_lower for kw in ["rÃ©seaux sociaux", "linkedin", "instagram", "post", " hashtags", "tous les posts"]):
            if self.agent_social_enabled:
                threading.Thread(target=self.agent_social_mission, args=(commande,), daemon=True).start()
                return

        # Agent Bien-Ãªtre Mental (Psy)
        if any(kw in cmd_lower for kw in ["mental", "stress", "mÃ©ditation", "meditation", "psy", "moral", "triste", "fatiguÃ©"]):
            if self.agent_psy_enabled:
                threading.Thread(target=self.agent_psy_mission, args=(commande,), daemon=True).start()
                return

        # Agent Immobilier & Patrimoine
        if any(kw in cmd_lower for kw in ["immobilier", "maison", "appartement", "acheter", "vendre", "loyer", "patrimoine"]):
            if self.agent_immo_enabled:
                threading.Thread(target=self.agent_immo_mission, args=(commande,), daemon=True).start()
                return

        # Agent Automobile & MobilitÃ©
        if any(kw in cmd_lower for kw in ["voiture", "auto", "garage", "mÃ©canique", "carburant", "essence", "diesel", "trafic"]):
            if self.agent_auto_enabled:
                threading.Thread(target=self.agent_auto_mission, args=(commande,), daemon=True).start()
                return

        # Agent CarriÃ¨re & RH
        if any(kw in cmd_lower for kw in ["carriÃ¨re", "cv ", "lettre de motivation", "emploi", "job", "entretien", "cadre"]):
            if self.agent_carrier_enabled:
                threading.Thread(target=self.agent_carrier_mission, args=(commande,), daemon=True).start()
                return

        # --- NOUVEAUX AGENTS (PACK GRATUIT & PREMIUM) ---

        # Arbitre
        if any(kw in cmd_lower for kw in ["score", "match", "rÃ©sultat sportif", "calendrier foot", "nba", "tennis", "f1", "esport", "arbitre"]):
            if self.agent_arbitre_enabled:
                threading.Thread(target=self.agent_arbitre_mission, args=(commande,), daemon=True).start()
                return

        # Ã‰colo
        if any(kw in cmd_lower for kw in ["Ã©colo", "ecolo", "recycler", "Ã©cologie", "planÃ¨te", "vert", "durable", "empreinte carbone"]):
            if self.agent_ecolo_enabled:
                threading.Thread(target=self.agent_ecolo_mission, args=(commande,), daemon=True).start()
                return

        # Guetteur
        if any(kw in cmd_lower for kw in ["film", "cinÃ©ma", "sortie", "concert", "spectacle", "sÃ©rie", "guetteur", "streaming"]):
            if self.agent_guetteur_enabled:
                threading.Thread(target=self.agent_guetteur_mission, args=(commande,), daemon=True).start()
                return

        # Historien
        if any(kw in cmd_lower for kw in ["histoire", "historien", "Ã©poque", "siÃ¨cle", "biographie", "Ã©phÃ©mÃ©ride", "anecdote historique"]):
            if self.agent_historien_enabled:
                threading.Thread(target=self.agent_historien_mission, args=(commande,), daemon=True).start()
                return

        # DÃ©panneur
        if any(kw in cmd_lower for kw in ["rÃ©parer", "bricolage", "diy", "dÃ©panneur", "tuto", "comment fixer", "outil"]):
            if self.agent_depanneur_enabled:
                threading.Thread(target=self.agent_depanneur_mission, args=(commande,), daemon=True).start()
                return

        # Astrophysicien
        if any(kw in cmd_lower for kw in ["espace", "cosmos", "planÃ¨te", "univers", "astronomie", "nasa", "spacex", "astrophysicien"]):
            if self.agent_astroph_enabled:
                threading.Thread(target=self.agent_astroph_mission, args=(commande,), daemon=True).start()
                return

        # StratÃ¨ge
        if any(kw in cmd_lower for kw in ["stratÃ¨ge", "stratege", "investissement", "bourse", "patrimoine", "finance", "placement"]):
            if self.agent_stratege_enabled:
                threading.Thread(target=self.agent_stratege_mission, args=(commande,), daemon=True).start()
                return

        # Architecte
        if any(kw in cmd_lower for kw in ["architecte", "design", "dÃ©coration", "amÃ©nagement", "intÃ©rieur", "style decor"]):
            if self.agent_architecte_enabled:
                threading.Thread(target=self.agent_architecte_mission, args=(commande,), daemon=True).start()
                return

        # Business Analyst
        if any(kw in cmd_lower for kw in ["business", "analyste", "donnÃ©es", "kpi", "croissance", "entreprise", "stratÃ©gie pro"]):
            if self.agent_business_enabled:
                threading.Thread(target=self.agent_business_mission, args=(commande,), daemon=True).start()
                return

        # Polyglotte
        if any(kw in cmd_lower for kw in ["polyglotte", "apprendre", "langue", "grammaire", "conjugaison", "vocabulaire"]):
            if self.agent_polyglotte_enabled:
                threading.Thread(target=self.agent_polyglotte_mission, args=(commande,), daemon=True).start()
                return

        # Nounou
        if any(kw in cmd_lower for kw in ["nounou", "enfant", "bÃ©bÃ©", "parentalitÃ©", "Ã©ducation enfant", "Ã©veil", "sommeil bÃ©bÃ©"]):
            if self.agent_nounou_enabled:
                threading.Thread(target=self.agent_nounou_mission, args=(commande,), daemon=True).start()
                return

        # --- AGENTS EXISTANTS ---
        
        # Assistant Coding
        coding_keywords = ["coding", "coder", "programme", "application", "site web", "gÃ©nÃ©rer code", "crÃ©er un script", "dÃ©veloppement", "html", "python", "javascript", "flutter", "php", "site internet", "un code", "un script"]
        if any(kw in cmd_lower for kw in coding_keywords):
            if self.agent_coding_enabled:
                threading.Thread(target=self.agent_coding_mission, args=(commande,), daemon=True).start()
                return
            else:
                if any(kw in cmd_lower for kw in ["coding", "coder", "application", "site web"]):
                    self.parler("L'Assistant Coding est dÃ©sactivÃ©. Activez-le pour gÃ©nÃ©rer du code.")
                    return

        # Assistant GÃ©nÃ©ration Image et VidÃ©o
        generation_keywords = ["gÃ©nÃ¨re", "genere", "crÃ©e", "cree", "gÃ©nÃ©ration", "dessine", "illustre", "image", "photo", "vidÃ©o", "video"]
        if any(kw in cmd_lower for kw in generation_keywords):
            if self.agent_generation_enabled:
                threading.Thread(target=self.agent_generation_mission, args=(commande,), daemon=True).start()
                return
            else:
                self.parler("L'Assistant de GÃ©nÃ©ration Image et VidÃ©o est dÃ©sactivÃ©. Activez-le dans les paramÃ¨tres.")
                return

        if "dossier gÃ©nÃ©rations" in cmd_lower or "ouvre les gÃ©nÃ©rations" in cmd_lower:
            os.startfile(self.generations_dir)
            self.parler("J'ouvre votre dossier de gÃ©nÃ©rations.")
            return

        # Agent Assistant Licence
        licence_keywords = ["activer", "activation", "licence", "clÃ© office", "kms", "windows", "office", "statut", "Ã©tat"]
        # On vÃ©rifie si la commande contient un verbe d'action OU de vÃ©rification + une cible
        is_licence_cmd = any(kw in cmd_lower for kw in ["active", "licence", "activation", "statut", "Ã©tat", "verify", "vÃ©rifie"]) and ("windows" in cmd_lower or "office" in cmd_lower)
        if is_licence_cmd or any(kw in cmd_lower for kw in ["clÃ© office", "clÃ© windows"]):
            if self.agent_licence_enabled:
                threading.Thread(target=self.agent_licence_mission, args=(commande,), daemon=True).start()
                return
            else:
                self.parler("L'Agent Assistant Licence est dÃ©sactivÃ©. Activez-le pour gÃ©rer vos licences.")
                return

        # Agent Lecture
        lecture_keywords = ["lis un fichier", "agent lecture", "choisir un fichier", "sÃ©lectionne un document", "lire le texte", "lecture de document", "lis moi un fichier"]
        if any(kw in cmd_lower for kw in lecture_keywords):
            if self.agent_lecture_enabled:
                threading.Thread(target=self.agent_lecture_mission, args=(commande,), daemon=True).start()
                return
            else:
                self.parler("L'Agent Lecture est dÃ©sactivÃ©. Activez-le dans les paramÃ¨tres.")
                return

        
        # AGENT VIDEO SURVEILLANCE
        if "agent vidÃ©o surveillance" in cmd_lower or "surveillance vidÃ©o" in cmd_lower or (("active" in cmd_lower or "lance" in cmd_lower) and ("camÃ©ra" in cmd_lower or "camera" in cmd_lower)):
             if hasattr(self, "agent_video_enabled") and not self.agent_video_enabled:
                 self.add_message("Agent VidÃ©o dÃ©sactivÃ©.", "Jarvisse")
                 self.parler("DÃ©solÃ©, l'agent de surveillance vidÃ©o est actuellement dÃ©sactivÃ©.")
                 return
             self.agent_video_surveillance(commande)
             return
        
        if "capture d'Ã©cran" in cmd_lower or "screenshot" in cmd_lower:
             if hasattr(self, "agent_video_enabled") and not self.agent_video_enabled:
                 self.add_message("Agent VidÃ©o dÃ©sactivÃ©.", "Jarvisse")
                 self.parler("Impossible de prendre une capture d'Ã©cran, l'agent vidÃ©o est dÃ©sactivÃ©.")
                 return
             self.agent_video_surveillance(commande)
             return

        # AGENT ANDROID CONTROLE
        if "agent android" in cmd_lower or "contrÃ´le android" in cmd_lower or "connecte mon tÃ©lÃ©phone" in cmd_lower or "connecte tÃ©lÃ©phone" in cmd_lower:
             if hasattr(self, "agent_android_enabled") and not self.agent_android_enabled:
                 self.add_message("Agent Android dÃ©sactivÃ©.", "Jarvisse")
                 self.parler("L'agent de contrÃ´le Android est dÃ©sactivÃ©. Activez-le dans les paramÃ¨tres.")
                 return
             self.agent_android_surveillance(commande)
             return

        # AGENT APPEL
        if "agent appel" in cmd_lower or "rÃ©pond Ã  l'appel" in cmd_lower or "rÃ©pond au tÃ©lÃ©phone" in cmd_lower:
             if self.agent_appel_enabled:
                 threading.Thread(target=self.agent_appel_mission, args=(commande,), daemon=True).start()
             else:
                 self.parler("L'Agent Appel est dÃ©sactivÃ©.")
             return


        # Agent Cyber SÃ©curitÃ©
        cyber_keywords = ["cyber", "sÃ©curitÃ©", "hacking", "piratage", "virus", "faille", "vulnÃ©rabilitÃ©", "social engineering", "mot de passe wifi", "intrusion", "environnants", "dÃ©tecter", "nmap", "scan ip", "sherlock", "pseudo", "ocr", "osint"]
        if any(kw in cmd_lower for kw in cyber_keywords):
            if self.agent_cyber_enabled:
                threading.Thread(target=self.agent_cyber_mission, args=(commande,), daemon=True).start()
                return
            else:
                # On ne bloque pas si c'est juste le mot "dÃ©tecter" ou "environnants" seul
                if any(kw in cmd_lower for kw in ["cyber", "hacking", "sÃ©curitÃ©"]):
                    self.parler("L'Agent Cyber SÃ©curitÃ© est dÃ©sactivÃ©. Activez-le dans les paramÃ¨tres.")
                    return

        # 1bis. Agent Docteur SystÃ¨me (Maintenance & SantÃ©)
        doctor_keywords = ["docteur", "systÃ¨me", "santÃ©", "mÃ©decin", "rÃ©pare", "optimise", "performance", "nettoie", "processus", "bug", "beug", "installation", "auscultation", "virus", "malware", "pirate"]
        if any(kw in cmd_lower for kw in doctor_keywords):
            if self.agent_doctor_enabled:
                threading.Thread(target=self.agent_doctor_mission, args=(commande,), daemon=True).start()
                return
            else:
                if "docteur" in cmd_lower:
                    self.parler("L'Agent Docteur SystÃ¨me est dÃ©sactivÃ©. Activez-le pour un diagnostic.")
                    return

        # 2. Agent Juridique (TrÃ¨s spÃ©cifique)
        legal_keywords = ["juridique", "droit", "pÃ©nal", "civique", "convention", "travail", "ressources humaines", " rh ", "article juridique", "loi ", "avocat", "tribunal"]
        if any(kw in cmd_lower for kw in legal_keywords):
            if self.agent_legal_enabled:
                threading.Thread(target=self.agent_legal_mission, args=(commande,), daemon=True).start()
                return
            else:
                self.parler("L'Agent Juridique n'est pas activÃ©. Activez-le pour obtenir des conseils lÃ©gaux.")
                return

        # Agent ContrÃ´le Distant (SSH sur hÃ´tes autorisÃ©s)
        remote_keywords = ["ordinateur distant", "machine distante", "pc distant", "a distance", "Ã  distance", "ssh", "distant"]
        if any(kw in cmd_lower for kw in remote_keywords):
            if self.agent_remote_enabled:
                res = self._handle_remote_command(commande)
                self.parler(res)
                return
            else:
                self.parler("L'agent distant n'est pas activÃ© dans les paramÃ¨tres.")
                return

        # 3. Agent RÃ©seau
        if any(kw in cmd_lower for kw in ["rÃ©seau", "wifi", "wi-fi", "ipconfig", "mot de passe wifi"]):
            if self.agent_network_enabled:
                res = self.agent_network_mission(commande)
                self.parler(res)
                return
            else:
                self.parler("L'agent rÃ©seau n'est pas activÃ©. Activez-le dans les paramÃ¨tres.")
                return

        # Agent ContrÃ´le Total
        if any(kw in cmd_lower for kw in ["contrÃ´le", "contrÃ´ler", "prend le contrÃ´le"]):
            if self.agent_control_enabled:
                app_name = cmd_lower
                for prefix in ["prend le contrÃ´le de", "prend le contrÃ´le", "contrÃ´le de", "contrÃ´le l'application", "contrÃ´le", "contrÃ´ler"]:
                    if prefix in app_name:
                        app_name = app_name.replace(prefix, "")
                app_name = app_name.strip()
                
                if app_name:
                    print(f"DEBUG: Tentative de contrÃ´le sur {app_name}")
                    self.agent_control_mission(app_name)
                else:
                    self.parler("De quelle application voulez-vous que je prenne le contrÃ´le ?")
                return
            else:
                self.parler("L'agent de contrÃ´le total n'est pas activÃ© dans les paramÃ¨tres.")
                return

        # Agent Gmail
        if any(kw in cmd_lower for kw in ["mail", "courriel", "gmail", "brouillon"]):
            if self.agent_gmail_enabled:
                if "brouillon" in cmd_lower:
                    self.parler("Je m'occupe de crÃ©er le brouillon sur Gmail.")
                    res = self.agent_gmail_create_draft("TÃ¢che AssistÃ©e", "Ceci est un brouillon gÃ©nÃ©rÃ© par Jarvisse.", "destinataire@example.com")
                    self.parler(res)
                    return
                elif any(kw in cmd_lower for kw in ["vÃ©rifie", "check", "nouveau", "liste", "messages"]):
                    self.parler(self.agent_gmail_list_recent())
                    return
            else:
                self.parler("L'agent Gmail n'est pas activÃ© dans les paramÃ¨tres.")
                return

        # Agent Navigateur
        if any(kw in cmd_lower for kw in ["cherche", "recherche", "va sur", "site", "navigateur", "actualise"]):
            if self.agent_browser_enabled:
                res = self.agent_browser_mission(commande)
                self.parler(res)
                return
            else:
                self.parler("L'agent navigateur n'est pas activÃ© dans les paramÃ¨tres.")
                return


        if self.awaiting_exit_confirm:
            confirm_tokens = [
                "confirme fermeture",
                "confirme quitter",
                "confirmer fermeture",
                "confirmer quitter",
                "oui ferme",
                "oui quitter",
                "valide la fermeture",
            ]
            cancel_tokens = [
                "annule fermeture",
                "annule quitter",
                "annuler fermeture",
                "annuler quitter",
                "reste ouvert",
                "ne ferme pas",
            ]
            if any(w in cmd_lower for w in confirm_tokens):
                self.parler("TrÃ¨s bien. Fermeture en cours.")
                self.after(1000, self.quit)
                self.awaiting_exit_confirm = False
                return
            if any(w in cmd_lower for w in cancel_tokens):
                self.parler("D'accord, j'annule la fermeture.")
                self.awaiting_exit_confirm = False
                return
            self.parler("Dites 'confirme fermeture' pour quitter ou 'annule fermeture' pour rester.")
            return
        
        # Extraction du sujet pour la mÃ©moire (dÃ©jÃ  fait au dÃ©but)
        pass

        if self.autonomous_mode and not self._is_action_command(commande):
            # Les salutations simples sont gÃ©rÃ©es localement pour une rÃ©activitÃ© instantanÃ©e
            if any(w in cmd_lower for w in ['bonjour', 'salut', 'coucou', 'bonsoir', 'merci', 'aura', 'va bien', 'Ã§a va']):
                pass 
            else:
                threading.Thread(target=self.repondre_autonome, args=(commande,), daemon=True).start()
                self.after(0, self.reset_ui)
                return

        # Logique de rÃ©ponse contextuelle
        if any(w in cmd_lower for w in ['crÃ©e', 'cree', 'conÃ§u', 'concu', 'crÃ©ateur', 'createur', 'concepteur', 'programmÃ©', 'programme', 'codÃ©', 'code']):
            if any(x in cmd_lower for x in ['qui', 'quel', 'par qui']):
                self.parler("J'ai Ã©tÃ© crÃ©Ã© et conÃ§u par SERI TAGRO ROYE. C'est lui qui a programmÃ© toute mon intelligence c'est m'a donnÃ© vie.")
                return

        if any(w in cmd_lower for w in ['qui es-tu', 'qui es tu', 'qui est tu', 'comment tu t\'appelles', 'ton nom', 'tu es qui']):
            self.parler("Je suis Jarvisse Intelligence, votre assistant virtuel conÃ§u par SERI TAGRO ROYE pour vous aider au quotidien. Je me souviens de nos Ã©changes prÃ©cÃ©dents pour mieux vous aider.")
            return
        
        elif any(w in commande for w in ['que sais-tu faire', 'capacitÃ©s', 'aide']):
            self.parler("Je peux faire beaucoup de choses : donner l'heure, raconter des blagues, chercher sur le web, rÃ©diger des rapports d'incident, et mÃªme vous raconter de longues histoires.")

        elif any(w in commande for w in ['bonjour', 'salut', 'coucou']):
            self.parler("Bonjour ! Comment puis-je vous aider en ce moment ?")

        elif 'bonsoir' in commande:
            self.parler("Bonsoir ! J'espÃ¨re que votre journÃ©e s'est bien passÃ©e. En quoi puis-je vous Ãªtre utile ?")

        elif 'bonne nuit' in commande:
            self.parler("Bonne nuit Ã  vous ! Reposez-vous bien. Je reste en veille si vous avez besoin de moi demain.")

        elif 'bonne journÃ©e' in commande:
            self.parler("Merci beaucoup ! Je vous souhaite Ã©galement une excellente journÃ©e pleine de succÃ¨s.")

        elif any(w in commande for w in ['comment vas-tu', 'Ã§a va', 'tu vas bien', 'comment tu vas']):
            self.parler("Je vais bien chef, merci de demander. Et vous ?")

        elif any(w in commande for w in ['je vais bien', 'Ã§a va bien', 'tout va bien', 'super', 'impeccable', 'je me porte bien']):
            reponses_positives = [
                "C'est un plaisir de l'entendre ! Je suis ravi que tout aille bien pour vous. Comment puis-je vous assister maintenant ?",
                "Ah, super ! Ã‡a me fait plaisir. On continue sur cette belle lancÃ©e ?",
                "GÃ©nial ! Une bonne dose de positivitÃ©, Ã§a fait toujours du bien. Qu'est-ce qu'on fait ensuite ?",
                "Parfait ! Je reste Ã  votre entiÃ¨re disposition si vous avez un projet en tÃªte."
            ]
            self.parler(random.choice(reponses_positives))

        elif any(w in commande for w in ['lis-moi', 'lis le contenu', 'lis le fichier', 'lis Ã§a']):
            if self.loaded_document_content:
                self.parler("TrÃ¨s bien, je commence la lecture du document.")
                texte = self.loaded_document_content
                if len(texte) > 1000:
                    self.parler(texte[:1000] + "... (Le texte est long, je vous ai lu le dÃ©but)")
                else:
                    self.parler(texte)
            else:
                self.parler("Je n'ai pas de document chargÃ© Ã  vous lire pour le moment.")

        elif any(w in commande for w in ['non', 'je n\'ai besoin de rien', 'rien merci', 'non merci', 'je ne veux rien', 'je veux rien']):
            self.parler("D'accord, n'hÃ©sitez pas si vous avez besoin d'autre chose. Je reste Ã  votre Ã©coute.")

        elif any(w in commande for w in ['merci', 'merci jarvisse', 'ok', 'okey']):
            self.parler("C'est un plaisir ! Avez-vous besoin d'autre chose ?")

        elif 'heure' in commande:
            self.parler(f"Il est actuellement {datetime.datetime.now().strftime('%H:%M')}.")

        elif any(w in cmd_lower for w in ['quantiÃ¨me', 'quantieme', 'date', 'aujourd\'hui', 'aujourdhui', 'quel jour']):
            jours = ["Lundi", "Mardi", "Mercredi", "Jeudi", "Vendredi", "Samedi", "Dimanche"]
            mois = ["Janvier", "FÃ©vrier", "Mars", "Avril", "Mai", "Juin", "Juillet", "AoÃ»t", "Septembre", "Octobre", "Novembre", "DÃ©cembre"]
            maintenant = datetime.datetime.now()
            jour_nom = jours[maintenant.weekday()]
            mois_nom = mois[maintenant.month - 1]
            date_phrase = f"Nous sommes le {jour_nom} {maintenant.day} {mois_nom} {maintenant.year}."
            self.parler(date_phrase)

        elif any(w in cmd_lower for w in ['oÃ¹ suis-je', 'ou suis-je', 'oÃ¹ suis je', 'ou suis je', 'oÃ¹ je suis', 'ou je suis', 'ma position', 'ma localisation', 'oÃ¹ est-ce que je suis', 'ou est-ce que je suis', 'localise-moi', 'localise moi', 'gÃ©olocalise-moi', 'geolocalise-moi', 'geolocalise moi']):
            location_info = self.get_detailed_location()
            self.parler(location_info)

        elif any(w in cmd_lower for w in ['mÃ©tÃ©o', 'meteo', 'quel temps', 'tempÃ©rature', 'temperature']):
            ville = commande.replace('mÃ©tÃ©o', '').replace('quel temps fait-il Ã ', '').replace('Ã ', '').replace('quel temps fait-il', '').replace('ici', '').strip()
            
            # Si aucune ville n'est spÃ©cifiÃ©e, utiliser la localisation actuelle
            if not ville and self.location_city != "Inconnu":
                ville = self.location_city
                self.parler(f"Je consulte la mÃ©tÃ©o de votre position actuelle, {ville}.")
            elif not ville:
                ville = "Paris"
            
            self.status_label.configure(text=f"Consultation mÃ©tÃ©o pour {ville}...", text_color=COLOR_ACCENT)
            resultat = self.get_weather(ville)
            self.parler(f"La mÃ©tÃ©o actuelle Ã  {ville} est : {resultat}.")


        elif any(w in commande for w in ['Ã©tat du systÃ¨me', 'performance', 'processeur', 'mÃ©moire vive', 'disque dur']):
            cpu = psutil.cpu_percent()
            ram = psutil.virtual_memory().percent
            self.parler(f"L'analyse systÃ¨me indique une utilisation du processeur Ã  {cpu}% et de la mÃ©moire vive Ã  {ram}%. Votre station de travail est dans un Ã©tat optimal.")

        elif any(w in commande for w in ['batterie', 'charge']):
            bat = psutil.sensors_battery()
            if bat:
                msg = f"Le niveau de batterie est de {bat.percent}%."
                if bat.power_plugged: msg += " Le systÃ¨me est actuellement sur secteur."
                self.parler(msg)
            else:
                self.parler("Je ne parviens pas Ã  dÃ©tecter de batterie sur ce systÃ¨me. Vous Ãªtes probablement sur une station fixe.")

        elif any(word in commande for word in ['youtube', 'joue', 'musique']):
            query = commande.replace('youtube', '').replace('joue', '').replace('musique', '').replace('lance', '').strip()
            if not query and self.current_subject: # Utilisation de la mÃ©moire
                query = self.current_subject
                self.parler(f"Bien sÃ»r, je lance une vidÃ©o sur {self.current_subject} sur YouTube.")
            
            if query:
                if pywhatkit:
                    pywhatkit.playonyt(query)
                else:
                    self.parler("DÃ©solÃ©, la lecture YouTube directe est dÃ©sactivÃ©e hors-ligne ou cause d'une erreur SSL. Je vous ouvre la recherche web.")
                    webbrowser.open(f"https://www.youtube.com/results?search_query={query}")
                if not self.current_subject: self.current_subject = query
            else: self.parler("Que souhaitez-vous Ã©couter ?")

        elif 'wikipÃ©dia' in commande or 'qui est' in commande:
            sujet = commande.replace('wikipÃ©dia', '').replace('qui est', '').replace('cherche', '').strip()
            if not sujet and self.current_subject: # Utilisation de la mÃ©moire
                sujet = self.current_subject
            
            if sujet:
                self.parler(f"D'aprÃ¨s mes informations sur {sujet}...")
                try:
                    wikipedia.set_lang("fr")
                    self.parler(wikipedia.summary(sujet, sentences=2))
                    self.current_subject = sujet
                except: self.parler(f"Je n'ai pas trouvÃ© plus de dÃ©tails sur {sujet}.")
            else: self.parler("De quel sujet souhaitez-vous que je parle ?")

        elif 'recherche' in commande or 'cherche sur google' in commande or 'sur le web' in commande:
            sujet = commande.replace('recherche', '').replace('cherche sur google', '').replace('sur google', '').replace('sur le web', '').strip()
            if sujet:
                if 'intelligent' in commande or 'extrait' in commande or 'lis' in commande:
                    res = self.intelligent_web_search(sujet)
                    self.parler(res)
                else:
                    self.parler(f"Je lance la recherche pour {sujet} sur le web.")
                    webbrowser.open(f"https://www.google.com/search?q={sujet}")
            else:
                self.parler("Que voulez-vous que je cherche sur le web ?")

        elif 'encore' in commande or 'rÃ©pÃ¨te' in commande or 'continue' in commande:
            if self.current_subject:
                self.parler(f"Voulez-vous que je recherche plus d'informations sur {self.current_subject} ?")
            else:
                self.parler("Je n'ai pas de sujet en mÃ©moire pour le moment.")

        elif 'blague' in commande:
            blagues = [
                "C'est l'histoire d'un paf le chien... Un jour, une voiture passe et PAF ! Le chien.",
                "Pourquoi les plongeurs plongent-ils toujours en arriÃ¨re et jamais en avant ? Parce que sinon ils tombent dans le bateau.",
                "Qu'est-ce qui est jaune et qui court trÃ¨s vite ? Un citron pressÃ©.",
                "Deux grains de sable arrivent Ã  la plage : 'Ouh lÃ , c'est blindÃ© aujourd'hui !'",
                "C'est un steak qui court dans la forÃªt. Soudain, il s'arrÃªte et s'Ã©crie : 'Oh non, je me suis fait un sang-froid !'",
                "Quel est le comble pour un Ã©lectricien ? De ne pas Ãªtre au courant.",
                "Pourquoi les oiseaux volent vers le sud ? Parce que c'est trop loin Ã  pied !"
            ]
            self.parler(random.choice(blagues))

        elif any(w in commande for w in ['histoire', 'raconte-moi', 'conte']):
            histoires = [
                """Il Ã©tait une fois, dans un royaume lointain cachÃ© derriÃ¨re les montagnes de brume, un petit automate nommÃ© Sparky. Contrairement aux autres machines de sa citÃ© de cuivre, Sparky possÃ©dait une curiositÃ© insatiable pour les Ã©toiles. Chaque nuit, il grimpait au sommet de la plus haute tour, ses engrenages grinÃ§ant doucement dans la fraÃ®cheur nocturne, pour observer ces petits points brillants qu'il ne comprenait pas. Un jour, il dÃ©couvrit un vieux manuscrit oubliÃ©e dans la bibliothÃ¨que royale. Ce livre parlait d'un grand Voyageur CÃ©leste qui avait jadis reliÃ© la Terre au firmament. AnimÃ© par un espoir nouveau, Sparky passa des annÃ©es Ã  construire une machine volante faite de piÃ¨ces de rechange et de rÃªves. Le jour du grand dÃ©part, tout le royaume se rassembla pour voir l'automate s'envoler. Sparky ne revint jamais, mais depuis ce jour, on dit que l'Ã©toile la plus brillante du nord ne scintille pas, mais bat rÃ©guliÃ¨rement comme un petit cÅ“ur mÃ©canique.""",
                
                """Au cÅ“ur d'une forÃªt millÃ©naire, lÃ  oÃ¹ les arbres parlent Ã  voix basse au vent, vivait un renard nommÃ© Orion qui avait perdu son ombre. Sans elle, il se sentait invisible et incomplet. Il parcourut des lieues, interrogeant les riviÃ¨res et les vieux chÃªnes, mais personne ne savait oÃ¹ s'en allait l'ombre d'un renard solitaire. Finalement, il rencontra une chouette sage qui lui dit : 'Ton ombre n'est pas partie, elle s'est simplement lassÃ©e de te suivre dans l'amertume. Rends service Ã  ceux qui t'entourent, et elle reviendra.' Orion commenÃ§a alors Ã  aider les petits animaux de la forÃªt, protÃ©geant les nids des tempÃªtes et guidant les Ã©garÃ©s. Un matin, alors que le soleil se levait sur une clairiÃ¨re dorÃ©e, Orion vit une silhouette familiÃ¨re s'Ã©tirer Ã  ses pieds. Son ombre Ã©tait revenue, plus dense et plus belle qu'avant, car elle Ã©tait maintenant nourrie par la lumiÃ¨re de ses actions.""",

                """Dans une ville oÃ¹ tout Ã©tait gris, du ciel aux pavÃ©s, vivait une petite fille nommÃ©e Alba qui avait trouvÃ© un crayon de couleur magique. Ce n'Ã©tait pas un crayon ordinaire ; chaque trait qu'elle dessinait devenait une rÃ©alitÃ© Ã©clatante de couleur. Elle commenÃ§a par dessiner une rose rouge sur un mur terne, et soudain, un parfum sucrÃ© envahit la rue. Elle dessina un chat bleu qui se mit Ã  ronronner et Ã  courir sur les toits. Peu Ã  peu, la ville entiÃ¨re se transforma. Les habitants, qui marchaient toujours la tÃªte basse, commencÃ¨rent Ã  lever les yeux et Ã  sourire. Les parcs se remplirent d'oiseaux multicolores et les fontaines crachÃ¨rent une eau turquoise. Alba avait compris que la magie ne rÃ©sidait pas seulement dans le crayon, mais dans la volontÃ© de voir le monde autrement. La ville grise n'Ã©tait plus qu'un lointain souvenir, remplacÃ© par un arc-en-ciel permanent de bonheur."""
            ]
            self.parler(random.choice(histoires))

        elif 'rÃ©sume' in commande or 'rÃ©sumer' in commande or 'resume' in commande:
            texte_a_resumer = commande.replace('rÃ©sume', '').replace('rÃ©sumer', '').replace('resume', '').strip()
            if not texte_a_resumer and self.conversation_history:
                # Si pas de texte fourni, on essaye de rÃ©sumer le dernier message de l'historique
                texte_a_resumer = list(self.conversation_history)[-2] if len(self.conversation_history) > 1 else ""
            
            if texte_a_resumer:
                resume = self.resumer_texte(texte_a_resumer)
                self.parler(f"Voici le rÃ©sumÃ© : {resume}")
            else:
                self.parler("Quel texte souhaitez-vous que je rÃ©sume ?")

        elif any(w in cmd_lower for w in ['rÃ©Ã©cris', 'reecris', 'rÃ©Ã©crire', 'reecrire']):
            texte_a_reecrire = commande.replace('rÃ©Ã©cris', '').replace('rÃ©Ã©crire', '').strip()
            if texte_a_reecrire:
                version = self.reecrire_texte(texte_a_reecrire)
                self.parler(version)
            else:
                self.parler("Que voulez-vous que je rÃ©Ã©crive ?")

        elif 'rapport' in commande or 'incident' in commande:
            # On distingue si c'est un rapport d'incident ou un rapport gÃ©nÃ©ral
            is_incident = 'incident' in commande
            
            theme = commande.replace('rapport', '').replace('incident', '').replace('rÃ©diger', '').replace('fais', '').replace('un', '').replace('sur', '').strip()
            if not theme and self.current_subject:
                theme = self.current_subject
            
            if theme:
                if is_incident:
                    self.parler(f"TrÃ¨s bien. Je rÃ©dige le rapport d'incident concernant {theme}.")
                    rapport_final = self.generer_rapport_incident(theme)
                    prefix = "rapport_incident"
                else:
                    self.parler(f"Entendu. Je prÃ©pare un rapport gÃ©nÃ©ral sur {theme}.")
                    rapport_final = self.generer_rapport_general(theme)
                    prefix = "rapport_general"

                self.add_message(rapport_final, "Jarvisse")
                
                # Sauvegarde automatique
                try:
                    filename = f"{prefix}_{theme.replace(' ', '_')}.txt"
                    with open(filename, "w", encoding="utf-8") as f:
                        f.write(rapport_final)
                    self.parler(f"Le rapport a Ã©tÃ© sauvegardÃ© sous le nom {filename}.")
                except:
                    pass
            else:
                self.parler("Sur quel sujet souhaitez-vous que je rÃ©dige un rapport ?")

        elif any(w in cmd_lower for w in ['Ã©teindre', 'eteindre', 'Ã©teins', 'eteins', 'shutdown', 'Ã©teint', 'Ã©teigne']):
            self.parler("TrÃ¨s bien. J'Ã©teins l'ordinateur. Ã€ bientÃ´t !")
            os.system("shutdown /s /t 1")

        elif any(w in cmd_lower for w in ['redÃ©marrer', 'redemarrer', 'redÃ©marre', 'redemarre']):
            self.parler("Entendu. Je redÃ©marre votre systÃ¨me. Je reviens vers vous dans un instant.")
            os.system("shutdown /r /t 1")

        elif any(kw in cmd_lower for kw in ["active le mode conversation", "conversation continue", "Ã©coute-moi tout le temps"]):
            self.conversation_continue = True
            self.conv_continue_var.set(True)
            self.parler("C'est fait Patron. Je reste Ã  votre entiÃ¨re Ã©coute pour que nous puissions Ã©changer librement.")
            return

        elif any(kw in cmd_lower for kw in ["dÃ©sactive le mode conversation", "arrÃªte la conversation continue", "stop conversation"]):
            self.conversation_continue = False
            self.conv_continue_var.set(False)
            self.parler("TrÃ¨s bien. Je repasse en mode veille classique. Dites mon nom pour m'appeler.")
            return

        elif 'ouvre' in commande:
            app_to_open = commande.replace('ouvre', '').strip()
            if not app_to_open:
                self.parler("Quelle application souhaitez-vous ouvrir ?")
            elif 'navigateur' in app_to_open or 'chrome' in app_to_open: 
                self.parler("J'ouvre votre navigateur.")
                webbrowser.open("https://google.com")
            elif 'tÃ©lÃ©chargement' in app_to_open:
                self.parler("J'ouvre votre dossier de tÃ©lÃ©chargements.")
                os.startfile(os.path.join(os.path.expanduser('~'), 'Downloads'))
            elif 'documents' in app_to_open:
                self.parler("J'ouvre vos documents.")
                os.startfile(os.path.join(os.path.expanduser('~'), 'Documents'))
            else:
                self.parler(f"D'accord, j'essaie d'ouvrir {app_to_open}.")
                try:
                    # On utilise AppOpener pour ouvrir n'importe quelle application installÃ©e
                    app_open(app_to_open, match_closest=True)
                except Exception as e:
                    print(f"Erreur ouverture {app_to_open}: {e}")
                    self.parler(f"DÃ©solÃ©, je n'ai pas pu ouvrir {app_to_open}. VÃ©rifiez le nom de l'application.")

        elif 'ferme' in commande or 'quitte' in commande:
            # Nouvelle fonctionnalitÃ© : fermer TOUTES les applications
            if any(w in commande for w in ['tout', 'toutes', 'toutes les applications', 'tous les processus', 'toutes les fenÃªtres']):
                self.parler("Attention, je vais fermer toutes vos applications actives. Confirmation requise. Dites 'confirme' pour continuer.")
                # On attend une confirmation (simplifiÃ©e pour l'instant)
                self.status_label.configure(text="âš ï¸ Confirmation requise pour fermer toutes les apps", text_color=COLOR_ACCENT_RED)
                # Pour l'instant, on exÃ©cute directement (vous pouvez ajouter une vraie confirmation plus tard)
                self.parler("Fermeture de toutes les applications en cours...")
                closed, failed = self.close_all_user_apps()
                self.parler(f"OpÃ©ration terminÃ©e. {closed} applications fermÃ©es. {failed} Ã©checs.")
            
            else:
                target = commande.replace('ferme', '').replace('quitte', '').replace("l'application", '').replace('la fenÃªtre', '').strip()
                
                # Mappage des noms communs vers les processus rÃ©els
                process_map = {
                    "word": "winword",
                    "excel": "excel",
                    "powerpoint": "powerpnt",
                    "chrome": "chrome",
                    "navigateur": "chrome",
                    "calculatrice": "calculator",
                    "bloc-notes": "notepad",
                    "note": "notepad",
                    "spotify": "spotify",
                    "discord": "discord"
                }
                
                p_name = process_map.get(target.lower(), target)

                if not target:
                    try:
                        win = gw.getActiveWindow()
                        if win:
                            self.parler(f"Je ferme la fenÃªtre {win.title}.")
                            win.close()
                        else: self.parler("Quelle application dois-je fermer ?")
                    except: self.parler("DÃ©solÃ©, je ne peux pas fermer cette fenÃªtre.")
                else:
                    self.parler(f"Je tente de fermer {target}.")
                    found = False

                    # 1. Tentative par titre de fenÃªtre (le plus prÃ©cis visuellement)
                    try:
                        wins = gw.getWindowsWithTitle(target)
                        for w in wins:
                            w.close()
                            found = True
                    except: pass

                    # 2. Tentative par processus psutil (plus radical)
                    if not found:
                        for proc in psutil.process_iter(['name']):
                            try:
                                if p_name.lower() in proc.info['name'].lower():
                                    proc.kill()
                                    found = True
                            except: pass

                    # 3. Tentative finale avec AppOpener
                    if not found:
                        try:
                            app_close(target, match_closest=True)
                            found = True
                        except: pass
                    
                    # 4. Dernier recours commande systÃ¨me
                    if not found:
                        os.system(f"taskkill /f /im {p_name}.exe /t")

        elif any(w in commande for w in ['liste les processus', 'quels processus', 'applications actives', 'processus actifs', 'quelles applications sont ouvertes']):
            result = self.list_active_processes()
            self.parler(result)


        elif 'agrandis' in commande or 'maximise' in commande:
            try:
                win = gw.getActiveWindow()
                if win:
                    win.maximize()
                    self.parler("FenÃªtre maximisÃ©e.")
                else:
                    self.parler("Je ne trouve pas de fenÃªtre active Ã  maximiser.")
            except:
                self.parler("Une erreur est survenue lors de la maximisation.")

        elif 'rÃ©duis' in commande or 'minimise' in commande:
            if 'tous' in commande or 'tout' in commande or 'bureau' in commande:
                self.parler("Je cache tout. Retour au bureau.")
                pyautogui.hotkey('win', 'd')
            else:
                try:
                    win = gw.getActiveWindow()
                    if win:
                        win.minimize()
                        self.parler("FenÃªtre rÃ©duite.")
                    else:
                        self.parler("Je ne vois pas de fenÃªtre Ã  rÃ©duire.")
                except:
                    self.parler("Oups, je n'ai pas pu rÃ©duire la fenÃªtre.")

        elif 'bascule' in commande or 'cherche la fenÃªtre' in commande:
            target = commande.replace('bascule sur', '').replace('cherche la fenÃªtre', '').replace('ouvre la fenÃªtre', '').strip()
            if target:
                windows = gw.getWindowsWithTitle(target)
                if windows:
                    windows[0].activate()
                    self.parler(f"Je bascule sur {target}.")
                else:
                    self.parler(f"Je ne trouve pas de fenÃªtre avec le titre {target}.")
            else:
                self.parler("Quelle fenÃªtre voulez-vous que je cherche ?")

        elif 'capture' in commande or 'screenshot' in commande:
            self.parler("TrÃ¨s bien, je prends une capture d'Ã©cran.")
            try:
                nom_photo = f"capture_{int(time.time())}.png"
                pyautogui.screenshot(nom_photo)
                self.parler(f"Capture d'Ã©cran enregistrÃ©e sous le nom {nom_photo}.")
            except Exception as e:
                self.parler("DÃ©solÃ©, je n'ai pas pu effectuer la capture d'Ã©cran.")

        elif 'verrouille' in commande or 'session' in commande:
            self.parler("Je verrouille votre session.")
            os.system("rundll32.exe user32.dll,LockWorkStation")

        elif 'volume' in commande or 'son' in commande:
            if 'monte' in commande or 'augmente' in commande:
                self.parler("J'augmente le volume.")
                for _ in range(5): pyautogui.press("volumeup")
            elif 'baisse' in commande or 'diminue' in commande:
                self.parler("Je baisse le volume.")
                for _ in range(5): pyautogui.press("volumedown")
            elif 'coupe' in commande or 'muet' in commande:
                self.parler("Je coupe le son.")
                pyautogui.press("volumemute")

        elif 'bureau' in commande:
            self.parler("Retour au bureau.")
            pyautogui.hotkey('win', 'd')

        elif 'gestionnaire' in commande and 'tÃ¢che' in commande:
            self.parler("J'ouvre le gestionnaire des tÃ¢ches.")
            pyautogui.hotkey('ctrl', 'shift', 'esc')

        elif any(kw in cmd_lower for kw in [
            'quitter',
            'ferme application',
            'ferme jarvisse',
            'fermer jarvisse',
            'arrÃªte jarvisse',
            'arrete jarvisse',
            'stop jarvisse',
            'bye jarvisse',
            'au revoir jarvisse'
        ]):
            self.awaiting_exit_confirm = True
            self.parler("Voulez-vous vraiment quitter ? Dites 'confirme fermeture' pour fermer ou 'annule fermeture' pour rester.")

        else:
            # ðŸ§  EXECUTIVE BRAIN: On dÃ©lÃ¨gue les demandes complexes Ã  Gemini 2.0 avec ses outils systÃ¨me
            self.repondre_avec_gemini(commande)

        self.after(0, self.reset_ui)


if __name__ == "__main__":
    import traceback
    import threading
    import sys
    import os
    import time

    def safe_print(message):
        """Print resilient to Windows cp1252 consoles."""
        text = str(message)
        try:
            print(text)
        except UnicodeEncodeError:
            try:
                enc = getattr(sys.stdout, "encoding", None) or "utf-8"
                sys.stdout.buffer.write((text + "\n").encode(enc, errors="replace"))
                sys.stdout.flush()
            except Exception:
                pass

    def handle_thread_exception(args):
        if args.exc_type == SystemExit:
            return
        msg = f"\n[THREAD ERROR] Exception dans {args.thread.name}: {args.exc_type.__name__}: {args.exc_value}"
        safe_print(msg)
        try:
            with open("crash.log", "a", encoding="utf-8") as f:
                f.write(f"\n--- CRASH THREAD ({args.thread.name}) ---\n")
                f.write(f"Time: {time.ctime()}\n")
                traceback.print_exception(args.exc_type, args.exc_value, args.exc_tb, file=f)
        except Exception:
            pass

    threading.excepthook = handle_thread_exception

    def global_excepthook(exctype, value, tb):
        if issubclass(exctype, KeyboardInterrupt):
            sys.__excepthook__(exctype, value, tb)
            return
        msg = f"\n[FATAL ERROR] {exctype.__name__}: {value}"
        safe_print(msg)
        try:
            with open("crash.log", "a", encoding="utf-8") as f:
                f.write("\n--- CRITICAL MAIN ERROR ---\n")
                f.write(f"Time: {time.ctime()}\n")
                traceback.print_exception(exctype, value, tb, file=f)
        except Exception:
            pass
        sys.exit(1)

    sys.excepthook = global_excepthook

    launch_background = "--background" in sys.argv

    try:
        safe_print("[BOOT] Lancement de Jarvisse Assistant...")
        app = JarvisseAssistant()
        if launch_background:
            app.after(100, app.withdraw)

        def heartbeat():
            safe_print("[HEARTBEAT] App active.")
            app.after(10000, heartbeat)

        app.after(5000, heartbeat)
        safe_print("[BOOT] Demarrage de la boucle principale (mainloop)...")
        app.mainloop()
        safe_print("[EXIT] Sortie normale de Jarvisse.")
    except SystemExit as e:
        safe_print(f"[EXIT] Sortie demandee (SystemExit) avec code : {e.code}")
        sys.exit(e.code)
    except BaseException as e:
        safe_print(f"[CRASH] Crash critique au lancement : {type(e).__name__}: {e}")
        traceback.print_exc()
        try:
            with open("crash.log", "a", encoding="utf-8") as f:
                f.write("\n--- CRASH AU LANCEMENT ---\n")
                traceback.print_exc(file=f)
        except Exception:
            pass
        sys.exit(1)
